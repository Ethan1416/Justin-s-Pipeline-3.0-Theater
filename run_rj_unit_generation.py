#!/usr/bin/env python3
"""
Romeo and Juliet Complete Unit Generator
=========================================

Generates the complete 6-week Romeo and Juliet unit (30 days) from scratch
using hardcoded agents for all generation, formatting, and validation.

Usage:
    python run_rj_unit_generation.py                    # Generate all 30 days
    python run_rj_unit_generation.py --week 1           # Generate Week 1 only
    python run_rj_unit_generation.py --day 1            # Generate Day 1 only
    python run_rj_unit_generation.py --enhance-pptx     # Also add trivia banners
"""

import sys
import json
import argparse
from pathlib import Path
from datetime import datetime
from typing import Any, Dict, List, Optional

# Add project to path
sys.path.insert(0, str(Path(__file__).parent))

from agents import (
    # Romeo & Juliet Generation Agents
    SceneCutterAgent,
    SceneSummaryGeneratorAgent,
    ReadingDayGeneratorAgent,
    ActivityDayGeneratorAgent,
    DifferentiationSelectorAgent,
    WeekPlannerAgent,
    TextExcerptSelectorAgent,
    RJUnitValidatorAgent,
    # Daily Generation Agents
    LessonPlanGeneratorAgent,
    WarmupGeneratorAgent,
    PowerPointGeneratorAgent,
    HandoutGeneratorAgent,
    JournalExitGeneratorAgent,
    # Validation Agents
    TruncationValidatorAgent,
    StructureValidatorAgent,
    TimingValidatorAgent,
    # Slide Enhancement Agents
    SlideContentEnhancerAgent,
    SlideEnhancementFormatterAgent,
    SlideEnhancementValidatorAgent,
    # Content Optimization Agents
    ContentBalanceOrchestratorAgent,
    TruncationDetectorAgent,
    SlideContentCondensationAgent,
)

# Import PowerPoint generator
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Pipeline root
PIPELINE_ROOT = Path(__file__).parent
PRODUCTION_DIR = PIPELINE_ROOT / "production" / "Unit_3_Romeo_and_Juliet"


# =============================================================================
# HARDCODED UNIT STRUCTURE
# =============================================================================

# HARDCODED: Complete 6-week unit structure (30 days)
UNIT_STRUCTURE = {
    1: {  # Week 1: Act 1 - Introduction & Setup
        "act": 1,
        "theme": "Introduction to Shakespeare & Romeo and Juliet",
        "days": [
            {"day": 1, "type": "activity", "topic": "Meet the Bard: Shakespeare's World", "scenes": []},
            {"day": 2, "type": "reading", "topic": "The Prologue: 14 Lines of Fate", "scenes": ["Prologue"]},
            {"day": 3, "type": "reading", "topic": "Act 1 Scene 1: The Feud Begins", "scenes": ["1.1"]},
            {"day": 4, "type": "reading", "topic": "Act 1 Scenes 2-3: Party Preparations", "scenes": ["1.2", "1.3"]},
            {"day": 5, "type": "activity", "topic": "Character Mapping & Predictions", "scenes": []},
        ]
    },
    2: {  # Week 2: Act 1 continued & Act 2 - Love Blooms
        "act": "1-2",
        "theme": "Love at First Sight",
        "days": [
            {"day": 6, "type": "reading", "topic": "Act 1 Scenes 4-5: The Masquerade Ball", "scenes": ["1.4", "1.5"]},
            {"day": 7, "type": "activity", "topic": "Queen Mab Speech Analysis", "scenes": []},
            {"day": 8, "type": "reading", "topic": "Act 2 Scene 2: The Balcony Scene", "scenes": ["2.2"]},
            {"day": 9, "type": "reading", "topic": "Act 2 Scenes 3-4: Friar's Plan", "scenes": ["2.3", "2.4"]},
            {"day": 10, "type": "activity", "topic": "Love Poetry Workshop", "scenes": []},
        ]
    },
    3: {  # Week 3: Act 2 continued & Act 3 - Turning Point
        "act": "2-3",
        "theme": "From Comedy to Tragedy",
        "days": [
            {"day": 11, "type": "reading", "topic": "Act 2 Scenes 5-6: The Secret Wedding", "scenes": ["2.5", "2.6"]},
            {"day": 12, "type": "reading", "topic": "Act 3 Scene 1: Mercutio & Tybalt", "scenes": ["3.1"]},
            {"day": 13, "type": "activity", "topic": "Stage Combat Workshop", "scenes": []},
            {"day": 14, "type": "reading", "topic": "Act 3 Scenes 2-3: Banishment", "scenes": ["3.2", "3.3"]},
            {"day": 15, "type": "reading", "topic": "Act 3 Scenes 4-5: Forced Marriage", "scenes": ["3.4", "3.5"]},
        ]
    },
    4: {  # Week 4: Act 4 - Desperate Measures
        "act": 4,
        "theme": "Desperate Times, Desperate Measures",
        "days": [
            {"day": 16, "type": "activity", "topic": "Character Hot Seat: Juliet", "scenes": []},
            {"day": 17, "type": "reading", "topic": "Act 4 Scene 1: The Friar's Potion", "scenes": ["4.1"]},
            {"day": 18, "type": "reading", "topic": "Act 4 Scene 3: Juliet's Soliloquy", "scenes": ["4.3"]},
            {"day": 19, "type": "activity", "topic": "Soliloquy Performance Workshop", "scenes": []},
            {"day": 20, "type": "reading", "topic": "Act 4 Scene 5: The Discovery", "scenes": ["4.5"]},
        ]
    },
    5: {  # Week 5: Act 5 - The Tragedy Unfolds
        "act": 5,
        "theme": "Fate vs. Free Will",
        "days": [
            {"day": 21, "type": "reading", "topic": "Act 5 Scene 1: Romeo's Dream", "scenes": ["5.1"]},
            {"day": 22, "type": "activity", "topic": "Dramatic Irony Analysis", "scenes": []},
            {"day": 23, "type": "reading", "topic": "Act 5 Scene 3 Part 1: The Tomb", "scenes": ["5.3a"]},
            {"day": 24, "type": "reading", "topic": "Act 5 Scene 3 Part 2: The End", "scenes": ["5.3b"]},
            {"day": 25, "type": "activity", "topic": "Theme Discussion: Who's to Blame?", "scenes": []},
        ]
    },
    6: {  # Week 6: Performance & Assessment
        "act": "performance",
        "theme": "Performance & Reflection",
        "days": [
            {"day": 26, "type": "activity", "topic": "Scene Selection & Rehearsal", "scenes": []},
            {"day": 27, "type": "activity", "topic": "Blocking & Staging Workshop", "scenes": []},
            {"day": 28, "type": "activity", "topic": "Dress Rehearsal", "scenes": []},
            {"day": 29, "type": "activity", "topic": "Performance Day", "scenes": []},
            {"day": 30, "type": "activity", "topic": "Reflection & Unit Assessment", "scenes": []},
        ]
    }
}

# HARDCODED: Learning objectives by day type
LEARNING_OBJECTIVES = {
    "reading": [
        "Analyze Shakespeare's language and dramatic techniques",
        "Identify key plot developments and character motivations",
        "Connect themes to real-world experiences"
    ],
    "activity": [
        "Apply theatrical concepts through hands-on practice",
        "Demonstrate understanding through creative expression",
        "Collaborate effectively with peers"
    ]
}

# HARDCODED: Vocabulary by week
VOCABULARY_BY_WEEK = {
    1: [
        {"term": "Prologue", "definition": "An introduction to a play that sets up the story"},
        {"term": "Feud", "definition": "A prolonged and bitter quarrel between families or groups"},
        {"term": "Iambic Pentameter", "definition": "A line of verse with five metrical feet, each with one unstressed syllable followed by one stressed"},
        {"term": "Sonnet", "definition": "A 14-line poem with a specific rhyme scheme"},
    ],
    2: [
        {"term": "Soliloquy", "definition": "A speech in which a character reveals their thoughts to the audience"},
        {"term": "Aside", "definition": "A remark spoken to the audience that other characters do not hear"},
        {"term": "Metaphor", "definition": "A comparison without using 'like' or 'as'"},
        {"term": "Imagery", "definition": "Descriptive language that appeals to the senses"},
    ],
    3: [
        {"term": "Dramatic Irony", "definition": "When the audience knows something the characters do not"},
        {"term": "Foil", "definition": "A character who contrasts with another to highlight qualities"},
        {"term": "Climax", "definition": "The turning point or most intense moment in a story"},
        {"term": "Tragedy", "definition": "A serious drama with an unhappy ending"},
    ],
    4: [
        {"term": "Monologue", "definition": "A long speech by one character"},
        {"term": "Foreshadowing", "definition": "Hints about what will happen later in the story"},
        {"term": "Catharsis", "definition": "Emotional release experienced by the audience"},
        {"term": "Hubris", "definition": "Excessive pride that leads to downfall"},
    ],
    5: [
        {"term": "Denouement", "definition": "The final resolution of a story's plot"},
        {"term": "Tragic Hero", "definition": "A protagonist whose downfall is caused by a fatal flaw"},
        {"term": "Theme", "definition": "The central message or insight about life in a work"},
        {"term": "Motif", "definition": "A recurring element that has symbolic significance"},
    ],
    6: [
        {"term": "Blocking", "definition": "The precise staging of actors for a scene"},
        {"term": "Stage Presence", "definition": "The ability to command attention while performing"},
        {"term": "Projection", "definition": "Speaking loudly and clearly enough to be heard"},
        {"term": "Ensemble", "definition": "A group of performers working together"},
    ]
}


# =============================================================================
# GENERATOR CLASS
# =============================================================================

class RomeoJulietUnitGenerator:
    """
    HARDCODED generator for the complete Romeo and Juliet unit.

    Uses hardcoded agents for all generation, validation, and formatting.
    """

    def __init__(self, enhance_pptx: bool = True, optimize_content: bool = True, verbose: bool = False):
        self.enhance_pptx = enhance_pptx
        self.optimize_content = optimize_content
        self.verbose = verbose

        # Initialize agents
        self.scene_cutter = SceneCutterAgent()
        self.scene_summary = SceneSummaryGeneratorAgent()
        self.reading_day = ReadingDayGeneratorAgent()
        self.activity_day = ActivityDayGeneratorAgent()
        self.differentiation = DifferentiationSelectorAgent()
        self.week_planner = WeekPlannerAgent()
        self.unit_validator = RJUnitValidatorAgent()

        # Validation agents
        self.truncation_validator = TruncationValidatorAgent("truncation_validator")
        self.structure_validator = StructureValidatorAgent("structure_validator")

        # Slide enhancement agents
        if enhance_pptx:
            self.slide_enhancer = SlideContentEnhancerAgent()
            self.slide_formatter = SlideEnhancementFormatterAgent()
            self.slide_validator = SlideEnhancementValidatorAgent()

        # Content optimization agents (for condensing slides and elaborating notes)
        if optimize_content:
            self.content_optimizer = ContentBalanceOrchestratorAgent()
            self.truncation_detector = TruncationDetectorAgent()
            self.content_condenser = SlideContentCondensationAgent()

        # Statistics
        self.stats = {
            "days_generated": 0,
            "pptx_generated": 0,
            "handouts_generated": 0,
            "validation_passed": 0,
            "validation_failed": 0,
            "trivia_added": 0,
            "content_optimized": 0,
            "truncations_fixed": 0,
        }

    def generate_unit(self, weeks: List[int] = None, days: List[int] = None):
        """Generate the complete unit or specified portions."""
        print("=" * 70)
        print("ROMEO AND JULIET UNIT GENERATOR")
        print("=" * 70)
        print(f"Started: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        print(f"Trivia Enhancement: {'Enabled (20pt font)' if self.enhance_pptx else 'Disabled'}")
        print(f"Content Optimization: {'Enabled (condensed slides, elaborated notes)' if self.optimize_content else 'Disabled'}")
        print()

        # Determine what to generate
        if days:
            target_days = days
        elif weeks:
            target_days = []
            for week in weeks:
                if week in UNIT_STRUCTURE:
                    for day_info in UNIT_STRUCTURE[week]["days"]:
                        target_days.append(day_info["day"])
        else:
            target_days = list(range(1, 31))  # All 30 days

        print(f"Generating {len(target_days)} days: {target_days}")
        print()

        # Create production directory
        PRODUCTION_DIR.mkdir(parents=True, exist_ok=True)

        # Generate each day
        for day_num in target_days:
            self._generate_day(day_num)

        # Run unit validation
        print()
        print("-" * 70)
        print("UNIT VALIDATION")
        print("-" * 70)
        self._validate_unit()

        # Print summary
        self._print_summary()

        return self.stats

    def _get_day_info(self, day_num: int) -> Optional[Dict]:
        """Get day info from unit structure."""
        for week_num, week_data in UNIT_STRUCTURE.items():
            for day_info in week_data["days"]:
                if day_info["day"] == day_num:
                    return {
                        "week": week_num,
                        "act": week_data["act"],
                        "theme": week_data["theme"],
                        **day_info
                    }
        return None

    def _generate_day(self, day_num: int):
        """Generate all materials for a single day."""
        day_info = self._get_day_info(day_num)
        if not day_info:
            print(f"  Day {day_num}: [ERROR] Not found in unit structure")
            return

        print(f"Day {day_num:02d}: {day_info['topic']}")

        # Create day directory
        day_dir = PRODUCTION_DIR / f"Day_{day_num:02d}"
        day_dir.mkdir(parents=True, exist_ok=True)

        # Generate day input data
        day_data = self._create_day_input(day_info)

        # Save input JSON
        input_path = day_dir / f"day{day_num:02d}_input.json"
        with open(input_path, 'w', encoding='utf-8') as f:
            json.dump(day_data, f, indent=2)

        # Generate lesson plan markdown
        lesson_path = day_dir / f"Day{day_num:02d}_LessonPlan.md"
        self._generate_lesson_plan(day_data, lesson_path)
        print(f"  [OK] Lesson plan: {lesson_path.name}")

        # Generate handout
        handout_path = day_dir / f"Day{day_num:02d}_Handout.md"
        self._generate_handout(day_data, handout_path)
        print(f"  [OK] Handout: {handout_path.name}")
        self.stats["handouts_generated"] += 1

        # Generate PowerPoint
        if PPTX_AVAILABLE:
            pptx_path = day_dir / f"Unit3_Day{day_num:02d}_{day_info['topic'].replace(' ', '_').replace(':', '')[:30]}.pptx"
            self._generate_pptx(day_data, pptx_path)
            print(f"  [OK] PowerPoint: {pptx_path.name}")
            self.stats["pptx_generated"] += 1

        # Validate day
        validation = self._validate_day(day_data)
        if validation["valid"]:
            print(f"  [PASS] Validation")
            self.stats["validation_passed"] += 1
        else:
            print(f"  [FAIL] Validation: {validation.get('errors', [])}")
            self.stats["validation_failed"] += 1

        self.stats["days_generated"] += 1
        print()

    def _create_day_input(self, day_info: Dict) -> Dict:
        """Create input data structure for a day."""
        week = day_info["week"]

        # Get vocabulary for this week
        vocab = VOCABULARY_BY_WEEK.get(week, [])

        # Get learning objectives based on day type
        objectives = LEARNING_OBJECTIVES.get(day_info["type"], [])

        # Build standards
        standards = [
            {"code": "CCSS.ELA-LITERACY.RL.9-10.3", "full_text": "Analyze how complex characters develop over the course of a text"},
            {"code": "CCSS.ELA-LITERACY.RL.9-10.5", "full_text": "Analyze how an author's choices concerning structure create effects"},
        ]

        if day_info["type"] == "activity":
            standards.append({
                "code": "NCAS.TH.Cr1.1.I",
                "full_text": "Apply basic research to construct ideas about theatrical works"
            })

        # Build warmup based on day type
        if day_info["type"] == "reading":
            warmup = {
                "name": "Quote of the Day",
                "instructions": f"Display a key quote from {day_info['scenes'][0] if day_info['scenes'] else 'the play'}. Students write 2-3 sentences predicting its meaning.",
                "duration_minutes": 5,
                "connection_to_lesson": f"Prepares students for {day_info['topic']}"
            }
        else:
            warmup = {
                "name": "Theater Tableau",
                "instructions": "In groups of 3-4, create a frozen image representing a key moment from the play so far.",
                "duration_minutes": 5,
                "connection_to_lesson": "Activates prior knowledge and gets students moving"
            }

        # Build activity
        activity = {
            "name": day_info["topic"],
            "instructions": f"Today we focus on {day_info['topic']}.",
            "duration_minutes": 25,
            "differentiation_options": []
        }

        # Build content points
        content_points = []
        if day_info["scenes"]:
            for scene in day_info["scenes"]:
                content_points.append({
                    "point": f"Scene {scene}: Key moments and language analysis",
                    "expanded": [
                        f"Reading and analysis of {scene}",
                        "Character motivation discussion",
                        "Language and imagery exploration"
                    ]
                })
        else:
            content_points.append({
                "point": day_info["topic"],
                "expanded": [
                    "Introduction to the activity",
                    "Guided practice",
                    "Student work time"
                ]
            })

        return {
            "unit": {
                "number": 3,
                "name": "Romeo and Juliet"
            },
            "week": week,
            "day": day_info["day"],
            "day_type": day_info["type"],
            "topic": day_info["topic"],
            "theme": day_info["theme"],
            "scenes": day_info["scenes"],
            "standards": standards,
            "learning_objectives": objectives,
            "vocabulary": vocab[:4],  # First 4 vocab words
            "prior_knowledge": f"Students have completed Days 1-{day_info['day']-1}" if day_info['day'] > 1 else "None required",
            "warmup": warmup,
            "activity": activity,
            "content_points": content_points,
            "journal_prompt": f"Reflect on today's exploration of {day_info['topic']}. What surprised you?",
            "exit_tickets": [
                f"What is one key takeaway from {day_info['topic']}?",
                "What question do you still have?"
            ],
            "materials_list": [
                "Folger Shakespeare edition of Romeo and Juliet",
                "Student journals",
                "Projector for PowerPoint"
            ]
        }

    def _generate_lesson_plan(self, day_data: Dict, output_path: Path):
        """Generate lesson plan markdown."""
        content = f"""# Lesson Plan: Romeo and Juliet - Day {day_data['day']}

## {day_data['topic']}

### Basic Information
| Field | Value |
|-------|-------|
| Unit | {day_data['unit']['name']} |
| Week | {day_data['week']} |
| Day | {day_data['day']} |
| Type | {day_data['day_type'].title()} Day |
| Duration | 56 minutes |

### Theme
{day_data['theme']}

### Standards Addressed
"""
        for std in day_data['standards']:
            content += f"- **{std['code']}**: {std['full_text']}\n"

        content += "\n### Learning Objectives\n"
        content += "By the end of this lesson, students will be able to:\n"
        for i, obj in enumerate(day_data['learning_objectives'], 1):
            content += f"{i}. {obj}\n"

        content += "\n### Vocabulary\n"
        for vocab in day_data['vocabulary']:
            content += f"- **{vocab['term']}**: {vocab['definition']}\n"

        content += "\n### Materials Needed\n"
        for material in day_data['materials_list']:
            content += f"- [ ] {material}\n"

        content += f"""
### Lesson Procedure

#### Opening (5 minutes)
**Journal Prompt:** {day_data['journal_prompt']}

#### Warmup (5 minutes)
**{day_data['warmup']['name']}**
{day_data['warmup']['instructions']}

#### Direct Instruction (15 minutes)
"""
        for cp in day_data['content_points']:
            point = cp['point'] if isinstance(cp, dict) else cp
            content += f"- {point}\n"

        content += f"""
#### Guided Practice (20 minutes)
**{day_data['activity']['name']}**
{day_data['activity']['instructions']}

#### Closure (10 minutes)
**Exit Ticket Questions:**
"""
        for i, q in enumerate(day_data['exit_tickets'], 1):
            content += f"{i}. {q}\n"

        content += "\n---\n*Generated by Romeo and Juliet Unit Generator*\n"

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)

    def _generate_handout(self, day_data: Dict, output_path: Path):
        """Generate student handout markdown."""
        content = f"""# Day {day_data['day']}: {day_data['topic']}

## Romeo and Juliet - Week {day_data['week']}

---

### Today's Focus
{day_data['theme']}

### Learning Objectives
"""
        for obj in day_data['learning_objectives']:
            content += f"- {obj}\n"

        content += "\n### Key Vocabulary\n"
        content += "| Term | Definition |\n|------|------------|\n"
        for vocab in day_data['vocabulary']:
            content += f"| {vocab['term']} | {vocab['definition']} |\n"

        if day_data['scenes']:
            content += "\n### Scenes We're Reading Today\n"
            for scene in day_data['scenes']:
                content += f"- {scene}\n"

        content += f"""
### Journal Prompt
{day_data['journal_prompt']}

---

### Notes Space
_Use this space to take notes during class:_

\\
\\
\\
\\
\\

### Exit Ticket
"""
        for i, q in enumerate(day_data['exit_tickets'], 1):
            content += f"{i}. {q}\n\n   _Answer:_ \n\n"

        content += "\n---\n*Romeo and Juliet Unit - Theater Education*\n"

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)

    def _generate_pptx(self, day_data: Dict, output_path: Path):
        """Generate PowerPoint presentation with trivia banners."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Reset enhancer for new presentation
        if self.enhance_pptx:
            self.slide_enhancer.reset()

        slides_data = []

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        self._add_title_slide(slide, day_data)
        self._add_presenter_notes(slide, self._generate_title_notes(day_data))
        slides_data.append({"title": day_data['topic'], "content": day_data['topic']})

        # Slide 2: Learning Objectives
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_objectives_slide(slide, day_data)
        self._add_presenter_notes(slide, self._generate_objectives_notes(day_data))
        slides_data.append({"title": "Learning Objectives", "content": " ".join(day_data['learning_objectives'])})

        # Slide 3: Vocabulary
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_vocabulary_slide(slide, day_data)
        self._add_presenter_notes(slide, self._generate_vocabulary_notes(day_data))
        slides_data.append({"title": "Vocabulary", "content": " ".join([v['term'] for v in day_data['vocabulary']])})

        # Slide 4: Warmup
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_warmup_slide(slide, day_data)
        self._add_presenter_notes(slide, self._generate_warmup_notes(day_data))
        slides_data.append({"title": "Warmup", "content": day_data['warmup']['instructions']})

        # Content slides
        for i, cp in enumerate(day_data['content_points']):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            point = cp['point'] if isinstance(cp, dict) else cp
            expanded = cp.get('expanded', []) if isinstance(cp, dict) else []
            self._add_content_slide(slide, point, expanded, i + 5)
            self._add_presenter_notes(slide, self._generate_content_notes(point, expanded, day_data))
            slides_data.append({"title": point[:30], "content": point + " " + " ".join(expanded)})

        # Activity slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_activity_slide(slide, day_data)
        self._add_presenter_notes(slide, self._generate_activity_notes(day_data))
        slides_data.append({"title": "Activity", "content": day_data['activity']['instructions']})

        # Exit ticket slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_exit_slide(slide, day_data)
        self._add_presenter_notes(slide, self._generate_exit_notes(day_data))
        slides_data.append({"title": "Exit Ticket", "content": " ".join(day_data['exit_tickets'])})

        # Add trivia banners if enabled
        if self.enhance_pptx:
            for i, slide in enumerate(prs.slides):
                if i == 0:  # Skip title slide
                    continue
                slide_info = slides_data[i] if i < len(slides_data) else {"title": "", "content": ""}
                self._add_trivia_banner(slide, slide_info)

        prs.save(output_path)

    def _add_title_slide(self, slide, day_data):
        """Add title slide content."""
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Day {day_data['day']}: {day_data['topic']}"
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Romeo and Juliet - Week {day_data['week']}"
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER

        # Theme
        theme_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
        tf = theme_box.text_frame
        p = tf.paragraphs[0]
        p.text = day_data['theme']
        p.font.size = Pt(20)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    def _add_objectives_slide(self, slide, day_data):
        """Add learning objectives slide."""
        # Header
        header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = "Learning Objectives"
        p.font.size = Pt(36)
        p.font.bold = True

        # Objectives
        obj_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
        tf = obj_box.text_frame
        for i, obj in enumerate(day_data['learning_objectives']):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = f"• {obj}"
            p.font.size = Pt(24)
            p.space_after = Pt(18)

    def _add_vocabulary_slide(self, slide, day_data):
        """Add vocabulary slide."""
        header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = "Key Vocabulary"
        p.font.size = Pt(36)
        p.font.bold = True

        vocab_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
        tf = vocab_box.text_frame
        for i, vocab in enumerate(day_data['vocabulary']):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = f"{vocab['term']}: {vocab['definition']}"
            p.font.size = Pt(22)
            p.space_after = Pt(14)

    def _add_warmup_slide(self, slide, day_data):
        """Add warmup slide."""
        header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = f"Warmup: {day_data['warmup']['name']}"
        p.font.size = Pt(36)
        p.font.bold = True

        inst_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(3))
        tf = inst_box.text_frame
        p = tf.paragraphs[0]
        p.text = day_data['warmup']['instructions']
        p.font.size = Pt(24)

    def _optimize_content_for_slide(self, content: str, content_type: str = "bullet") -> tuple:
        """Optimize content using content optimization agents.

        Returns:
            tuple: (condensed_content, elaborated_notes)
        """
        if not self.optimize_content:
            return content, ""

        result = self.content_optimizer.execute({
            "content": content,
            "slide_type": content_type,
            "additional_context": {}
        })

        output = result.output
        if output.get("success"):
            self.stats["content_optimized"] += 1
            if output.get("truncation_fixes_applied", 0) > 0:
                self.stats["truncations_fixed"] += output["truncation_fixes_applied"]

            return (
                output.get("optimized_slide_content", content),
                output.get("presenter_notes", "")
            )

        return content, ""

    def _add_content_slide(self, slide, point, expanded, slide_num):
        """Add content slide with optimized content."""
        # Optimize the title/point
        optimized_point, _ = self._optimize_content_for_slide(point, "title")
        title_text = optimized_point if isinstance(optimized_point, str) else point
        if len(title_text) > 50:
            title_text = title_text[:50] + "..."

        header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True

        if expanded:
            # Optimize each bullet point for condensed display
            optimized_bullets = []
            for exp in expanded:
                optimized, _ = self._optimize_content_for_slide(exp, "bullet")
                optimized_text = optimized if isinstance(optimized, str) else exp
                optimized_bullets.append(optimized_text)

            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4))
            tf = content_box.text_frame
            for i, exp in enumerate(optimized_bullets):
                if i > 0:
                    p = tf.add_paragraph()
                else:
                    p = tf.paragraphs[0]
                p.text = f"• {exp}"
                p.font.size = Pt(22)
                p.space_after = Pt(12)

    def _add_activity_slide(self, slide, day_data):
        """Add activity slide."""
        header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = f"Activity: {day_data['activity']['name']}"
        p.font.size = Pt(36)
        p.font.bold = True

        inst_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(3))
        tf = inst_box.text_frame
        p = tf.paragraphs[0]
        p.text = day_data['activity']['instructions']
        p.font.size = Pt(24)

    def _add_exit_slide(self, slide, day_data):
        """Add exit ticket slide."""
        header = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = "Exit Ticket"
        p.font.size = Pt(36)
        p.font.bold = True

        exit_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4))
        tf = exit_box.text_frame
        for i, q in enumerate(day_data['exit_tickets']):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = f"{i+1}. {q}"
            p.font.size = Pt(24)
            p.space_after = Pt(20)

    def _add_trivia_banner(self, slide, slide_info):
        """Add trivia banner to slide using formatter and validator agents."""
        # Get enhancement from content enhancer
        result = self.slide_enhancer.execute({
            "slide_content": slide_info["content"],
            "slide_title": slide_info["title"],
        })

        if not result.output.get("success"):
            return

        label = result.output["label"]
        content = result.output["content"]

        # Format using formatter agent
        formatted = self.slide_formatter.execute({
            "label": label,
            "content": content,
        })

        if not formatted.output.get("success"):
            return

        # Validate
        validation = self.slide_validator.execute({"formatted_data": formatted.output})
        if not validation.output.get("valid"):
            return

        # Add banner
        box_config = formatted.output["box"]
        label_config = formatted.output["label"]
        content_config = formatted.output["content"]

        pos = box_config["position"]
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(pos["left"]),
            Inches(pos["top"]),
            Inches(pos["width"]),
            Inches(pos["height"])
        )

        shape.fill.solid()
        fill_rgb = box_config["fill_color_rgb"]
        shape.fill.fore_color.rgb = RGBColor(fill_rgb[0], fill_rgb[1], fill_rgb[2])
        border_rgb = box_config["border_color_rgb"]
        shape.line.color.rgb = RGBColor(border_rgb[0], border_rgb[1], border_rgb[2])
        shape.line.width = Pt(box_config["border_width"])

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        # Label - 20pt
        run = p.add_run()
        run.text = label_config["text"]
        run.font.bold = label_config["bold"]
        run.font.size = Pt(label_config["font_size"])  # 20pt
        label_rgb = label_config["color_rgb"]
        run.font.color.rgb = RGBColor(label_rgb[0], label_rgb[1], label_rgb[2])

        # Content - 20pt
        run2 = p.add_run()
        run2.text = content_config["text"]
        run2.font.bold = content_config["bold"]
        run2.font.size = Pt(content_config["font_size"])  # 20pt
        content_rgb = content_config["color_rgb"]
        run2.font.color.rgb = RGBColor(content_rgb[0], content_rgb[1], content_rgb[2])

        self.stats["trivia_added"] += 1

    # =========================================================================
    # PRESENTER NOTES GENERATION METHODS
    # =========================================================================

    def _add_presenter_notes(self, slide, notes_text: str):
        """Add presenter notes to a slide."""
        if not notes_text:
            return

        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

    def _generate_title_notes(self, day_data: Dict) -> str:
        """Generate presenter notes for title slide."""
        notes = f"""PRESENTER NOTES - Day {day_data['day']}: {day_data['topic']}

Welcome everyone to Day {day_data['day']} of our Romeo and Juliet unit.

Today's focus: {day_data['topic']}

This is Week {day_data['week']}, and our theme is: {day_data['theme']}

[PAUSE] Take a moment to let students settle in and get their materials ready.

Before we begin, remind students that all our work today connects to our larger exploration of Shakespeare's tragic love story and what it teaches us about the human experience.

[TRANSITION] When ready, advance to the learning objectives slide.
"""
        return notes

    def _generate_objectives_notes(self, day_data: Dict) -> str:
        """Generate presenter notes for learning objectives slide."""
        objectives_list = "\n".join([f"  - {obj}" for obj in day_data['learning_objectives']])
        notes = f"""LEARNING OBJECTIVES

Read through each objective with the class. [PAUSE between each]

Today's objectives:
{objectives_list}

[CHECK FOR UNDERSTANDING] Ask: "Based on these objectives, what do you think we'll be doing today?"

Call on 2-3 students to share their predictions. This activates prior knowledge and builds anticipation.

[TRANSITION] "Before we dive in, let's review some key vocabulary that will help us today."
"""
        return notes

    def _generate_vocabulary_notes(self, day_data: Dict) -> str:
        """Generate presenter notes for vocabulary slide."""
        vocab_script = ""
        for vocab in day_data['vocabulary']:
            vocab_script += f"""
{vocab['term'].upper()}:
Say the word with me: "{vocab['term']}"
Definition: {vocab['definition']}
[PAUSE] Can anyone use this word in a sentence related to what we've read so far?
"""

        notes = f"""KEY VOCABULARY

Introduce each term using the "Say-Define-Use" method:

{vocab_script}

[ASSESSMENT] Quick check - Have students give a thumbs up if they feel confident with these terms, thumbs sideways if they need more practice.

[TRANSITION] "Now let's warm up our minds and bodies for today's work."
"""
        return notes

    def _generate_warmup_notes(self, day_data: Dict) -> str:
        """Generate presenter notes for warmup slide."""
        warmup = day_data['warmup']
        notes = f"""WARMUP: {warmup['name']}

Duration: {warmup.get('duration_minutes', 5)} minutes

Instructions for students:
{warmup['instructions']}

[TEACHER ACTIONS]
1. Give clear, step-by-step directions
2. Set a visible timer
3. Circulate and observe student engagement
4. Note interesting responses to share in discussion

Connection to today's lesson:
{warmup.get('connection_to_lesson', 'This warmup prepares students for the main content.')}

[DEBRIEF] After time is up, ask 2-3 students to share. Validate responses and connect to today's topic.

[TRANSITION] "Excellent work! Now let's dig into our main content for today."
"""
        return notes

    def _generate_content_notes(self, point: str, expanded: List[str], day_data: Dict) -> str:
        """Generate presenter notes for content slides."""
        expanded_text = "\n".join([f"  - {exp}" for exp in expanded]) if expanded else "  - Discuss this point in depth with students"

        notes = f"""CONTENT: {point}

Key Points to Cover:
{expanded_text}

[DELIVERY TIPS]
- Speak clearly and check for understanding frequently
- Use specific examples from the text when possible
- Encourage students to take notes
- Pause for questions

[ENGAGEMENT STRATEGIES]
- "Turn and talk to your neighbor about this for 30 seconds"
- "Can someone give me an example of this from what we've read?"
- "Why do you think Shakespeare chose to include this?"

[CONNECTION] Relate this content to:
- Previous scenes/events in the play
- Students' own experiences
- Broader themes of love, fate, and conflict

[PACING] Don't rush - ensure students grasp this before moving on.
"""
        return notes

    def _generate_activity_notes(self, day_data: Dict) -> str:
        """Generate presenter notes for activity slide."""
        activity = day_data['activity']
        notes = f"""ACTIVITY: {activity['name']}

Duration: {activity.get('duration_minutes', 20)} minutes

Instructions:
{activity['instructions']}

[SETUP]
1. Ensure all materials are ready before beginning
2. Group students strategically if working in teams
3. Post/display clear expectations

[DURING THE ACTIVITY]
- Circulate and provide support
- Ask probing questions to deepen thinking
- Keep track of time and give warnings (10 min, 5 min, 2 min)
- Note exemplary work to highlight later

[DIFFERENTIATION TIPS]
- For struggling students: Provide sentence starters or graphic organizers
- For advanced students: Add complexity or leadership roles
- For ELL students: Pair with supportive partners, provide visuals

[CLOSURE] Save 3-5 minutes at the end for sharing and reflection.
"""
        return notes

    def _generate_exit_notes(self, day_data: Dict) -> str:
        """Generate presenter notes for exit ticket slide."""
        questions = "\n".join([f"  {i+1}. {q}" for i, q in enumerate(day_data['exit_tickets'])])
        notes = f"""EXIT TICKET

Questions:
{questions}

[INSTRUCTIONS]
1. Distribute exit tickets or have students use their journals
2. Give students 3-5 minutes of quiet writing time
3. Remind them: "This is not graded for correctness - I want to see your genuine thinking"

[COLLECTION]
- Collect exit tickets as students leave
- Or have students submit digitally

[TEACHER FOLLOW-UP]
After class:
- Review responses to gauge understanding
- Identify misconceptions to address next class
- Note students who may need additional support

[CLOSING]
"Thank you for your hard work today. Tomorrow we'll build on what we learned about {day_data['topic']}. Don't forget to [any homework/reminder]."

[DISMISS] Students may pack up and prepare to leave when the bell rings.
"""
        return notes

    def _validate_day(self, day_data: Dict) -> Dict:
        """Validate a day's generated content."""
        errors = []

        # Structure validation
        required_fields = ["topic", "learning_objectives", "vocabulary", "warmup", "activity"]
        for field in required_fields:
            if not day_data.get(field):
                errors.append(f"Missing required field: {field}")

        # Content validation
        if len(day_data.get("learning_objectives", [])) < 1:
            errors.append("At least 1 learning objective required")

        if len(day_data.get("vocabulary", [])) < 2:
            errors.append("At least 2 vocabulary terms required")

        return {"valid": len(errors) == 0, "errors": errors}

    def _validate_unit(self):
        """Run unit-level validation."""
        # Check all days exist
        days_found = list(PRODUCTION_DIR.glob("Day_*/"))
        print(f"Days generated: {len(days_found)}/30")

        # Check PowerPoints
        pptx_found = list(PRODUCTION_DIR.glob("Day_*/*.pptx"))
        print(f"PowerPoints: {len(pptx_found)}")

        # Check handouts
        handouts_found = list(PRODUCTION_DIR.glob("Day_*/*Handout.md"))
        print(f"Handouts: {len(handouts_found)}")

    def _print_summary(self):
        """Print generation summary."""
        print()
        print("=" * 70)
        print("GENERATION COMPLETE")
        print("=" * 70)
        print(f"Days generated:     {self.stats['days_generated']}")
        print(f"PowerPoints:        {self.stats['pptx_generated']}")
        print(f"Handouts:           {self.stats['handouts_generated']}")
        print(f"Validation passed:  {self.stats['validation_passed']}")
        print(f"Validation failed:  {self.stats['validation_failed']}")
        print(f"Trivia banners:     {self.stats['trivia_added']} (20pt font)")
        print(f"Content optimized:  {self.stats['content_optimized']} items")
        print(f"Truncations fixed:  {self.stats['truncations_fixed']}")
        print(f"Output directory:   {PRODUCTION_DIR}")
        print()


# =============================================================================
# MAIN
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Generate complete Romeo and Juliet 6-week unit",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s                      Generate all 30 days
  %(prog)s --week 1             Generate Week 1 (Days 1-5)
  %(prog)s --week 1 2           Generate Weeks 1 and 2
  %(prog)s --day 1              Generate Day 1 only
  %(prog)s --day 1 2 3          Generate Days 1, 2, and 3
  %(prog)s --no-enhance         Skip trivia banner enhancement
  %(prog)s --no-optimize        Skip content optimization (condensation)
        """
    )

    parser.add_argument('--week', type=int, nargs='+', choices=[1, 2, 3, 4, 5, 6],
                       help='Generate specific week(s)')
    parser.add_argument('--day', type=int, nargs='+',
                       help='Generate specific day(s) (1-30)')
    parser.add_argument('--no-enhance', action='store_true',
                       help='Skip trivia banner enhancement')
    parser.add_argument('--no-optimize', action='store_true',
                       help='Skip content optimization (slides will be verbose)')
    parser.add_argument('--verbose', '-v', action='store_true',
                       help='Verbose output')

    args = parser.parse_args()

    # Validate day numbers
    if args.day:
        for d in args.day:
            if d < 1 or d > 30:
                parser.error(f"Day must be between 1 and 30, got {d}")

    # Run generator
    generator = RomeoJulietUnitGenerator(
        enhance_pptx=not args.no_enhance,
        optimize_content=not args.no_optimize,
        verbose=args.verbose
    )

    generator.generate_unit(weeks=args.week, days=args.day)

    return 0


if __name__ == "__main__":
    sys.exit(main())
