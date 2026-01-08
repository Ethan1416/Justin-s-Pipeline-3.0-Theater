"""
Step 12: PowerPoint Population
Template-based slide generation with visual integration

Requirements:
- python-pptx (pip install python-pptx)
- Pillow (pip install Pillow)
"""

import os
import re
import json
import shutil
from pathlib import Path
from datetime import datetime
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_FILL_TYPE

# ============================================
# CONFIGURATION - LOADED FROM pipeline_config.json
# ============================================

def load_config():
    """Load configuration from pipeline_config.json

    Config v3.0+ uses relative paths from the repository root.
    Templates are in templates/ folder within the repo.
    """
    # Repository root is the same directory as this script
    repo_root = Path(__file__).parent
    config_path = repo_root / "pipeline_config.json"

    if not config_path.exists():
        raise FileNotFoundError(
            f"Config file not found at: {config_path}\n"
            f"Please ensure pipeline_config.json exists in the repository root."
        )

    with open(config_path, 'r', encoding='utf-8') as f:
        config = json.load(f)

    # Add repo_root to config for path resolution
    config['_repo_root'] = str(repo_root)

    # Validate production_folder is set (required for output)
    if not config.get('paths', {}).get('production_folder'):
        raise ValueError(
            f"Missing 'production_folder' in pipeline_config.json.\n"
            f"Please set paths.production_folder to your production folder path.\n"
            f"Example: C:\\Users\\mcdan\\Desktop\\Domain_Production_2024-12-22"
        )

    return config

# Load configuration
CONFIG = load_config()

# Extract paths from config - using repository-relative paths
REPO_ROOT = Path(CONFIG['_repo_root'])
TEMPLATE_PATH = REPO_ROOT / CONFIG['templates']['content_master']
VISUAL_TEMPLATE_PATH = REPO_ROOT / CONFIG['templates']['visual_organizer']
PRODUCTION_FOLDER = CONFIG['paths']['production_folder']
DOMAIN_NAME = CONFIG['domain'].get('display_name') or CONFIG['domain'].get('name') or "Unknown_Domain"

# Shape name mappings (from template requirements in config)
shape_mappings = CONFIG.get('template_requirements', {}).get('shape_mappings', {})
SHAPE_TITLE = next((k for k, v in shape_mappings.items() if v == 'title'), "TextBox 2")
SHAPE_BODY = next((k for k, v in shape_mappings.items() if v == 'body'), "TextBox 19")
SHAPE_TIP = next((k for k, v in shape_mappings.items() if v == 'tip'), "TextBox 18")

# Font settings from config (with defaults)
FONT_SETTINGS = CONFIG.get('template_requirements', {}).get('font_settings', {})

# Table settings from config
TABLE_SETTINGS = CONFIG.get('template_requirements', {}).get('table_settings', {})
TABLE_TEMPLATE_PATH = TEMPLATE_PATH  # Use content_master for tables

# Visual settings
DIAGRAM_MAX_WIDTH = Inches(11)
DIAGRAM_MAX_HEIGHT = Inches(5)
TABLE_FONT_SIZE = Pt(TABLE_SETTINGS.get('font_size_pt', 20))
HEADER_FONT_SIZE = Pt(28)

# Minimum font size for all visual aids (graphic organizers, diagrams, etc.)
# This ensures readability and professional appearance
VISUAL_MIN_FONT_SIZE = Pt(18)

# ============================================
# HELPER FUNCTIONS
# ============================================

def sanitize_filename(name):
    """Convert section name to valid filename."""
    name = name.replace(" ", "_")
    name = re.sub(r'[^\w\-_]', '', name)
    return name

def find_shape_by_name(slide, shape_name):
    """Find a shape by its exact name."""
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape
    return None

def clear_shape_text(shape):
    """Clear all text from a shape while preserving the shape."""
    if shape and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            paragraph.clear()
        if shape.text_frame.paragraphs:
            shape.text_frame.paragraphs[0].text = ""

def set_shape_text(shape, text, shape_name=None):
    """Set text in a shape, applying font settings from config."""
    if not shape or not shape.has_text_frame:
        return False

    # Get font settings from config, or use shape-specific defaults
    font_config = FONT_SETTINGS.get(shape_name, {}) if shape_name else {}
    font_name = font_config.get('font_name', 'Aptos')
    font_size = Pt(font_config.get('font_size_pt', 20))
    font_bold = font_config.get('bold', False)

    # Shape-specific default colors: title=white, body/tip=black
    if shape_name == SHAPE_TITLE:
        default_color = [255, 255, 255]  # White for title
    else:
        default_color = [0, 0, 0]  # Black for body and tip
    font_color_rgb = font_config.get('font_color_rgb', default_color)

    # Clear existing text
    shape.text_frame.clear()

    # For titles: remove newlines and truncate to 28 characters
    if shape_name == SHAPE_TITLE:
        text = text.replace('\n', ' ').replace('\r', ' ')
        if len(text) > 28:
            text = text[:28]

    # Add text with configured formatting
    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = font_bold
    run.font.color.rgb = RGBColor(font_color_rgb[0], font_color_rgb[1], font_color_rgb[2])

    return True

# ============================================
# SHAPE DRAWING HELPER FUNCTIONS
# ============================================
# IMPORTANT: These functions implement correct rotation positioning.
# See SHAPE_ROTATION_GUIDELINES.md for full documentation.
# NEVER create rotated lines directly - always use these helpers!

def draw_line(slide, x1, y1, x2, y2, color, width=Pt(2)):
    """
    Draw a line using a thin rectangle (avoids connector corruption).

    PowerPoint rotates shapes around their CENTER, so we need to:
    1. Calculate the midpoint between start/end
    2. Position rectangle center at midpoint
    3. Rotate around that center point

    See SHAPE_ROTATION_GUIDELINES.md for detailed explanation.

    Args:
        slide: PowerPoint slide object
        x1, y1: Start coordinates (Inches or Emu)
        x2, y2: End coordinates (Inches or Emu)
        color: RGBColor object
        width: Line thickness (default Pt(2))

    Returns:
        The rectangle shape representing the line
    """
    import math

    # Convert all to Emu for calculation
    x1_emu = x1.emu if hasattr(x1, 'emu') else Emu(x1)
    y1_emu = y1.emu if hasattr(y1, 'emu') else Emu(y1)
    x2_emu = x2.emu if hasattr(x2, 'emu') else Emu(x2)
    y2_emu = y2.emu if hasattr(y2, 'emu') else Emu(y2)
    width_emu = width.emu if hasattr(width, 'emu') else Emu(width)

    # Calculate length and angle
    dx = x2_emu - x1_emu
    dy = y2_emu - y1_emu
    length = math.sqrt(dx**2 + dy**2)
    angle = math.atan2(dy, dx)

    # Calculate midpoint (where rectangle center should be)
    mid_x = (x1_emu + x2_emu) / 2
    mid_y = (y1_emu + y2_emu) / 2

    # Calculate top-left position such that center is at midpoint
    # For a rectangle: center_x = left + width/2, center_y = top + height/2
    # So: left = center_x - width/2, top = center_y - height/2
    # Before rotation, the "width" is the length and "height" is the line thickness
    left = mid_x - length / 2
    top = mid_y - width_emu / 2

    # Create rectangle positioned so its center is at the midpoint
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(left)), Emu(int(top)),
        Emu(int(length)), Emu(width_emu)
    )

    # Style the line
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = color
    line_shape.line.fill.background()  # No border

    # Rotate around center (which is now at the midpoint between start/end)
    if angle != 0:
        line_shape.rotation = math.degrees(angle)

    return line_shape

def draw_arrow(slide, x1, y1, x2, y2, color, line_width=Pt(2), arrow_size=Inches(0.15)):
    """
    Draw an arrow from (x1,y1) to (x2,y2) using rectangles and a triangle.

    Args:
        slide: PowerPoint slide object
        x1, y1: Start coordinates
        x2, y2: End coordinates (arrowhead location)
        color: RGBColor object
        line_width: Thickness of line
        arrow_size: Size of arrowhead triangle

    Returns:
        Tuple of (line_shape, arrowhead_shape)
    """
    import math

    # Convert to Emu
    x1_emu = x1.emu if hasattr(x1, 'emu') else Emu(x1)
    y1_emu = y1.emu if hasattr(y1, 'emu') else Emu(y1)
    x2_emu = x2.emu if hasattr(x2, 'emu') else Emu(x2)
    y2_emu = y2.emu if hasattr(y2, 'emu') else Emu(y2)
    arrow_size_emu = arrow_size.emu if hasattr(arrow_size, 'emu') else Emu(arrow_size)

    # Calculate angle
    dx = x2_emu - x1_emu
    dy = y2_emu - y1_emu
    total_length = math.sqrt(dx**2 + dy**2)
    angle = math.atan2(dy, dx)

    # Shorten line to make room for arrowhead (account for arrowhead length)
    # The arrowhead extends from its center by arrow_size/2 in the direction it points
    line_end_offset = arrow_size_emu * 0.6  # Arrowhead extends this far from line end
    line_length = total_length - line_end_offset

    # Calculate where the line should end (before arrowhead)
    line_end_x = x1_emu + line_length * math.cos(angle)
    line_end_y = y1_emu + line_length * math.sin(angle)

    # Draw line
    line_shape = draw_line(slide,
                           Emu(x1_emu), Emu(y1_emu),
                           Emu(int(line_end_x)), Emu(int(line_end_y)),
                           color, line_width)

    # Position arrowhead
    # For isosceles triangle pointing right (default), the tip is at the right center
    # We want the tip at (x2, y2) after rotation
    # Triangle's center should be offset back from tip by half its width
    arrow_center_offset = arrow_size_emu / 3  # Center is 1/3 back from tip
    arrow_center_x = x2_emu - arrow_center_offset * math.cos(angle)
    arrow_center_y = y2_emu - arrow_center_offset * math.sin(angle)

    # Calculate top-left position so center is at calculated position
    arrow_left = arrow_center_x - arrow_size_emu / 2
    arrow_top = arrow_center_y - arrow_size_emu / 2

    # Draw arrowhead (isosceles triangle)
    arrowhead = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Emu(int(arrow_left)), Emu(int(arrow_top)),
        Emu(arrow_size_emu), Emu(arrow_size_emu)
    )

    # Style arrowhead
    arrowhead.fill.solid()
    arrowhead.fill.fore_color.rgb = color
    arrowhead.line.fill.background()

    # Rotate arrowhead to point in direction of line
    # Adjust angle by 90 degrees because isosceles triangle points up by default in PowerPoint
    arrowhead.rotation = math.degrees(angle) + 90

    return line_shape, arrowhead

# ============================================
# BLUEPRINT PARSING
# ============================================

def parse_blueprint(filepath):
    """Parse a Step 10 integrated blueprint into slide data."""
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    slides = []

    # Extract section name from header
    section_match = re.search(r'Section:\s*\d+\s*-\s*(.+)', content)
    if not section_match:
        section_match = re.search(r'Section:\s*(.+)', content)
    section_name = section_match.group(1).strip() if section_match else "Unknown_Section"

    # Split by slide markers
    slide_pattern = r'={40,}\s*SLIDE\s+(\d+[AB]?):\s*(.+?)\s*={40,}'
    slide_blocks = re.split(slide_pattern, content)

    i = 1
    while i < len(slide_blocks) - 2:
        slide_num = slide_blocks[i].strip()
        slide_title = slide_blocks[i + 1].strip()
        slide_content = slide_blocks[i + 2] if i + 2 < len(slide_blocks) else ""

        slide_data = parse_slide_content(slide_num, slide_title, slide_content)
        slides.append(slide_data)

        i += 3

    return section_name, slides

def parse_slide_content(num, title, content):
    """Parse individual slide content."""
    slide = {
        'number': num,
        'title': title,
        'type': 'Content',
        'visual': False,
        'visual_type': None,
        'header': title,
        'body': '',
        'tip': '',
        'notes': '',
        'table_data': None,
        'key_differentiator_data': None,
        'decision_tree_data': None,
        'flowchart_data': None,
        'hierarchy_data': None,
        'timeline_data': None,
        'spectrum_data': None,
        'visual_spec': None
    }

    # Extract type
    type_match = re.search(r'Type:\s*(.+)', content)
    if type_match:
        slide['type'] = type_match.group(1).strip()

    # Check for visual type in Type field (e.g., "Type: VISUAL" or "Synthesis Visual Aid - TABLE")
    if slide['type'].upper() == 'VISUAL':
        slide['visual'] = True
    elif 'TABLE' in slide['type'].upper():
        slide['visual'] = True
        slide['visual_type'] = 'table'
    elif 'DECISION TREE' in slide['type'].upper():
        slide['visual'] = True
        slide['visual_type'] = 'decision_tree'

    # Extract visual type from "Visual Type:" field (e.g., "Visual Type: TABLE")
    visual_type_match = re.search(r'Visual Type:\s*(.+)', content)
    if visual_type_match:
        slide['visual'] = True
        slide['visual_type'] = visual_type_match.group(1).strip().lower()

    # Extract visual info from Visual: field (legacy format)
    visual_match = re.search(r'^Visual:\s*(.+)', content, re.MULTILINE)
    if visual_match:
        visual_text = visual_match.group(1).strip()
        if visual_text.lower().startswith('yes'):
            slide['visual'] = True
            type_parts = visual_text.split('-')
            if len(type_parts) > 1:
                slide['visual_type'] = type_parts[1].strip().lower()

    # Extract header
    header_match = re.search(r'HEADER[^:]*:\s*\n(.*?)(?=\n\s*\n|\nBODY)', content, re.DOTALL)
    if header_match:
        slide['header'] = header_match.group(1).strip()

    # Extract body - handle "Line N:" format
    body_match = re.search(r'BODY[^:]*:\s*\n(.*?)(?=\nVISUAL SPECIFICATION:|\nNCLEX TIP:|\nPRESENTER NOTES:|\n-{20,}|\n={20,})', content, re.DOTALL)
    if body_match:
        body_text = body_match.group(1).strip()
        # Remove "Line N:" prefixes if present
        lines = body_text.split('\n')
        cleaned_lines = []
        for line in lines:
            line_match = re.match(r'Line\s*\d+:\s*(.*)', line)
            if line_match:
                cleaned_lines.append(line_match.group(1))
            else:
                cleaned_lines.append(line)
        slide['body'] = '\n'.join(cleaned_lines)

    # Extract visual specification (text format)
    visual_spec_match = re.search(r'VISUAL SPECIFICATION:\s*\n(.*?)(?=\nNCLEX TIP:|\nPRESENTER NOTES:)', content, re.DOTALL)
    if visual_spec_match:
        slide['visual_spec'] = visual_spec_match.group(1).strip()
        if slide['visual_type'] == 'table':
            slide['table_data'] = parse_table_spec(slide['visual_spec'])

    # Also try VISUAL_DATA JSON format if VISUAL SPECIFICATION not found
    if not slide['visual_spec']:
        visual_data_match = re.search(r'VISUAL_DATA:\s*\n(\{[\s\S]*?\n\})', content)
        if visual_data_match:
            try:
                visual_json = json.loads(visual_data_match.group(1))
                slide['visual_spec'] = visual_data_match.group(1)  # Keep raw for debugging
                # Parse JSON into appropriate data structure based on type
                vtype = visual_json.get('type', '').lower()
                if vtype == 'table' and 'columns' in visual_json and 'rows' in visual_json:
                    # Convert JSON table to table_data format: [headers, row1, row2, ...]
                    slide['table_data'] = [visual_json['columns']] + visual_json['rows']
                elif vtype == 'hierarchy' and 'levels' in visual_json:
                    # Convert JSON levels array to expected format
                    h_data = {'layout': visual_json.get('layout', 'A'), 'level1': None, 'level2': [], 'level3': []}
                    for lvl in visual_json['levels']:
                        lvl_num = lvl.get('level', 1)
                        nodes = lvl.get('nodes', [])
                        if lvl_num == 1 and nodes:
                            h_data['level1'] = {'name': nodes[0], 'description': ''}
                        elif lvl_num == 2:
                            h_data['level2'] = [{'name': n, 'description': ''} for n in nodes]
                        elif lvl_num == 3:
                            # Try to assign parents based on position
                            l2_names = [n['name'] for n in h_data['level2']] if h_data['level2'] else []
                            for i, n in enumerate(nodes):
                                parent = l2_names[i % len(l2_names)] if l2_names else ''
                                h_data['level3'].append({'name': n, 'description': '', 'parent': parent})
                    slide['hierarchy_data'] = h_data
                elif vtype == 'timeline' and 'events' in visual_json:
                    # Timeline expects 'events' list with 'time', 'label', 'description'
                    slide['timeline_data'] = visual_json
                elif vtype == 'flowchart' and 'steps' in visual_json:
                    # Convert steps to expected format with 'title' and 'description'
                    fc_data = {'layout': visual_json.get('layout', 'B'), 'steps': [], 'connectors': []}
                    for step in visual_json['steps']:
                        fc_data['steps'].append({
                            'title': step.get('header', step.get('title', '')),
                            'description': step.get('body', step.get('description', ''))
                        })
                    slide['flowchart_data'] = fc_data
                elif vtype == 'spectrum':
                    # Convert 'dimensions' to 'points' format expected by add_spectrum_content
                    if 'dimensions' in visual_json:
                        s_data = {'layout': visual_json.get('layout', 'C'), 'points': []}
                        for dim in visual_json['dimensions']:
                            s_data['points'].append({
                                'name': dim.get('trait', ''),
                                'description': dim.get('direction', '')
                            })
                        slide['spectrum_data'] = s_data
                    elif 'points' in visual_json:
                        slide['spectrum_data'] = visual_json
                elif vtype == 'key_differentiators' and 'concepts' in visual_json:
                    # Convert concepts with color names to RGB
                    color_map = {
                        'blue': {'header': (21, 101, 192), 'light': (227, 242, 253)},
                        'green': {'header': (46, 125, 50), 'light': (232, 245, 233)},
                        'purple': {'header': (74, 20, 140), 'light': (243, 229, 245)},
                        'orange': {'header': (230, 81, 0), 'light': (255, 243, 224)},
                        'teal': {'header': (0, 121, 107), 'light': (224, 242, 241)},
                        'red': {'header': (183, 28, 28), 'light': (255, 235, 238)}
                    }
                    kd_data = {'concepts': [], 'colors': {}, 'key_differentiator': visual_json.get('key_differentiator', '')}
                    for i, concept in enumerate(visual_json['concepts'], 1):
                        color_name = concept.get('color', 'blue').lower()
                        colors = color_map.get(color_name, color_map['blue'])
                        kd_data['concepts'].append({
                            'number': i,
                            'name': concept.get('name', ''),
                            'features': concept.get('features', [])
                        })
                        kd_data['colors'][i] = colors
                    slide['key_differentiator_data'] = kd_data
                elif vtype == 'decision_tree':
                    slide['decision_tree_data'] = visual_json
            except json.JSONDecodeError:
                pass  # If JSON parsing fails, continue with text-based parsing

    # Parse table from body if it contains ASCII table format
    if slide['visual_type'] == 'table' and not slide['table_data']:
        slide['table_data'] = parse_table_from_body(slide['body'])

    # Parse key differentiator data if visual type is key_differentiators
    if slide['visual_type'] == 'key_differentiators' and slide['visual_spec']:
        slide['key_differentiator_data'] = parse_key_differentiator_spec(slide['visual_spec'])

    # Parse decision tree data if visual type is decision_tree
    if slide['visual_type'] == 'decision_tree' and slide['visual_spec']:
        slide['decision_tree_data'] = parse_decision_tree_spec(slide['visual_spec'])

    # Parse flowchart data if visual type is flowchart
    if slide['visual_type'] == 'flowchart' and slide['visual_spec']:
        slide['flowchart_data'] = parse_flowchart_spec(slide['visual_spec'])

    # Parse hierarchy data if visual type is hierarchy
    if slide['visual_type'] == 'hierarchy' and slide['visual_spec']:
        slide['hierarchy_data'] = parse_hierarchy_spec(slide['visual_spec'])

    # Parse timeline data if visual type is timeline
    if slide['visual_type'] == 'timeline' and slide['visual_spec']:
        slide['timeline_data'] = parse_timeline_spec(slide['visual_spec'])

    # Parse spectrum data if visual type is spectrum
    if slide['visual_type'] == 'spectrum' and slide['visual_spec']:
        slide['spectrum_data'] = parse_spectrum_spec(slide['visual_spec'])

    # Extract NCLEX tip
    tip_match = re.search(r'NCLEX TIP:\s*\n(.*?)(?=\nPRESENTER NOTES:|\n-{20,}|\n={20,}|$)', content, re.DOTALL)
    if tip_match:
        tip_text = tip_match.group(1).strip()
        if tip_text.lower() not in ['none', 'n/a', '[none]', '[n/a]', '']:
            slide['tip'] = tip_text

    # Extract presenter notes
    notes_match = re.search(r'PRESENTER NOTES:\s*\n(.*?)(?=\n-{20,}|\n={40,}|$)', content, re.DOTALL)
    if notes_match:
        slide['notes'] = notes_match.group(1).strip()

    # Validate 8-line limit for content slides
    if not slide['visual'] and slide['body']:
        # Count non-empty lines in body
        body_lines = [line for line in slide['body'].split('\n') if line.strip()]
        if len(body_lines) > 8:
            print(f"\nWARNING: Slide {slide['number']} exceeds 8-line limit for content slides!")
            print(f"   Title: {slide['header']}")
            print(f"   Lines: {len(body_lines)} (maximum: 8)")
            print(f"   This may cause text overflow on the PowerPoint slide.")
            print(f"   Please condense the BODY section to 8 lines or fewer.\n")

    return slide

def parse_table_spec(spec_text):
    """Parse markdown table from visual specification."""
    lines = spec_text.strip().split('\n')
    table_lines = [l for l in lines if '|' in l and not l.strip().startswith('|---') and not l.strip() == '|']

    if not table_lines:
        return None

    table_data = []
    for line in table_lines:
        cells = [c.strip() for c in line.split('|') if c.strip()]
        if cells:
            table_data.append(cells)

    return table_data if table_data else None

def parse_table_from_body(body_text):
    """Parse ASCII table from body text (handles box-drawing characters)."""
    if not body_text:
        return None

    lines = body_text.split('\n')
    table_data = []

    for line in lines:
        # Check for pipe character or box-drawing characters
        if '|' in line or '│' in line:
            # Skip border/separator lines
            if set(line.strip()) <= set('─┌┐└┘├┤┬┴┼|-+'):
                continue
            if re.match(r'^[\s\-─┌┐└┘├┤┬┴┼|+]+$', line):
                continue

            # Split by pipe or box-drawing vertical
            cells = re.split(r'[|│]', line)
            cells = [c.strip() for c in cells if c.strip()]

            if cells:
                table_data.append(cells)

    return table_data if len(table_data) >= 2 else None

def parse_decision_tree_spec(spec_text):
    """Parse decision tree specification from visual spec."""
    if not spec_text:
        return None

    dt_data = {
        'layout': 'A',
        'level1': {},
        'level2a': {},
        'level2b': {},
        'outcomes': []
    }

    # Extract layout
    layout_match = re.search(r'Layout:\s*([A-F]|AUTO)', spec_text, re.IGNORECASE)
    if layout_match:
        dt_data['layout'] = layout_match.group(1).upper()

    # Parse LEVEL 1
    level1_match = re.search(r'LEVEL 1:(.*?)(?=LEVEL 2A:|LEVEL 2B:|OUTCOMES:|$)', spec_text, re.DOTALL)
    if level1_match:
        level1_text = level1_match.group(1)
        header_match = re.search(r'-\s*Header:\s*["\']?([^"\'\n]+)["\']?', level1_text)
        question_match = re.search(r'-\s*Question:\s*["\']?([^"\'\n]+)["\']?', level1_text)
        paths_match = re.search(r'-\s*Paths:\s*([^\n]+)', level1_text)

        if header_match and question_match and paths_match:
            paths = [p.strip() for p in paths_match.group(1).split(',')]
            dt_data['level1'] = {
                'header': header_match.group(1).strip(),
                'question': question_match.group(1).strip(),
                'paths': paths
            }

    # Parse LEVEL 2A
    level2a_match = re.search(r'LEVEL 2A:(.*?)(?=LEVEL 2B:|OUTCOMES:|$)', spec_text, re.DOTALL)
    if level2a_match:
        level2a_text = level2a_match.group(1)
        header_match = re.search(r'-\s*Header:\s*["\']?([^"\'\n]+)["\']?', level2a_text)
        question_match = re.search(r'-\s*Question:\s*["\']?([^"\'\n]+)["\']?', level2a_text)
        paths_match = re.search(r'-\s*Paths:\s*([^\n]+)', level2a_text)
        parent_match = re.search(r'-\s*Parent Path:\s*([^\n]+)', level2a_text)

        if header_match and question_match and paths_match:
            paths = [p.strip() for p in paths_match.group(1).split(',')]
            dt_data['level2a'] = {
                'header': header_match.group(1).strip(),
                'question': question_match.group(1).strip(),
                'paths': paths,
                'parent_path': parent_match.group(1).strip() if parent_match else ''
            }

    # Parse LEVEL 2B
    level2b_match = re.search(r'LEVEL 2B:(.*?)(?=OUTCOMES:|$)', spec_text, re.DOTALL)
    if level2b_match:
        level2b_text = level2b_match.group(1)
        header_match = re.search(r'-\s*Header:\s*["\']?([^"\'\n]+)["\']?', level2b_text)
        question_match = re.search(r'-\s*Question:\s*["\']?([^"\'\n]+)["\']?', level2b_text)
        paths_match = re.search(r'-\s*Paths:\s*([^\n]+)', level2b_text)
        parent_match = re.search(r'-\s*Parent Path:\s*([^\n]+)', level2b_text)

        if header_match and question_match and paths_match:
            paths = [p.strip() for p in paths_match.group(1).split(',')]
            dt_data['level2b'] = {
                'header': header_match.group(1).strip(),
                'question': question_match.group(1).strip(),
                'paths': paths,
                'parent_path': parent_match.group(1).strip() if parent_match else ''
            }

    # Parse OUTCOMES
    outcomes_match = re.search(r'OUTCOMES:(.*?)(?=NCLEX TIP:|PRESENTER NOTES:|$)', spec_text, re.DOTALL)
    if outcomes_match:
        outcomes_text = outcomes_match.group(1)
        # Format: - O1: "Name" | Color: green | Parent: L2A-YES
        for line in outcomes_text.split('\n'):
            line = line.strip()
            if not line or not line.startswith('-'):
                continue

            outcome_match = re.match(r'-\s*O\d+:\s*["\']?([^"\'\|]+)["\']?\s*\|\s*Color:\s*(\w+)\s*\|\s*Parent:\s*([^\n]+)', line)
            if outcome_match:
                dt_data['outcomes'].append({
                    'name': outcome_match.group(1).strip(),
                    'color': outcome_match.group(2).strip().lower(),
                    'parent': outcome_match.group(3).strip()
                })

    return dt_data if dt_data['level1'] and dt_data['outcomes'] else None

def parse_flowchart_spec(spec_text):
    """Parse flowchart specification from visual spec."""
    if not spec_text:
        return None

    fc_data = {
        'layout': 'A',
        'steps': [],
        'connectors': []
    }

    # Extract layout
    layout_match = re.search(r'Layout:\s*([A-F]|AUTO)', spec_text, re.IGNORECASE)
    if layout_match:
        fc_data['layout'] = layout_match.group(1).upper()

    # Parse STEPS
    steps_match = re.search(r'STEPS:(.*?)(?=CONNECTORS:|NCLEX TIP:|PRESENTER NOTES:|$)', spec_text, re.DOTALL)
    if steps_match:
        steps_text = steps_match.group(1)
        for line in steps_text.split('\n'):
            line = line.strip()
            if not line or not re.match(r'^\d+\.', line):
                continue

            # Format: "1. Title → Description"
            step_match = re.match(r'^\d+\.\s*([^→]+)→\s*(.+)', line)
            if step_match:
                fc_data['steps'].append({
                    'title': step_match.group(1).strip(),
                    'description': step_match.group(2).strip()
                })

    # Parse CONNECTORS (optional)
    connectors_match = re.search(r'CONNECTORS:(.*?)(?=NCLEX TIP:|PRESENTER NOTES:|$)', spec_text, re.DOTALL)
    if connectors_match:
        connectors_text = connectors_match.group(1)
        for line in connectors_text.split('\n'):
            line = line.strip()
            if not line:
                continue

            # Format: "1→2: [Label]"
            conn_match = re.match(r'(\d+)→(\d+):\s*\[([^\]]+)\]', line)
            if conn_match:
                fc_data['connectors'].append({
                    'from': int(conn_match.group(1)),
                    'to': int(conn_match.group(2)),
                    'label': conn_match.group(3).strip()
                })

    return fc_data if fc_data['steps'] else None

def parse_hierarchy_spec(spec_text):
    """Parse hierarchy specification from visual spec."""
    if not spec_text:
        return None

    h_data = {
        'layout': 'A',
        'levels': 3,
        'level1': None,
        'level2': [],
        'level3': []
    }

    # Extract layout and levels
    layout_match = re.search(r'Layout:\s*([A-F]|AUTO)', spec_text, re.IGNORECASE)
    if layout_match:
        h_data['layout'] = layout_match.group(1).upper()

    levels_match = re.search(r'Levels:\s*(\d+)', spec_text)
    if levels_match:
        h_data['levels'] = int(levels_match.group(1))

    # Parse LEVEL 1
    level1_match = re.search(r'LEVEL 1:(.*?)(?=LEVEL 2:|NCLEX TIP:|PRESENTER NOTES:|$)', spec_text, re.DOTALL)
    if level1_match:
        level1_text = level1_match.group(1)
        # Format: - "Name" | "Description"
        item_match = re.search(r'-\s*["\']([^"\']+)["\']\s*\|\s*["\']([^"\']+)["\']', level1_text)
        if item_match:
            h_data['level1'] = {
                'name': item_match.group(1).strip(),
                'description': item_match.group(2).strip()
            }

    # Parse LEVEL 2
    level2_match = re.search(r'LEVEL 2:(.*?)(?=LEVEL 3:|NCLEX TIP:|PRESENTER NOTES:|$)', spec_text, re.DOTALL)
    if level2_match:
        level2_text = level2_match.group(1)
        for line in level2_text.split('\n'):
            item_match = re.search(r'-\s*["\']([^"\']+)["\']\s*\|\s*["\']([^"\']+)["\']', line)
            if item_match:
                h_data['level2'].append({
                    'name': item_match.group(1).strip(),
                    'description': item_match.group(2).strip()
                })

    # Parse LEVEL 3
    level3_match = re.search(r'LEVEL 3:(.*?)(?=NCLEX TIP:|PRESENTER NOTES:|$)', spec_text, re.DOTALL)
    if level3_match:
        level3_text = level3_match.group(1)
        for line in level3_text.split('\n'):
            # Format: - "Name" | "Description" → Parent: ParentName
            item_match = re.search(r'-\s*["\']([^"\']+)["\']\s*\|\s*["\']([^"\']+)["\']\s*→\s*Parent:\s*([^\n]+)', line)
            if item_match:
                h_data['level3'].append({
                    'name': item_match.group(1).strip(),
                    'description': item_match.group(2).strip(),
                    'parent': item_match.group(3).strip()
                })

    return h_data if h_data['level1'] else None

def parse_timeline_spec(spec_text):
    """Parse timeline specification from visual spec."""
    if not spec_text:
        return None

    t_data = {
        'layout': 'A',
        'events': []
    }

    # Extract layout
    layout_match = re.search(r'Layout:\s*([A-F]|AUTO)', spec_text, re.IGNORECASE)
    if layout_match:
        t_data['layout'] = layout_match.group(1).upper()

    # Parse EVENTS
    events_match = re.search(r'EVENTS:(.*?)(?=NCLEX TIP:|PRESENTER NOTES:|$)', spec_text, re.DOTALL)
    if events_match:
        events_text = events_match.group(1)
        for line in events_text.split('\n'):
            line = line.strip()
            if not line or not re.match(r'^\d+\.', line):
                continue

            # Format: "1. 1953 → Title | Description"
            event_match = re.match(r'^\d+\.\s*([^→]+)→\s*([^\|]+)\|\s*(.+)', line)
            if event_match:
                t_data['events'].append({
                    'year': event_match.group(1).strip(),
                    'title': event_match.group(2).strip(),
                    'description': event_match.group(3).strip()
                })

    return t_data if t_data['events'] else None

def parse_spectrum_spec(spec_text):
    """Parse spectrum specification from visual spec."""
    if not spec_text:
        return None

    s_data = {
        'layout': 'A',
        'segments': 5,
        'points': []
    }

    # Extract layout and segments
    layout_match = re.search(r'Layout:\s*([A-F]|AUTO)', spec_text, re.IGNORECASE)
    if layout_match:
        s_data['layout'] = layout_match.group(1).upper()

    segments_match = re.search(r'Segments:\s*(\d+)', spec_text)
    if segments_match:
        s_data['segments'] = int(segments_match.group(1))

    # Parse SPECTRUM POINTS
    points_match = re.search(r'SPECTRUM POINTS:(.*?)(?=NCLEX TIP:|PRESENTER NOTES:|$)', spec_text, re.DOTALL)
    if points_match:
        points_text = points_match.group(1)
        for line in points_text.split('\n'):
            line = line.strip()
            if not line or not re.match(r'^\d+\.', line):
                continue

            # Format: "1. Name → Description"
            point_match = re.match(r'^\d+\.\s*([^→]+)→\s*(.+)', line)
            if point_match:
                s_data['points'].append({
                    'name': point_match.group(1).strip(),
                    'description': point_match.group(2).strip()
                })

    return s_data if s_data['points'] else None

def parse_key_differentiator_spec(spec_text):
    """Parse key differentiator specification from visual spec."""
    if not spec_text:
        return None

    kd_data = {
        'layout': 'A',  # Default to side-by-side
        'concepts': [],
        'key_differences': [],
        'colors': {
            1: {'header': (21, 101, 192), 'light': (227, 242, 253)},  # Blue
            2: {'header': (183, 28, 28), 'light': (255, 235, 238)}    # Red
        }
    }

    # Try NEW format first (ITEM A/B with DIMENSIONS)
    item_a_match = re.search(r'ITEM A:\s*(.+)', spec_text)
    item_b_match = re.search(r'ITEM B:\s*(.+)', spec_text)

    if item_a_match and item_b_match:
        # Parse ITEM A/B format
        concept_a_name = item_a_match.group(1).strip()
        concept_b_name = item_b_match.group(1).strip()

        # Extract dimensions
        dimensions_match = re.search(r'DIMENSIONS:(.*?)(?=NCLEX TIP:|PRESENTER NOTES:|$)', spec_text, re.DOTALL)
        if dimensions_match:
            dimensions_text = dimensions_match.group(1)

            # Parse each dimension line
            features_a = []
            features_b = []
            for line in dimensions_text.split('\n'):
                line = line.strip()
                if not line or not re.match(r'^\d+\.', line):
                    continue

                # Format: "1. Name → A: feature | B: feature"
                parts = line.split('→', 1)
                if len(parts) == 2:
                    dimension_name = re.sub(r'^\d+\.\s*', '', parts[0]).strip()
                    comparison = parts[1]

                    # Split by | to get A and B parts
                    if '|' in comparison:
                        a_part, b_part = comparison.split('|', 1)
                        a_value = re.sub(r'^A:\s*', '', a_part.strip())
                        b_value = re.sub(r'^B:\s*', '', b_part.strip())

                        features_a.append(f"{dimension_name}: {a_value}")
                        features_b.append(f"{dimension_name}: {b_value}")

            kd_data['concepts'].append({
                'number': 1,
                'name': concept_a_name,
                'features': features_a
            })

            kd_data['concepts'].append({
                'number': 2,
                'name': concept_b_name,
                'features': features_b
            })

            return kd_data

    # Try OLD format (Concept/Features format)
    layout_match = re.search(r'Layout:\s*([A-Z])\s*\((.+?)\)', spec_text)
    if layout_match:
        kd_data['layout'] = layout_match.group(1)

    # Extract concepts
    concept_pattern = r'-\s*Concept\s*(\d+):\s*"(.+?)"\s*\n\s*-\s*Features:\s*"(.+?)"'
    for match in re.finditer(concept_pattern, spec_text):
        concept_num = int(match.group(1))
        concept_name = match.group(2)
        features = match.group(3).split('", "')
        features = [f.strip('"') for f in features]
        kd_data['concepts'].append({
            'number': concept_num,
            'name': concept_name,
            'features': features
        })

    # Extract key differences
    kd_pattern = r'-\s*KD\d+:\s*"(.+?)"'
    kd_data['key_differences'] = re.findall(kd_pattern, spec_text)

    # Extract colors from old format
    color_pattern = r'-\s*Concept\s*(\d+)\s*(header|light):\s*RGB\((\d+),\s*(\d+),\s*(\d+)\)'
    for match in re.finditer(color_pattern, spec_text):
        concept_num = int(match.group(1))
        color_type = match.group(2)
        rgb = (int(match.group(3)), int(match.group(4)), int(match.group(5)))

        if concept_num not in kd_data['colors']:
            kd_data['colors'][concept_num] = {}
        kd_data['colors'][concept_num][color_type] = rgb

    return kd_data if kd_data['concepts'] else None

def add_key_differentiator_content(slide, kd_data):
    """Add key differentiator concept boxes to an existing slide."""
    # Clear the body textbox placeholder
    body_shape = find_shape_by_name(slide, 'TextBox 3') or find_shape_by_name(slide, 'TextBox 19')
    if body_shape:
        clear_shape_text(body_shape)

    # Layout concepts horizontally
    num_concepts = len(kd_data['concepts'])
    if num_concepts == 2:
        # Two-column layout with dynamic height based on content
        concept_width = Inches(5.5)  # Slightly wider for better text fit
        spacing = Inches(0.4)
        start_x = Inches(0.65)  # Adjusted for wider boxes
        start_y = Inches(1.3)

        # Calculate dynamic height based on number of features
        # Each concept has features, calculate max features count
        max_features = max(len(c['features']) for c in kd_data['concepts'])
        # Base height + additional height per feature (0.7" per feature line for 18pt font + bullet + spacing)
        features_height = Inches(0.7 + max_features * 0.7)
        concept_height = Inches(0.7) + features_height  # Header + features

        for i, concept in enumerate(kd_data['concepts']):
            x_pos = start_x + i * (concept_width + spacing)

            # Get colors for this concept
            header_color = kd_data['colors'].get(concept['number'], {}).get('header', (21, 101, 192))
            light_color = kd_data['colors'].get(concept['number'], {}).get('light', (227, 242, 253))

            # Create concept header box
            header_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, start_y,
                concept_width, Inches(0.7)
            )
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = RGBColor(*header_color)
            header_shape.line.color.rgb = RGBColor(*header_color)

            # Add concept name text
            header_tf = header_shape.text_frame
            header_tf.word_wrap = True
            header_p = header_tf.paragraphs[0]
            header_p.text = concept['name']
            header_p.font.size = Pt(24)
            header_p.font.bold = True
            header_p.font.name = "Aptos"
            header_p.font.color.rgb = RGBColor(255, 255, 255)
            header_p.alignment = PP_ALIGN.CENTER

            # Create features box (using dynamic height)
            features_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, start_y + Inches(0.7),
                concept_width, features_height
            )
            features_shape.fill.solid()
            features_shape.fill.fore_color.rgb = RGBColor(*light_color)
            features_shape.line.color.rgb = RGBColor(*header_color)
            features_shape.line.width = Pt(2)

            # Add features text
            features_tf = features_shape.text_frame
            features_tf.word_wrap = True
            features_tf.margin_left = Inches(0.2)
            features_tf.margin_right = Inches(0.2)
            features_tf.margin_top = Inches(0.2)
            features_tf.margin_bottom = Inches(0.2)

            # Add each feature as a bullet
            for j, feature in enumerate(concept['features']):
                if j > 0:
                    p = features_tf.add_paragraph()
                else:
                    p = features_tf.paragraphs[0]
                p.text = f"• {feature}"
                p.font.size = Pt(18)
                p.font.name = "Aptos"
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.space_before = Pt(6)
                p.space_after = Pt(6)

        # Add key differences box in center bottom (positioned dynamically)
        if kd_data['key_differences']:
            # Position below the concept boxes with some spacing
            kd_y = start_y + concept_height + Inches(0.3)
            kd_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(2.5), kd_y,
                Inches(8), Inches(1.2)
            )
            kd_box.fill.solid()
            kd_box.fill.fore_color.rgb = RGBColor(255, 235, 59)  # Yellow
            kd_box.line.color.rgb = RGBColor(245, 124, 0)  # Orange
            kd_box.line.width = Pt(3)

            kd_tf = kd_box.text_frame
            kd_tf.word_wrap = True
            kd_tf.margin_left = Inches(0.3)
            kd_tf.margin_right = Inches(0.3)

            # Add title
            title_p = kd_tf.paragraphs[0]
            title_p.text = "KEY DIFFERENCE:"
            title_p.font.size = Pt(18)
            title_p.font.bold = True
            title_p.font.name = "Aptos"
            title_p.font.color.rgb = RGBColor(0, 0, 0)
            title_p.alignment = PP_ALIGN.CENTER

            # Add differences
            for kd_text in kd_data['key_differences']:
                p = kd_tf.add_paragraph()
                p.text = kd_text
                p.font.size = Pt(18)  # Minimum 18pt for visual aid template slides
                p.font.name = "Aptos"
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER

def add_decision_tree_content(slide, dt_data):
    """Add decision tree diagram to an existing slide."""
    # Clear the body textbox placeholder
    body_shape = find_shape_by_name(slide, 'TextBox 3') or find_shape_by_name(slide, 'TextBox 19')
    if body_shape:
        clear_shape_text(body_shape)

    # Color mapping
    color_map = {
        'green': RGBColor(0, 102, 68),
        'red': RGBColor(153, 0, 0),
        'blue': RGBColor(0, 71, 133),
        'purple': RGBColor(75, 0, 110)
    }

    # Layout for 3-level binary tree (Layout A)
    # Level 1: Top center
    # Level 2A: Left, Level 2B: Right
    # Outcomes: O1, O2 below L2A; O3, O4 below L2B

    # Dimensions for nodes
    decision_node_width = Inches(2.5)
    decision_node_height = Inches(1.2)
    outcome_node_width = Inches(2.2)
    outcome_node_height = Inches(1.0)

    # Positions (Layout A: 1-2-4 structure)
    # Level 1 (top center)
    l1_x = Inches(5.4)
    l1_y = Inches(1.0)

    # Level 2 (middle row)
    l2a_x = Inches(2.5)
    l2a_y = Inches(2.8)
    l2b_x = Inches(8.3)
    l2b_y = Inches(2.8)

    # Outcomes (bottom row)
    o1_x = Inches(1.0)
    o1_y = Inches(4.8)
    o2_x = Inches(4.0)
    o2_y = Inches(4.8)
    o3_x = Inches(7.0)
    o3_y = Inches(4.8)
    o4_x = Inches(10.0)
    o4_y = Inches(4.8)

    # Create Level 1 decision node
    if dt_data['level1']:
        l1 = dt_data['level1']
        # Header box
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            l1_x, l1_y,
            decision_node_width, Inches(0.4)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(0, 51, 102)
        header.line.color.rgb = RGBColor(0, 51, 102)

        tf = header.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = l1['header']
        p.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Body box
        body = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            l1_x, l1_y + Inches(0.4),
            decision_node_width, Inches(0.8)
        )
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(240, 244, 248)
        body.line.color.rgb = RGBColor(0, 51, 102)
        body.line.width = Pt(2)

        tf = body.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = l1['question']
        p.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(33, 37, 41)
        p.alignment = PP_ALIGN.CENTER

    # Create Level 2A decision node
    if dt_data['level2a']:
        l2a = dt_data['level2a']
        # Header box
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            l2a_x, l2a_y,
            decision_node_width, Inches(0.4)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(0, 51, 102)
        header.line.color.rgb = RGBColor(0, 51, 102)

        tf = header.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = l2a['header']
        p.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Body box
        body = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            l2a_x, l2a_y + Inches(0.4),
            decision_node_width, Inches(0.8)
        )
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(240, 244, 248)
        body.line.color.rgb = RGBColor(0, 51, 102)
        body.line.width = Pt(2)

        tf = body.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = l2a['question']
        p.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(33, 37, 41)
        p.alignment = PP_ALIGN.CENTER

    # Create Level 2B decision node
    if dt_data['level2b']:
        l2b = dt_data['level2b']
        # Header box
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            l2b_x, l2b_y,
            decision_node_width, Inches(0.4)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(0, 51, 102)
        header.line.color.rgb = RGBColor(0, 51, 102)

        tf = header.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = l2b['header']
        p.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Body box
        body = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            l2b_x, l2b_y + Inches(0.4),
            decision_node_width, Inches(0.8)
        )
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(240, 244, 248)
        body.line.color.rgb = RGBColor(0, 51, 102)
        body.line.width = Pt(2)

        tf = body.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = l2b['question']
        p.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(33, 37, 41)
        p.alignment = PP_ALIGN.CENTER

    # Create outcome nodes
    outcome_positions = [
        (o1_x, o1_y), (o2_x, o2_y), (o3_x, o3_y), (o4_x, o4_y)
    ]

    for i, outcome in enumerate(dt_data['outcomes']):
        if i >= len(outcome_positions):
            break

        x, y = outcome_positions[i]
        color = color_map.get(outcome['color'], RGBColor(0, 71, 133))

        # Outcome box
        outcome_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y,
            outcome_node_width, outcome_node_height
        )
        outcome_box.fill.solid()
        outcome_box.fill.fore_color.rgb = color
        outcome_box.line.color.rgb = color
        outcome_box.line.width = Pt(2)

        tf = outcome_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = outcome['name']
        p.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Add connector lines using rectangle-based lines (not python-pptx connectors)
    line_color = RGBColor(100, 100, 100)

    # Lines from L1 to L2A and L2B
    if dt_data['level1'] and dt_data['level2a']:
        # L1 to L2A
        draw_line(slide,
                 l1_x + decision_node_width / 2, l1_y + decision_node_height,
                 l2a_x + decision_node_width / 2, l2a_y,
                 line_color, Pt(2))

    if dt_data['level1'] and dt_data['level2b']:
        # L1 to L2B
        draw_line(slide,
                 l1_x + decision_node_width / 2, l1_y + decision_node_height,
                 l2b_x + decision_node_width / 2, l2b_y,
                 line_color, Pt(2))

    # Lines from L2A to outcomes O1 and O2
    if dt_data['level2a'] and len(dt_data['outcomes']) >= 1:
        # L2A to O1
        draw_line(slide,
                 l2a_x + decision_node_width / 2, l2a_y + decision_node_height,
                 o1_x + outcome_node_width / 2, o1_y,
                 line_color, Pt(2))

    if dt_data['level2a'] and len(dt_data['outcomes']) >= 2:
        # L2A to O2
        draw_line(slide,
                 l2a_x + decision_node_width / 2, l2a_y + decision_node_height,
                 o2_x + outcome_node_width / 2, o2_y,
                 line_color, Pt(2))

    # Lines from L2B to outcomes O3 and O4
    if dt_data['level2b'] and len(dt_data['outcomes']) >= 3:
        # L2B to O3
        draw_line(slide,
                 l2b_x + decision_node_width / 2, l2b_y + decision_node_height,
                 o3_x + outcome_node_width / 2, o3_y,
                 line_color, Pt(2))

    if dt_data['level2b'] and len(dt_data['outcomes']) >= 4:
        # L2B to O4
        draw_line(slide,
                 l2b_x + decision_node_width / 2, l2b_y + decision_node_height,
                 o4_x + outcome_node_width / 2, o4_y,
                 line_color, Pt(2))

def add_flowchart_content(slide, fc_data):
    """Add flowchart diagram to an existing slide."""
    # Clear the body textbox placeholder
    body_shape = find_shape_by_name(slide, 'TextBox 3') or find_shape_by_name(slide, 'TextBox 19')
    if body_shape:
        clear_shape_text(body_shape)

    # Colors
    step_header_color = RGBColor(0, 102, 102)
    step_body_color = RGBColor(224, 242, 241)
    start_node_color = RGBColor(46, 125, 50)
    connector_color = RGBColor(0, 102, 102)

    # Layout B: Vertical flowchart
    step_width = Inches(3.0)
    step_height = Inches(1.0)
    spacing_y = Inches(0.4)

    # Calculate positions for vertical layout
    num_steps = len(fc_data['steps'])
    start_x = Inches(5.15)  # Center horizontally
    start_y = Inches(1.2)

    # Create step boxes
    for i, step in enumerate(fc_data['steps']):
        y_pos = start_y + i * (step_height + spacing_y)

        # Determine color (first step is green/start, rest are teal)
        if i == 0:
            box_color = start_node_color
            text_color = RGBColor(255, 255, 255)
        else:
            box_color = step_header_color
            text_color = RGBColor(255, 255, 255)

        # Create step box
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            start_x, y_pos,
            step_width, step_height
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = box_color
        step_box.line.color.rgb = box_color
        step_box.line.width = Pt(2)

        # Add text
        tf = step_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)

        # Title (bold)
        p = tf.paragraphs[0]
        p.text = step['title']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

        # Description (if room)
        if step['description']:
            p2 = tf.add_paragraph()
            p2.text = step['description']
            p2.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
            p2.font.name = "Aptos"
            p2.font.color.rgb = text_color
            p2.alignment = PP_ALIGN.CENTER

        # Add down arrow between steps using rectangle-based line
        if i < num_steps - 1:
            # Calculate arrow start and end positions (center of boxes)
            arrow_start_x = start_x + step_width / 2
            arrow_start_y = y_pos + step_height
            arrow_end_x = start_x + step_width / 2
            arrow_end_y = y_pos + step_height + spacing_y

            # Draw arrow
            draw_arrow(slide,
                      arrow_start_x, arrow_start_y,
                      arrow_end_x, arrow_end_y,
                      connector_color,
                      line_width=Pt(3),
                      arrow_size=Inches(0.12))

def add_hierarchy_content(slide, h_data):
    """Add hierarchy diagram to an existing slide."""
    # Clear the body textbox placeholder
    body_shape = find_shape_by_name(slide, 'TextBox 3') or find_shape_by_name(slide, 'TextBox 19')
    if body_shape:
        clear_shape_text(body_shape)

    # Colors - gradient from dark to light blue
    level1_color = RGBColor(1, 87, 155)
    level2_color = RGBColor(2, 119, 189)
    level3_color = RGBColor(3, 155, 229)

    # Layout A: Top-down tree
    # Level 1: 1 node at top center
    # Level 2: 2 nodes in middle row
    # Level 3: 4 nodes in bottom row

    l1_width = Inches(3.5)
    l1_height = Inches(0.8)
    l2_width = Inches(3.0)
    l2_height = Inches(0.7)
    l3_width = Inches(2.5)
    l3_height = Inches(0.7)

    # Level 1 position (top center)
    l1_x = Inches(4.9)
    l1_y = Inches(1.0)

    # Create Level 1 node
    if h_data['level1']:
        l1_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            l1_x, l1_y, l1_width, l1_height
        )
        l1_box.fill.solid()
        l1_box.fill.fore_color.rgb = level1_color
        l1_box.line.color.rgb = level1_color

        tf = l1_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = h_data['level1']['name']
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Level 2 positions (middle row, spread horizontally)
    num_l2 = len(h_data['level2'])
    if num_l2 == 2:
        l2_positions = [Inches(2.5), Inches(8.1)]
    else:
        # Distribute evenly
        l2_positions = [Inches(1.5 + i * 3.5) for i in range(num_l2)]

    l2_y = Inches(2.6)

    # Create Level 2 nodes
    for i, l2_node in enumerate(h_data['level2']):
        if i >= len(l2_positions):
            break

        l2_x = l2_positions[i]

        l2_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            l2_x, l2_y, l2_width, l2_height
        )
        l2_box.fill.solid()
        l2_box.fill.fore_color.rgb = level2_color
        l2_box.line.color.rgb = level2_color

        tf = l2_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = l2_node['name']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Draw lines from L1 to all L2 nodes
    line_color = RGBColor(100, 100, 100)
    if h_data['level1']:
        for i in range(min(len(h_data['level2']), len(l2_positions))):
            l2_x = l2_positions[i]
            draw_line(slide,
                     l1_x + l1_width / 2, l1_y + l1_height,
                     l2_x + l2_width / 2, l2_y,
                     line_color, Pt(2))

    # Level 3 positions and nodes
    l3_y = Inches(4.2)

    for l3_node in h_data['level3']:
        # Find parent L2 node index
        parent_name = l3_node['parent']
        parent_idx = -1
        for i, l2_node in enumerate(h_data['level2']):
            if l2_node['name'] == parent_name:
                parent_idx = i
                break

        if parent_idx == -1 or parent_idx >= len(l2_positions):
            continue

        # Position L3 node below its parent
        # If multiple L3 nodes share same parent, offset them
        parent_x = l2_positions[parent_idx]

        # Count how many L3 nodes have this parent
        siblings = [n for n in h_data['level3'] if n['parent'] == parent_name]
        sibling_idx = siblings.index(l3_node)
        num_siblings = len(siblings)

        if num_siblings == 1:
            # Single child: center under parent
            l3_x = parent_x + (l2_width - l3_width) / 2
        elif num_siblings == 2:
            # Two children: spread them out significantly to prevent overlap
            # L3 child width is 2.5", need at least 0.5" gap between them
            # Parent is 3.0" wide, so extend 1.25" on each side
            if sibling_idx == 0:
                l3_x = parent_x - Inches(1.25)  # Left child: extends left significantly
            else:
                l3_x = parent_x + l2_width - l3_width + Inches(1.25)  # Right child: extends right significantly
        else:
            # More than 2 children: distribute evenly
            horizontal_span = l2_width + Inches(1.0)
            l3_x = parent_x - Inches(0.5) + sibling_idx * (horizontal_span / (num_siblings - 1))

        l3_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            l3_x, l3_y, l3_width, l3_height
        )
        l3_box.fill.solid()
        l3_box.fill.fore_color.rgb = level3_color
        l3_box.line.color.rgb = level3_color

        tf = l3_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = l3_node['name']
        p.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Draw line from parent L2 to this L3 node
        draw_line(slide,
                 parent_x + l2_width / 2, l2_y + l2_height,
                 l3_x + l3_width / 2, l3_y,
                 line_color, Pt(2))

def add_timeline_content(slide, t_data):
    """Add timeline diagram to an existing slide."""
    # Clear the body textbox placeholder
    body_shape = find_shape_by_name(slide, 'TextBox 3') or find_shape_by_name(slide, 'TextBox 19')
    if body_shape:
        clear_shape_text(body_shape)

    # Colors
    bar_color = RGBColor(74, 20, 140)
    marker_color = RGBColor(123, 31, 162)
    card_background = RGBColor(243, 229, 245)
    card_border = RGBColor(74, 20, 140)

    # Layout A: Horizontal timeline
    num_events = len(t_data['events'])

    # Timeline bar
    bar_x = Inches(1.0)
    bar_y = Inches(3.5)
    bar_width = Inches(11.3)
    bar_height = Inches(0.2)

    timeline_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        bar_x, bar_y, bar_width, bar_height
    )
    timeline_bar.fill.solid()
    timeline_bar.fill.fore_color.rgb = bar_color
    timeline_bar.line.fill.background()

    # Event markers and cards
    event_spacing = bar_width / (num_events - 1) if num_events > 1 else 0

    for i, event in enumerate(t_data['events']):
        marker_x = bar_x + i * event_spacing

        # Event marker (circle)
        marker_size = Inches(0.3)
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            marker_x - marker_size/2, bar_y - marker_size/2 + bar_height/2,
            marker_size, marker_size
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = marker_color
        marker.line.color.rgb = RGBColor(255, 255, 255)
        marker.line.width = Pt(2)

        # Event card (alternating above/below bar)
        card_width = Inches(2.0)
        card_height = Inches(1.2)

        if i % 2 == 0:
            # Above the bar
            card_y = bar_y - card_height - Inches(0.4)
        else:
            # Below the bar
            card_y = bar_y + bar_height + Inches(0.4)

        card_x = marker_x - card_width/2

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            card_x, card_y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = card_background
        card.line.color.rgb = card_border
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)

        # Year
        p1 = tf.paragraphs[0]
        p1.text = event['year']
        p1.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p1.font.bold = True
        p1.font.name = "Aptos"
        p1.font.color.rgb = bar_color
        p1.alignment = PP_ALIGN.CENTER

        # Title
        p2 = tf.add_paragraph()
        p2.text = event['title']
        p2.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p2.font.bold = True
        p2.font.name = "Aptos"
        p2.font.color.rgb = RGBColor(0, 0, 0)
        p2.alignment = PP_ALIGN.CENTER

def add_spectrum_content(slide, s_data):
    """Add spectrum diagram to an existing slide."""
    # Clear the body textbox placeholder
    body_shape = find_shape_by_name(slide, 'TextBox 3') or find_shape_by_name(slide, 'TextBox 19')
    if body_shape:
        clear_shape_text(body_shape)

    # Layout A: Horizontal bar spectrum
    num_points = len(s_data['points'])

    bar_x = Inches(1.5)
    bar_y = Inches(2.5)
    bar_width = Inches(10.3)
    bar_height = Inches(1.0)

    # Create gradient spectrum bar using multiple segments
    segment_width = bar_width / num_points

    # Color gradient from light blue to dark blue
    for i in range(num_points):
        # Calculate color gradient
        ratio = i / (num_points - 1) if num_points > 1 else 0
        r = int(227 - ratio * 214)  # 227 -> 13
        g = int(242 - ratio * 171)  # 242 -> 71
        b = int(253 - ratio * 92)   # 253 -> 161

        segment = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bar_x + i * segment_width, bar_y,
            segment_width, bar_height
        )
        segment.fill.solid()
        segment.fill.fore_color.rgb = RGBColor(r, g, b)
        segment.line.fill.background()

    # Add labels for each point
    for i, point in enumerate(s_data['points']):
        label_x = bar_x + i * segment_width
        label_y = bar_y + bar_height + Inches(0.2)
        label_width = segment_width
        label_height = Inches(1.2)

        label_box = slide.shapes.add_textbox(
            label_x, label_y, label_width, label_height
        )

        tf = label_box.text_frame
        tf.word_wrap = True

        # Name (bold)
        p1 = tf.paragraphs[0]
        p1.text = point['name']
        p1.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p1.font.bold = True
        p1.font.name = "Aptos"
        p1.font.color.rgb = RGBColor(0, 0, 0)
        p1.alignment = PP_ALIGN.CENTER

        # Description
        p2 = tf.add_paragraph()
        p2.text = point['description'][:40]  # Truncate if too long
        p2.font.size = VISUAL_MIN_FONT_SIZE  # Minimum 18pt for visual aids
        p2.font.name = "Aptos"
        p2.font.color.rgb = RGBColor(60, 60, 60)
        p2.alignment = PP_ALIGN.CENTER

# ============================================
# VISUAL SLIDE GENERATION
# ============================================

def create_key_differentiator_slide(prs, slide_data, kd_data, visual_layout):
    """Create a key differentiator slide with colored concept boxes."""
    try:
        slide = prs.slides.add_slide(visual_layout)
        title_shape = find_shape_by_name(slide, SHAPE_TITLE)
        if title_shape and slide_data['header']:
            set_shape_text(title_shape, slide_data['header'], SHAPE_TITLE)

        num_concepts = len(kd_data['concepts'])
        if num_concepts >= 2:
            concept_width = Inches(5) if num_concepts == 2 else Inches(3.8)
            spacing = Inches(0.5) if num_concepts == 2 else Inches(0.3)
            start_x = Inches(0.75) if num_concepts == 2 else Inches(0.5)
            start_y = Inches(1.5)

            for i, concept in enumerate(kd_data['concepts']):
                x_pos = start_x + i * (concept_width + spacing)
                header_color = kd_data['colors'].get(concept['number'], {}).get('header', (21, 101, 192))
                light_color = kd_data['colors'].get(concept['number'], {}).get('light', (227, 242, 253))

                header_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, start_y, concept_width, Inches(0.7))
                header_box.fill.solid()
                header_box.fill.fore_color.rgb = RGBColor(*header_color)
                header_box.line.color.rgb = RGBColor(*header_color)
                header_tf = header_box.text_frame
                header_tf.word_wrap = True
                header_p = header_tf.paragraphs[0]
                header_p.text = concept['name']
                header_p.font.size = Pt(24 if num_concepts == 2 else 20)
                header_p.font.bold = True
                header_p.font.name = "Aptos"
                header_p.font.color.rgb = RGBColor(255, 255, 255)
                header_p.alignment = PP_ALIGN.CENTER

                features_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, start_y + Inches(0.7), concept_width, Inches(2.8 if num_concepts == 2 else 2.4))
                features_shape.fill.solid()
                features_shape.fill.fore_color.rgb = RGBColor(*light_color)
                features_shape.line.color.rgb = RGBColor(*header_color)
                features_shape.line.width = Pt(2)
                features_tf = features_shape.text_frame
                features_tf.word_wrap = True
                features_tf.margin_left = Inches(0.2)
                features_tf.margin_right = Inches(0.2)
                features_tf.margin_top = Inches(0.2)
                features_tf.margin_bottom = Inches(0.2)

                for j, feature in enumerate(concept['features']):
                    p = features_tf.add_paragraph() if j > 0 else features_tf.paragraphs[0]
                    p.text = f"* {feature}"
                    p.font.size = Pt(18 if num_concepts == 2 else 16)
                    p.font.name = "Aptos"
                    p.font.color.rgb = RGBColor(0, 0, 0)

            if kd_data.get('key_differences'):
                kd_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.2), Inches(8), Inches(1.2))
                kd_box.fill.solid()
                kd_box.fill.fore_color.rgb = RGBColor(255, 235, 59)
                kd_box.line.color.rgb = RGBColor(245, 124, 0)
                kd_box.line.width = Pt(3)
                kd_tf = kd_box.text_frame
                kd_tf.word_wrap = True
                kd_tf.margin_left = Inches(0.3)
                kd_tf.margin_right = Inches(0.3)
                title_p = kd_tf.paragraphs[0]
                title_p.text = "KEY DIFFERENCE:"
                title_p.font.size = Pt(18)
                title_p.font.bold = True
                title_p.font.name = "Aptos"
                title_p.font.color.rgb = RGBColor(0, 0, 0)
                title_p.alignment = PP_ALIGN.CENTER
                for kd_text in kd_data['key_differences']:
                    p = kd_tf.add_paragraph()
                    p.text = kd_text
                    p.font.size = Pt(18)
                    p.font.name = "Aptos"
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment = PP_ALIGN.CENTER

        if slide_data['notes']:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']
        return slide
    except Exception as e:
        print(f"  Warning: Could not create key differentiator slide: {e}")
        return None


def create_table_slide_from_template(prs, slide_data, visual_layout):
    """Create a table slide using visual layout inheritance."""

    # Create slide from visual layout (automatically inherits all background shapes)
    slide = prs.slides.add_slide(visual_layout)

    # Find and populate the title shape (already exists from layout)
    title_shape = find_shape_by_name(slide, SHAPE_TITLE)
    if title_shape and slide_data['header']:
        set_shape_text(title_shape, slide_data['header'], SHAPE_TITLE)

    # Add the table
    if slide_data.get('table_data'):
        create_table_on_slide(slide, slide_data['table_data'])

    # Add presenter notes
    if slide_data['notes']:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['notes']

    return slide


def create_vignette_answer_slide(prs, slide_data, visual_layout):
    """Create a Vignette or Answer slide using visual layout inheritance."""
    try:
        # Create slide from visual layout (automatically inherits all background shapes)
        slide = prs.slides.add_slide(visual_layout)

        # Find and populate the title shape (already exists from layout)
        title_shape = find_shape_by_name(slide, SHAPE_TITLE)
        if title_shape and slide_data['header']:
            set_shape_text(title_shape, slide_data['header'], SHAPE_TITLE)

        # Find and populate the body shape (already exists from layout)
        # Visual layout has TextBox 3 for body content
        body_shape = find_shape_by_name(slide, 'TextBox 3')
        if not body_shape:
            # Fallback to TextBox 19 if TextBox 3 not found
            body_shape = find_shape_by_name(slide, SHAPE_BODY)

        if body_shape and slide_data['body']:
            body_shape.text_frame.clear()
            p = body_shape.text_frame.paragraphs[0]
            p.text = slide_data['body']
            p.font.name = 'Aptos'
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0, 0, 0)

        # Add presenter notes
        if slide_data['notes']:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']

        return slide

    except Exception as e:
        print(f"  Warning: Could not create vignette/answer slide: {e}")
        return None


def create_visual_slide(prs, slide_data, diagrams_folder, section_num, slide_idx, visual_layout):
    """Create a custom slide for visual content using the visual layout."""

    # For Vignette/Answer slides, use visual aid template
    if slide_data.get('type') in ['Vignette', 'Answer']:
        vignette_slide = create_vignette_answer_slide(prs, slide_data, visual_layout)
        if vignette_slide:
            return vignette_slide
        # Fall through to default if it fails

    visual_type = slide_data.get('visual_type', '')

    # For table slides, use the table template
    if visual_type == 'table' and slide_data.get('table_data'):
        table_slide = create_table_slide_from_template(prs, slide_data, visual_layout)
        if table_slide:
            return table_slide
        # Fall through to default approach if table template method fails

    # For key differentiator slides, use the custom key differentiator function
    if visual_type == 'key_differentiators' and slide_data.get('key_differentiator_data'):
        kd_slide = create_key_differentiator_slide(prs, slide_data, slide_data['key_differentiator_data'], visual_layout)
        if kd_slide:
            return kd_slide
        # Fall through to default if it fails

    # Use visual layout for all other visual slides (fallback)
    slide = prs.slides.add_slide(visual_layout)

    # Add header text box (matching template style)
    header_left = Inches(0.5)
    header_top = Inches(0.3)
    header_width = Inches(12)
    header_height = Inches(0.8)

    header_box = slide.shapes.add_textbox(header_left, header_top, header_width, header_height)
    header_frame = header_box.text_frame
    header_para = header_frame.paragraphs[0]
    header_para.text = slide_data['header']
    header_para.font.size = HEADER_FONT_SIZE
    header_para.font.bold = True
    header_para.font.name = "Aptos"

    # Add visual content based on type
    if visual_type == 'table' and slide_data.get('table_data'):
        # Fallback if table template method failed
        create_table_on_slide(slide, slide_data['table_data'])
    elif visual_type == 'decision_tree' and slide_data.get('decision_tree_data'):
        # Decision trees use dedicated generator
        # See step12_decision_tree_generation.txt for full implementation
        create_decision_tree_on_slide(slide, slide_data['decision_tree_data'])
    else:
        # Look for diagram image
        diagram_pattern = f"section_{section_num}_slide_{slide_idx + 1}_*.png"
        diagram_files = list(Path(diagrams_folder).glob(diagram_pattern))

        if diagram_files:
            insert_diagram_on_slide(slide, str(diagram_files[0]))

    # Add presenter notes (required for all slides including visuals)
    if slide_data['notes']:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['notes']

    return slide

def create_table_on_slide(slide, table_data):
    """Create a native PowerPoint table using config settings."""
    if not table_data or len(table_data) < 1:
        return None

    rows = len(table_data)
    cols = len(table_data[0]) if table_data else 0

    if cols == 0:
        return None

    # Table position from config (anchored to bottom, above copyright)
    table_pos = TABLE_SETTINGS.get('position', {})
    left = Inches(table_pos.get('left_inches', 0.5))
    width = Inches(table_pos.get('width_inches', 12.3))
    height = Inches(min(rows * 0.5, table_pos.get('max_height_inches', 4.5)))
    top = Inches(table_pos.get('top_inches', 1.8))

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Get formatting from config
    font_name = TABLE_SETTINGS.get('font_name', 'Aptos')
    font_size = Pt(TABLE_SETTINGS.get('font_size_pt', 20))
    font_color = TABLE_SETTINGS.get('font_color_rgb', [0, 0, 0])
    header_color = TABLE_SETTINGS.get('header_font_color_rgb', [255, 255, 255])
    header_bg = TABLE_SETTINGS.get('header_background_rgb', [0, 51, 102])
    header_bold = TABLE_SETTINGS.get('header_bold', True)

    # Populate and format cells
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            if col_idx < len(table.columns):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)

                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = font_size
                    paragraph.font.name = font_name

                    if row_idx == 0:
                        # Header row formatting
                        paragraph.font.bold = header_bold
                        paragraph.font.color.rgb = RGBColor(header_color[0], header_color[1], header_color[2])
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(header_bg[0], header_bg[1], header_bg[2])
                    else:
                        # Data row formatting
                        paragraph.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])

    return table

def insert_diagram_on_slide(slide, image_path):
    """Insert a diagram image onto the slide."""
    try:
        from PIL import Image

        if not Path(image_path).exists():
            return None

        with Image.open(image_path) as img:
            img_width_px, img_height_px = img.size

        # Convert pixels to EMUs (914400 EMUs per inch, assuming 96 DPI)
        img_width_emu = Emu(img_width_px * 914400 / 96)
        img_height_emu = Emu(img_height_px * 914400 / 96)

        # Scale to fit within max dimensions
        scale_w = DIAGRAM_MAX_WIDTH / img_width_emu if img_width_emu > DIAGRAM_MAX_WIDTH else 1
        scale_h = DIAGRAM_MAX_HEIGHT / img_height_emu if img_height_emu > DIAGRAM_MAX_HEIGHT else 1
        scale = min(scale_w, scale_h, 1.0)

        final_width = int(img_width_emu * scale)
        final_height = int(img_height_emu * scale)

        # Center horizontally
        slide_width = Inches(13.33)
        left = int((slide_width - final_width) / 2)
        top = Inches(1.5)

        picture = slide.shapes.add_picture(image_path, left, top, final_width, final_height)
        return picture

    except Exception as e:
        print(f"  Warning: Could not insert diagram: {e}")
        return None


def create_decision_tree_on_slide(slide, tree_data):
    """
    Create a decision tree using native PowerPoint shapes.
    Placeholder implementation - see step12_decision_tree_generation.txt for full version.
    """
    # This is a simplified placeholder
    # Full implementation requires the dedicated decision tree generator
    print("  Note: Decision tree generation requires step12_decision_tree_generation.txt")
    return slide


# ============================================
# MAIN POPULATION LOGIC
# ============================================

def hide_duplicate_tier1_textboxes(prs):
    """Hide duplicate TIER 1 textboxes by clearing their text (no deletion to avoid corruption)."""
    CORRECT_LEFT_EMU = 10506456
    CORRECT_TOP_EMU = 64008

    for slide in prs.slides:
        # Clear text from TIER 1 textboxes NOT at correct position
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text_frame'):
                    left_inches = shape.left / 914400
                    top_inches = shape.top / 914400

                    # TIER 1 area (right side, top)
                    if left_inches > 11.0 and top_inches < 1.0:
                        # If NOT at correct position, clear text to hide it
                        if shape.left != CORRECT_LEFT_EMU or shape.top != CORRECT_TOP_EMU:
                            shape.text = ""  # Clear text only - no deletion to avoid corruption
            except:
                pass


def ensure_tier1_on_all_slides(prs):
    """Ensure all slides have TIER 1 text at correct position."""
    CORRECT_LEFT_EMU = 10506456
    CORRECT_TOP_EMU = 64008
    CORRECT_LEFT_IN = CORRECT_LEFT_EMU / 914400
    CORRECT_TOP_IN = CORRECT_TOP_EMU / 914400

    for slide in prs.slides:
        # Check if slide has TIER 1 text at correct position
        has_tier1_text = False
        tier1_shape = None

        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text_frame'):
                    if shape.left == CORRECT_LEFT_EMU and shape.top == CORRECT_TOP_EMU:
                        tier1_shape = shape
                        if shape.text and shape.text.strip():
                            has_tier1_text = True
                        break
            except:
                pass

        # Add text if missing
        if tier1_shape and not has_tier1_text:
            # Shape exists but has no text - add text
            tier1_shape.text = "TIER 1"
            if tier1_shape.text_frame.paragraphs:
                p = tier1_shape.text_frame.paragraphs[0]
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.name = 'Arial'
        elif not tier1_shape:
            # Shape doesn't exist - create it
            txBox = slide.shapes.add_textbox(
                Inches(CORRECT_LEFT_IN),
                Inches(CORRECT_TOP_IN),
                Inches(0.5),
                Inches(0.3)
            )
            tf = txBox.text_frame
            tf.text = "TIER 1"
            p = tf.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.name = 'Arial'


def duplicate_slide(prs, slide_index):
    """Duplicate a slide within the same presentation by cloning its XML."""
    template_slide = prs.slides[slide_index]

    # Get the slide layout
    blank_layout = template_slide.slide_layout

    # Create new slide with same layout
    new_slide = prs.slides.add_slide(blank_layout)

    # Copy all shapes from template slide to new slide
    for shape in template_slide.shapes:
        el = shape.element
        newel = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return new_slide

def populate_section(template_path, blueprint_path, output_path, diagrams_folder, section_num, log_entries):
    """Populate a single section PowerPoint from blueprint."""

    # Parse blueprint
    section_name, slides_data = parse_blueprint(blueprint_path)

    if not slides_data:
        log_entries.append(f"  ERROR: No slides parsed from blueprint")
        return False, section_name

    log_entries.append(f"  Section: {section_name}")
    log_entries.append(f"  Total slides in blueprint: {len(slides_data)}")

    # Separate content slides from visual slides
    # Treat Vignette/Answer slides as visual slides (custom generated)
    content_slides = [(idx, s) for idx, s in enumerate(slides_data) if not s['visual'] and s['type'] not in ['Vignette', 'Answer']]
    visual_slides = [(idx, s) for idx, s in enumerate(slides_data) if s['visual'] or s['type'] in ['Vignette', 'Answer']]

    log_entries.append(f"  Content slides (using template): {len(content_slides)}")
    log_entries.append(f"  Visual slides (custom generated): {len(visual_slides)}")

    # Load template master
    prs = Presentation(template_path)

    # Template master has 2 slides:
    # Slide 0 = content template (with NCLEX tip)
    # Slide 1 = visual aid template (without NCLEX tip)

    log_entries.append(f"  Loaded template master with {len(prs.slides)} template slides")

    # Keep track of which slides to populate
    all_slides_data = []

    # Create content slides by duplicating slide 0
    for content_idx, (orig_idx, slide_data) in enumerate(content_slides):
        # Duplicate the content template slide (slide 0)
        slide = duplicate_slide(prs, 0)
        all_slides_data.append(('content', slide, slide_data))
        log_entries.append(f"  Created content slide {content_idx + 1} (duplicated from template slide 1)")

    # Create visual slides by duplicating slide 1
    for visual_idx, (orig_idx, slide_data) in enumerate(visual_slides):
        # Duplicate the visual template slide (slide 1)
        slide = duplicate_slide(prs, 1)
        all_slides_data.append(('visual', slide, slide_data))
        log_entries.append(f"  Created visual slide {visual_idx + 1} (duplicated from template slide 2)")

    # Remove the original 2 template slides
    # Remove in reverse order to maintain indices
    rId = prs.slides._sldIdLst[1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[1]

    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

    log_entries.append(f"  Removed original template slides")

    # Now populate all slides
    for slide_type, slide, slide_data in all_slides_data:
        if slide_type == 'content':
            # Find shapes by exact name (already exist from duplication)
            title_shape = find_shape_by_name(slide, SHAPE_TITLE)
            body_shape = find_shape_by_name(slide, SHAPE_BODY)
            tip_shape = find_shape_by_name(slide, SHAPE_TIP)

            # Populate title/header (36pt white bold from config)
            if title_shape and slide_data['header']:
                set_shape_text(title_shape, slide_data['header'], SHAPE_TITLE)

            # Populate body (20pt from config)
            if body_shape and slide_data['body']:
                set_shape_text(body_shape, slide_data['body'], SHAPE_BODY)

            # Handle NCLEX tip based on slide type
            if tip_shape:
                if slide_data['tip'] and slide_data['type'] in ['Content', 'Core Content', 'Connection Slide']:
                    set_shape_text(tip_shape, slide_data['tip'], SHAPE_TIP)
                else:
                    # Clear tip for non-content slides
                    clear_shape_text(tip_shape)

            # Add presenter notes
            if slide_data['notes']:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data['notes']

            log_entries.append(f"    Content: {slide_data['type']} - {slide_data['header'][:35]}...")

        else:  # visual slide
            # Find and populate title
            title_shape = find_shape_by_name(slide, SHAPE_TITLE)
            if title_shape and slide_data['header']:
                set_shape_text(title_shape, slide_data['header'], SHAPE_TITLE)

            # Handle visual content based on type
            visual_type = slide_data.get('visual_type', '')

            # For table slides
            if visual_type == 'table' and slide_data.get('table_data'):
                create_table_on_slide(slide, slide_data['table_data'])

            # For key differentiator slides - add concept boxes
            elif visual_type == 'key_differentiators' and slide_data.get('key_differentiator_data'):
                kd_data = slide_data['key_differentiator_data']
                # Add the concept comparison boxes on top of the template
                add_key_differentiator_content(slide, kd_data)

            # For decision tree slides - add decision tree diagram
            elif visual_type == 'decision_tree' and slide_data.get('decision_tree_data'):
                add_decision_tree_content(slide, slide_data['decision_tree_data'])

            # For flowchart slides - add flowchart diagram
            elif visual_type == 'flowchart' and slide_data.get('flowchart_data'):
                add_flowchart_content(slide, slide_data['flowchart_data'])

            # For hierarchy slides - add hierarchy diagram
            elif visual_type == 'hierarchy' and slide_data.get('hierarchy_data'):
                add_hierarchy_content(slide, slide_data['hierarchy_data'])

            # For timeline slides - add timeline diagram
            elif visual_type == 'timeline' and slide_data.get('timeline_data'):
                add_timeline_content(slide, slide_data['timeline_data'])

            # For spectrum slides - add spectrum diagram
            elif visual_type == 'spectrum' and slide_data.get('spectrum_data'):
                add_spectrum_content(slide, slide_data['spectrum_data'])

            # For vignette/answer slides
            elif slide_data.get('type') in ['Vignette', 'Answer']:
                # Find body textbox
                body_shape = find_shape_by_name(slide, 'TextBox 3')
                if not body_shape:
                    body_shape = find_shape_by_name(slide, SHAPE_BODY)
                if body_shape and slide_data['body']:
                    body_shape.text_frame.clear()
                    p = body_shape.text_frame.paragraphs[0]
                    p.text = slide_data['body']
                    p.font.name = 'Aptos'
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(0, 0, 0)

            # Add presenter notes
            if slide_data['notes']:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data['notes']

            log_entries.append(f"    Visual: {visual_type or slide_data.get('type', 'unknown')} - {slide_data['header'][:35]}...")

    # Ensure all slides have TIER 1 text at correct position
    ensure_tier1_on_all_slides(prs)
    log_entries.append(f"  Ensured TIER 1 on all slides")

    # Hide duplicate TIER 1 textboxes (clear text without deleting to avoid corruption)
    hide_duplicate_tier1_textboxes(prs)
    log_entries.append(f"  Hidden duplicate TIER 1 textboxes")

    # Save presentation
    prs.save(output_path)
    log_entries.append(f"  Final slide count: {len(prs.slides)}")
    log_entries.append(f"  Saved: {output_path}")

    return True, section_name

def main():
    """Main execution function."""
    log_entries = []
    log_entries.append("=" * 60)
    log_entries.append("STEP 12: POWERPOINT POPULATION LOG")
    log_entries.append("=" * 60)
    log_entries.append(f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    log_entries.append(f"Template: {TEMPLATE_PATH}")
    log_entries.append(f"Domain: {DOMAIN_NAME}")
    log_entries.append("")

    # Validate paths
    if not Path(TEMPLATE_PATH).exists():
        print(f"ERROR: Template not found: {TEMPLATE_PATH}")
        return

    production_path = Path(PRODUCTION_FOLDER)
    if not production_path.exists():
        print(f"ERROR: Production folder not found: {PRODUCTION_FOLDER}")
        return

    # Use integrated/ subfolder for blueprints (Step 10 outputs)
    integrated_folder = production_path / "integrated"
    diagrams_folder = production_path / "diagrams"

    # Output to production folder's powerpoints/ subfolder
    powerpoints_folder = production_path / "powerpoints"
    logs_folder = production_path / "logs"

    powerpoints_folder.mkdir(parents=True, exist_ok=True)
    logs_folder.mkdir(parents=True, exist_ok=True)

    log_entries.append(f"Output folder: {powerpoints_folder}")
    log_entries.append("")

    # Find all integrated blueprint files (from Step 10)
    blueprint_files = sorted(integrated_folder.glob("step10_integrated_*.txt"))

    if not blueprint_files:
        log_entries.append("ERROR: No integrated blueprint files found in integrated/ subfolder")
        print("ERROR: No blueprint files found in integrated folder")
        return

    log_entries.append(f"Found {len(blueprint_files)} blueprint files")
    log_entries.append("")

    success_count = 0
    error_count = 0

    # Process each section
    for section_num, blueprint_file in enumerate(blueprint_files, 1):
        log_entries.append("-" * 60)
        log_entries.append(f"PROCESSING SECTION {section_num}")
        log_entries.append("-" * 60)
        log_entries.append(f"  Blueprint: {blueprint_file.name}")

        try:
            # Parse to get section name for output filename
            section_name, _ = parse_blueprint(str(blueprint_file))
            safe_name = sanitize_filename(section_name)
            output_path = powerpoints_folder / f"{safe_name}.pptx"

            success, _ = populate_section(
                TEMPLATE_PATH,
                str(blueprint_file),
                str(output_path),
                str(diagrams_folder),
                section_num,
                log_entries
            )

            if success:
                success_count += 1
            else:
                error_count += 1

        except Exception as e:
            log_entries.append(f"  ERROR: {str(e)}")
            error_count += 1
            import traceback
            log_entries.append(f"  Traceback: {traceback.format_exc()}")

        log_entries.append("")

    # Summary
    log_entries.append("=" * 60)
    log_entries.append("SUMMARY")
    log_entries.append("=" * 60)
    log_entries.append(f"Sections processed: {len(blueprint_files)}")
    log_entries.append(f"Successful: {success_count}")
    log_entries.append(f"Errors: {error_count}")
    log_entries.append("")

    if error_count == 0:
        log_entries.append("STATUS: ALL SECTIONS POPULATED SUCCESSFULLY")
    else:
        log_entries.append("STATUS: COMPLETED WITH ERRORS - Review log")

    log_entries.append("=" * 60)

    # Write log file to logs/ subfolder
    log_path = logs_folder / "population_log.txt"
    with open(log_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(log_entries))

    print(f"\nPopulation complete!")
    print(f"PowerPoints: {powerpoints_folder}")
    print(f"Log: {log_path}")
    print(f"Successful: {success_count}, Errors: {error_count}")


if __name__ == "__main__":
    main()
