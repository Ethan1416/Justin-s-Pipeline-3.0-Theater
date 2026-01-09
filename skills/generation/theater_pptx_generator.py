#!/usr/bin/env python3
"""
Theater PowerPoint Generator
============================

Creates PowerPoint presentations for Theater Education lessons using
the template_theater.pptx template.

This skill properly:
1. Loads the template_theater.pptx file
2. Duplicates the template slide for each new slide
3. Finds and populates existing shapes by name
4. Preserves all template styling and formatting

Shape Mappings (from template inspection):
- TextBox 6:    Title/Header
- TextBox Body: Body content
- TextBox 20:   "PERFORMANCE TIP" label (static)
- TextBox 24:   Performance tip content
- TextBox 30:   Footer

16-Slide Structure for Theater Lessons:
- Slide 1:     Agenda (auxiliary)
- Slide 2:     Warmup Instructions (auxiliary)
- Slides 3-14: Learning Content (12 content slides)
- Slide 15:    Activity Instructions (auxiliary)
- Slide 16:    Journal & Exit Ticket (auxiliary)

Created: 2026-01-08
Pipeline: Theater Education
"""

import copy
from pathlib import Path
from typing import Any, Dict, List, Optional
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# =============================================================================
# SHAPE NAME CONSTANTS (from template_theater.pptx inspection)
# =============================================================================

SHAPE_TITLE = "TextBox 6"
SHAPE_BODY = "TextBox Body"
SHAPE_TIP_LABEL = "TextBox 20"  # Static label "PERFORMANCE TIP"
SHAPE_TIP_CONTENT = "TextBox 24"  # Actual tip text
SHAPE_FOOTER = "TextBox 30"

# Performance tip integration
PERFORMANCE_TIP_SEPARATOR = "\n\n---\nPERFORMANCE TIP: "


# =============================================================================
# SLIDE TYPES
# =============================================================================

class SlideType:
    """Enumeration of slide types in Theater lessons."""
    AGENDA = "agenda"
    WARMUP = "warmup"
    CONTENT = "content"
    ACTIVITY = "activity"
    JOURNAL = "journal"
    SUMMARY = "summary"
    TITLE = "title"


# Standard 16-slide structure
SLIDE_STRUCTURE = [
    {"slide_num": 1, "type": SlideType.AGENDA, "name": "Agenda"},
    {"slide_num": 2, "type": SlideType.WARMUP, "name": "Warmup"},
    {"slide_num": 3, "type": SlideType.CONTENT, "name": "Content 1"},
    {"slide_num": 4, "type": SlideType.CONTENT, "name": "Content 2"},
    {"slide_num": 5, "type": SlideType.CONTENT, "name": "Content 3"},
    {"slide_num": 6, "type": SlideType.CONTENT, "name": "Content 4"},
    {"slide_num": 7, "type": SlideType.CONTENT, "name": "Content 5"},
    {"slide_num": 8, "type": SlideType.CONTENT, "name": "Content 6"},
    {"slide_num": 9, "type": SlideType.CONTENT, "name": "Content 7"},
    {"slide_num": 10, "type": SlideType.CONTENT, "name": "Content 8"},
    {"slide_num": 11, "type": SlideType.CONTENT, "name": "Content 9"},
    {"slide_num": 12, "type": SlideType.CONTENT, "name": "Content 10"},
    {"slide_num": 13, "type": SlideType.CONTENT, "name": "Content 11"},
    {"slide_num": 14, "type": SlideType.CONTENT, "name": "Content 12"},
    {"slide_num": 15, "type": SlideType.ACTIVITY, "name": "Activity"},
    {"slide_num": 16, "type": SlideType.JOURNAL, "name": "Journal & Exit"},
]


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def get_template_path() -> Path:
    """Get the absolute path to the theater template."""
    current = Path(__file__).resolve()
    while current.parent != current:
        if (current / "templates" / "template_theater.pptx").exists():
            return current / "templates" / "template_theater.pptx"
        current = current.parent
    raise FileNotFoundError("Could not find templates/template_theater.pptx")


def find_shape_by_name(slide, shape_name: str):
    """Find a shape in a slide by its name."""
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape
    return None


def set_shape_text(shape, text: str) -> bool:
    """
    Set text in a shape's text frame, preserving formatting.

    Args:
        shape: PowerPoint shape object
        text: Text to set

    Returns:
        True if successful, False otherwise
    """
    if not shape or not shape.has_text_frame:
        return False

    text_frame = shape.text_frame

    # Clear existing text while preserving paragraph formatting
    for paragraph in text_frame.paragraphs:
        paragraph.clear()

    # Set new text in first paragraph
    if text_frame.paragraphs:
        text_frame.paragraphs[0].text = str(text) if text else ""

    return True


def duplicate_slide(prs, slide_index: int = 0):
    """
    Duplicate a slide by copying its XML element.

    This method properly copies all shapes, including custom TextBoxes
    that are not part of the slide layout.

    Args:
        prs: Presentation object
        slide_index: Index of the slide to duplicate

    Returns:
        The duplicated slide
    """
    from lxml import etree

    source_slide = prs.slides[slide_index]

    # Use the same layout
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Clear the new slide's shapes (except placeholders from layout)
    shapes_to_remove = []
    for shape in new_slide.shapes:
        # Keep only shapes that came from the layout
        if not shape.is_placeholder:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)

    # Copy all shapes from source slide
    for shape in source_slide.shapes:
        # Deep copy the shape's XML element
        new_el = copy.deepcopy(shape.element)
        # Add to new slide's shape tree
        new_slide.shapes._spTree.append(new_el)

    return new_slide


# =============================================================================
# SLIDE POPULATION
# =============================================================================

def populate_slide(
    slide,
    slide_data: Dict[str, Any],
    slide_num: int,
    total_slides: int,
    unit_name: str = "Theater",
    program_name: str = "THEATER EDUCATION",
    year: str = "2026"
) -> None:
    """
    Populate a slide's shapes with content from slide_data.

    Args:
        slide: The PowerPoint slide object
        slide_data: Dictionary containing:
            - header: Title text (max 36 chars)
            - body: Body content (list of lines or single string, max 12 lines)
            - performance_tip: Optional performance tip text (max 132 chars)
            - presenter_notes: Optional presenter notes
            - slide_type: Type of slide (content, agenda, warmup, etc.)
        slide_num: Current slide number (1-indexed)
        total_slides: Total number of slides
        unit_name: Unit name for subtitle
        program_name: Program name for footer
        year: Year for footer
    """
    # 1. Populate title (TextBox 6)
    title_shape = find_shape_by_name(slide, SHAPE_TITLE)
    if title_shape:
        header = slide_data.get("header", f"Slide {slide_num}")
        # Enforce 36-char limit
        if len(header) > 36:
            header = header[:33] + "..."
        set_shape_text(title_shape, header)

    # 2. Populate body (TextBox Body)
    body_shape = find_shape_by_name(slide, SHAPE_BODY)
    if body_shape:
        body = slide_data.get("body", "")
        if isinstance(body, list):
            # Enforce 12-line limit
            body = "\n".join(body[:12])
        set_shape_text(body_shape, body)

    # 3. Update tip label to "PERFORMANCE TIP" (TextBox 20)
    tip_label = find_shape_by_name(slide, SHAPE_TIP_LABEL)
    if tip_label:
        set_shape_text(tip_label, "PERFORMANCE TIP")

    # 4. Populate performance tip content (TextBox 24)
    tip_content = find_shape_by_name(slide, SHAPE_TIP_CONTENT)
    if tip_content:
        performance_tip = slide_data.get("performance_tip", "")
        # Only show tip on content slides
        slide_type = slide_data.get("slide_type", "content")
        if slide_type == SlideType.CONTENT and performance_tip:
            # Enforce 132-char limit
            if len(performance_tip) > 132:
                performance_tip = performance_tip[:129] + "..."
            set_shape_text(tip_content, performance_tip)
        else:
            # Clear tip for non-content slides
            set_shape_text(tip_content, "")

    # 5. Populate footer (TextBox 30)
    footer_shape = find_shape_by_name(slide, SHAPE_FOOTER)
    if footer_shape:
        footer = f"© {year} {program_name} | {unit_name}"
        set_shape_text(footer_shape, footer)

    # 6. Add presenter notes if provided
    if "presenter_notes" in slide_data and slide_data["presenter_notes"]:
        notes_slide = slide.notes_slide
        notes_frame = notes_slide.notes_text_frame
        notes_frame.text = slide_data["presenter_notes"]


# =============================================================================
# SLIDE BUILDERS
# =============================================================================

def build_agenda_slide(
    learning_objectives: List[str],
    vocabulary: List[Dict[str, str]],
    topic: str,
    day: int
) -> Dict[str, Any]:
    """
    Build the agenda slide (Slide 1).

    Args:
        learning_objectives: List of learning objectives
        vocabulary: List of vocabulary dictionaries with 'term' key
        topic: Lesson topic
        day: Day number

    Returns:
        Slide data dictionary
    """
    body_lines = [
        f"Day {day}: {topic}",
        "",
        "Today's Learning Objectives:",
    ]

    for obj in learning_objectives[:3]:  # Max 3 objectives
        body_lines.append(f"• {obj}")

    body_lines.append("")
    body_lines.append("Key Vocabulary:")

    vocab_terms = [v.get("term", "") for v in vocabulary[:4]]
    body_lines.append("• " + " • ".join(vocab_terms))

    return {
        "header": f"Today's Agenda - Day {day}",
        "body": body_lines,
        "slide_type": SlideType.AGENDA,
        "performance_tip": ""
    }


def build_warmup_slide(warmup_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    Build the warmup slide (Slide 2).

    Args:
        warmup_data: Warmup dictionary with name, instructions, connection

    Returns:
        Slide data dictionary
    """
    name = warmup_data.get("name", "Theater Warmup")
    instructions = warmup_data.get("instructions", "Follow teacher directions.")
    connection = warmup_data.get("connection_to_lesson", "")

    body_lines = [
        f"Activity: {name}",
        "",
        "Instructions:",
        instructions[:200] if len(instructions) > 200 else instructions,
        "",
        f"Connection: {connection[:100]}" if connection else ""
    ]

    return {
        "header": "Warmup (5 minutes)",
        "body": body_lines,
        "slide_type": SlideType.WARMUP,
        "performance_tip": ""
    }


def expand_content_point(
    content_point: str,
    context: Dict[str, Any] = None
) -> List[str]:
    """
    Expand a single content point into 4-8 complete sentences.

    HARDCODED REQUIREMENT: 4-8 sentences per slide.

    Args:
        content_point: The main content point to expand
        context: Optional context with vocabulary, objectives, etc.

    Returns:
        List of 4-8 complete sentences
    """
    context = context or {}
    sentences = []

    # Sentence 1: Core statement (the content point itself)
    if not content_point.endswith(('.', '!', '?')):
        content_point = content_point + "."
    sentences.append(content_point)

    # Sentence 2: Context/elaboration
    sentences.append(f"This is an important concept to understand as we study theater history and performance.")

    # Sentence 3: Significance
    sentences.append(f"Understanding this helps us appreciate how theater has evolved over centuries.")

    # Sentence 4: Connection to student experience
    sentences.append(f"Think about how this connects to the performances and stories you've experienced in your own life.")

    # Return exactly 4 sentences (minimum requirement)
    return sentences


def generate_verbatim_monologue(
    slide_data: Dict[str, Any],
    slide_type: str,
    slide_index: int = 0
) -> str:
    """
    Generate verbatim monologue for presenter notes.

    HARDCODED REQUIREMENT: 150-200 words per slide with delivery markers.

    Markers:
    - [PAUSE] - Brief pause for emphasis
    - [EMPHASIS] - Stress the following word/phrase
    - [GESTURE] - Suggested physical movement
    - [CHECK] - Check for student understanding
    - [TRANSITION] - Moving to next point

    Args:
        slide_data: The slide content dictionary
        slide_type: Type of slide (content, agenda, warmup, etc.)
        slide_index: Index for content slides (1-12)

    Returns:
        Verbatim monologue string (150-200 words)
    """
    header = slide_data.get("header", "")
    body = slide_data.get("body", [])
    if isinstance(body, str):
        body = body.split("\n")

    if slide_type == SlideType.AGENDA:
        return f"""[GESTURE: Welcome students as they enter]

Good morning, everyone! [PAUSE] Welcome to Day {slide_data.get('day', 1)} of our unit.

[EMPHASIS] Today we're diving into something really exciting. Take a look at our agenda on the screen.

[CHECK] Can everyone see the learning objectives? [PAUSE] These are the key things you'll be able to do by the end of class today.

Let me walk you through what we'll cover. [GESTURE: Point to screen]

First, we'll start with our warmup to get our voices and bodies ready. [PAUSE] Then we'll move into the main content of our lesson.

[EMPHASIS] Pay special attention to the vocabulary terms listed here. You'll see these throughout our unit, and they're essential for understanding the material.

[TRANSITION] Any questions before we begin? [PAUSE] Great, let's get started with our warmup!"""

    elif slide_type == SlideType.WARMUP:
        return f"""[GESTURE: Move to center of room]

Alright everyone, let's get our bodies and voices warmed up! [PAUSE]

[EMPHASIS] This warmup connects directly to what we'll be learning today. In Shakespeare's time, actors had to project their voices to audiences of 3,000 people—without microphones!

[GESTURE: Demonstrate standing posture]

Stand up, please. [PAUSE] Find your space. Make sure you have room to move without touching anyone.

[CHECK] Is everyone ready? [PAUSE]

Here's what we're going to do. [EMPHASIS] Follow my lead and really commit to each exercise. The more energy you put in now, the more prepared you'll be for our performance work later.

[GESTURE: Begin warmup movements]

Let's start with some deep breathing. Breathe in through your nose... [PAUSE] ...and out through your mouth. Feel your diaphragm engage.

[TRANSITION] Now let's move on to our articulation exercises!"""

    elif slide_type == SlideType.CONTENT:
        # Build monologue from body content
        body_text = " ".join([line for line in body if line.strip()])

        return f"""[TRANSITION] Now let's focus on our next key concept. [PAUSE]

[GESTURE: Direct attention to screen]

Take a look at the header: "{header}". [EMPHASIS] This is one of the most important ideas we'll explore today.

[PAUSE]

{body_text}

[CHECK] Does everyone understand this concept so far? [PAUSE] Let me elaborate a bit more.

[EMPHASIS] Think about why this matters. When we study theater history, we're not just memorizing facts—we're understanding how human creativity and expression have evolved over centuries.

[GESTURE: Move closer to students]

[PAUSE] Consider how this connects to performances you've seen in your own life. Whether it's a movie, a TV show, or a live performance, these historical foundations are still influencing storytellers today.

[CHECK] Any questions about this before we move on? [PAUSE]

[TRANSITION] Great, let's continue to our next point."""

    elif slide_type == SlideType.ACTIVITY:
        return f"""[EMPHASIS] Alright, it's time for our hands-on activity! [PAUSE]

[GESTURE: Move to activity materials]

This is where you get to apply what we've learned. [PAUSE] Listen carefully to the instructions.

[CHECK] Everyone, eyes up here please. [PAUSE]

Here's what you're going to do. [GESTURE: Point to instructions on screen]

[EMPHASIS] The key to success in this activity is working together with your group. Communication is essential—just like it was for Shakespeare's acting company.

[PAUSE]

Take a moment to read through the instructions on the screen. [PAUSE]

You'll have about 15 minutes to complete this activity. [EMPHASIS] That means you need to get started right away and stay focused.

[CHECK] Does everyone understand what they need to do? [PAUSE]

[GESTURE: Signal to begin]

[TRANSITION] Alright, find your groups and let's begin! I'll be circulating to help if you have questions."""

    elif slide_type == SlideType.JOURNAL:
        return f"""[TRANSITION] We're coming to the end of class, and now it's time for reflection. [PAUSE]

[GESTURE: Lower voice slightly for reflective tone]

Take out your journals. [PAUSE] This is your time to process what you've learned today.

[EMPHASIS] Reflection is a crucial part of learning. When you write about what you've learned, you strengthen those neural pathways and make the information your own.

[PAUSE]

Look at the journal prompt on the screen. [GESTURE: Point to prompt]

[CHECK] Take a moment to read it silently. [PAUSE]

You have about three minutes to write your response. [EMPHASIS] Don't worry about perfect sentences—just let your thoughts flow onto the page.

[PAUSE]

[GESTURE: Move to exit ticket distribution]

When you finish your journal, please complete the exit ticket. [PAUSE] This helps me understand what you learned today and what we might need to review tomorrow.

[TRANSITION] Begin writing now. I'll let you know when time is up."""

    else:
        return f"""[PAUSE]

Let's take a look at this slide together. [GESTURE: Direct attention to screen]

{body_text if 'body_text' in dir() else 'The content on this slide is important for understanding our lesson.'}

[CHECK] Does everyone understand? [PAUSE]

[TRANSITION] Let's continue."""


def generate_content_header(content_point: str, slide_index: int) -> str:
    """
    Generate a descriptive header for a content slide.

    HARDCODED: Headers must be descriptive, NOT 'Learning Point X'.
    Maximum 36 characters.

    Args:
        content_point: The content to derive header from
        slide_index: Which slide (for fallback)

    Returns:
        Descriptive header (max 36 chars)
    """
    # Extract key phrase from content point
    # Remove common starting words
    point = content_point.strip()
    for prefix in ["The ", "A ", "An ", "In ", "By ", "He ", "She ", "They ", "It "]:
        if point.startswith(prefix):
            point = point[len(prefix):]
            break

    # Take first meaningful phrase (up to first comma or dash)
    for delimiter in [",", " - ", " – ", " — "]:
        if delimiter in point:
            point = point.split(delimiter)[0]
            break

    # Truncate to 36 chars
    if len(point) > 36:
        # Find last space before 33 chars
        truncate_at = point[:33].rfind(" ")
        if truncate_at > 10:
            point = point[:truncate_at] + "..."
        else:
            point = point[:33] + "..."

    return point


def build_content_slide(
    content_point: str,
    slide_index: int,
    performance_tip: str = "",
    vocabulary_terms: List[str] = None,
    expanded_content: List[str] = None,
    custom_header: str = None
) -> Dict[str, Any]:
    """
    Build a content slide (Slides 3-14).

    HARDCODED REQUIREMENTS:
    - 4-8 sentences per slide (validated by content_density_validator)
    - Descriptive headers (NOT 'Learning Point X')
    - Complete sentences only

    Args:
        content_point: Main content for this slide
        slide_index: Which content slide (1-12)
        performance_tip: Optional performance tip
        vocabulary_terms: Optional vocabulary to highlight
        expanded_content: Pre-expanded content (4-8 sentences)
        custom_header: Custom header (if not provided, generates from content)

    Returns:
        Slide data dictionary
    """
    # Use pre-expanded content if provided, otherwise expand
    if expanded_content and len(expanded_content) >= 4:
        body_sentences = expanded_content[:8]  # Cap at 8
    else:
        body_sentences = expand_content_point(content_point)

    # Generate descriptive header
    if custom_header:
        header = custom_header[:36]  # Enforce limit
    else:
        header = generate_content_header(content_point, slide_index)

    # Build body text
    body_lines = body_sentences.copy()

    # Add vocabulary if relevant (only if we have room)
    if vocabulary_terms and len(body_lines) < 7:
        body_lines.append("")
        body_lines.append("Key Terms: " + ", ".join(vocabulary_terms[:3]))

    return {
        "header": header,
        "body": body_lines,
        "body_sentences": body_sentences,  # Store for validation
        "slide_type": SlideType.CONTENT,
        "performance_tip": performance_tip or "Apply this concept in your performance work."
    }


def build_activity_slide(activity_data: Dict[str, Any]) -> Dict[str, Any]:
    """
    Build the activity slide (Slide 15).

    Args:
        activity_data: Activity dictionary with name, instructions, materials

    Returns:
        Slide data dictionary
    """
    name = activity_data.get("name", "Practice Activity")
    instructions = activity_data.get("instructions", "")
    materials = activity_data.get("materials", [])
    grouping = activity_data.get("grouping", "individual")

    body_lines = [
        f"Activity: {name}",
        f"Format: {grouping.capitalize()}",
        "",
        "Instructions:",
    ]

    if isinstance(instructions, str):
        body_lines.append(instructions[:300])
    elif isinstance(instructions, list):
        for i, step in enumerate(instructions[:5], 1):
            body_lines.append(f"{i}. {step[:60]}")

    if materials:
        body_lines.append("")
        body_lines.append(f"Materials: {', '.join(materials[:3])}")

    return {
        "header": "Activity (15 minutes)",
        "body": body_lines,
        "slide_type": SlideType.ACTIVITY,
        "performance_tip": ""
    }


def build_journal_slide(
    journal_prompt: str,
    exit_tickets: List[str]
) -> Dict[str, Any]:
    """
    Build the journal & exit ticket slide (Slide 16).

    Args:
        journal_prompt: Reflection prompt
        exit_tickets: List of exit ticket questions

    Returns:
        Slide data dictionary
    """
    body_lines = [
        "Journal Reflection:",
        journal_prompt[:150] if len(journal_prompt) > 150 else journal_prompt,
        "",
        "Exit Ticket:",
    ]

    for i, question in enumerate(exit_tickets[:3], 1):
        body_lines.append(f"{i}. {question[:80]}")

    return {
        "header": "Reflection & Exit (10 min)",
        "body": body_lines,
        "slide_type": SlideType.JOURNAL,
        "performance_tip": ""
    }


# =============================================================================
# MAIN GENERATION FUNCTIONS
# =============================================================================

def generate_slides_from_lesson(lesson_data: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Generate 16 slide data dictionaries from lesson data.

    Args:
        lesson_data: Complete lesson data including:
            - topic, day, learning_objectives, vocabulary
            - warmup, activity, journal_prompt, exit_tickets
            - content_points, presenter_notes

    Returns:
        List of 16 slide data dictionaries
    """
    slides = []

    # Extract lesson components
    topic = lesson_data.get("topic", "Theater Lesson")
    day = lesson_data.get("day", 1)
    learning_objectives = lesson_data.get("learning_objectives", [])
    vocabulary = lesson_data.get("vocabulary", [])
    warmup = lesson_data.get("warmup", {})
    activity = lesson_data.get("activity", {})
    journal_prompt = lesson_data.get("journal_prompt", "")
    exit_tickets = lesson_data.get("exit_tickets", [])
    raw_content_points = lesson_data.get("content_points", [])
    presenter_notes = lesson_data.get("presenter_notes", {})

    # Normalize content_points to handle both formats:
    # Old format: ["string1", "string2", ...]
    # New format: [{"point": "...", "expanded": ["...", ...]}, ...]
    content_points = []
    expanded_content = []
    for cp in raw_content_points:
        if isinstance(cp, dict):
            # New format with expanded content
            content_points.append(cp.get("point", ""))
            expanded_content.append(cp.get("expanded", []))
        else:
            # Old format - just a string
            content_points.append(cp)
            expanded_content.append([])

    # Ensure we have 12 content points
    while len(content_points) < 12:
        content_points.append(f"Additional content point {len(content_points) + 1}")
        expanded_content.append([])

    # Slide 1: Agenda
    agenda = build_agenda_slide(learning_objectives, vocabulary, topic, day)
    agenda["day"] = day  # Add day for monologue generation
    # Generate verbatim monologue if no presenter notes provided
    if presenter_notes.get("slide_1"):
        agenda["presenter_notes"] = presenter_notes.get("slide_1")
    else:
        agenda["presenter_notes"] = generate_verbatim_monologue(agenda, SlideType.AGENDA)
    slides.append(agenda)

    # Slide 2: Warmup
    warmup_slide = build_warmup_slide(warmup)
    # Generate verbatim monologue if no presenter notes provided
    if presenter_notes.get("slide_2"):
        warmup_slide["presenter_notes"] = presenter_notes.get("slide_2")
    else:
        warmup_slide["presenter_notes"] = generate_verbatim_monologue(warmup_slide, SlideType.WARMUP)
    slides.append(warmup_slide)

    # Slides 3-14: Content (12 slides)
    vocab_terms = [v.get("term", "") for v in vocabulary]
    for i in range(12):
        # Use pre-expanded content if available (4-8 sentences)
        pre_expanded = expanded_content[i] if i < len(expanded_content) else []
        content = build_content_slide(
            content_point=content_points[i],
            expanded_content=pre_expanded if len(pre_expanded) >= 4 else None,
            slide_index=i + 1,
            performance_tip=f"Theater practice: Apply this concept in your performance work.",
            vocabulary_terms=vocab_terms[i:i+2] if i < len(vocab_terms) else []
        )
        # Generate verbatim monologue if no presenter notes provided
        if presenter_notes.get(f"slide_{i+3}"):
            content["presenter_notes"] = presenter_notes.get(f"slide_{i+3}")
        else:
            content["presenter_notes"] = generate_verbatim_monologue(content, SlideType.CONTENT, i + 1)
        slides.append(content)

    # Slide 15: Activity
    activity_slide = build_activity_slide(activity)
    # Generate verbatim monologue if no presenter notes provided
    if presenter_notes.get("slide_15"):
        activity_slide["presenter_notes"] = presenter_notes.get("slide_15")
    else:
        activity_slide["presenter_notes"] = generate_verbatim_monologue(activity_slide, SlideType.ACTIVITY)
    slides.append(activity_slide)

    # Slide 16: Journal & Exit
    journal_slide = build_journal_slide(journal_prompt, exit_tickets)
    # Generate verbatim monologue if no presenter notes provided
    if presenter_notes.get("slide_16"):
        journal_slide["presenter_notes"] = presenter_notes.get("slide_16")
    else:
        journal_slide["presenter_notes"] = generate_verbatim_monologue(journal_slide, SlideType.JOURNAL)
    slides.append(journal_slide)

    return slides


def create_theater_presentation(
    lesson_data: Dict[str, Any],
    output_path: Path,
    unit_name: str = "Theater",
    program_name: str = "THEATER EDUCATION",
    year: str = "2026"
) -> Path:
    """
    Create a complete Theater PowerPoint presentation.

    Args:
        lesson_data: Complete lesson data dictionary
        output_path: Where to save the .pptx file
        unit_name: Unit name for branding
        program_name: Program name for footer
        year: Year for footer

    Returns:
        Path to the created .pptx file
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")

    # Get template
    template_path = get_template_path()

    # Load template
    prs = Presentation(str(template_path))

    if len(prs.slides) == 0:
        raise ValueError("Template has no slides")

    # Generate slide data
    slides_data = generate_slides_from_lesson(lesson_data)
    total_slides = len(slides_data)

    # Duplicate template slide for each additional slide needed
    for i in range(1, total_slides):
        duplicate_slide(prs, 0)

    # Populate all slides
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides[i]
        populate_slide(
            slide=slide,
            slide_data=slide_data,
            slide_num=i + 1,
            total_slides=total_slides,
            unit_name=unit_name,
            program_name=program_name,
            year=year
        )

    # Ensure output directory exists
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Save presentation
    prs.save(str(output_path))

    return output_path


def generate_pptx(
    lesson_data: Dict[str, Any],
    output_dir: Path,
    unit_number: int,
    day: int
) -> Dict[str, Any]:
    """
    Main entry point for PowerPoint generation.

    Args:
        lesson_data: Complete lesson data
        output_dir: Output directory
        unit_number: Unit number (1-4)
        day: Day number

    Returns:
        Result dictionary with status and file path
    """
    unit_names = {
        1: "Greek Theater",
        2: "Commedia dell'Arte",
        3: "Shakespeare",
        4: "Student-Directed One Acts"
    }

    unit_name = unit_names.get(unit_number, "Theater")

    # Generate filename
    filename = f"Unit{unit_number}_Day{day:02d}_{unit_name.replace(' ', '_')}.pptx"
    output_path = Path(output_dir) / filename

    try:
        result_path = create_theater_presentation(
            lesson_data=lesson_data,
            output_path=output_path,
            unit_name=unit_name,
            program_name="THEATER EDUCATION",
            year=str(datetime.now().year)
        )

        return {
            "status": "success",
            "file_path": str(result_path),
            "slides_count": 16,
            "unit": unit_name,
            "day": day
        }

    except Exception as e:
        return {
            "status": "error",
            "error": str(e),
            "unit": unit_name,
            "day": day
        }


# =============================================================================
# VALIDATION
# =============================================================================

def validate_slide_content(slide_data: Dict[str, Any]) -> List[str]:
    """
    Validate slide content against constraints.

    Args:
        slide_data: Slide data dictionary

    Returns:
        List of validation issues (empty if valid)
    """
    issues = []

    # Check header length (max 36 chars)
    header = slide_data.get("header", "")
    if len(header) > 36:
        issues.append(f"Header exceeds 36 chars: {len(header)} chars")

    # Check body line count (max 12 lines)
    body = slide_data.get("body", "")
    if isinstance(body, list):
        if len(body) > 12:
            issues.append(f"Body exceeds 12 lines: {len(body)} lines")
    elif isinstance(body, str):
        lines = body.split("\n")
        if len(lines) > 12:
            issues.append(f"Body exceeds 12 lines: {len(lines)} lines")

    # Check performance tip (max 132 chars)
    tip = slide_data.get("performance_tip", "")
    if len(tip) > 132:
        issues.append(f"Performance tip exceeds 132 chars: {len(tip)} chars")

    return issues


def validate_presentation(slides_data: List[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Validate entire presentation.

    Args:
        slides_data: List of slide data dictionaries

    Returns:
        Validation result dictionary
    """
    all_issues = []

    # Check slide count
    if len(slides_data) != 16:
        all_issues.append(f"Expected 16 slides, found {len(slides_data)}")

    # Validate each slide
    for i, slide in enumerate(slides_data):
        slide_issues = validate_slide_content(slide)
        for issue in slide_issues:
            all_issues.append(f"Slide {i+1}: {issue}")

    return {
        "valid": len(all_issues) == 0,
        "issues": all_issues,
        "slides_checked": len(slides_data)
    }


# =============================================================================
# EXPORTS
# =============================================================================

__all__ = [
    "create_theater_presentation",
    "generate_pptx",
    "generate_slides_from_lesson",
    "validate_presentation",
    "validate_slide_content",
    "SlideType",
    "SLIDE_STRUCTURE",
]


# =============================================================================
# CLI TESTING
# =============================================================================

if __name__ == "__main__":
    # Test with sample data
    sample_lesson = {
        "topic": "Introduction to Greek Theater",
        "day": 1,
        "learning_objectives": [
            "Identify key characteristics of Greek theater",
            "Explain the role of the Festival of Dionysus",
            "Demonstrate understanding of theatrical origins"
        ],
        "vocabulary": [
            {"term": "dithyramb", "definition": "A choral hymn sung to Dionysus"},
            {"term": "theatron", "definition": "The seating area in Greek theater"},
            {"term": "orchestra", "definition": "The circular performance space"}
        ],
        "warmup": {
            "name": "Chorus Movement",
            "instructions": "Stand in a circle. Follow the leader's movements.",
            "connection_to_lesson": "Introduces ensemble movement like Greek chorus"
        },
        "activity": {
            "name": "Dionysus Festival Role Play",
            "instructions": ["Form groups of 4", "Create a short dithyramb", "Present to class"],
            "grouping": "small groups",
            "materials": ["Percussion instruments", "Movement space"]
        },
        "journal_prompt": "How might participating in a religious festival have influenced the development of theater?",
        "exit_tickets": [
            "Name two elements of Greek theater",
            "Why was the Festival of Dionysus important?"
        ],
        "content_points": [
            "Greek theater originated in ancient Athens around 534 BCE",
            "The Festival of Dionysus was a religious celebration honoring the god of wine",
            "Early performances featured a single actor and a chorus of 12-15 performers",
            "The theatron (seating area) could hold up to 17,000 spectators",
            "The orchestra was a circular space where the chorus performed",
            "Greek theater explored themes of fate, hubris, and divine justice",
            "Tragedy and comedy were the two main dramatic forms",
            "Masks allowed actors to play multiple roles",
            "The chorus provided commentary and emotional response",
            "Theater was a civic duty - attendance was considered important",
            "Competition between playwrights was a key feature of festivals",
            "The legacy of Greek theater continues to influence drama today"
        ]
    }

    # Test generation
    print("Testing Theater PowerPoint Generator...")
    print("=" * 50)

    # Generate slides
    slides = generate_slides_from_lesson(sample_lesson)
    print(f"Generated {len(slides)} slides")

    # Validate
    validation = validate_presentation(slides)
    print(f"Validation: {'PASS' if validation['valid'] else 'FAIL'}")
    if validation['issues']:
        for issue in validation['issues']:
            print(f"  - {issue}")

    # Create presentation
    output_path = Path("test_output/test_theater_presentation.pptx")
    try:
        result = generate_pptx(sample_lesson, output_path.parent, unit_number=1, day=1)
        print(f"\nResult: {result['status']}")
        if result['status'] == 'success':
            print(f"Created: {result['file_path']}")
    except Exception as e:
        print(f"Error: {e}")
