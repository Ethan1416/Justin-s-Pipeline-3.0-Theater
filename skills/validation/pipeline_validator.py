"""
Pipeline Validator - Comprehensive Validation Suite for NCLEX Lecture Pipeline
==============================================================================

This script provides validation functions for all pipeline steps:
1. Template Integrity Check
2. Character Limit Validation
3. Pre-Flight Validation (input file checks)
4. Presenter Notes Word Counter
5. Visual-Diagram Alignment Check
6. Slide Count Constraint Validation
7. Master Validation Report

Requirements:
- python-pptx (pip install python-pptx)

Usage:
- Run directly: python pipeline_validator.py
- Or import functions: from pipeline_validator import validate_template
"""

import os
import re
import json
import sys
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Tuple, Optional

# Ensure UTF-8 output
if sys.stdout:
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except:
        pass

# ============================================
# CONFIGURATION LOADER
# ============================================

def load_config(config_path: str = None) -> dict:
    """Load pipeline configuration from JSON file."""
    if config_path is None:
        config_path = Path(__file__).parent / "pipeline_config.json"

    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except FileNotFoundError:
        print(f"WARNING: Config file not found at {config_path}")
        return get_default_config()
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid JSON in config file: {e}")
        return get_default_config()


def get_default_config() -> dict:
    """Return default configuration if file not found."""
    return {
        "template_requirements": {
            "required_shapes": ["TextBox 2", "TextBox 19", "TextBox 18"]
        },
        "constraints": {
            "character_limits": {
                "header_chars_per_line": 32,
                "header_max_lines": 2,
                "body_chars_per_line": 66,
                "body_max_lines": 8,
                "tip_chars_per_line": 66,
                "tip_max_lines": 2
            },
            "presenter_notes": {
                "max_words": 450
            },
            "slides": {
                "section_min": 12,
                "section_max": 35,
                "session_min": 72,
                "session_max": 90
            }
        }
    }


# ============================================
# 1. TEMPLATE INTEGRITY CHECK
# ============================================

def validate_template(template_path: str, config: dict = None) -> Tuple[bool, List[str]]:
    """
    Verify PowerPoint template has required shapes and structure.

    Returns:
        Tuple of (is_valid, list_of_issues)
    """
    issues = []

    try:
        from pptx import Presentation
    except ImportError:
        return False, ["ERROR: python-pptx not installed. Run: pip install python-pptx"]

    if config is None:
        config = load_config()

    required_shapes = config.get("template_requirements", {}).get(
        "required_shapes", ["TextBox 2", "TextBox 19", "TextBox 18"]
    )

    # Check file exists
    if not Path(template_path).exists():
        return False, [f"Template file not found: {template_path}"]

    try:
        prs = Presentation(template_path)
    except Exception as e:
        return False, [f"Cannot open template: {e}"]

    # Check we have at least one slide
    if len(prs.slides) == 0:
        issues.append("Template has no slides")
        return False, issues

    # Check first slide for required shapes
    first_slide = prs.slides[0]
    found_shapes = set()

    for shape in first_slide.shapes:
        if shape.name in required_shapes:
            found_shapes.add(shape.name)

    missing_shapes = set(required_shapes) - found_shapes

    if missing_shapes:
        for shape in missing_shapes:
            issues.append(f"Missing required shape: {shape}")

    # Check slide dimensions
    expected_width = config.get("template_requirements", {}).get(
        "slide_dimensions", {}
    ).get("width_inches", 13.33)

    actual_width = prs.slide_width.inches
    if abs(actual_width - expected_width) > 0.1:
        issues.append(f"Slide width mismatch: expected {expected_width}\", got {actual_width:.2f}\"")

    is_valid = len(issues) == 0

    if is_valid:
        print(f"[PASS] Template validation: All {len(required_shapes)} required shapes found")
    else:
        print(f"[FAIL] Template validation: {len(issues)} issue(s) found")

    return is_valid, issues


# ============================================
# 2. CHARACTER LIMIT VALIDATION
# ============================================

def count_line_characters(text: str) -> List[int]:
    """Count characters per line in text."""
    if not text:
        return [0]
    lines = text.strip().split('\n')
    return [len(line) for line in lines]


def validate_character_limits(content: dict, config: dict = None) -> Tuple[bool, List[str]]:
    """
    Validate header, body, and tip text against character limits.

    Args:
        content: Dict with keys 'header', 'body', 'tip', 'slide_number'
        config: Pipeline configuration

    Returns:
        Tuple of (is_valid, list_of_issues)
    """
    if config is None:
        config = load_config()

    limits = config.get("constraints", {}).get("character_limits", {})
    issues = []
    slide_id = content.get('slide_number', 'Unknown')

    # Check header
    header = content.get('header', '')
    if header:
        header_chars = count_line_characters(header)
        header_lines = len(header_chars)
        max_header_chars = limits.get('header_chars_per_line', 32)
        max_header_lines = limits.get('header_max_lines', 2)

        if header_lines > max_header_lines:
            issues.append(f"Slide {slide_id}: Header has {header_lines} lines (max {max_header_lines})")

        for i, char_count in enumerate(header_chars):
            if char_count > max_header_chars:
                issues.append(f"Slide {slide_id}: Header line {i+1} has {char_count} chars (max {max_header_chars})")

    # Check body
    body = content.get('body', '')
    if body:
        body_chars = count_line_characters(body)
        body_lines = len(body_chars)
        max_body_chars = limits.get('body_chars_per_line', 66)
        max_body_lines = limits.get('body_max_lines', 8)

        if body_lines > max_body_lines:
            issues.append(f"Slide {slide_id}: Body has {body_lines} lines (max {max_body_lines})")

        for i, char_count in enumerate(body_chars):
            if char_count > max_body_chars:
                issues.append(f"Slide {slide_id}: Body line {i+1} has {char_count} chars (max {max_body_chars})")

    # Check tip
    tip = content.get('tip', '')
    if tip and tip.lower() not in ['none', 'n/a', '']:
        tip_chars = count_line_characters(tip)
        tip_lines = len(tip_chars)
        max_tip_chars = limits.get('tip_chars_per_line', 66)
        max_tip_lines = limits.get('tip_max_lines', 2)

        if tip_lines > max_tip_lines:
            issues.append(f"Slide {slide_id}: NCLEX Tip has {tip_lines} lines (max {max_tip_lines})")

        for i, char_count in enumerate(tip_chars):
            if char_count > max_tip_chars:
                issues.append(f"Slide {slide_id}: NCLEX Tip line {i+1} has {char_count} chars (max {max_tip_chars})")

    return len(issues) == 0, issues


def validate_blueprint_character_limits(blueprint_path: str, config: dict = None) -> Tuple[bool, List[str]]:
    """
    Validate all slides in a blueprint file for character limits.
    """
    if config is None:
        config = load_config()

    if not Path(blueprint_path).exists():
        return False, [f"Blueprint file not found: {blueprint_path}"]

    with open(blueprint_path, 'r', encoding='utf-8') as f:
        content = f.read()

    all_issues = []
    slides = parse_slides_from_blueprint(content)

    for slide in slides:
        is_valid, issues = validate_character_limits(slide, config)
        all_issues.extend(issues)

    if all_issues:
        print(f"[FAIL] Character limit validation: {len(all_issues)} violation(s)")
    else:
        print(f"[PASS] Character limit validation: All {len(slides)} slides within limits")

    return len(all_issues) == 0, all_issues


# ============================================
# 3. PRE-FLIGHT VALIDATION
# ============================================

def validate_preflight(step_number: int, pipeline_folder: str, config: dict = None) -> Tuple[bool, List[str]]:
    """
    Validate that all required input files exist for a given step.

    Args:
        step_number: The step to validate inputs for (6-12)
        pipeline_folder: Path to pipeline folder
        config: Pipeline configuration

    Returns:
        Tuple of (all_inputs_valid, list_of_issues)
    """
    if config is None:
        config = load_config()

    issues = []
    pipeline_path = Path(pipeline_folder)

    # Define required inputs per step
    required_inputs = {
        6: ["step4_output_*.txt", "step5_presentation_standards.txt"],
        7: ["step6_blueprint_*.txt"],
        8: ["step7_revised_*.txt", "step4_output_*.txt"],
        9: ["step7_revised_*.txt", "step8_validation_score_*.txt"],
        10: ["step9_visual_specs_*.txt", "step7_revised_*.txt"],
        11: ["step10_integrated_blueprint_*.txt"],
        12: ["visual template.pptx", "step10_integrated_blueprint_*.txt"]
    }

    if step_number not in required_inputs:
        return True, []

    for pattern in required_inputs[step_number]:
        matches = list(pipeline_path.glob(pattern))
        if not matches:
            issues.append(f"Step {step_number}: Missing required input matching '{pattern}'")

    if issues:
        print(f"[FAIL] Pre-flight validation for Step {step_number}: {len(issues)} missing input(s)")
    else:
        print(f"[PASS] Pre-flight validation for Step {step_number}: All inputs present")

    return len(issues) == 0, issues


# ============================================
# 4. PRESENTER NOTES WORD COUNTER
# ============================================

def count_words(text: str) -> int:
    """Count words in text."""
    if not text:
        return 0
    # Remove markers like [PAUSE], [EMPHASIS: ...]
    cleaned = re.sub(r'\[PAUSE\]|\[EMPHASIS:[^\]]*\]', '', text)
    words = cleaned.split()
    return len(words)


def validate_presenter_notes(notes: str, slide_id: str = "Unknown", config: dict = None) -> Tuple[bool, List[str]]:
    """
    Validate presenter notes word count.

    Returns:
        Tuple of (is_valid, list_of_issues)
    """
    if config is None:
        config = load_config()

    max_words = config.get("constraints", {}).get("presenter_notes", {}).get("max_words", 450)
    issues = []

    word_count = count_words(notes)

    if word_count > max_words:
        issues.append(f"Slide {slide_id}: Presenter notes have {word_count} words (max {max_words})")
        estimated_seconds = int(word_count / 2.5)  # ~150 wpm
        issues.append(f"  Estimated duration: ~{estimated_seconds} seconds (max 180)")

    return len(issues) == 0, issues


def validate_blueprint_presenter_notes(blueprint_path: str, config: dict = None) -> Tuple[bool, List[str]]:
    """
    Validate all presenter notes in a blueprint file.
    """
    if config is None:
        config = load_config()

    if not Path(blueprint_path).exists():
        return False, [f"Blueprint file not found: {blueprint_path}"]

    with open(blueprint_path, 'r', encoding='utf-8') as f:
        content = f.read()

    all_issues = []
    slides = parse_slides_from_blueprint(content)

    for slide in slides:
        notes = slide.get('notes', '')
        if notes:
            is_valid, issues = validate_presenter_notes(notes, slide.get('slide_number', 'Unknown'), config)
            all_issues.extend(issues)

    if all_issues:
        print(f"[FAIL] Presenter notes validation: {len(all_issues)} issue(s)")
    else:
        print(f"[PASS] Presenter notes validation: All {len(slides)} slides within word limits")

    return len(all_issues) == 0, all_issues


# ============================================
# 5. VISUAL-DIAGRAM ALIGNMENT CHECK
# ============================================

def validate_visual_alignment(blueprint_path: str, diagrams_folder: str) -> Tuple[bool, List[str]]:
    """
    Check that visual specifications in blueprint have corresponding diagram files.

    Returns:
        Tuple of (all_aligned, list_of_issues)
    """
    if not Path(blueprint_path).exists():
        return False, [f"Blueprint file not found: {blueprint_path}"]

    diagrams_path = Path(diagrams_folder)
    if not diagrams_path.exists():
        # Not necessarily an error - may have no Python-generated diagrams
        return True, ["Note: Diagrams folder does not exist (may be OK if only tables)"]

    with open(blueprint_path, 'r', encoding='utf-8') as f:
        content = f.read()

    issues = []

    # Find all visual slides that need Python-generated diagrams
    visual_pattern = r'SLIDE\s+(\d+[AB]?):[^\n]+\n.*?Visual:\s*Yes\s*-\s*(Decision Tree|Flowchart|Hierarchy|Concept Map)'
    matches = re.findall(visual_pattern, content, re.DOTALL | re.IGNORECASE)

    for slide_num, visual_type in matches:
        # Look for corresponding diagram file
        diagram_patterns = [
            f"*slide_{slide_num}*.png",
            f"*slide{slide_num}*.png",
            f"*{slide_num}*.png"
        ]

        found = False
        for pattern in diagram_patterns:
            if list(diagrams_path.glob(pattern)):
                found = True
                break

        if not found:
            issues.append(f"Slide {slide_num}: Missing diagram image for {visual_type}")

    if issues:
        print(f"[FAIL] Visual alignment: {len(issues)} missing diagram(s)")
    else:
        print(f"[PASS] Visual alignment: All visual specifications have corresponding diagrams")

    return len(issues) == 0, issues


# ============================================
# 6. SLIDE COUNT CONSTRAINT VALIDATION
# ============================================

def validate_slide_counts(blueprint_path: str, config: dict = None) -> Tuple[bool, List[str]]:
    """
    Validate slide counts against section/session constraints.
    """
    if config is None:
        config = load_config()

    if not Path(blueprint_path).exists():
        return False, [f"Blueprint file not found: {blueprint_path}"]

    with open(blueprint_path, 'r', encoding='utf-8') as f:
        content = f.read()

    constraints = config.get("constraints", {}).get("slides", {})
    section_min = constraints.get("section_min", 12)
    section_max = constraints.get("section_max", 35)

    issues = []

    # Count slides
    slide_matches = re.findall(r'SLIDE\s+(\d+[AB]?):', content)
    slide_count = len(slide_matches)

    # Extract section name
    section_match = re.search(r'Section:\s*(.+)', content)
    section_name = section_match.group(1).strip() if section_match else "Unknown"

    if slide_count < section_min:
        issues.append(f"Section '{section_name}': {slide_count} slides (below minimum {section_min})")
    elif slide_count > section_max:
        issues.append(f"Section '{section_name}': {slide_count} slides (exceeds maximum {section_max})")

    if issues:
        print(f"[FAIL] Slide count validation: {len(issues)} constraint violation(s)")
    else:
        print(f"[PASS] Slide count validation: {slide_count} slides within {section_min}-{section_max} range")

    return len(issues) == 0, issues


# ============================================
# HELPER: PARSE SLIDES FROM BLUEPRINT
# ============================================

def parse_slides_from_blueprint(content: str) -> List[dict]:
    """Parse blueprint content into list of slide dictionaries."""
    slides = []

    # Split by slide markers
    slide_pattern = r'={40,}\s*SLIDE\s+(\d+[AB]?):\s*(.+?)\s*={40,}'
    parts = re.split(slide_pattern, content)

    i = 1
    while i < len(parts) - 2:
        slide_num = parts[i].strip()
        slide_title = parts[i + 1].strip()
        slide_content = parts[i + 2] if i + 2 < len(parts) else ""

        slide = {
            'slide_number': slide_num,
            'title': slide_title,
            'header': '',
            'body': '',
            'tip': '',
            'notes': ''
        }

        # Extract header
        header_match = re.search(r'HEADER:\s*\n(.*?)(?=\n\s*\n|\nBODY:)', slide_content, re.DOTALL)
        if header_match:
            slide['header'] = header_match.group(1).strip()

        # Extract body
        body_match = re.search(r'BODY:\s*\n(.*?)(?=\nVISUAL|\nNCLEX TIP:|\nPRESENTER NOTES:)', slide_content, re.DOTALL)
        if body_match:
            slide['body'] = body_match.group(1).strip()

        # Extract NCLEX tip
        tip_match = re.search(r'NCLEX TIP:\s*\n(.*?)(?=\nPRESENTER NOTES:)', slide_content, re.DOTALL)
        if tip_match:
            slide['tip'] = tip_match.group(1).strip()

        # Extract presenter notes
        notes_match = re.search(r'PRESENTER NOTES:\s*\n(.*?)(?=\n={40,}|$)', slide_content, re.DOTALL)
        if notes_match:
            slide['notes'] = notes_match.group(1).strip()

        slides.append(slide)
        i += 3

    return slides


# ============================================
# 7. MASTER VALIDATION REPORT
# ============================================

def run_master_validation(pipeline_folder: str, step: int = None) -> dict:
    """
    Run comprehensive validation and generate report.

    Args:
        pipeline_folder: Path to pipeline folder
        step: Specific step to validate (None = validate all)

    Returns:
        Dictionary with validation results
    """
    config = load_config(Path(pipeline_folder) / "pipeline_config.json")
    pipeline_path = Path(pipeline_folder)

    results = {
        'timestamp': datetime.now().isoformat(),
        'pipeline_folder': str(pipeline_folder),
        'validations': {},
        'overall_status': 'PASS',
        'total_issues': 0
    }

    print("=" * 60)
    print("PIPELINE MASTER VALIDATION REPORT")
    print("=" * 60)
    print(f"Timestamp: {results['timestamp']}")
    print(f"Folder: {pipeline_folder}")
    print()

    # 1. Template validation
    print("-" * 40)
    print("1. TEMPLATE INTEGRITY")
    print("-" * 40)
    template_path = pipeline_path / config.get("paths", {}).get("template_file", "visual template.pptx")
    is_valid, issues = validate_template(str(template_path), config)
    results['validations']['template'] = {'valid': is_valid, 'issues': issues}
    if not is_valid:
        results['overall_status'] = 'FAIL'
    results['total_issues'] += len(issues)
    print()

    # 2. Pre-flight validation for specified step
    if step:
        print("-" * 40)
        print(f"2. PRE-FLIGHT (Step {step})")
        print("-" * 40)
        is_valid, issues = validate_preflight(step, str(pipeline_path), config)
        results['validations']['preflight'] = {'valid': is_valid, 'issues': issues}
        if not is_valid:
            results['overall_status'] = 'FAIL'
        results['total_issues'] += len(issues)
        print()

    # 3. Blueprint validations (find all integrated blueprints)
    blueprints = list(pipeline_path.glob("step10_integrated_blueprint_*.txt"))
    if not blueprints:
        blueprints = list(pipeline_path.glob("step7_revised_*.txt"))

    if blueprints:
        print("-" * 40)
        print("3. BLUEPRINT VALIDATIONS")
        print("-" * 40)

        for bp in blueprints:
            print(f"\nValidating: {bp.name}")

            # Character limits
            is_valid, issues = validate_blueprint_character_limits(str(bp), config)
            results['validations'][f'char_limits_{bp.stem}'] = {'valid': is_valid, 'issues': issues}
            if not is_valid:
                results['overall_status'] = 'FAIL'
            results['total_issues'] += len(issues)

            # Presenter notes
            is_valid, issues = validate_blueprint_presenter_notes(str(bp), config)
            results['validations'][f'notes_{bp.stem}'] = {'valid': is_valid, 'issues': issues}
            if not is_valid:
                results['overall_status'] = 'FAIL'
            results['total_issues'] += len(issues)

            # Slide counts
            is_valid, issues = validate_slide_counts(str(bp), config)
            results['validations'][f'slides_{bp.stem}'] = {'valid': is_valid, 'issues': issues}
            if not is_valid:
                results['overall_status'] = 'FAIL'
            results['total_issues'] += len(issues)

    # Summary
    print()
    print("=" * 60)
    print("VALIDATION SUMMARY")
    print("=" * 60)
    print(f"Overall Status: {results['overall_status']}")
    print(f"Total Issues: {results['total_issues']}")

    if results['total_issues'] > 0:
        print("\nIssues by Category:")
        for key, val in results['validations'].items():
            if val['issues']:
                print(f"  {key}: {len(val['issues'])} issue(s)")
                for issue in val['issues'][:3]:  # Show first 3
                    print(f"    - {issue}")
                if len(val['issues']) > 3:
                    print(f"    ... and {len(val['issues']) - 3} more")

    print("=" * 60)

    return results


# ============================================
# MAIN EXECUTION
# ============================================

def main():
    """Main entry point for validation."""
    # Get pipeline folder from script location
    pipeline_folder = Path(__file__).parent

    print("\nNCLEX Lecture Pipeline Validator")
    print("================================\n")

    # Run master validation
    results = run_master_validation(str(pipeline_folder))

    # Save results to file
    report_path = pipeline_folder / f"validation_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
    with open(report_path, 'w', encoding='utf-8') as f:
        json.dump(results, f, indent=2)

    print(f"\nReport saved to: {report_path}")

    return 0 if results['overall_status'] == 'PASS' else 1


if __name__ == "__main__":
    sys.exit(main())
