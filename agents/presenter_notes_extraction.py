"""
Presenter Notes Extraction Agents (HARDCODED)
==============================================

Specialized agents for extracting presenter notes/monologues from
PowerPoint slides and generating Word documents.

These agents handle:
- Extracting verbatim presenter notes from PPTX files
- Formatting notes for Word document output
- Generating professional Word documents
- Validating extraction completeness
"""

from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple
from pathlib import Path
from dataclasses import dataclass, field
import re

from .base import Agent, AgentResult, AgentStatus

# Check for required libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    from docx.shared import Inches as DocxInches, Pt as DocxPt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False


# =============================================================================
# HARDCODED CONFIGURATION
# =============================================================================

# HARDCODED: Word document styling configuration
WORD_DOCUMENT_CONFIG = {
    "title_font_size": 24,
    "heading_font_size": 14,
    "body_font_size": 12,
    "slide_label_font_size": 11,
    "font_name": "Calibri",
    "title_font_name": "Calibri Light",
    "line_spacing": 1.15,
    "paragraph_spacing_after": 12,
    "slide_separator": True,
    "include_slide_titles": True,
    "include_slide_numbers": True,
    "page_margins": {
        "top": 1.0,
        "bottom": 1.0,
        "left": 1.0,
        "right": 1.0,
    }
}

# HARDCODED: Notes extraction configuration
EXTRACTION_CONFIG = {
    "min_note_length": 10,           # Minimum characters to consider valid note
    "max_note_length": 10000,        # Maximum characters per slide
    "strip_whitespace": True,
    "normalize_line_breaks": True,
    "remove_empty_slides": False,    # Keep empty slides as placeholders
    "include_slide_title": True,
    "placeholder_for_empty": "[No presenter notes for this slide]",
}

# HARDCODED: Validation thresholds
VALIDATION_CONFIG = {
    "min_slides_with_notes": 0.5,    # At least 50% of slides should have notes
    "min_total_word_count": 100,     # Minimum total words across all notes
    "warn_short_note_threshold": 20, # Warn if note is shorter than this
}


# =============================================================================
# DATA CLASSES
# =============================================================================

@dataclass
class SlideNotes:
    """Container for extracted slide notes."""
    slide_number: int
    slide_title: str
    notes_text: str
    word_count: int
    has_content: bool
    extraction_status: str  # "success", "empty", "error"


@dataclass
class ExtractionResult:
    """Result of notes extraction from a presentation."""
    source_file: str
    total_slides: int
    slides_with_notes: int
    slides_empty: int
    total_word_count: int
    slide_notes: List[SlideNotes]
    extraction_time: datetime
    errors: List[str]


# =============================================================================
# PresenterNotesExtractorAgent
# =============================================================================

class PresenterNotesExtractorAgent(Agent):
    """
    HARDCODED agent for extracting presenter notes from PowerPoint files.

    Extracts verbatim monologue/presenter notes from every slide and
    provides structured output for further processing.
    """

    def __init__(self):
        super().__init__(name="PresenterNotesExtractorAgent")
        self.config = EXTRACTION_CONFIG

    def _process(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Extract presenter notes from a PowerPoint file."""
        pptx_path = context.get("pptx_path")

        if not pptx_path:
            return {
                "success": False,
                "error": "No PowerPoint path provided"
            }

        if not PPTX_AVAILABLE:
            return {
                "success": False,
                "error": "python-pptx library not available"
            }

        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            return {
                "success": False,
                "error": f"File not found: {pptx_path}"
            }

        try:
            return self._extract_notes(pptx_path)
        except Exception as e:
            return {
                "success": False,
                "error": f"Extraction failed: {str(e)}"
            }

    def _extract_notes(self, pptx_path: Path) -> Dict[str, Any]:
        """Perform the actual notes extraction."""
        prs = Presentation(str(pptx_path))

        slide_notes_list = []
        total_word_count = 0
        slides_with_notes = 0
        errors = []

        for i, slide in enumerate(prs.slides, 1):
            # Get slide title
            slide_title = self._get_slide_title(slide, i)

            # Get presenter notes
            notes_text = ""
            try:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    notes_text = notes_slide.notes_text_frame.text
            except Exception as e:
                errors.append(f"Slide {i}: Error reading notes - {str(e)}")

            # Process notes text
            if self.config["strip_whitespace"]:
                notes_text = notes_text.strip()

            if self.config["normalize_line_breaks"]:
                notes_text = self._normalize_line_breaks(notes_text)

            # Calculate word count
            word_count = len(notes_text.split()) if notes_text else 0
            total_word_count += word_count

            # Determine if slide has meaningful content
            has_content = len(notes_text) >= self.config["min_note_length"]
            if has_content:
                slides_with_notes += 1

            # Handle empty notes
            if not has_content and not self.config["remove_empty_slides"]:
                display_text = self.config["placeholder_for_empty"]
            else:
                display_text = notes_text

            slide_notes = SlideNotes(
                slide_number=i,
                slide_title=slide_title,
                notes_text=display_text,
                word_count=word_count,
                has_content=has_content,
                extraction_status="success" if has_content else "empty"
            )
            slide_notes_list.append(slide_notes)

        # Build result
        result = ExtractionResult(
            source_file=str(pptx_path),
            total_slides=len(prs.slides),
            slides_with_notes=slides_with_notes,
            slides_empty=len(prs.slides) - slides_with_notes,
            total_word_count=total_word_count,
            slide_notes=slide_notes_list,
            extraction_time=datetime.now(),
            errors=errors
        )

        return {
            "success": True,
            "source_file": result.source_file,
            "total_slides": result.total_slides,
            "slides_with_notes": result.slides_with_notes,
            "slides_empty": result.slides_empty,
            "total_word_count": result.total_word_count,
            "slide_notes": [
                {
                    "slide_number": sn.slide_number,
                    "slide_title": sn.slide_title,
                    "notes_text": sn.notes_text,
                    "word_count": sn.word_count,
                    "has_content": sn.has_content,
                    "extraction_status": sn.extraction_status
                }
                for sn in result.slide_notes
            ],
            "extraction_time": result.extraction_time.isoformat(),
            "errors": result.errors
        }

    def _get_slide_title(self, slide, slide_number: int) -> str:
        """Extract slide title or generate default."""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                return slide.shapes.title.text.strip()
        except:
            pass

        # Try to find any text that looks like a title
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if len(text) < 100 and text:  # Short text likely a title
                    return text[:50]

        return f"Slide {slide_number}"

    def _normalize_line_breaks(self, text: str) -> str:
        """Normalize line breaks in text."""
        # Replace multiple newlines with double newline
        text = re.sub(r'\n{3,}', '\n\n', text)
        # Replace Windows line endings
        text = text.replace('\r\n', '\n')
        return text

    def extract_from_directory(self, directory: Path) -> List[Dict[str, Any]]:
        """Extract notes from all PowerPoint files in a directory."""
        results = []
        pptx_files = list(directory.glob("**/*.pptx"))

        for pptx_file in sorted(pptx_files):
            result = self.execute({"pptx_path": str(pptx_file)})
            results.append(result.output)

        return results


# =============================================================================
# MonologueFormatterAgent
# =============================================================================

class MonologueFormatterAgent(Agent):
    """
    HARDCODED agent for formatting extracted notes for Word output.

    Transforms raw extracted notes into a structured format suitable
    for professional Word document generation.
    """

    def __init__(self):
        super().__init__(name="MonologueFormatterAgent")
        self.config = WORD_DOCUMENT_CONFIG

    def _process(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Format extracted notes for Word document generation."""
        extraction_result = context.get("extraction_result")
        presentation_title = context.get("presentation_title", "Presenter Notes")
        include_metadata = context.get("include_metadata", True)

        if not extraction_result:
            return {
                "success": False,
                "error": "No extraction result provided"
            }

        slide_notes = extraction_result.get("slide_notes", [])

        # Format document structure
        formatted = {
            "title": presentation_title,
            "subtitle": f"Presenter Notes Monologue",
            "metadata": {},
            "sections": []
        }

        # Add metadata if requested
        if include_metadata:
            formatted["metadata"] = {
                "source_file": extraction_result.get("source_file", "Unknown"),
                "total_slides": extraction_result.get("total_slides", 0),
                "slides_with_notes": extraction_result.get("slides_with_notes", 0),
                "total_word_count": extraction_result.get("total_word_count", 0),
                "extraction_date": extraction_result.get("extraction_time", ""),
            }

        # Format each slide's notes
        for slide in slide_notes:
            section = {
                "slide_number": slide["slide_number"],
                "slide_title": slide["slide_title"],
                "has_content": slide["has_content"],
                "word_count": slide["word_count"],
                "formatted_header": self._format_slide_header(slide),
                "formatted_body": self._format_slide_body(slide),
            }
            formatted["sections"].append(section)

        # Calculate reading time estimate
        total_words = extraction_result.get("total_word_count", 0)
        reading_time_minutes = round(total_words / 150, 1)  # ~150 WPM speaking

        formatted["summary"] = {
            "total_sections": len(formatted["sections"]),
            "sections_with_content": sum(1 for s in formatted["sections"] if s["has_content"]),
            "total_word_count": total_words,
            "estimated_speaking_time_minutes": reading_time_minutes,
        }

        return {
            "success": True,
            "formatted_document": formatted,
            "config_used": self.config
        }

    def _format_slide_header(self, slide: Dict) -> str:
        """Format the header for a slide section."""
        if self.config["include_slide_numbers"]:
            header = f"Slide {slide['slide_number']}"
            if self.config["include_slide_titles"] and slide["slide_title"]:
                header += f": {slide['slide_title']}"
        else:
            header = slide["slide_title"] or f"Slide {slide['slide_number']}"

        return header

    def _format_slide_body(self, slide: Dict) -> str:
        """Format the body text for a slide section."""
        text = slide["notes_text"]

        # Clean up formatting
        text = text.strip()

        # Add word count annotation if configured
        if slide["has_content"] and slide["word_count"] > 0:
            # Could add annotation here if needed
            pass

        return text


# =============================================================================
# NotesToWordGeneratorAgent
# =============================================================================

class NotesToWordGeneratorAgent(Agent):
    """
    HARDCODED agent for generating Word documents from formatted notes.

    Creates professional Word documents with proper styling, headers,
    and formatting for presenter monologues.
    """

    def __init__(self):
        super().__init__(name="NotesToWordGeneratorAgent")
        self.config = WORD_DOCUMENT_CONFIG

    def _process(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Generate a Word document from formatted notes."""
        formatted_document = context.get("formatted_document")
        output_path = context.get("output_path")

        if not formatted_document:
            return {
                "success": False,
                "error": "No formatted document provided"
            }

        if not output_path:
            return {
                "success": False,
                "error": "No output path provided"
            }

        if not DOCX_AVAILABLE:
            return {
                "success": False,
                "error": "python-docx library not available"
            }

        try:
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)

            doc = self._create_document(formatted_document)
            doc.save(str(output_path))

            return {
                "success": True,
                "output_path": str(output_path),
                "document_title": formatted_document.get("title", ""),
                "sections_written": len(formatted_document.get("sections", [])),
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Document generation failed: {str(e)}"
            }

    def _create_document(self, formatted_document: Dict) -> Document:
        """Create the Word document with proper formatting."""
        doc = Document()

        # Set page margins
        for section in doc.sections:
            section.top_margin = DocxInches(self.config["page_margins"]["top"])
            section.bottom_margin = DocxInches(self.config["page_margins"]["bottom"])
            section.left_margin = DocxInches(self.config["page_margins"]["left"])
            section.right_margin = DocxInches(self.config["page_margins"]["right"])

        # Add title
        title = doc.add_heading(formatted_document.get("title", "Presenter Notes"), 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add subtitle
        subtitle = doc.add_paragraph(formatted_document.get("subtitle", ""))
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add metadata section if present
        metadata = formatted_document.get("metadata", {})
        if metadata:
            doc.add_paragraph()  # Spacer
            meta_para = doc.add_paragraph()
            meta_para.add_run("Document Information").bold = True
            doc.add_paragraph(f"Source: {metadata.get('source_file', 'N/A')}")
            doc.add_paragraph(f"Total Slides: {metadata.get('total_slides', 0)}")
            doc.add_paragraph(f"Slides with Notes: {metadata.get('slides_with_notes', 0)}")
            doc.add_paragraph(f"Total Words: {metadata.get('total_word_count', 0)}")
            doc.add_paragraph()  # Spacer

        # Add horizontal line
        doc.add_paragraph("─" * 50)

        # Add each slide section
        for section in formatted_document.get("sections", []):
            self._add_slide_section(doc, section)

        # Add summary at end
        summary = formatted_document.get("summary", {})
        if summary:
            doc.add_paragraph()
            doc.add_paragraph("─" * 50)
            summary_para = doc.add_paragraph()
            summary_para.add_run("Summary").bold = True
            doc.add_paragraph(f"Total Sections: {summary.get('total_sections', 0)}")
            doc.add_paragraph(f"Sections with Content: {summary.get('sections_with_content', 0)}")
            doc.add_paragraph(f"Total Word Count: {summary.get('total_word_count', 0)}")
            doc.add_paragraph(f"Estimated Speaking Time: {summary.get('estimated_speaking_time_minutes', 0)} minutes")

        return doc

    def _add_slide_section(self, doc: Document, section: Dict):
        """Add a slide section to the document."""
        # Add slide header
        header = doc.add_heading(section["formatted_header"], level=2)

        # Add word count annotation
        if section["has_content"]:
            annotation = doc.add_paragraph()
            run = annotation.add_run(f"[{section['word_count']} words]")
            run.italic = True
            run.font.size = DocxPt(10)

        # Add notes body
        body_text = section["formatted_body"]
        if body_text:
            # Split by paragraphs and add each
            paragraphs = body_text.split('\n\n')
            for para_text in paragraphs:
                if para_text.strip():
                    para = doc.add_paragraph(para_text.strip())
                    para.paragraph_format.line_spacing = self.config["line_spacing"]

        # Add separator if configured
        if self.config["slide_separator"]:
            doc.add_paragraph()  # Space before separator


# =============================================================================
# NotesExtractionValidatorAgent
# =============================================================================

class NotesExtractionValidatorAgent(Agent):
    """
    HARDCODED validator agent for notes extraction quality.

    Validates that:
    - Extraction was successful
    - Minimum coverage thresholds are met
    - Notes quality meets requirements
    """

    def __init__(self):
        super().__init__(name="NotesExtractionValidatorAgent")
        self.config = VALIDATION_CONFIG

    def _process(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Validate extraction results."""
        extraction_result = context.get("extraction_result")

        if not extraction_result:
            return {
                "valid": False,
                "errors": ["No extraction result provided"],
                "warnings": []
            }

        if not extraction_result.get("success"):
            return {
                "valid": False,
                "errors": [extraction_result.get("error", "Extraction failed")],
                "warnings": []
            }

        errors = []
        warnings = []

        # Check coverage ratio
        total_slides = extraction_result.get("total_slides", 0)
        slides_with_notes = extraction_result.get("slides_with_notes", 0)

        if total_slides > 0:
            coverage_ratio = slides_with_notes / total_slides
            if coverage_ratio < self.config["min_slides_with_notes"]:
                warnings.append(
                    f"Low notes coverage: {coverage_ratio:.1%} "
                    f"(threshold: {self.config['min_slides_with_notes']:.1%})"
                )

        # Check total word count
        total_words = extraction_result.get("total_word_count", 0)
        if total_words < self.config["min_total_word_count"]:
            warnings.append(
                f"Low total word count: {total_words} "
                f"(threshold: {self.config['min_total_word_count']})"
            )

        # Check individual slide notes
        slide_notes = extraction_result.get("slide_notes", [])
        short_notes_count = 0
        for slide in slide_notes:
            if slide["has_content"] and slide["word_count"] < self.config["warn_short_note_threshold"]:
                short_notes_count += 1

        if short_notes_count > 0:
            warnings.append(f"{short_notes_count} slides have very short notes")

        # Check for extraction errors
        extraction_errors = extraction_result.get("errors", [])
        if extraction_errors:
            for err in extraction_errors:
                warnings.append(f"Extraction warning: {err}")

        return {
            "valid": len(errors) == 0,
            "errors": errors,
            "warnings": warnings,
            "statistics": {
                "total_slides": total_slides,
                "slides_with_notes": slides_with_notes,
                "coverage_ratio": coverage_ratio if total_slides > 0 else 0,
                "total_word_count": total_words,
                "short_notes_count": short_notes_count,
            }
        }


# =============================================================================
# PresentationNotesOrchestratorAgent
# =============================================================================

class PresentationNotesOrchestratorAgent(Agent):
    """
    HARDCODED orchestrator agent that coordinates the full notes extraction
    and Word document generation pipeline.

    Orchestrates:
    1. Notes extraction from PowerPoint
    2. Formatting for Word output
    3. Word document generation
    4. Validation of results
    """

    def __init__(self):
        super().__init__(name="PresentationNotesOrchestratorAgent")
        self.extractor = PresenterNotesExtractorAgent()
        self.formatter = MonologueFormatterAgent()
        self.generator = NotesToWordGeneratorAgent()
        self.validator = NotesExtractionValidatorAgent()

    def _process(self, context: Dict[str, Any]) -> Dict[str, Any]:
        """Orchestrate the full extraction and generation pipeline."""
        pptx_path = context.get("pptx_path")
        output_path = context.get("output_path")
        presentation_title = context.get("presentation_title")

        if not pptx_path:
            return {"success": False, "error": "No PowerPoint path provided"}

        # Auto-generate output path if not provided
        if not output_path:
            pptx_path_obj = Path(pptx_path)
            output_path = pptx_path_obj.parent / f"{pptx_path_obj.stem}_PresenterNotes.docx"

        # Auto-generate title if not provided
        if not presentation_title:
            presentation_title = Path(pptx_path).stem.replace("_", " ")

        # Step 1: Extract notes
        extraction_result = self.extractor.execute({"pptx_path": pptx_path})
        if not extraction_result.output.get("success"):
            return {
                "success": False,
                "step_failed": "extraction",
                "error": extraction_result.output.get("error")
            }

        # Step 2: Validate extraction
        validation_result = self.validator.execute({
            "extraction_result": extraction_result.output
        })

        # Step 3: Format for Word
        format_result = self.formatter.execute({
            "extraction_result": extraction_result.output,
            "presentation_title": presentation_title,
            "include_metadata": True
        })
        if not format_result.output.get("success"):
            return {
                "success": False,
                "step_failed": "formatting",
                "error": format_result.output.get("error")
            }

        # Step 4: Generate Word document
        generation_result = self.generator.execute({
            "formatted_document": format_result.output["formatted_document"],
            "output_path": str(output_path)
        })
        if not generation_result.output.get("success"):
            return {
                "success": False,
                "step_failed": "generation",
                "error": generation_result.output.get("error")
            }

        return {
            "success": True,
            "output_path": str(output_path),
            "extraction_stats": {
                "total_slides": extraction_result.output.get("total_slides"),
                "slides_with_notes": extraction_result.output.get("slides_with_notes"),
                "total_word_count": extraction_result.output.get("total_word_count"),
            },
            "validation": {
                "valid": validation_result.output.get("valid"),
                "warnings": validation_result.output.get("warnings", []),
            }
        }

    def process_directory(self, directory: Path, output_dir: Path = None) -> Dict[str, Any]:
        """Process all PowerPoint files in a directory."""
        directory = Path(directory)
        if output_dir is None:
            output_dir = directory

        pptx_files = list(directory.glob("**/*.pptx"))
        results = []
        success_count = 0
        error_count = 0

        for pptx_file in sorted(pptx_files):
            # Generate output path
            relative_path = pptx_file.relative_to(directory)
            output_path = output_dir / relative_path.parent / f"{pptx_file.stem}_PresenterNotes.docx"

            result = self.execute({
                "pptx_path": str(pptx_file),
                "output_path": str(output_path)
            })

            if result.output.get("success"):
                success_count += 1
            else:
                error_count += 1

            results.append({
                "source": str(pptx_file),
                "output": str(output_path) if result.output.get("success") else None,
                "success": result.output.get("success"),
                "error": result.output.get("error") if not result.output.get("success") else None
            })

        return {
            "total_files": len(pptx_files),
            "success_count": success_count,
            "error_count": error_count,
            "results": results
        }
