"""
PowerPoint Enhancement Script
=============================

Updates PowerPoint slides with fun facts or performance tips
using the SlideContentEnhancerAgent.
"""

import sys
from pathlib import Path

# Add project to path
sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from agents import (
    SlideContentEnhancerAgent,
    FunFactGeneratorAgent,
    SlideEnhancementFormatterAgent,
    SlideEnhancementValidatorAgent,
)

def get_slide_text(slide):
    """Extract all text from a slide."""
    text_parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_parts.append(shape.text)
    return " ".join(text_parts)

def get_slide_title(slide):
    """Get the title of a slide."""
    if slide.shapes.title:
        return slide.shapes.title.text
    return ""

def add_trivia_banner(slide, fact_text, label="Did You Know?", formatter=None, validator=None):
    """
    Add a trivia banner to the bottom of a slide using formatter and validator agents.

    Uses SlideEnhancementFormatterAgent for consistent 20pt font formatting
    and SlideEnhancementValidatorAgent to ensure requirements are met.
    """
    # Use formatter agent if provided, otherwise create one
    if formatter is None:
        formatter = SlideEnhancementFormatterAgent()

    # Get formatted configuration from agent
    formatted = formatter.execute({
        "label": label,
        "content": fact_text,
    })

    if not formatted.output.get("success"):
        print(f"    WARNING: Formatting failed - {formatted.output.get('error')}")
        return None

    # Validate if validator provided
    if validator:
        validation = validator.execute({"formatted_data": formatted.output})
        if not validation.output.get("valid"):
            for error in validation.output.get("errors", []):
                print(f"    VALIDATION ERROR: {error}")
            return None
        for warning in validation.output.get("warnings", []):
            print(f"    WARNING: {warning}")

    # Extract configuration from formatted output
    box_config = formatted.output["box"]
    label_config = formatted.output["label"]
    content_config = formatted.output["content"]

    # Position: bottom of slide (from formatter config)
    pos = box_config["position"]
    left = Inches(pos["left"])
    top = Inches(pos["top"])
    width = Inches(pos["width"])
    height = Inches(pos["height"])

    # Add shape
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )

    # Style the shape - from formatter config
    shape.fill.solid()
    fill_rgb = box_config["fill_color_rgb"]
    shape.fill.fore_color.rgb = RGBColor(fill_rgb[0], fill_rgb[1], fill_rgb[2])
    border_rgb = box_config["border_color_rgb"]
    shape.line.color.rgb = RGBColor(border_rgb[0], border_rgb[1], border_rgb[2])
    shape.line.width = Pt(box_config["border_width"])

    # Add text
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    # Label paragraph - 20pt font from formatter
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label_config["text"]
    run.font.bold = label_config["bold"]
    run.font.size = Pt(label_config["font_size"])  # 20pt
    label_rgb = label_config["color_rgb"]
    run.font.color.rgb = RGBColor(label_rgb[0], label_rgb[1], label_rgb[2])

    # Fact text - 20pt font from formatter
    run2 = p.add_run()
    run2.text = content_config["text"]
    run2.font.bold = content_config["bold"]
    run2.font.size = Pt(content_config["font_size"])  # 20pt
    content_rgb = content_config["color_rgb"]
    run2.font.color.rgb = RGBColor(content_rgb[0], content_rgb[1], content_rgb[2])

    return shape


# Keep old function name as alias for backwards compatibility
def add_fun_fact_box(slide, fact_text, label="Did You Know?"):
    """Backwards compatible wrapper for add_trivia_banner."""
    return add_trivia_banner(slide, fact_text, label)

def add_performance_tip_box(slide, tip_text):
    """Add a performance tip box to the bottom of a slide."""
    left = Inches(0.5)
    top = Inches(6.5)
    width = Inches(9)
    height = Inches(0.8)

    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )

    # Style - light blue for performance tips
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 243, 255)  # Light blue
    shape.line.color.rgb = RGBColor(0, 102, 204)  # Blue
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "🎭 Performance Tip: "
    run.font.bold = True
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0, 51, 102)

    run2 = p.add_run()
    run2.text = tip_text
    run2.font.size = Pt(10)
    run2.font.color.rgb = RGBColor(51, 51, 51)

    return shape

def enhance_presentation(pptx_path, output_path=None):
    """Enhance a PowerPoint presentation with fun facts and tips using hardcoded agents."""
    print(f"Loading presentation: {pptx_path}")
    prs = Presentation(pptx_path)

    # Initialize hardcoded agents
    enhancer = SlideContentEnhancerAgent()
    formatter = SlideEnhancementFormatterAgent()
    validator = SlideEnhancementValidatorAgent()

    enhancer.reset()

    slides_enhanced = 0
    facts_added = 0
    tips_added = 0
    validation_errors = 0

    print(f"Processing {len(prs.slides)} slides...")
    print(f"Font size: {formatter.get_config()['content_font_size']}pt (validated)")
    print()

    for i, slide in enumerate(prs.slides, 1):
        title = get_slide_title(slide)
        content = get_slide_text(slide)

        # Skip title slides or very short content
        if len(content) < 20:
            print(f"  Slide {i}: [SKIP] '{title[:30]}...' (too short)")
            continue

        # Get enhancement from content enhancer agent
        result = enhancer.execute({
            "slide_content": content,
            "slide_title": title,
        })

        if result.output.get("success"):
            enhancement_type = result.output["enhancement_type"]
            enhancement_content = result.output["content"]
            label = result.output["label"]

            # Use formatter and validator agents for 20pt font
            banner_result = add_trivia_banner(
                slide,
                enhancement_content,
                label,
                formatter=formatter,
                validator=validator
            )

            if banner_result:
                if enhancement_type == "fun_fact":
                    facts_added += 1
                else:
                    tips_added += 1
                slides_enhanced += 1
                print(f"  Slide {i}: [{label}] '{title[:30]}...' (20pt validated)")
            else:
                validation_errors += 1
                print(f"  Slide {i}: [ERROR] '{title[:30]}...' (validation failed)")
        else:
            print(f"  Slide {i}: [NONE] '{title[:30]}...'")

    # Save
    if output_path is None:
        output_path = pptx_path

    prs.save(output_path)

    print()
    print("=" * 60)
    print("ENHANCEMENT COMPLETE")
    print("=" * 60)
    print(f"Total slides:      {len(prs.slides)}")
    print(f"Slides enhanced:   {slides_enhanced}")
    print(f"Fun facts added:   {facts_added}")
    print(f"Trivia added:      {tips_added}")
    print(f"Validation errors: {validation_errors}")
    print(f"Font size:         {formatter.get_config()['content_font_size']}pt")
    print(f"Saved to:          {output_path}")

    return output_path

def main():
    # Default to Romeo and Juliet presentation
    pptx_path = Path(__file__).parent / "production" / "Unit_3_Romeo_and_Juliet" / "Day_01" / "Unit3_Day01_Shakespeare.pptx"

    if not pptx_path.exists():
        print(f"ERROR: PowerPoint not found: {pptx_path}")
        return 1

    print("=" * 60)
    print("POWERPOINT ENHANCEMENT: Adding Fun Facts & Performance Tips")
    print("=" * 60)
    print()

    output_path = enhance_presentation(str(pptx_path))

    return 0

if __name__ == "__main__":
    sys.exit(main())
