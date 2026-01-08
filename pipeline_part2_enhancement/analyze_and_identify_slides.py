"""
Slide Analysis and Graphic Organizer Identification
Automatically analyzes PowerPoint slides to identify candidates for graphic organizers
"""
from pptx import Presentation
import re
import json
from collections import defaultdict

# =============================================================================
# CONDITION TESTING - Identifies which graphic organizer type fits best
# =============================================================================

def test_table_condition(slide_content, notes):
    """Test if slide content is suitable for a TABLE"""
    score = 0
    text = f"{slide_content} {notes}".lower()

    # Strong indicators (each worth 3 points)
    strong_keywords = [
        'compared to', 'comparison', 'versus', 'vs', 'differs from',
        'in contrast', 'characteristics', 'features include',
        'properties', 'attributes'
    ]
    for keyword in strong_keywords:
        if keyword in text:
            score += 3

    # List structure indicators (worth 2 points)
    if text.count('•') >= 4 or text.count('\n-') >= 4:
        score += 2

    # Multiple categories (worth 2 points)
    category_words = ['type', 'class', 'category', 'group', 'kind']
    if sum(word in text for word in category_words) >= 2:
        score += 2

    # Numerical or statistical data (worth 3 points)
    if any(word in text for word in ['percent', '%', 'mg', 'dose', 'range', 'score']):
        score += 3

    # Multiple items with parallel structure (worth 2 points)
    if text.count(':') >= 3:
        score += 2

    return min(score, 15)  # Cap at 15

def test_flowchart_condition(slide_content, notes):
    """Test if slide content is suitable for a FLOWCHART"""
    score = 0
    text = f"{slide_content} {notes}".lower()

    # Sequential process indicators (worth 4 points each)
    sequential_keywords = [
        'first', 'then', 'next', 'finally', 'step',
        'process', 'mechanism', 'pathway', 'sequence'
    ]
    for keyword in sequential_keywords:
        if keyword in text:
            score += 4
            break  # Only count once

    # Numbered steps (worth 4 points)
    if re.search(r'\d+\.\s+', text) or 'step 1' in text or 'stage 1' in text:
        score += 4

    # Cause-effect or temporal (worth 3 points)
    if any(word in text for word in ['leads to', 'results in', 'causes', 'produces', 'triggers']):
        score += 3

    # Procedural language (worth 2 points)
    if any(word in text for word in ['procedure', 'protocol', 'method', 'technique']):
        score += 2

    return min(score, 15)

def test_decision_tree_condition(slide_content, notes):
    """Test if slide content is suitable for a DECISION TREE"""
    score = 0
    text = f"{slide_content} {notes}".lower()

    # Decision language (worth 4 points each)
    decision_keywords = [
        'if', 'whether', 'diagnosis', 'diagnostic', 'classify',
        'determine', 'distinguish', 'differentiate'
    ]
    for keyword in decision_keywords:
        if keyword in text:
            score += 4
            break

    # Binary choices (worth 3 points)
    if 'yes' in text and 'no' in text:
        score += 3
    if 'present' in text and 'absent' in text:
        score += 3

    # Multiple outcomes (worth 2 points)
    if text.count('disorder') >= 2 or text.count('condition') >= 2:
        score += 2

    # Rule-out language (worth 3 points)
    if any(word in text for word in ['rule out', 'exclude', 'criteria', 'assessment']):
        score += 3

    return min(score, 15)

def test_hierarchy_condition(slide_content, notes):
    """Test if slide content is suitable for a HIERARCHY"""
    score = 0
    text = f"{slide_content} {notes}".lower()

    # Hierarchical structure indicators (worth 4 points)
    hierarchy_keywords = [
        'classification', 'taxonomy', 'category', 'subcategory',
        'type', 'subtype', 'level', 'organization'
    ]
    for keyword in hierarchy_keywords:
        if keyword in text:
            score += 4
            break

    # Multiple levels evident (worth 3 points)
    if text.count('include') >= 2 or text.count('consist of') >= 1:
        score += 3

    # Nested structure (worth 3 points)
    if '  •' in slide_content or '    -' in slide_content:  # Indented bullets
        score += 3

    # System or family words (worth 2 points)
    if any(word in text for word in ['system', 'family', 'group', 'cluster']):
        score += 2

    return min(score, 15)

def test_timeline_condition(slide_content, notes):
    """Test if slide content is suitable for a TIMELINE"""
    score = 0
    text = f"{slide_content} {notes}".lower()

    # Temporal indicators (worth 4 points each)
    temporal_keywords = [
        'stage', 'phase', 'period', 'development', 'progression',
        'chronological', 'history', 'evolution', 'course'
    ]
    for keyword in temporal_keywords:
        if keyword in text:
            score += 4
            break

    # Time markers (worth 3 points)
    if re.search(r'\d+\s*(year|month|week|day|minute|hour)', text):
        score += 3
    if any(word in text for word in ['early', 'late', 'initial', 'final', 'onset']):
        score += 3

    # Sequential events (worth 2 points)
    if any(word in text for word in ['begins', 'ends', 'during', 'after', 'before']):
        score += 2

    return min(score, 15)

def test_spectrum_condition(slide_content, notes):
    """Test if slide content is suitable for a SPECTRUM"""
    score = 0
    text = f"{slide_content} {notes}".lower()

    # Continuum indicators (worth 4 points)
    spectrum_keywords = [
        'range', 'continuum', 'spectrum', 'gradient',
        'severity', 'intensity', 'degree'
    ]
    for keyword in spectrum_keywords:
        if keyword in text:
            score += 4
            break

    # Severity levels (worth 3 points)
    if 'mild' in text and 'severe' in text:
        score += 3
    if any(combo in text for combo in ['mild to severe', 'low to high', 'minimal to maximal']):
        score += 3

    # Bipolar dimensions (worth 3 points)
    if any(word in text for word in ['opposite', 'extreme', 'pole', 'end']):
        score += 3

    return min(score, 15)

def test_key_differentiators_condition(slide_content, notes):
    """Test if slide content is suitable for KEY DIFFERENTIATORS"""
    score = 0
    text = f"{slide_content} {notes}".lower()

    # High-yield discrimination language (worth 5 points)
    differentiator_keywords = [
        'commonly confused', 'distinguish', 'discrimination',
        'key difference', 'critical distinction', 'vs', 'versus'
    ]
    for keyword in differentiator_keywords:
        if keyword in text:
            score += 5
            break

    # Comparison of exactly 2-3 items (worth 4 points)
    concepts_count = sum(1 for marker in [':', '\n'] if marker in slide_content)
    if 2 <= concepts_count <= 4:
        score += 4

    # NCLEX relevance mentioned (worth 3 points)
    if 'nclex' in text or 'exam' in text or 'test' in text:
        score += 3

    # Differential diagnosis language (worth 3 points)
    if any(word in text for word in ['differential', 'discriminate', 'differentiate']):
        score += 3

    return min(score, 15)

# =============================================================================
# SLIDE SCORING AND RANKING
# =============================================================================

def score_slide_for_visual_types(slide_data):
    """Score a slide for all visual types and return ranked recommendations"""
    content = ' '.join(slide_data['content'])
    notes = slide_data['presenter_notes']

    scores = {
        'TABLE': test_table_condition(content, notes),
        'FLOWCHART': test_flowchart_condition(content, notes),
        'DECISION_TREE': test_decision_tree_condition(content, notes),
        'HIERARCHY': test_hierarchy_condition(content, notes),
        'TIMELINE': test_timeline_condition(content, notes),
        'SPECTRUM': test_spectrum_condition(content, notes),
        'KEY_DIFFERENTIATORS': test_key_differentiators_condition(content, notes),
    }

    # Sort by score
    ranked = sorted(scores.items(), key=lambda x: x[1], reverse=True)

    return {
        'scores': scores,
        'top_recommendation': ranked[0][0],
        'top_score': ranked[0][1],
        'ranked_recommendations': ranked
    }

# =============================================================================
# VARIETY ALGORITHM
# =============================================================================

def apply_variety_algorithm(candidates, min_quota=0.2, max_quota=0.4):
    """
    Apply variety algorithm to prevent repetition
    Enforces 20-40% quota and maximum 2 of same type before switching
    """
    if not candidates:
        return []

    selected = []
    type_counts = defaultdict(int)
    recent_types = []

    for candidate in candidates:
        visual_type = candidate['recommended_type']

        # Check if we've used this type twice in a row
        if len(recent_types) >= 2 and recent_types[-1] == recent_types[-2] == visual_type:
            # Try to use second-best recommendation
            alt_type = candidate['scores']['ranked_recommendations'][1][0]
            if type_counts[alt_type] < 3:  # Max 3 of any type per section
                visual_type = alt_type
            else:
                # Skip this slide to enforce variety
                continue

        # Add to selection
        selected.append({
            **candidate,
            'final_type': visual_type
        })
        type_counts[visual_type] += 1
        recent_types.append(visual_type)

        # Keep only last 2 types in memory
        if len(recent_types) > 2:
            recent_types.pop(0)

    return selected


# =============================================================================
# COMPOSITE SCORING SYSTEM (Option B: Weighted Selection)
# =============================================================================

def calculate_anchor_point_density(slide_content, notes):
    """
    Calculate anchor point density based on high-value NCLEX keywords

    Returns score 0-10 based on presence of key concepts that appear
    in anchor point summaries (neurotransmitters, disorders, treatments, etc.)
    """
    text = f"{slide_content} {notes}".lower()
    score = 0

    # High-frequency NCLEX anchor point keywords (each worth points)
    neurotransmitter_keywords = [
        'dopamine', 'serotonin', 'norepinephrine', 'gaba', 'glutamate',
        'acetylcholine', 'endorphin', 'substance p'
    ]

    brain_structure_keywords = [
        'hippocampus', 'amygdala', 'prefrontal cortex', 'basal ganglia',
        'thalamus', 'hypothalamus', 'cerebellum', 'substantia nigra',
        'caudate', 'putamen', 'nucleus accumbens'
    ]

    disorder_keywords = [
        'schizophrenia', 'depression', 'bipolar', 'anxiety', 'ptsd',
        'ocd', 'adhd', 'autism', "alzheimer's", "parkinson's", "huntington's",
        'dementia', 'panic disorder', 'phobia'
    ]

    treatment_keywords = [
        'ssri', 'snri', 'maoi', 'benzodiazepine', 'antipsychotic',
        'lithium', 'carbamazepine', 'valproate', 'lamotrigine',
        'levodopa', 'dopamine agonist', 'cholinesterase inhibitor'
    ]

    mechanism_keywords = [
        'reuptake', 'receptor', 'agonist', 'antagonist', 'inhibitor',
        'neurotransmission', 'synapse', 'pathway', 'mechanism of action',
        'blood-brain barrier', 'half-life', 'bioavailability'
    ]

    all_keywords = (
        neurotransmitter_keywords +
        brain_structure_keywords +
        disorder_keywords +
        treatment_keywords +
        mechanism_keywords
    )

    # Count unique keyword categories present (max 5 categories)
    categories_present = 0
    if any(kw in text for kw in neurotransmitter_keywords):
        categories_present += 1
    if any(kw in text for kw in brain_structure_keywords):
        categories_present += 1
    if any(kw in text for kw in disorder_keywords):
        categories_present += 1
    if any(kw in text for kw in treatment_keywords):
        categories_present += 1
    if any(kw in text for kw in mechanism_keywords):
        categories_present += 1

    # Base score from categories (0-5 points)
    score += categories_present

    # Additional points for multiple keywords in same category (depth)
    keyword_count = sum(1 for kw in all_keywords if kw in text)
    if keyword_count >= 3:
        score += 2
    if keyword_count >= 5:
        score += 2
    if keyword_count >= 8:
        score += 1

    return min(score, 10)  # Cap at 10


def calculate_position_bonus(slide_number, total_slides, section_size):
    """
    Calculate position bonus to encourage spread throughout section

    Gives slight preference to slides distributed across:
    - Beginning (first 20%)
    - Middle (30-70%)
    - End (last 20%)

    Returns 0-3 points
    """
    if section_size <= 3:
        return 1  # Small section, position doesn't matter much

    position_in_section = (slide_number % section_size) / section_size

    # Prefer slides in strategic positions
    if position_in_section <= 0.2:  # First 20% (intro/foundation)
        return 2
    elif position_in_section >= 0.8:  # Last 20% (synthesis/comparison)
        return 2
    elif 0.3 <= position_in_section <= 0.7:  # Middle 40% (core content)
        return 3
    else:
        return 1  # Transitions


def calculate_composite_score(candidate, section_size, total_slides):
    """
    Calculate composite score using weighted factors

    Weights:
    - Base quality score: 60% (how well slide fits graphic organizer type)
    - Anchor point density: 30% (NCLEX concept density)
    - Position bonus: 10% (strategic placement in section)

    Returns composite score (0-100 scale)
    """
    # Extract data
    base_score = candidate['score']  # Already 0-15 scale
    slide_number = candidate['slide_number']
    slide_content = candidate.get('content_preview', '')
    notes = candidate.get('notes', '')

    # Calculate components
    anchor_density = calculate_anchor_point_density(slide_content, notes)  # 0-10
    position_bonus = calculate_position_bonus(slide_number, total_slides, section_size)  # 0-3

    # Normalize to 0-100 scale
    normalized_base = (base_score / 15.0) * 100  # Base score max 15
    normalized_anchor = (anchor_density / 10.0) * 100  # Anchor max 10
    normalized_position = (position_bonus / 3.0) * 100  # Position max 3

    # Apply weights
    composite = (
        normalized_base * 0.6 +      # 60% base quality
        normalized_anchor * 0.3 +     # 30% anchor density
        normalized_position * 0.1     # 10% position
    )

    return round(composite, 2)


# =============================================================================
# SECTION DETECTION AND ANALYSIS
# =============================================================================

def detect_sections(slides):
    """Detect section boundaries in slides"""
    sections = []
    current_section = {'start': 0, 'title': 'Section 1', 'slides': []}

    for idx, slide in enumerate(slides):
        # Section title indicators
        title_text = ' '.join(slide['content']).lower()

        # Check if this is a section intro slide
        if any(indicator in title_text for indicator in [
            'section', 'part', 'module', 'chapter'
        ]) or (len(slide['content']) <= 2 and idx > 0):
            # Save previous section
            if current_section['slides']:
                sections.append(current_section)
            # Start new section
            current_section = {
                'start': idx,
                'title': slide['content'][0] if slide['content'] else f"Section {len(sections)+1}",
                'slides': []
            }
        else:
            current_section['slides'].append(idx)

    # Add final section
    if current_section['slides']:
        sections.append(current_section)

    # If no sections detected, treat whole presentation as one section
    if not sections:
        sections = [{'start': 0, 'title': 'Main', 'slides': list(range(len(slides)))}]

    return sections

def analyze_powerpoint_for_graphic_organizers(pptx_path, threshold=6):
    """
    Main analysis function
    Returns list of slides that should have graphic organizers
    """
    prs = Presentation(pptx_path)

    # Extract slide data
    slides = []
    for idx, slide in enumerate(prs.slides):
        slide_data = {
            'slide_number': idx + 1,
            'content': [],
            'presenter_notes': ''
        }

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    slide_data['content'].append(text)

        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                slide_data['presenter_notes'] = notes_slide.notes_text_frame.text.strip()

        slides.append(slide_data)

    # Detect sections
    sections = detect_sections(slides)

    # Analyze each section
    all_recommendations = []

    for section in sections:
        print(f"\n{'='*80}")
        print(f"SECTION: {section['title']}")
        print(f"Slides: {len(section['slides'])}")
        print(f"{'='*80}\n")

        # Score slides in this section
        candidates = []
        section_size = len(section['slides'])
        total_slides = len(slides)

        for slide_idx in section['slides']:
            slide = slides[slide_idx]
            scoring = score_slide_for_visual_types(slide)

            # Only consider if top score meets threshold
            if scoring['top_score'] >= threshold:
                candidate = {
                    'slide_number': slide_idx + 1,
                    'recommended_type': scoring['top_recommendation'],
                    'score': scoring['top_score'],
                    'scores': scoring,
                    'content_preview': slide['content'][0][:80] if slide['content'] else '',
                    'notes': slide['presenter_notes']  # Include notes for anchor density
                }
                candidates.append(candidate)

        # COMPOSITE SCORING: Calculate weighted score for each candidate
        print(f"\nCalculating composite scores (Base 60% + Anchor Density 30% + Position 10%)...")
        for candidate in candidates:
            composite = calculate_composite_score(candidate, section_size, total_slides)
            candidate['composite_score'] = composite

            # Debug output
            base = candidate['score']
            anchor = calculate_anchor_point_density(
                candidate['content_preview'],
                candidate['notes']
            )
            position = calculate_position_bonus(
                candidate['slide_number'],
                total_slides,
                section_size
            )
            print(f"  Slide {candidate['slide_number']}: "
                  f"Base={base:.1f}/15 | Anchor={anchor:.1f}/10 | Position={position}/3 "
                  f"-> Composite={composite:.1f}/100")

        # SORT BY COMPOSITE SCORE (best slides first)
        candidates.sort(key=lambda x: x['composite_score'], reverse=True)
        print(f"\n[OK] Sorted {len(candidates)} candidates by composite score (highest first)")

        # Apply variety algorithm (now processes best candidates first)
        selected = apply_variety_algorithm(candidates, min_quota=0.2, max_quota=0.4)

        # Enforce quota (20-40% of section)
        section_size = len(section['slides'])
        min_required = int(section_size * 0.2)
        max_allowed = int(section_size * 0.4)

        if len(selected) < min_required:
            print(f"WARNING: Only {len(selected)} slides selected, but minimum is {min_required}")
            print(f"Lowering threshold or reviewing section content recommended.")

        if len(selected) > max_allowed:
            selected = selected[:max_allowed]
            print(f"Limiting to {max_allowed} slides (40% quota)")

        print(f"\n[OK] Selected {len(selected)} slides for graphic organizers ({len(selected)/section_size*100:.1f}%)")

        for item in selected:
            composite = item.get('composite_score', item['score'])
            print(f"  Slide {item['slide_number']}: {item['final_type']} "
                  f"(composite: {composite:.1f}/100, base: {item['score']:.1f}/15)")
            print(f"    Preview: {item['content_preview']}")

        all_recommendations.extend(selected)

    return all_recommendations, slides

# =============================================================================
# SAVE ANALYSIS REPORT
# =============================================================================

def save_analysis_report(pptx_path, recommendations):
    """Save analysis report to JSON"""
    output_path = pptx_path.replace('.pptx', '_graphic_organizer_recommendations.json')

    report = {
        'source_file': pptx_path,
        'total_slides_analyzed': 0,
        'total_recommended': len(recommendations),
        'recommendations': recommendations
    }

    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(report, f, indent=2, ensure_ascii=False)

    print(f"\n{'='*80}")
    print(f"Analysis report saved to:")
    print(f"{output_path}")
    print(f"{'='*80}\n")

    return output_path

# =============================================================================
# MAIN EXECUTION
# =============================================================================

if __name__ == "__main__":
    # Configure for your production folder
    import glob
    # UPDATE THIS PATH to your production powerpoints folder
    pptx_dir = r"C:\path\to\production\powerpoints"
    files = glob.glob(pptx_dir + r"\*.pptx")

    for file_path in files:
        print(f"\n{'#'*80}")
        print(f"ANALYZING: {file_path.split('\\')[-1]}")
        print(f"{'#'*80}")

        recommendations, slides = analyze_powerpoint_for_graphic_organizers(file_path, threshold=6)

        save_analysis_report(file_path, recommendations)
