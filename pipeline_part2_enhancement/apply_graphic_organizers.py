"""
Apply Graphic Organizers to PowerPoint Slides - VERSION 3
Updated based on Section 9 success patterns

Key Improvements:
- No arbitrary feature limits (4 → 10)
- No character truncation (40 → 100 chars)
- Professional table styling by default
- Blueprint-compatible functions
- Better fallback labels
- Two-panel layout support
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import json
import copy
import re
import os

# =============================================================================
# GLOBAL CONSTANTS
# =============================================================================

VISUAL_MIN_FONT_SIZE = Pt(18)

COLORS = {
    'accent_red': RGBColor(220, 38, 38),
    'dark_header': RGBColor(26, 26, 26),
    'light_gray_bg': RGBColor(245, 245, 245),
    'white_content': RGBColor(255, 255, 255),
    'medium_gray': RGBColor(128, 128, 128),
    'light_blue': RGBColor(96, 165, 250),
    'light_green': RGBColor(74, 222, 128),
    'light_yellow': RGBColor(250, 204, 21),
    'black': RGBColor(0, 0, 0),
    'white': RGBColor(255, 255, 255),
    'dark_navy': RGBColor(0, 51, 102),
    'light_gray_blue': RGBColor(240, 244, 248),
    'dark_gray': RGBColor(33, 37, 41),
    'table_header_blue': RGBColor(31, 73, 125),  # NEW: Professional table header
    'table_alt_row': RGBColor(242, 242, 242),    # NEW: Alternating rows
    'table_first_col': RGBColor(220, 220, 220),  # NEW: First column
}

CONTENT_AREA = {
    'left': Inches(0.54),
    'top': Inches(0.72),
    'width': Inches(12.26),
    'height': Inches(6.52),
}

HEADER_ZONE_BOTTOM = Inches(0.7)
FOOTER_ZONE_TOP = Inches(6.5)  # Tightened from 7.3 to exclude supplementary text
FAR_RIGHT_ZONE = Inches(11.0)

# =============================================================================
# NEW: PROFESSIONAL TABLE CREATION
# =============================================================================

def create_professional_table(slide, rows, cols, left, top, width, height):
    """
    Create a table with professional styling like Section 9

    Features:
    - Dark blue header row with white bold text
    - Light gray first column with bold text
    - Alternating row shading
    - Proper cell margins
    - All borders visible
    """
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Apply professional styling
    for i, row in enumerate(table.rows):
        for j, cell in enumerate(row.cells):
            # Set cell margins
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)

            # Header row styling
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['table_header_blue']
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.size = VISUAL_MIN_FONT_SIZE  # 18pt minimum
                            run.font.color.rgb = COLORS['black']  # Always black text
            else:
                # Regular cell styling
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.LEFT
                        for run in paragraph.runs:
                            run.font.size = VISUAL_MIN_FONT_SIZE  # 18pt minimum
                            run.font.color.rgb = COLORS['black']

                # First column bold + light gray
                if j == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['table_first_col']
                    if cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = True

                # Alternating row shading
                elif i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['table_alt_row']

    return table


def apply_table_from_blueprint(slide, blueprint_data):
    """
    Apply table from structured blueprint data

    blueprint_data format:
    {
        'title': 'Table Title',  # Optional - for reference
        'headers': ['Column 1', 'Column 2', 'Column 3'],
        'rows': [
            {'Column 1': 'value', 'Column 2': 'value', 'Column 3': 'value'},
            ...
        ],
        'column_widths': [0.25, 0.375, 0.375],  # Optional proportions
        'footer': 'Footer text'  # Optional
    }
    """
    headers = blueprint_data.get('headers', [])
    rows = blueprint_data.get('rows', [])
    column_widths = blueprint_data.get('column_widths', None)
    footer = blueprint_data.get('footer', None)

    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    # Calculate dimensions
    table_left = CONTENT_AREA['left']
    table_top = CONTENT_AREA['top']
    table_width = CONTENT_AREA['width']
    table_height = Inches(5.8) if not footer else Inches(5.5)

    # Create table with professional styling
    table = create_professional_table(
        slide, num_rows, num_cols,
        table_left, table_top, table_width, table_height
    )

    # Set column widths if specified
    if column_widths and len(column_widths) == num_cols:
        for i, width_prop in enumerate(column_widths):
            table.columns[i].width = Inches(table_width.inches * width_prop)

    # Fill headers
    for j, header in enumerate(headers):
        table.cell(0, j).text = header

    # Fill data rows
    for i, row_data in enumerate(rows):
        for j, header in enumerate(headers):
            cell_value = row_data.get(header, '')
            table.cell(i + 1, j).text = str(cell_value)

    # Add footer if present
    if footer:
        footer_box = slide.shapes.add_textbox(
            table_left,
            table_top + table_height + Inches(0.1),
            table_width,
            Inches(0.5)
        )
        tf = footer_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = VISUAL_MIN_FONT_SIZE  # 18pt minimum
        p.font.color.rgb = COLORS['black']
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return table


# =============================================================================
# NEW: TWO-PANEL LAYOUT (Like Slide 5)
# =============================================================================

def apply_two_panel_layout(slide, panel_a_data, panel_b_data, connection_text=None):
    """
    Create two-panel side-by-side layout like Section 9 Slide 5

    panel_a_data/panel_b_data format:
    {
        'title': 'Panel Title',
        'features': ['Feature 1', 'Feature 2', ...],
        'color_light': RGBColor(...),  # Optional
        'color_dark': RGBColor(...)    # Optional
    }
    """
    # Default colors
    default_colors_a = {
        'color_light': RGBColor(220, 230, 241),  # Light blue
        'color_dark': RGBColor(31, 73, 125)      # Dark blue
    }
    default_colors_b = {
        'color_light': RGBColor(242, 220, 219),  # Light red
        'color_dark': RGBColor(192, 0, 0)        # Dark red
    }

    panel_width = Inches(5.8)
    panel_height = Inches(5.5)
    panel_gap = Inches(0.66)
    panel_left_A = CONTENT_AREA['left']
    panel_left_B = panel_left_A + panel_width + panel_gap

    # Panel A
    panel_a = slide.shapes.add_shape(
        1,  # Rectangle
        panel_left_A,
        CONTENT_AREA['top'],
        panel_width,
        panel_height
    )
    panel_a.fill.solid()
    panel_a.fill.fore_color.rgb = panel_a_data.get('color_light', default_colors_a['color_light'])
    panel_a.line.color.rgb = panel_a_data.get('color_dark', default_colors_a['color_dark'])
    panel_a.line.width = Pt(2)

    tf_a = panel_a.text_frame
    tf_a.word_wrap = True
    tf_a.margin_left = Inches(0.2)
    tf_a.margin_right = Inches(0.2)
    tf_a.margin_top = Inches(0.15)

    # Panel A title
    p = tf_a.paragraphs[0]
    p.text = panel_a_data['title']
    p.font.bold = True
    p.font.size = VISUAL_MIN_FONT_SIZE  # 18pt minimum
    p.font.color.rgb = COLORS['black']  # Always black text
    p.alignment = PP_ALIGN.CENTER

    # Panel A features
    for feature in panel_a_data.get('features', []):
        p = tf_a.add_paragraph()
        p.text = "• " + feature
        p.font.size = VISUAL_MIN_FONT_SIZE  # 18pt minimum
        p.font.color.rgb = COLORS['black']
        p.space_before = Pt(6)
        p.level = 0

    # Panel B (similar structure)
    panel_b = slide.shapes.add_shape(
        1,  # Rectangle
        panel_left_B,
        CONTENT_AREA['top'],
        panel_width,
        panel_height
    )
    panel_b.fill.solid()
    panel_b.fill.fore_color.rgb = panel_b_data.get('color_light', default_colors_b['color_light'])
    panel_b.line.color.rgb = panel_b_data.get('color_dark', default_colors_b['color_dark'])
    panel_b.line.width = Pt(2)

    tf_b = panel_b.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = Inches(0.2)
    tf_b.margin_right = Inches(0.2)
    tf_b.margin_top = Inches(0.15)

    # Panel B title
    p = tf_b.paragraphs[0]
    p.text = panel_b_data['title']
    p.font.bold = True
    p.font.size = VISUAL_MIN_FONT_SIZE  # 18pt minimum
    p.font.color.rgb = COLORS['black']  # Always black text
    p.alignment = PP_ALIGN.CENTER

    # Panel B features
    for feature in panel_b_data.get('features', []):
        p = tf_b.add_paragraph()
        p.text = "• " + feature
        p.font.size = VISUAL_MIN_FONT_SIZE  # 18pt minimum
        p.font.color.rgb = COLORS['black']
        p.space_before = Pt(6)
        p.level = 0

    # Connection statement
    if connection_text:
        connection = slide.shapes.add_textbox(
            panel_left_A,
            CONTENT_AREA['top'] + panel_height + Inches(0.1),
            panel_width * 2 + panel_gap,
            Inches(0.4)
        )
        tf_conn = connection.text_frame
        p = tf_conn.paragraphs[0]
        p.text = connection_text
        p.font.size = VISUAL_MIN_FONT_SIZE  # 18pt minimum
        p.font.color.rgb = COLORS['black']
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER


# =============================================================================
# SMART CONTENT EXTRACTION (From Section 8 - Already Working)
# =============================================================================

def extract_slide_content(slide):
    """
    Extract ONLY educational content for graphic organizers using hybrid detection.
    Excludes: titles, copyright, NCLEX boxes, tier labels, elements outside content area.
    Returns: List of individual lines/bullet points (not multi-line blocks)
    """
    content = []

    for shape in slide.shapes:
        # Skip if no text
        if not shape.has_text_frame or not shape.text or not shape.text.strip():
            continue

        text = shape.text.strip()
        text_lower = text.lower()

        # EXCLUSION MODE 1: Keyword-based detection
        exclusion_keywords = [
            'copyright', '©', '�', 'prepjet', 'all rights reserved',
            '📚 nclex application', 'nclex application',
            'tier 1', 'tier 2', 'tier 3',
            'licensure preparation', 'education service'
        ]
        if any(keyword in text_lower for keyword in exclusion_keywords):
            continue

        # EXCLUSION MODE 2: Position-based detection
        try:
            if shape.top < HEADER_ZONE_BOTTOM:
                continue
            if shape.top > FOOTER_ZONE_TOP:
                continue
            if shape.left > FAR_RIGHT_ZONE:
                continue

            # INCLUSION: Shape must be within or overlapping main content area
            shape_left = shape.left
            shape_top = shape.top
            shape_right = shape.left + shape.width
            shape_bottom = shape.top + shape.height

            content_left = CONTENT_AREA['left']
            content_top = CONTENT_AREA['top']
            content_right = CONTENT_AREA['left'] + CONTENT_AREA['width']
            content_bottom = CONTENT_AREA['top'] + CONTENT_AREA['height']

            overlaps = (
                shape_right > content_left and
                shape_left < content_right and
                shape_bottom > content_top and
                shape_top < content_bottom
            )

            if overlaps:
                # CRITICAL: Split multi-line text into individual lines
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        content.append(line)

        except Exception as e:
            pass

    return content


def clear_slide_content_area(slide):
    """
    Remove shapes from content area, but preserve:
    - Title (header zone)
    - TIER labels (far right)
    - Copyright (footer zone)
    - NCLEX APPLICATION boxes (detected by keyword)
    """
    shapes_to_delete = []

    for shape in slide.shapes:
        try:
            # PRESERVE MODE 1: Keyword-based preservation
            if shape.has_text_frame and shape.text:
                text_lower = shape.text.lower()
                preservation_keywords = [
                    'copyright', '©', '�', 'prepjet', 'all rights reserved',
                    '📚 nclex application', 'nclex application',
                    'tier 1', 'tier 2', 'tier 3',
                    'licensure preparation', 'education service'
                ]
                if any(keyword in text_lower for keyword in preservation_keywords):
                    continue  # Don't delete

            # PRESERVE MODE 2: Position-based preservation
            if shape.top < HEADER_ZONE_BOTTOM:
                continue
            if shape.top > FOOTER_ZONE_TOP:
                continue
            if shape.left > FAR_RIGHT_ZONE:
                continue

            # DELETE: Shapes in content area that aren't preserved
            if (shape.top >= CONTENT_AREA['top'] and
                shape.left >= CONTENT_AREA['left'] and
                shape.top + shape.height <= CONTENT_AREA['top'] + CONTENT_AREA['height']):
                shapes_to_delete.append(shape)

        except Exception as e:
            pass

    for shape in reversed(shapes_to_delete):
        sp = shape.element
        sp.getparent().remove(sp)


# =============================================================================
# UPDATED: KEY DIFFERENTIATORS (With Improvements)
# =============================================================================

def apply_key_differentiators_to_slide(slide, content_list, notes):
    """
    Apply KEY DIFFERENTIATORS graphic organizer

    IMPROVED:
    - Feature limit increased: 4 → 10
    - Character limit increased: 40 → 100
    - Better fallback labels (use content instead of "Concept A/B")
    """
    items = [line.strip(' •-') for line in content_list if line.strip() and not line.startswith('TIER')]

    # Extract two main concepts
    concepts = []
    current_concept = {'title': '', 'features': []}

    for line in items:
        if ':' in line and not current_concept['title']:
            current_concept['title'] = line.split(':')[0]
        elif line and current_concept['title']:
            # IMPROVED: Increased from 4 to 10, and 40 to 100 chars
            if len(current_concept['features']) < 10:
                current_concept['features'].append('• ' + line[:100])
            else:
                concepts.append(current_concept)
                current_concept = {'title': '', 'features': []}

    if current_concept['title'] and current_concept['features']:
        concepts.append(current_concept)

    # IMPROVED: Better fallback labels
    if len(concepts) < 2:
        # Try to use actual content as titles instead of generic "Concept A/B"
        title_a = items[0] if items else 'Primary Concept'
        title_b = items[int(len(items)/2)] if len(items) > 6 else 'Secondary Concept'

        concepts = [
            {'title': title_a[:50], 'features': ['• ' + line[:100] for line in items[1:6]]},
            {'title': title_b[:50], 'features': ['• ' + line[:100] for line in items[6:11]]}
        ]

    # Two comparison boxes
    config = [
        {'color_light': RGBColor(224, 242, 254), 'color_dark': RGBColor(59, 130, 246)},
        {'color_light': RGBColor(220, 252, 231), 'color_dark': RGBColor(34, 197, 94)}
    ]

    box_w = 4.8
    box_h = 4.0
    box_y = 1.5
    spacing = 1.0
    total_w = 2 * box_w + spacing
    start_x = (CONTENT_AREA['width'].inches - total_w) / 2 + CONTENT_AREA['left'].inches

    for idx, (concept, cfg) in enumerate(zip(concepts[:2], config)):
        x = start_x + idx * (box_w + spacing)

        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(box_y),
            Inches(box_w), Inches(box_h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = cfg['color_light']
        box.line.color.rgb = cfg['color_dark']
        box.line.width = Pt(2)

        # Text
        tf = box.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)
        tf.word_wrap = True

        # Title
        p = tf.paragraphs[0]
        p.text = concept['title']
        p.font.size = VISUAL_MIN_FONT_SIZE  # 18pt minimum
        p.font.bold = True
        p.font.color.rgb = COLORS['black']  # Always black text
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)

        # Features
        for feature in concept['features']:
            p = tf.add_paragraph()
            p.text = feature
            p.font.size = VISUAL_MIN_FONT_SIZE  # 18pt minimum
            p.font.color.rgb = COLORS['black']
            p.space_before = Pt(4)
            p.level = 0


# =============================================================================
# UPDATED: TABLE APPLICATION (With Professional Styling)
# =============================================================================

def extract_table_data_from_content(content_list):
    """Extract tabular data from bullet points"""
    # Look for patterns like "Item: description" or "Category - details"
    rows = []
    headers = ['Item', 'Description']

    for line in content_list:
        # Split on : or -
        if ':' in line:
            parts = line.split(':', 1)
            if len(parts) == 2:
                rows.append([parts[0].strip(' •-'), parts[1].strip()])
        elif ' - ' in line:
            parts = line.split(' - ', 1)
            if len(parts) == 2:
                rows.append([parts[0].strip(' •-'), parts[1].strip()])

    if not rows:
        # Fallback: just make a simple list
        headers = ['Item']
        rows = [[line.strip(' •-')] for line in content_list if line.strip()]

    # IMPROVED: Increased from 6 to 10 rows
    return [headers] + rows[:10]


def apply_table_to_slide(slide, content_list, notes):
    """
    Apply TABLE graphic organizer

    IMPROVED: Uses professional styling by default
    """
    # Extract data
    table_data = extract_table_data_from_content(content_list)

    if len(table_data) < 2:
        return

    num_rows = len(table_data)
    num_cols = len(table_data[0])

    # IMPROVED: Use professional table creation
    table = create_professional_table(
        slide, num_rows, num_cols,
        CONTENT_AREA['left'],
        CONTENT_AREA['top'],
        CONTENT_AREA['width'],
        Inches(min(6.0, num_rows * 0.6))
    )

    # Fill table
    for i, row in enumerate(table_data):
        for j, cell_text in enumerate(row):
            table.cell(i, j).text = str(cell_text)


# =============================================================================
# TIMELINE GRAPHIC ORGANIZER
# =============================================================================

def apply_timeline_to_slide(slide, content_list, notes):
    """
    Apply TIMELINE graphic organizer - horizontal timeline with stages
    """
    items = [line.strip(' •-') for line in content_list if line.strip() and not line.startswith('TIER')]

    if len(items) < 2:
        items = ['Stage 1', 'Stage 2', 'Stage 3']

    # Limit to 5 stages for readability
    stages = items[:5]
    num_stages = len(stages)

    # Timeline dimensions
    timeline_y = 2.5
    timeline_height = 0.15
    stage_box_height = 1.8

    # Calculate positions
    content_left = CONTENT_AREA['left'].inches
    content_width = CONTENT_AREA['width'].inches
    stage_width = (content_width - 0.5) / num_stages

    # Draw timeline bar
    timeline_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(content_left), Inches(timeline_y),
        Inches(content_width), Inches(timeline_height)
    )
    timeline_bar.fill.solid()
    timeline_bar.fill.fore_color.rgb = COLORS['dark_gray']
    timeline_bar.line.fill.background()

    # Add stage boxes
    colors = [
        RGBColor(224, 242, 254),  # Light blue
        RGBColor(220, 252, 231),  # Light green
        RGBColor(254, 243, 199),  # Light yellow
        RGBColor(254, 226, 226),  # Light red
        RGBColor(243, 232, 255),  # Light purple
    ]

    for idx, stage in enumerate(stages):
        x = content_left + idx * stage_width + 0.1

        # Stage box above timeline
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(timeline_y - stage_box_height - 0.2),
            Inches(stage_width - 0.2), Inches(stage_box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[idx % len(colors)]
        box.line.color.rgb = COLORS['dark_gray']
        box.line.width = Pt(1)

        # Stage text
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = stage[:80]  # Limit length
        p.font.size = VISUAL_MIN_FONT_SIZE
        p.font.color.rgb = COLORS['black']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Stage number circle on timeline
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + (stage_width - 0.2) / 2 - 0.2), Inches(timeline_y - 0.1),
            Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['black']
        circle.line.fill.background()

        # Number in circle
        ctf = circle.text_frame
        ctf.paragraphs[0].text = str(idx + 1)
        ctf.paragraphs[0].font.size = Pt(14)
        ctf.paragraphs[0].font.color.rgb = COLORS['white']
        ctf.paragraphs[0].font.bold = True
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER


# =============================================================================
# FLOWCHART GRAPHIC ORGANIZER
# =============================================================================

def apply_flowchart_to_slide(slide, content_list, notes):
    """
    Apply FLOWCHART graphic organizer - process flow with arrows
    """
    items = [line.strip(' •-') for line in content_list if line.strip() and not line.startswith('TIER')]

    if len(items) < 2:
        items = ['Start', 'Process', 'End']

    # Limit to 4 steps for readability
    steps = items[:4]
    num_steps = len(steps)

    # Flowchart dimensions
    box_width = 2.5
    box_height = 1.2
    arrow_width = 0.8
    start_y = 2.0

    content_left = CONTENT_AREA['left'].inches
    content_width = CONTENT_AREA['width'].inches
    total_width = num_steps * box_width + (num_steps - 1) * arrow_width
    start_x = content_left + (content_width - total_width) / 2

    for idx, step in enumerate(steps):
        x = start_x + idx * (box_width + arrow_width)

        # Process box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(start_y),
            Inches(box_width), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(224, 242, 254)  # Light blue
        box.line.color.rgb = COLORS['dark_gray']
        box.line.width = Pt(2)

        # Step text
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = step[:60]
        p.font.size = VISUAL_MIN_FONT_SIZE
        p.font.color.rgb = COLORS['black']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Arrow to next step
        if idx < num_steps - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + box_width + 0.1), Inches(start_y + box_height / 2 - 0.15),
                Inches(arrow_width - 0.2), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['dark_gray']
            arrow.line.fill.background()


# =============================================================================
# DECISION TREE GRAPHIC ORGANIZER
# =============================================================================

def apply_decision_tree_to_slide(slide, content_list, notes):
    """
    Apply DECISION_TREE graphic organizer - if-then branching
    """
    items = [line.strip(' •-') for line in content_list if line.strip() and not line.startswith('TIER')]

    if len(items) < 3:
        items = ['Decision Point?', 'Yes: Outcome A', 'No: Outcome B']

    content_left = CONTENT_AREA['left'].inches
    content_width = CONTENT_AREA['width'].inches
    center_x = content_left + content_width / 2

    # Decision diamond at top
    diamond_size = 1.5
    diamond = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        Inches(center_x - diamond_size / 2), Inches(1.2),
        Inches(diamond_size), Inches(diamond_size)
    )
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = RGBColor(254, 243, 199)  # Light yellow
    diamond.line.color.rgb = COLORS['dark_gray']
    diamond.line.width = Pt(2)

    tf = diamond.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = items[0][:40] if items else 'Decision?'
    p.font.size = VISUAL_MIN_FONT_SIZE
    p.font.color.rgb = COLORS['black']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Two outcome boxes
    box_width = 3.5
    box_height = 2.0
    box_y = 3.5
    spacing = 1.5

    outcomes = [
        {'label': 'YES', 'text': items[1] if len(items) > 1 else 'Outcome A', 'color': RGBColor(220, 252, 231)},
        {'label': 'NO', 'text': items[2] if len(items) > 2 else 'Outcome B', 'color': RGBColor(254, 226, 226)}
    ]

    for idx, outcome in enumerate(outcomes):
        x = center_x - box_width - spacing / 2 + idx * (box_width + spacing)

        # Outcome box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(box_y),
            Inches(box_width), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = outcome['color']
        box.line.color.rgb = COLORS['dark_gray']
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)

        # Label
        p = tf.paragraphs[0]
        p.text = outcome['label']
        p.font.size = VISUAL_MIN_FONT_SIZE
        p.font.color.rgb = COLORS['black']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Content
        p2 = tf.add_paragraph()
        p2.text = outcome['text'][:100]
        p2.font.size = VISUAL_MIN_FONT_SIZE
        p2.font.color.rgb = COLORS['black']
        p2.alignment = PP_ALIGN.CENTER


# =============================================================================
# HIERARCHY GRAPHIC ORGANIZER
# =============================================================================

def apply_hierarchy_to_slide(slide, content_list, notes):
    """
    Apply HIERARCHY graphic organizer - top-down structure
    """
    items = [line.strip(' •-') for line in content_list if line.strip() and not line.startswith('TIER')]

    if len(items) < 2:
        items = ['Top Level', 'Sub Level 1', 'Sub Level 2', 'Sub Level 3']

    content_left = CONTENT_AREA['left'].inches
    content_width = CONTENT_AREA['width'].inches
    center_x = content_left + content_width / 2

    # Top level box
    top_width = 4.0
    top_height = 1.0
    top_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(center_x - top_width / 2), Inches(1.0),
        Inches(top_width), Inches(top_height)
    )
    top_box.fill.solid()
    top_box.fill.fore_color.rgb = RGBColor(31, 73, 125)  # Dark blue
    top_box.line.fill.background()

    tf = top_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = items[0][:50]
    p.font.size = VISUAL_MIN_FONT_SIZE
    p.font.color.rgb = COLORS['black']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Sub-level boxes (up to 4)
    sub_items = items[1:5]
    num_subs = len(sub_items)

    if num_subs > 0:
        sub_width = 2.5
        sub_height = 1.5
        sub_y = 2.8
        total_sub_width = num_subs * sub_width + (num_subs - 1) * 0.3
        start_x = center_x - total_sub_width / 2

        colors = [
            RGBColor(224, 242, 254),
            RGBColor(220, 252, 231),
            RGBColor(254, 243, 199),
            RGBColor(243, 232, 255),
        ]

        for idx, sub in enumerate(sub_items):
            x = start_x + idx * (sub_width + 0.3)

            sub_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(sub_y),
                Inches(sub_width), Inches(sub_height)
            )
            sub_box.fill.solid()
            sub_box.fill.fore_color.rgb = colors[idx % len(colors)]
            sub_box.line.color.rgb = COLORS['dark_gray']
            sub_box.line.width = Pt(1)

            tf = sub_box.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)

            p = tf.paragraphs[0]
            p.text = sub[:60]
            p.font.size = VISUAL_MIN_FONT_SIZE
            p.font.color.rgb = COLORS['black']
            p.alignment = PP_ALIGN.CENTER


# =============================================================================
# SPECTRUM GRAPHIC ORGANIZER
# =============================================================================

def apply_spectrum_to_slide(slide, content_list, notes):
    """
    Apply SPECTRUM graphic organizer - range/gradient visualization
    """
    items = [line.strip(' •-') for line in content_list if line.strip() and not line.startswith('TIER')]

    if len(items) < 2:
        items = ['Low End', 'Middle', 'High End']

    # Use first and last as endpoints, middle items as points
    endpoints = [items[0], items[-1]] if len(items) >= 2 else ['Low', 'High']
    midpoints = items[1:-1] if len(items) > 2 else []

    content_left = CONTENT_AREA['left'].inches
    content_width = CONTENT_AREA['width'].inches

    # Spectrum bar
    bar_height = 0.6
    bar_y = 2.5
    bar_left = content_left + 0.5
    bar_width = content_width - 1.0

    # Gradient bar (use multiple rectangles to simulate gradient)
    num_segments = 5
    segment_width = bar_width / num_segments
    gradient_colors = [
        RGBColor(254, 226, 226),  # Light red
        RGBColor(254, 243, 199),  # Light yellow
        RGBColor(254, 249, 195),  # Lighter yellow
        RGBColor(220, 252, 231),  # Light green
        RGBColor(187, 247, 208),  # Greener
    ]

    for i in range(num_segments):
        segment = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bar_left + i * segment_width), Inches(bar_y),
            Inches(segment_width + 0.01), Inches(bar_height)
        )
        segment.fill.solid()
        segment.fill.fore_color.rgb = gradient_colors[i]
        segment.line.fill.background()

    # Endpoint labels
    label_width = 2.5
    label_height = 1.2

    # Left endpoint
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(bar_left), Inches(bar_y + bar_height + 0.3),
        Inches(label_width), Inches(label_height)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(254, 226, 226)
    left_box.line.color.rgb = COLORS['dark_gray']
    left_box.line.width = Pt(1)

    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = endpoints[0][:50]
    p.font.size = VISUAL_MIN_FONT_SIZE
    p.font.color.rgb = COLORS['black']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Right endpoint
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(bar_left + bar_width - label_width), Inches(bar_y + bar_height + 0.3),
        Inches(label_width), Inches(label_height)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(220, 252, 231)
    right_box.line.color.rgb = COLORS['dark_gray']
    right_box.line.width = Pt(1)

    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = endpoints[1][:50]
    p.font.size = VISUAL_MIN_FONT_SIZE
    p.font.color.rgb = COLORS['black']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Middle label if exists
    if midpoints:
        mid_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bar_left + bar_width / 2 - label_width / 2), Inches(bar_y - label_height - 0.3),
            Inches(label_width), Inches(label_height)
        )
        mid_box.fill.solid()
        mid_box.fill.fore_color.rgb = RGBColor(254, 249, 195)
        mid_box.line.color.rgb = COLORS['dark_gray']
        mid_box.line.width = Pt(1)

        tf = mid_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = midpoints[0][:50]
        p.font.size = VISUAL_MIN_FONT_SIZE
        p.font.color.rgb = COLORS['black']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER


# =============================================================================
# MAIN APPLICATION WITH BLUEPRINT SUPPORT
# =============================================================================

def apply_from_blueprint(slide, blueprint):
    """Apply graphic organizer from blueprint data"""
    org_type = blueprint.get('type', '').upper()

    clear_slide_content_area(slide)

    if org_type == 'TABLE':
        apply_table_from_blueprint(slide, blueprint)
    elif org_type == 'TWO_PANEL' or org_type == 'KEY_DIFFERENTIATORS':
        if 'panel_a' in blueprint and 'panel_b' in blueprint:
            apply_two_panel_layout(
                slide,
                blueprint['panel_a'],
                blueprint['panel_b'],
                blueprint.get('connection', None)
            )
        else:
            # Fallback to regular KEY_DIFFERENTIATORS
            content = extract_slide_content(slide)
            apply_key_differentiators_to_slide(slide, content, '')
    # Add other types as needed


def apply_graphic_organizers(pptx_path, recommendations_json, blueprints_dir=None, output_path=None):
    """
    Main function to apply graphic organizers

    NEW: Checks for blueprints first, falls back to extraction
    """
    # Load presentation
    prs = Presentation(pptx_path)

    # Load recommendations
    with open(recommendations_json, 'r', encoding='utf-8') as f:
        recommendations_data = json.load(f)

    # Handle wrapper structure from analyze_and_identify_slides.py
    if isinstance(recommendations_data, dict) and 'recommendations' in recommendations_data:
        recommendations = recommendations_data['recommendations']
    else:
        recommendations = recommendations_data

    # Process each slide
    for rec in recommendations:
        slide_num = rec['slide_number']
        # Support both field names from different JSON formats
        organizer_type = rec.get('final_type') or rec.get('recommended_type') or rec.get('recommended_organizer')
        slide_index = slide_num - 1

        if slide_index >= len(prs.slides):
            continue

        slide = prs.slides[slide_index]

        # CHECK FOR BLUEPRINT FIRST
        if blueprints_dir and os.path.exists(blueprints_dir):
            blueprint_file = os.path.join(blueprints_dir, f'slide_{slide_num}_blueprint.json')
            if os.path.exists(blueprint_file):
                print(f"Slide {slide_num}: Using blueprint")
                with open(blueprint_file, 'r', encoding='utf-8') as f:
                    blueprint = json.load(f)
                apply_from_blueprint(slide, blueprint)
                continue

        # OTHERWISE, USE EXTRACTION APPROACH
        print(f"Slide {slide_num}: Using extraction for {organizer_type}")
        content = extract_slide_content(slide)
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ''

        clear_slide_content_area(slide)

        # Apply appropriate organizer
        if organizer_type == 'TABLE':
            apply_table_to_slide(slide, content, notes)
        elif organizer_type == 'KEY_DIFFERENTIATORS':
            apply_key_differentiators_to_slide(slide, content, notes)
        elif organizer_type == 'TIMELINE':
            apply_timeline_to_slide(slide, content, notes)
        elif organizer_type == 'FLOWCHART':
            apply_flowchart_to_slide(slide, content, notes)
        elif organizer_type == 'DECISION_TREE':
            apply_decision_tree_to_slide(slide, content, notes)
        elif organizer_type == 'HIERARCHY':
            apply_hierarchy_to_slide(slide, content, notes)
        elif organizer_type == 'SPECTRUM':
            apply_spectrum_to_slide(slide, content, notes)
        else:
            print(f"  WARNING: Unknown organizer type '{organizer_type}' - skipping")

    # Save
    if not output_path:
        base = os.path.splitext(pptx_path)[0]
        output_path = f"{base}_ENHANCED.pptx"

    prs.save(output_path)
    print(f"\n[SUCCESS] Enhanced PowerPoint saved to: {output_path}")
    return output_path


# =============================================================================
# EXAMPLE USAGE
# =============================================================================

if __name__ == "__main__":
    import glob

    # Configure for your production folder
    pptx_dir = r"C:\Users\mcdan\Desktop\Lifespan_Development_Production_2026-01-01\powerpoints"
    blueprints_dir = r"C:\Users\mcdan\Desktop\Lifespan_Development_Production_2026-01-01\integrated"

    pptx_files = glob.glob(pptx_dir + r"\*.pptx")
    # Exclude already enhanced files
    pptx_files = [f for f in pptx_files if '_ENHANCED' not in f]

    print(f"Found {len(pptx_files)} PowerPoint files to enhance\n")

    for pptx_path in pptx_files:
        # Find corresponding recommendations JSON
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        recommendations_json = os.path.join(pptx_dir, f"{base_name}_graphic_organizer_recommendations.json")

        if os.path.exists(recommendations_json):
            print(f"\n{'='*60}")
            print(f"Processing: {base_name}")
            print(f"{'='*60}")
            try:
                apply_graphic_organizers(pptx_path, recommendations_json, blueprints_dir)
            except Exception as e:
                print(f"ERROR processing {base_name}: {e}")
        else:
            print(f"Skipping {base_name} - no recommendations JSON found")
