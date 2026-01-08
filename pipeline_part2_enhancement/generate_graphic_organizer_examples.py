"""
Generate Example Graphic Organizers
Creates a PowerPoint with examples of all 7 graphic organizer types
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy

# =============================================================================
# GLOBAL CONSTANTS - From visual_specs.txt
# =============================================================================

# Minimum font size (18pt requirement)
VISUAL_MIN_FONT_SIZE = Pt(18)

# Color Palette from template
COLORS = {
    'accent_red': RGBColor(220, 38, 38),
    'dark_header': RGBColor(26, 26, 26),
    'light_gray_bg': RGBColor(245, 245, 245),
    'white_content': RGBColor(255, 255, 255),
    'medium_gray': RGBColor(128, 128, 128),
    'light_blue': RGBColor(96, 165, 250),
    'light_green': RGBColor(74, 222, 128),
    'light_yellow': RGBColor(250, 204, 21),
    'black': RGBColor(0, 0, 0),
    'white': RGBColor(255, 255, 255),
    # Additional colors for specific organizers
    'dark_navy': RGBColor(0, 51, 102),
    'light_gray_blue': RGBColor(240, 244, 248),
    'dark_gray': RGBColor(33, 37, 41),
}

# Content area boundaries
CONTENT_AREA = {
    'left': Inches(0.54),
    'top': Inches(0.72),
    'width': Inches(12.26),
    'height': Inches(6.52),
}

# =============================================================================
# TABLE GENERATION
# =============================================================================

def set_cell_borders(cell, border_color_rgb=(0, 51, 102), width_emu=12700):
    """Set borders on a table cell"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    hex_color = '{:02X}{:02X}{:02X}'.format(*border_color_rgb)

    for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
        for existing in tcPr.findall(qn(f'a:{border_name}')):
            tcPr.remove(existing)

        border_xml = (
            f'<a:{border_name} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            f'w="{width_emu}" cap="flat" cmpd="sng" algn="ctr">'
            f'<a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
            f'<a:prstDash val="solid"/>'
            f'<a:round/>'
            f'<a:headEnd type="none" w="med" len="med"/>'
            f'<a:tailEnd type="none" w="med" len="med"/>'
            f'</a:{border_name}>'
        )
        border_elem = parse_xml(border_xml)
        tcPr.append(border_elem)

def create_table_example(slide, title_text):
    """Create a table example comparing antipsychotics"""
    # Add title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.5)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['dark_header']

    # Table data
    table_data = [
        ['Drug Class', 'Mechanism', 'Key Side Effects'],
        ['Conventional', 'D2 antagonist', 'High EPS risk, TD'],
        ['Atypical', '5-HT2A + D2', 'Low EPS, metabolic'],
        ['Clozapine', 'Multi-receptor', 'Agranulocytosis'],
    ]

    rows = len(table_data)
    cols = len(table_data[0])

    # Create table
    left = CONTENT_AREA['left']
    top = Inches(1.2)
    width = CONTENT_AREA['width']
    height = Inches(3.0)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    col_width = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_width

    # Populate table
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = "Calibri"
            paragraph.font.size = Pt(20)

            if row_idx == 0:
                # Header row
                paragraph.font.bold = True
                paragraph.font.color.rgb = COLORS['white']
                paragraph.alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['dark_navy']
            else:
                # Data rows
                paragraph.font.color.rgb = COLORS['black']
                paragraph.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['white']

            set_cell_borders(cell, (0, 51, 102), 12700)

# =============================================================================
# FLOWCHART GENERATION
# =============================================================================

def create_flowchart_example(slide, title_text):
    """Create a flowchart showing sleep stages"""
    # Add title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.5)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['dark_header']

    # Flowchart steps
    steps = [
        'Stage N1: Light Sleep',
        'Stage N2: Sleep Spindles',
        'Stage N3: Deep Sleep',
        'REM Sleep: Dreaming'
    ]

    start_y = 1.5
    box_width = 3.5
    box_height = 0.9
    spacing = 0.4
    center_x = (CONTENT_AREA['width'].inches - box_width) / 2 + CONTENT_AREA['left'].inches

    prev_y = start_y
    for idx, step_text in enumerate(steps):
        # Create box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(center_x), Inches(prev_y),
            Inches(box_width), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['light_blue']
        box.line.color.rgb = RGBColor(59, 130, 246)
        box.line.width = Pt(2.5)

        # Add text
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_text
        p.font.size = VISUAL_MIN_FONT_SIZE
        p.font.name = "Calibri"
        p.font.color.rgb = COLORS['black']
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = 1  # Middle

        # Add arrow if not last step
        if idx < len(steps) - 1:
            arrow_start_y = prev_y + box_height
            arrow_end_y = prev_y + box_height + spacing

            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(center_x + box_width/2), Inches(arrow_start_y),
                Inches(center_x + box_width/2), Inches(arrow_end_y)
            )
            connector.line.color.rgb = RGBColor(30, 64, 175)
            connector.line.width = Pt(3)

        prev_y += box_height + spacing

# =============================================================================
# DECISION TREE GENERATION
# =============================================================================

def create_decision_tree_example(slide, title_text):
    """Create a decision tree for antipsychotic selection"""
    # Add title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.5)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['dark_header']

    # Level 1 decision
    l1_x, l1_y = 4.9, 1.0
    l1_w, l1_h = 3.5, 1.1

    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(l1_x), Inches(l1_y),
        Inches(l1_w), Inches(l1_h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(254, 249, 195)
    box.line.color.rgb = RGBColor(234, 179, 8)
    box.line.width = Pt(2.5)

    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Movement Disorder Risk?"
    p.font.size = VISUAL_MIN_FONT_SIZE
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['black']
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1

    # Outcomes
    outcomes = [
        ('High Risk', 'Use Atypical', 2.5, 3.5, COLORS['light_green']),
        ('Low Risk', 'Either Type OK', 7.5, 3.5, COLORS['light_blue']),
    ]

    for label, text, x, y, color in outcomes:
        # Create outcome box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(2.5), Inches(0.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = COLORS['medium_gray']
        box.line.width = Pt(2)

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = VISUAL_MIN_FONT_SIZE
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = COLORS['black']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = 1

        # Add connector from L1
        start_x = l1_x + (0 if x < l1_x else l1_w)
        start_y = l1_y + l1_h
        end_x = x + 1.25
        end_y = y

        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(start_x), Inches(start_y),
            Inches(end_x), Inches(end_y)
        )
        connector.line.color.rgb = COLORS['medium_gray']
        connector.line.width = Pt(2)

# =============================================================================
# HIERARCHY GENERATION
# =============================================================================

def create_hierarchy_example(slide, title_text):
    """Create a hierarchy of neurotransmitter systems"""
    # Add title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.5)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['dark_header']

    # Level 1 - Parent
    l1_w, l1_h = 4.5, 0.9
    l1_x = (CONTENT_AREA['width'].inches - l1_w) / 2 + CONTENT_AREA['left'].inches
    l1_y = 1.2

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l1_x), Inches(l1_y),
        Inches(l1_w), Inches(l1_h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['accent_red']
    box.line.color.rgb = RGBColor(153, 27, 27)
    box.line.width = Pt(3)

    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Neurotransmitter Systems"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1

    # Level 2 - Children
    l2_items = ['Dopamine', 'Serotonin', 'GABA']
    l2_w, l2_h = 2.8, 0.75
    l2_y = l1_y + l1_h + 1.0
    l2_spacing = 0.8
    l2_total_width = len(l2_items) * l2_w + (len(l2_items) - 1) * l2_spacing
    l2_start_x = (CONTENT_AREA['width'].inches - l2_total_width) / 2 + CONTENT_AREA['left'].inches

    for idx, item in enumerate(l2_items):
        x = l2_start_x + idx * (l2_w + l2_spacing)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(l2_y),
            Inches(l2_w), Inches(l2_h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['light_blue']
        box.line.color.rgb = RGBColor(59, 130, 246)
        box.line.width = Pt(2)

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = VISUAL_MIN_FONT_SIZE
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = COLORS['black']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = 1

        # Add connector
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(l1_x + l1_w/2), Inches(l1_y + l1_h),
            Inches(x + l2_w/2), Inches(l2_y)
        )
        connector.line.color.rgb = COLORS['medium_gray']
        connector.line.width = Pt(2.5)

# =============================================================================
# TIMELINE GENERATION
# =============================================================================

def create_timeline_example(slide, title_text):
    """Create a timeline of sleep stages"""
    # Add title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.5)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['dark_header']

    # Timeline bar
    bar_w = 10.5
    bar_h = 0.3
    bar_x = (CONTENT_AREA['width'].inches - bar_w) / 2 + CONTENT_AREA['left'].inches
    bar_y = 3.5

    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(bar_x), Inches(bar_y),
        Inches(bar_w), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['medium_gray']
    bar.line.color.rgb = RGBColor(75, 85, 99)
    bar.line.width = Pt(2)

    # Events
    events = [
        ('0 min', 'Sleep Onset', 0),
        ('10 min', 'Stage N2', 0.25),
        ('30 min', 'Stage N3', 0.5),
        ('90 min', 'REM', 1.0)
    ]

    for time, label, position in events:
        x = bar_x + position * (bar_w - 0.5)

        # Marker
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(bar_y - 0.1),
            Inches(0.5), Inches(0.5)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = COLORS['accent_red']
        marker.line.color.rgb = RGBColor(153, 27, 27)
        marker.line.width = Pt(2.5)

        # Event box above
        event_box = slide.shapes.add_textbox(
            Inches(x - 0.5), Inches(bar_y - 1.3),
            Inches(1.5), Inches(0.8)
        )
        tf = event_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{time}\n{label}"
        p.font.size = VISUAL_MIN_FONT_SIZE
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = COLORS['black']
        p.alignment = PP_ALIGN.CENTER

# =============================================================================
# SPECTRUM GENERATION
# =============================================================================

def create_spectrum_example(slide, title_text):
    """Create a spectrum of symptom severity"""
    # Add title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.5)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['dark_header']

    # Spectrum bar
    bar_w = 10.0
    bar_h = 1.2
    bar_x = (CONTENT_AREA['width'].inches - bar_w) / 2 + CONTENT_AREA['left'].inches
    bar_y = 2.5

    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(bar_x), Inches(bar_y),
        Inches(bar_w), Inches(bar_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['light_blue']
    bar.line.color.rgb = COLORS['medium_gray']
    bar.line.width = Pt(3)

    # Endpoint labels
    labels = [
        ('Mild', bar_x - 0.3, bar_y + bar_h/2 - 0.2),
        ('Moderate', bar_x + bar_w/2 - 0.6, bar_y + bar_h/2 - 0.2),
        ('Severe', bar_x + bar_w + 0.3, bar_y + bar_h/2 - 0.2)
    ]

    for text, x, y in labels:
        label_box = slide.shapes.add_textbox(
            Inches(x), Inches(y),
            Inches(1.2), Inches(0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = COLORS['black']
        p.alignment = PP_ALIGN.CENTER

# =============================================================================
# KEY DIFFERENTIATORS GENERATION
# =============================================================================

def create_key_differentiators_example(slide, title_text):
    """Create key differentiators comparison"""
    # Add title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.5)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['dark_header']

    # Two comparison boxes
    items = [
        {
            'title': 'Tardive Dyskinesia',
            'color_light': RGBColor(224, 242, 254),
            'color_dark': RGBColor(59, 130, 246),
            'features': [
                '• Involuntary movements',
                '• Face, tongue, jaw',
                '• Often irreversible',
                '• Older adults, women'
            ]
        },
        {
            'title': 'Akathisia',
            'color_light': RGBColor(220, 252, 231),
            'color_dark': RGBColor(34, 197, 94),
            'features': [
                '• Inner restlessness',
                '• Can\'t sit still',
                '• Subjective distress',
                '• Dose-related'
            ]
        }
    ]

    box_w = 4.5
    box_h = 3.5
    box_y = 1.5
    spacing = 1.0
    total_w = 2 * box_w + spacing
    start_x = (CONTENT_AREA['width'].inches - total_w) / 2 + CONTENT_AREA['left'].inches

    for idx, item in enumerate(items):
        x = start_x + idx * (box_w + spacing)

        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(box_y),
            Inches(box_w), Inches(box_h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = item['color_light']
        box.line.color.rgb = item['color_dark']
        box.line.width = Pt(3)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(box_y),
            Inches(box_w), Inches(0.6)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = item['color_dark']
        title_bar.line.fill.background()

        tf = title_bar.text_frame
        p = tf.paragraphs[0]
        p.text = item['title']
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = 1

        # Features
        features_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(box_y + 0.8),
            Inches(box_w - 0.4), Inches(box_h - 1.0)
        )
        tf = features_box.text_frame
        tf.word_wrap = True

        for feature in item['features']:
            p = tf.add_paragraph()
            p.text = feature
            p.font.size = VISUAL_MIN_FONT_SIZE
            p.font.name = "Calibri"
            p.font.color.rgb = COLORS['black']
            p.level = 0

    # VS divider
    vs_x = start_x + box_w + spacing/2
    vs_y = box_y + box_h/2

    vs_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(vs_x - 0.5), Inches(vs_y - 0.5),
        Inches(1.0), Inches(1.0)
    )
    vs_circle.fill.solid()
    vs_circle.fill.fore_color.rgb = COLORS['white']
    vs_circle.line.color.rgb = COLORS['medium_gray']
    vs_circle.line.width = Pt(3)

    tf = vs_circle.text_frame
    p = tf.paragraphs[0]
    p.text = "VS"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['black']
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1

# =============================================================================
# MAIN GENERATION
# =============================================================================

def generate_examples():
    """Generate all example graphic organizers"""
    # Load template
    template_path = r"C:\Users\mcdan\OneDrive\Desktop\pipeline_part2_enhancement\template master.pptx"
    prs = Presentation(template_path)

    # Get blank layout
    blank_layout = prs.slide_layouts[6]  # Typically blank layout

    # Create example slides
    examples = [
        ('TABLE: Antipsychotic Comparison', create_table_example),
        ('FLOWCHART: Sleep Stages', create_flowchart_example),
        ('DECISION TREE: Antipsychotic Selection', create_decision_tree_example),
        ('HIERARCHY: Neurotransmitter Systems', create_hierarchy_example),
        ('TIMELINE: Sleep Cycle Progression', create_timeline_example),
        ('SPECTRUM: Symptom Severity', create_spectrum_example),
        ('KEY DIFFERENTIATORS: TD vs Akathisia', create_key_differentiators_example)
    ]

    for title, create_func in examples:
        slide = prs.slides.add_slide(blank_layout)
        create_func(slide, title)
        print(f"Created: {title}")

    # Save
    output_path = r"C:\Users\mcdan\OneDrive\Desktop\pipeline_part2_enhancement\GRAPHIC_ORGANIZER_EXAMPLES.pptx"
    prs.save(output_path)
    print(f"\nExample PowerPoint saved to:\n{output_path}")

if __name__ == "__main__":
    generate_examples()
