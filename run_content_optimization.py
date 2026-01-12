#!/usr/bin/env python3
"""
Slide Content Optimization Script
==================================

Optimizes PowerPoint slide content for student note-taking by:
- Detecting and fixing truncated statements
- Condensing verbose content to concise bullets
- Elaborating presenter notes for teacher delivery
- Validating all content meets quality standards

Uses hardcoded agents for all optimization, validation, and formatting tasks.

Usage:
    python run_content_optimization.py                           # Process all PPTX in production/
    python run_content_optimization.py --path /path/to/file.pptx # Process single file
    python run_content_optimization.py --dir /path/to/folder     # Process directory
    python run_content_optimization.py --unit 3                  # Process Unit 3 only
    python run_content_optimization.py --test                    # Run test samples
"""

import sys
import argparse
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any
import yaml

# Add project to path
sys.path.insert(0, str(Path(__file__).parent))

from agents import (
    TruncationDetectorAgent,
    TruncationFixerAgent,
    SlideContentCondensationAgent,
    PresenterNotesElaboratorAgent,
    SlideContentValidatorAgent,
    ContentBalanceOrchestratorAgent,
)

# Pipeline root
PIPELINE_ROOT = Path(__file__).parent
PRODUCTION_DIR = PIPELINE_ROOT / "production"
SAMPLES_DIR = PIPELINE_ROOT / "samples"


class ContentOptimizationProcessor:
    """
    Processor for optimizing slide content.

    Uses hardcoded agents for all tasks:
    - TruncationDetectorAgent: Detects truncated statements
    - TruncationFixerAgent: Fixes truncated content
    - SlideContentCondensationAgent: Condenses verbose content
    - PresenterNotesElaboratorAgent: Elaborates presenter notes
    - SlideContentValidatorAgent: Validates content quality
    - ContentBalanceOrchestratorAgent: Orchestrates full pipeline
    """

    def __init__(self, verbose: bool = False):
        self.verbose = verbose

        # Initialize hardcoded agents
        self.detector = TruncationDetectorAgent()
        self.fixer = TruncationFixerAgent()
        self.condenser = SlideContentCondensationAgent()
        self.elaborator = PresenterNotesElaboratorAgent()
        self.validator = SlideContentValidatorAgent()
        self.orchestrator = ContentBalanceOrchestratorAgent()

        # Statistics
        self.stats = {
            "files_processed": 0,
            "files_success": 0,
            "files_failed": 0,
            "slides_processed": 0,
            "truncations_detected": 0,
            "truncations_fixed": 0,
            "content_condensed": 0,
            "validation_passed": 0,
            "validation_failed": 0,
            "total_words_reduced": 0,
        }

    def process_content(self, content: str, content_type: str = "bullet") -> Dict[str, Any]:
        """Process a single piece of content through the optimization pipeline."""
        result = self.orchestrator.execute({
            "content": content,
            "slide_type": content_type,
            "additional_context": {}
        })

        output = result.output

        # Update stats
        if output.get("truncation_detected"):
            self.stats["truncations_detected"] += 1
        if output.get("truncation_fixes_applied", 0) > 0:
            self.stats["truncations_fixed"] += output["truncation_fixes_applied"]

        condensation = output.get("condensation_stats", {})
        if condensation.get("reduction_percentage", 0) > 0:
            self.stats["content_condensed"] += 1
            self.stats["total_words_reduced"] += (
                condensation.get("original_words", 0) - condensation.get("condensed_words", 0)
            )

        if output.get("validation_passed"):
            self.stats["validation_passed"] += 1
        else:
            self.stats["validation_failed"] += 1

        self.stats["slides_processed"] += 1

        return output

    def run_test_samples(self) -> Dict[str, Any]:
        """Run optimization on test samples."""
        samples_file = SAMPLES_DIR / "slide_content_test_samples.yaml"

        if not samples_file.exists():
            print(f"ERROR: Sample file not found: {samples_file}")
            return {"success": False, "error": "Sample file not found"}

        with open(samples_file, 'r') as f:
            samples = yaml.safe_load(f)

        results = {
            "truncation_tests": [],
            "condensation_tests": [],
            "validation_tests": [],
            "summary": {}
        }

        print("\n" + "=" * 70)
        print("TRUNCATION DETECTION TESTS")
        print("=" * 70)

        truncation_pass = 0
        truncation_total = 0

        for test in samples.get("truncation_test_cases", []):
            truncation_total += 1
            detection = self.detector.execute({
                "content": test["input"],
                "content_type": "bullet"
            })

            detected = detection.output.get("has_truncation", False)
            expected = test.get("expected_truncation", False)
            passed = detected == expected

            if passed:
                truncation_pass += 1
                status = "[PASS]"
            else:
                status = "[FAIL]"

            print(f"  {status} {test['id']}: {test['description']}")
            if self.verbose:
                print(f"         Input: {test['input'][:60]}...")
                print(f"         Expected: {expected}, Got: {detected}")

            results["truncation_tests"].append({
                "id": test["id"],
                "passed": passed,
                "expected": expected,
                "actual": detected
            })

        print(f"\n  Truncation Tests: {truncation_pass}/{truncation_total} passed")

        print("\n" + "=" * 70)
        print("CONTENT CONDENSATION TESTS")
        print("=" * 70)

        condense_pass = 0
        condense_total = 0

        for test in samples.get("condensation_test_cases", []):
            condense_total += 1
            condensation = self.condenser.execute({
                "content": test["input"],
                "content_type": "bullet"
            })

            reduction = condensation.output.get("reduction_percentage", 0)
            expected_reduction = test.get("expected_word_reduction", False)
            min_reduction = test.get("min_reduction_percentage", 0)

            if expected_reduction:
                passed = reduction >= min_reduction
            else:
                passed = reduction < 10  # No significant reduction expected

            if passed:
                condense_pass += 1
                status = "[PASS]"
            else:
                status = "[FAIL]"

            print(f"  {status} {test['id']}: {test['description']}")
            if self.verbose:
                print(f"         Reduction: {reduction}%, Expected min: {min_reduction}%")

            results["condensation_tests"].append({
                "id": test["id"],
                "passed": passed,
                "reduction_percentage": reduction
            })

        print(f"\n  Condensation Tests: {condense_pass}/{condense_total} passed")

        print("\n" + "=" * 70)
        print("FULL OPTIMIZATION TESTS")
        print("=" * 70)

        optimize_pass = 0
        optimize_total = 0

        for test in samples.get("full_optimization_test_cases", []):
            optimize_total += 1
            content = test.get("input", {}).get("content", [])

            result = self.orchestrator.execute({
                "content": content,
                "slide_type": test.get("input", {}).get("type", "content"),
                "additional_context": {"title": test.get("input", {}).get("title", "")}
            })

            expected = test.get("expected_output", {})
            actual = {
                "truncation_fixed": result.output.get("truncation_fixes_applied", 0) > 0,
                "content_condensed": result.output.get("condensation_stats", {}).get("reduction_percentage", 0) > 0,
                "notes_elaborated": len(result.output.get("presenter_notes", "")) > 100,
                "validation_passed": result.output.get("validation_passed", False)
            }

            # Check if all expected conditions are met
            passed = True
            for key, expected_value in expected.items():
                if actual.get(key) != expected_value:
                    passed = False
                    break

            if passed:
                optimize_pass += 1
                status = "[PASS]"
            else:
                status = "[FAIL]"

            print(f"  {status} {test['id']}: {test['description']}")
            if self.verbose:
                print(f"         Expected: {expected}")
                print(f"         Actual: {actual}")

        print(f"\n  Full Optimization Tests: {optimize_pass}/{optimize_total} passed")

        # Summary
        total_pass = truncation_pass + condense_pass + optimize_pass
        total_tests = truncation_total + condense_total + optimize_total

        results["summary"] = {
            "total_tests": total_tests,
            "total_passed": total_pass,
            "pass_rate": round(total_pass / max(total_tests, 1) * 100, 1),
            "truncation_accuracy": round(truncation_pass / max(truncation_total, 1) * 100, 1),
            "condensation_accuracy": round(condense_pass / max(condense_total, 1) * 100, 1),
            "optimization_accuracy": round(optimize_pass / max(optimize_total, 1) * 100, 1),
        }

        return results

    def process_pptx_file(self, pptx_path: Path) -> Dict[str, Any]:
        """Process a single PowerPoint file."""
        try:
            from pptx import Presentation
        except ImportError:
            print("ERROR: python-pptx not installed. Run: pip install python-pptx")
            return {"success": False, "error": "python-pptx not installed"}

        print(f"Processing: {pptx_path.name}")

        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            self.stats["files_failed"] += 1
            return {"success": False, "error": str(e)}

        slide_results = []
        truncations_found = 0
        truncations_fixed = 0
        words_reduced = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text.append(text)

            if not slide_text:
                continue

            # Process each text element
            for i, text in enumerate(slide_text):
                result = self.process_content(text, "bullet")

                if result.get("truncation_detected"):
                    truncations_found += 1
                if result.get("truncation_fixes_applied", 0) > 0:
                    truncations_fixed += result["truncation_fixes_applied"]

                condensation = result.get("condensation_stats", {})
                words_reduced += (
                    condensation.get("original_words", 0) -
                    condensation.get("condensed_words", 0)
                )

                # Update slide content if optimized
                optimized = result.get("optimized_slide_content", text)
                notes = result.get("presenter_notes", "")

                slide_results.append({
                    "slide": slide_num,
                    "original": text[:50] + "..." if len(text) > 50 else text,
                    "optimized": optimized[:50] + "..." if len(str(optimized)) > 50 else optimized,
                    "validation_passed": result.get("validation_passed", False)
                })

        self.stats["files_processed"] += 1
        self.stats["files_success"] += 1

        print(f"  [OK] Processed {len(slide_results)} text elements")
        print(f"       Truncations found: {truncations_found}, Fixed: {truncations_fixed}")
        print(f"       Words reduced: {words_reduced}")

        return {
            "success": True,
            "slides_processed": len(slide_results),
            "truncations_found": truncations_found,
            "truncations_fixed": truncations_fixed,
            "words_reduced": words_reduced,
            "results": slide_results
        }

    def process_directory(self, directory: Path) -> List[Dict[str, Any]]:
        """Process all PowerPoint files in a directory."""
        pptx_files = sorted(directory.glob("**/*.pptx"))

        if not pptx_files:
            print(f"No PowerPoint files found in: {directory}")
            return []

        print(f"Found {len(pptx_files)} PowerPoint files")
        print()

        results = []
        for pptx_file in pptx_files:
            result = self.process_pptx_file(pptx_file)
            results.append(result)
            print()

        return results

    def process_unit(self, unit_number: int) -> List[Dict[str, Any]]:
        """Process all PowerPoint files for a specific unit."""
        unit_folders = list(PRODUCTION_DIR.glob(f"Unit_{unit_number}_*"))

        if not unit_folders:
            print(f"Unit {unit_number} folder not found in production/")
            return []

        unit_folder = unit_folders[0]
        print(f"Processing Unit {unit_number}: {unit_folder.name}")
        print()

        return self.process_directory(unit_folder)

    def print_summary(self):
        """Print processing summary."""
        print()
        print("=" * 70)
        print("CONTENT OPTIMIZATION COMPLETE")
        print("=" * 70)
        print(f"Files processed:      {self.stats['files_processed']}")
        print(f"Files successful:     {self.stats['files_success']}")
        print(f"Files failed:         {self.stats['files_failed']}")
        print()
        print(f"Slides processed:     {self.stats['slides_processed']}")
        print(f"Truncations detected: {self.stats['truncations_detected']}")
        print(f"Truncations fixed:    {self.stats['truncations_fixed']}")
        print(f"Content condensed:    {self.stats['content_condensed']}")
        print()
        print(f"Validation passed:    {self.stats['validation_passed']}")
        print(f"Validation failed:    {self.stats['validation_failed']}")
        print(f"Total words reduced:  {self.stats['total_words_reduced']}")
        print()


def main():
    parser = argparse.ArgumentParser(
        description="Optimize slide content for student note-taking",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s --test                    Run test samples
  %(prog)s --unit 3                  Process Unit 3 (Romeo and Juliet)
  %(prog)s --path file.pptx          Process single file
  %(prog)s --dir /path/to/folder     Process all PPTX in directory
        """
    )

    parser.add_argument('--path', type=str, help='Path to single PowerPoint file')
    parser.add_argument('--dir', type=str, help='Directory containing PowerPoint files')
    parser.add_argument('--unit', type=int, choices=[1, 2, 3, 4],
                       help='Process specific unit (1-4)')
    parser.add_argument('--test', action='store_true', help='Run test samples')
    parser.add_argument('--verbose', '-v', action='store_true', help='Verbose output')

    args = parser.parse_args()

    # Initialize processor
    processor = ContentOptimizationProcessor(verbose=args.verbose)

    print("=" * 70)
    print("SLIDE CONTENT OPTIMIZATION")
    print("=" * 70)
    print(f"Started: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    # Process based on arguments
    if args.test:
        # Run test samples
        results = processor.run_test_samples()
        print()
        print("=" * 70)
        print("TEST SUMMARY")
        print("=" * 70)
        summary = results.get("summary", {})
        print(f"Total Tests:         {summary.get('total_tests', 0)}")
        print(f"Tests Passed:        {summary.get('total_passed', 0)}")
        print(f"Pass Rate:           {summary.get('pass_rate', 0)}%")
        print()
        print(f"Truncation Accuracy: {summary.get('truncation_accuracy', 0)}%")
        print(f"Condensation Accuracy: {summary.get('condensation_accuracy', 0)}%")
        print(f"Optimization Accuracy: {summary.get('optimization_accuracy', 0)}%")
        print()
        return 0 if summary.get('pass_rate', 0) >= 80 else 1

    elif args.path:
        # Single file
        pptx_path = Path(args.path)
        if not pptx_path.exists():
            print(f"ERROR: File not found: {pptx_path}")
            return 1
        processor.process_pptx_file(pptx_path)

    elif args.dir:
        # Directory
        directory = Path(args.dir)
        if not directory.exists():
            print(f"ERROR: Directory not found: {directory}")
            return 1
        processor.process_directory(directory)

    elif args.unit:
        # Specific unit
        processor.process_unit(args.unit)

    else:
        # Default: process all production folders
        print("Processing all PowerPoint files in production/")
        print()
        processor.process_directory(PRODUCTION_DIR)

    # Print summary
    processor.print_summary()

    return 0 if processor.stats["files_failed"] == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
