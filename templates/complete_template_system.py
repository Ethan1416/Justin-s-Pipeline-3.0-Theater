#!/usr/bin/env python3
"""
=============================================================================
NURSING LECTURE COMPLETE TEMPLATE SYSTEM
All 21 Slide Types + Generators + Validators
=============================================================================

TEMPLATE TYPES (21 Total):

VOLUME 1 - STANDARD TEMPLATES:
  1.  title         - Opening/section dividers
  2.  content       - Standard lecture content
  3.  two_column    - Side-by-side comparisons
  4.  quote         - Impactful statements
  5.  process       - Step-by-step flows
  6.  statistics    - Data visualization
  7.  case_study    - Clinical scenarios
  8.  summary       - Key points recap
  9.  qa            - Discussion closing
  10. image         - Visual content placeholder

VOLUME 2 - PREMIUM TEMPLATES:
  11. decision_tree - Clinical decision pathways
  12. timeline      - Milestone progression
  13. table         - Adaptive data grid
  14. canvas        - Freeform content space
  15. checklist     - Task tracking
  16. matrix        - Feature comparison grid
  17. hierarchy     - Organizational structure
  18. definition    - Glossary/terminology
  19. before_after  - Transformation comparison
  20. alert         - Critical information highlight

VOLUME 3 - SPECIALTY TEMPLATES:
  21. nclex_tip     - Clean content with NCLEX tip banner

FEATURES:
- Unified color psychology (alpha-wave cyan + dopamine coral)
- Consistent typography hierarchy (Georgia + Arial)
- 3D shadow depth effects
- Professional medical aesthetic
- Adaptive content layouts

USAGE:
    python complete_template_system.py                    # Generate all templates
    python complete_template_system.py --type timeline    # Generate specific type
    python complete_template_system.py --json data.json   # Batch from JSON
    python complete_template_system.py --validate         # Run validation

Author: Generated for nursing/medical education programs

=============================================================================
ADAPTIVE CONTENT INSTRUCTIONS
=============================================================================

Each template adapts to content within defined ranges. When your content
exceeds or falls below these ranges, follow these guidelines:

TEMPLATE CAPACITY REFERENCE:
┌─────────────────┬──────────┬──────────┬─────────────────────────────────┐
│ Template        │ Min      │ Max      │ Overflow Strategy               │
├─────────────────┼──────────┼──────────┼─────────────────────────────────┤
│ content         │ 1 bullet │ 4 bullets│ Split into multiple slides      │
│ two_column      │ 1 point  │ 5 points │ Split columns across slides     │
│ process         │ 2 steps  │ 5 steps  │ Use multiple process slides     │
│ statistics      │ 1 stat   │ 4 stats  │ Group related stats per slide   │
│ case_study      │ 1 item   │ 3 consid.│ Use follow-up discussion slide  │
│ summary         │ 1 point  │ 4 points │ Prioritize or split slides      │
│ qa              │ 0 prompts│ 3 prompts│ Additional prompts on next slide│
│ decision_tree   │ 2 branch │ 3 branch │ Create sub-decision slides      │
│ timeline        │ 2 events │ 6 events │ Split into phases/parts         │
│ table           │ 1 row    │ 6 rows   │ Continue table on next slide    │
│ checklist       │ 1 item   │ 8 items  │ Split into procedure parts      │
│ matrix          │ 1 row    │ 5 rows   │ Split comparison categories     │
│ hierarchy       │ 1 L1     │ 4 L1     │ Zoom into sub-hierarchies       │
│ definition      │ 1 term   │ 4 terms  │ Continue glossary on next slide │
│ before_after    │ 1 point  │ 5 points │ Group by category               │
│ alert           │ 0 actions│ 3 actions│ Most critical actions only      │
│ nclex_tip       │ 1 line   │ 9 lines  │ Split content across slides     │
└─────────────────┴──────────┴──────────┴─────────────────────────────────┘

AUTOMATIC ADAPTATIONS:
The system automatically handles content within ranges:

1. SPACING ADJUSTMENT
   - Elements reposition based on item count
   - Fewer items = more whitespace (cleaner look)
   - More items = tighter spacing (information dense)

2. SIZE SCALING
   - Font sizes may reduce slightly for more content
   - Card/box widths adjust to fit available space
   - Timeline markers spread evenly regardless of count

3. TRUNCATION (when exceeding max)
   - Content beyond max is silently dropped
   - Use multiple slides for complete content
   - Example: 7 timeline events → only first 6 render

CONTENT OVERFLOW STRATEGIES:

Strategy 1: SPLIT ACROSS SLIDES
   If you have 8 process steps:
   - Slide 1: Steps 1-4 with title "Process Overview (Part 1)"
   - Slide 2: Steps 5-8 with title "Process Overview (Part 2)"

Strategy 2: ZOOM INTO DETAIL
   If hierarchy has 6 level-1 nodes:
   - Slide 1: Show top 4 with "See detailed breakdown →"
   - Slide 2-3: Individual hierarchy slides for remaining

Strategy 3: CATEGORY GROUPING
   If you have 12 definitions:
   - Group by theme: "Cardiac Terms" (4), "Respiratory Terms" (4), etc.
   - Each group gets its own definition slide

Strategy 4: PROGRESSIVE DISCLOSURE
   If checklist has 15 items:
   - Slide 1: "Pre-Procedure Checklist" (items 1-8)
   - Slide 2: "During Procedure Checklist" (items 9-15)

CONTENT UNDERFLOW HANDLING:
When you have fewer items than optimal:

- 1-2 bullet points: Consider using quote or alert template instead
- 1 statistic: Use larger display, add context/comparison
- 2 timeline events: Add context cards or milestone markers
- 1 hierarchy level: Flatten to simple list or expand detail

MULTI-SLIDE GENERATION EXAMPLE:

```python
# Handling 10 process steps across 2 slides
all_steps = [...10 steps...]

# Slide 1: First 5 steps
slide1 = ProcessSlide()
slide1.create(ProcessContent(
    title="Medication Administration (Steps 1-5)",
    steps=all_steps[:5],
    slide_number="12"
))

# Slide 2: Remaining 5 steps  
slide2 = ProcessSlide()
slide2.create(ProcessContent(
    title="Medication Administration (Steps 6-10)",
    steps=all_steps[5:],
    slide_number="13"
))
```

DYNAMIC TEMPLATE SELECTION:
Choose template based on content characteristics:

| Content Type              | Best Template    | Alternative       |
|---------------------------|------------------|-------------------|
| 1-2 key points            | quote, alert     | content           |
| 3-4 key points            | content, summary | two_column        |
| Sequential steps (2-5)    | process          | timeline          |
| Sequential steps (6+)     | timeline (split) | multiple process  |
| Yes/No decision           | decision_tree    | two_column        |
| Multi-branch decision     | decision_tree    | hierarchy         |
| Feature comparison        | matrix           | table             |
| Data with status          | table            | statistics        |
| Before/after (1-5 pts)    | before_after     | two_column        |
| Organizational structure  | hierarchy        | process           |
| Vocabulary/glossary       | definition       | content           |
| Critical safety info      | alert            | quote             |
| Task tracking             | checklist        | content           |
| Open-ended workspace      | canvas           | image             |

=============================================================================
"""

import os
import sys
import json
from dataclasses import dataclass, field
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# =============================================================================
# TEMPLATE CAPACITY SPECIFICATIONS
# =============================================================================

TEMPLATE_SPECS = {
    'title': {
        'description': 'Opening/section divider slide',
        'capacity': {'fields': ['main_title', 'subtitle', 'instructor', 'date', 'institution']},
        'min_content': 1,
        'max_content': 1,
        'overflow_strategy': 'Not applicable - single slide template',
        'adaptive_features': ['None - fixed layout']
    },
    'content': {
        'description': 'Standard lecture content with sidebar',
        'capacity': {'bullet_points': 4, 'key_takeaways': 3},
        'min_content': 1,
        'max_content': 4,
        'overflow_strategy': 'Split bullets across multiple content slides',
        'adaptive_features': ['Bullet spacing adjusts to count', 'Sidebar scales with takeaways']
    },
    'two_column': {
        'description': 'Side-by-side comparison layout',
        'capacity': {'left_points': 5, 'right_points': 5},
        'min_content': 1,
        'max_content': 5,
        'overflow_strategy': 'Split into multiple comparison slides by category',
        'adaptive_features': ['Point spacing adjusts', 'Columns balance automatically']
    },
    'quote': {
        'description': 'Impactful statement with attribution',
        'capacity': {'quote': 1, 'attribution': 1, 'context': 1},
        'min_content': 1,
        'max_content': 1,
        'overflow_strategy': 'Use multiple quote slides for different quotes',
        'adaptive_features': ['Quote text wraps automatically']
    },
    'process': {
        'description': 'Step-by-step workflow visualization',
        'capacity': {'steps': 5},
        'min_content': 2,
        'max_content': 5,
        'overflow_strategy': 'Split into "Part 1/Part 2" process slides',
        'adaptive_features': ['Step width adjusts to count', 'Connectors scale', 'Cards resize']
    },
    'statistics': {
        'description': 'Data visualization with stat cards',
        'capacity': {'stats': 4},
        'min_content': 1,
        'max_content': 4,
        'overflow_strategy': 'Group related statistics per slide',
        'adaptive_features': ['Card width adjusts to count', 'Spacing rebalances']
    },
    'case_study': {
        'description': 'Clinical scenario with considerations',
        'capacity': {'considerations': 3},
        'min_content': 1,
        'max_content': 3,
        'overflow_strategy': 'Use follow-up discussion slide for additional considerations',
        'adaptive_features': ['Text areas wrap content', 'Consideration list adjusts']
    },
    'summary': {
        'description': 'Key takeaways in card grid',
        'capacity': {'points': 4},
        'min_content': 1,
        'max_content': 4,
        'overflow_strategy': 'Prioritize top 4 or split across slides',
        'adaptive_features': ['2x2 grid layout', 'Cards resize to content']
    },
    'qa': {
        'description': 'Q&A and discussion closing',
        'capacity': {'prompts': 3},
        'min_content': 0,
        'max_content': 3,
        'overflow_strategy': 'Most important prompts only; others verbal',
        'adaptive_features': ['Prompts stack vertically', 'Contact info optional']
    },
    'image': {
        'description': 'Visual focus with caption',
        'capacity': {'caption': 1},
        'min_content': 1,
        'max_content': 1,
        'overflow_strategy': 'One image per slide; use multiple for galleries',
        'adaptive_features': ['Caption wraps automatically']
    },
    'decision_tree': {
        'description': 'Clinical decision pathway with branches',
        'capacity': {'branches': 3},
        'min_content': 2,
        'max_content': 3,
        'overflow_strategy': 'Create sub-decision slides for complex pathways',
        'adaptive_features': ['Branch width adjusts', 'Connectors reposition', 'Labels scale']
    },
    'timeline': {
        'description': 'Milestone progression visualization',
        'capacity': {'events': 6},
        'min_content': 2,
        'max_content': 6,
        'overflow_strategy': 'Split into phases: "Phase 1 Timeline", "Phase 2 Timeline"',
        'adaptive_features': ['Event spacing calculated dynamically', 'Cards resize', 'Markers reposition']
    },
    'table': {
        'description': 'Adaptive data grid with status styling',
        'capacity': {'headers': 5, 'rows': 6},
        'min_content': 1,
        'max_content': 6,
        'overflow_strategy': 'Continue table on next slide with repeated headers',
        'adaptive_features': ['Column widths calculate from count', 'Row heights fixed', 'Status auto-styled']
    },
    'canvas': {
        'description': 'Freeform content workspace',
        'capacity': {'placeholder_text': 1},
        'min_content': 0,
        'max_content': 1,
        'overflow_strategy': 'Not applicable - open workspace',
        'adaptive_features': ['Grid dots for alignment', 'Corner accents frame space']
    },
    'checklist': {
        'description': 'Task tracking with progress',
        'capacity': {'items': 8},
        'min_content': 1,
        'max_content': 8,
        'overflow_strategy': 'Split by procedure phase or category',
        'adaptive_features': ['Progress bar calculates from checked', 'Items stack dynamically', 'Summary updates']
    },
    'matrix': {
        'description': 'Feature comparison grid',
        'capacity': {'row_headers': 5, 'col_headers': 4},
        'min_content': 1,
        'max_content': 5,
        'overflow_strategy': 'Split by feature category across slides',
        'adaptive_features': ['Cell sizes calculate', 'Check/X auto-styled', 'Alternating rows']
    },
    'hierarchy': {
        'description': 'Organizational structure tree',
        'capacity': {'level1': 4, 'level2_per_l1': 2},
        'min_content': 1,
        'max_content': 4,
        'overflow_strategy': 'Zoom into sub-hierarchies on separate slides',
        'adaptive_features': ['L1 spacing calculates', 'L2 positions under parents', 'Connectors redraw']
    },
    'definition': {
        'description': 'Glossary/terminology cards',
        'capacity': {'terms': 4},
        'min_content': 1,
        'max_content': 4,
        'overflow_strategy': 'Continue glossary on next slide',
        'adaptive_features': ['Cards stack vertically', 'Definition text wraps', 'Accent colors rotate']
    },
    'before_after': {
        'description': 'Transformation comparison',
        'capacity': {'before_points': 5, 'after_points': 5},
        'min_content': 1,
        'max_content': 5,
        'overflow_strategy': 'Group by transformation category',
        'adaptive_features': ['Points stack in columns', 'Markers align', 'Arrow centers']
    },
    'alert': {
        'description': 'Critical information highlight',
        'capacity': {'action_items': 3},
        'min_content': 0,
        'max_content': 3,
        'overflow_strategy': 'Most critical actions only; additional in follow-up',
        'adaptive_features': ['Alert type changes colors/icon', 'Actions stack', 'Message wraps']
    },
    'nclex_tip': {
        'description': 'Clean content slide with NCLEX tip banner',
        'capacity': {'content_lines': 9, 'nclex_tip': 1},
        'min_content': 1,
        'max_content': 9,
        'overflow_strategy': 'Split content across multiple nclex_tip slides',
        'adaptive_features': ['Content lines stack vertically', 'Tip banner fixed at bottom', 'Minimal embellishments', 'Custom copyright footer', '20pt body text']
    }
}


def get_template_capacity(template_type: str) -> dict:
    """Get capacity info for a template type"""
    return TEMPLATE_SPECS.get(template_type, {})


def check_content_fit(template_type: str, content: dict) -> dict:
    """
    Check if content fits within template capacity.
    Returns: {'fits': bool, 'warnings': list, 'suggestions': list}
    """
    spec = TEMPLATE_SPECS.get(template_type)
    if not spec:
        return {'fits': True, 'warnings': ['Unknown template type'], 'suggestions': []}
    
    warnings = []
    suggestions = []
    fits = True
    
    capacity = spec.get('capacity', {})
    
    for field, max_count in capacity.items():
        if isinstance(max_count, int):
            content_value = content.get(field, [])
            if isinstance(content_value, list):
                actual_count = len(content_value)
                if actual_count > max_count:
                    fits = False
                    overflow = actual_count - max_count
                    warnings.append(f"{field}: {actual_count} items exceeds max {max_count} (overflow: {overflow})")
                    suggestions.append(f"Strategy: {spec.get('overflow_strategy', 'Split across slides')}")
                elif actual_count < spec.get('min_content', 0):
                    warnings.append(f"{field}: {actual_count} items below recommended minimum")
                    suggestions.append("Consider using a simpler template or adding content")
    
    return {'fits': fits, 'warnings': warnings, 'suggestions': suggestions}


def suggest_template(content_characteristics: dict) -> list:
    """
    Suggest best template(s) based on content characteristics.
    
    Args:
        content_characteristics: dict with keys like:
            - 'item_count': number of main items
            - 'has_comparison': bool
            - 'is_sequential': bool
            - 'is_hierarchical': bool
            - 'has_data': bool
            - 'is_critical': bool
            - 'content_type': 'text'|'data'|'process'|'decision'
    
    Returns:
        List of (template_type, confidence, reason) tuples
    """
    suggestions = []
    
    item_count = content_characteristics.get('item_count', 0)
    has_comparison = content_characteristics.get('has_comparison', False)
    is_sequential = content_characteristics.get('is_sequential', False)
    is_hierarchical = content_characteristics.get('is_hierarchical', False)
    has_data = content_characteristics.get('has_data', False)
    is_critical = content_characteristics.get('is_critical', False)
    content_type = content_characteristics.get('content_type', 'text')
    
    # Critical information
    if is_critical:
        suggestions.append(('alert', 0.95, 'Critical information needs high visibility'))
    
    # Decision/branching content
    if content_type == 'decision':
        if item_count <= 3:
            suggestions.append(('decision_tree', 0.9, 'Branching decisions fit decision tree'))
        else:
            suggestions.append(('decision_tree', 0.7, 'May need multiple decision slides'))
            suggestions.append(('hierarchy', 0.6, 'Complex decisions can use hierarchy'))
    
    # Sequential/process content
    if is_sequential:
        if item_count <= 5:
            suggestions.append(('process', 0.9, 'Sequential steps fit process template'))
        elif item_count <= 6:
            suggestions.append(('timeline', 0.9, 'More steps work well in timeline'))
        else:
            suggestions.append(('timeline', 0.7, 'Split into multiple timeline slides'))
            suggestions.append(('process', 0.6, 'Or multiple process slides'))
    
    # Comparison content
    if has_comparison:
        if item_count <= 5:
            suggestions.append(('two_column', 0.85, 'Direct comparison fits two-column'))
            suggestions.append(('before_after', 0.8, 'If showing transformation'))
        suggestions.append(('matrix', 0.75, 'Multi-feature comparison'))
    
    # Hierarchical content
    if is_hierarchical:
        if item_count <= 4:
            suggestions.append(('hierarchy', 0.9, 'Organizational structure'))
        else:
            suggestions.append(('hierarchy', 0.6, 'Split into sub-hierarchies'))
    
    # Data-heavy content
    if has_data:
        if item_count <= 4:
            suggestions.append(('statistics', 0.85, 'Key metrics visualization'))
        suggestions.append(('table', 0.8, 'Structured data display'))
    
    # Default text content
    if content_type == 'text' and not suggestions:
        if item_count <= 1:
            suggestions.append(('quote', 0.7, 'Single impactful statement'))
        elif item_count <= 4:
            suggestions.append(('content', 0.85, 'Standard content layout'))
        else:
            suggestions.append(('content', 0.7, 'Split across multiple slides'))
            suggestions.append(('summary', 0.65, 'Key points format'))
    
    # Sort by confidence
    suggestions.sort(key=lambda x: x[1], reverse=True)
    return suggestions[:3]


def split_content_for_template(template_type: str, content: dict) -> list:
    """
    Split content that exceeds template capacity into multiple slide contents.
    
    Returns: List of content dicts, each fitting template capacity
    """
    spec = TEMPLATE_SPECS.get(template_type)
    if not spec:
        return [content]
    
    result = []
    capacity = spec.get('capacity', {})
    
    # Find the main list field that might overflow
    list_fields = {k: v for k, v in capacity.items() if isinstance(v, int)}
    
    if not list_fields:
        return [content]
    
    # Get the primary list field (usually the one with most capacity)
    primary_field = max(list_fields.keys(), key=lambda k: list_fields[k])
    max_items = list_fields[primary_field]
    
    content_list = content.get(primary_field, [])
    
    if not isinstance(content_list, list) or len(content_list) <= max_items:
        return [content]
    
    # Split into chunks
    chunks = [content_list[i:i + max_items] for i in range(0, len(content_list), max_items)]
    
    for i, chunk in enumerate(chunks):
        new_content = content.copy()
        new_content[primary_field] = chunk
        
        # Update title to indicate part number if multiple
        if len(chunks) > 1 and 'title' in new_content:
            new_content['title'] = f"{content.get('title', 'Content')} (Part {i + 1}/{len(chunks)})"
        
        # Update slide number if present
        if 'slide_number' in new_content:
            base_num = int(content.get('slide_number', '01'))
            new_content['slide_number'] = f"{base_num + i:02d}"
        
        result.append(new_content)
    
    return result


# =============================================================================
# COLOR PALETTE - Color Psychology Optimized
# =============================================================================

class Colors:
    """
    Complete color palette with color psychology principles:
    - Cyan spectrum: Alpha-wave activation (480-520nm wavelength)
    - Coral spectrum: Dopamine trigger for attention anchoring
    - Navy spectrum: Authority and calm
    - Sand spectrum: Biophilic grounding
    """
    
    # Deep ocean navy (headers, dark panels)
    OCEAN_DEEP = RGBColor(0x0c, 0x19, 0x29)
    OCEAN_MID = RGBColor(0x13, 0x32, 0x46)
    OCEAN_LIGHT = RGBColor(0x1a, 0x4a, 0x5e)
    OCEAN_SUBTLE = RGBColor(0x1e, 0x40, 0x52)
    
    # Cyan spectrum (alpha-wave 480-520nm)
    CYAN_BRIGHT = RGBColor(0x22, 0xd3, 0xee)
    CYAN_PRIMARY = RGBColor(0x06, 0xb6, 0xd4)
    CYAN_SOFT = RGBColor(0x67, 0xe8, 0xf9)
    CYAN_PALE = RGBColor(0xa5, 0xf3, 0xfc)
    CYAN_MUTED = RGBColor(0x16, 0x4e, 0x63)
    CYAN_DARK = RGBColor(0x0e, 0x7a, 0x90)
    
    # Warm spectrum (dopamine activation)
    CORAL_BRIGHT = RGBColor(0xfb, 0x92, 0x3c)
    CORAL_PRIMARY = RGBColor(0xf9, 0x73, 0x16)
    CORAL_SOFT = RGBColor(0xfd, 0xba, 0x74)
    CORAL_PALE = RGBColor(0xfe, 0xd7, 0xaa)
    CORAL_DEEP = RGBColor(0xea, 0x58, 0x0c)
    
    # Biophilic neutrals (grounding)
    SAND_LIGHT = RGBColor(0xfe, 0xfc, 0xf8)
    SAND_WARM = RGBColor(0xfa, 0xf5, 0xed)
    SAND_MID = RGBColor(0xf0, 0xe9, 0xdd)
    SAND_DARK = RGBColor(0xe8, 0xe0, 0xd4)
    SAND_DEEP = RGBColor(0xd6, 0xd0, 0xc4)
    
    # Text hierarchy
    TEXT_DARK = RGBColor(0x0f, 0x17, 0x2a)
    TEXT_PRIMARY = RGBColor(0x1e, 0x3a, 0x4c)
    TEXT_SECONDARY = RGBColor(0x3d, 0x5a, 0x6c)
    TEXT_MUTED = RGBColor(0x6b, 0x84, 0x94)
    TEXT_LIGHT = RGBColor(0x94, 0xa3, 0xb8)
    TEXT_PALE = RGBColor(0xcb, 0xd5, 0xe1)
    
    # Status colors
    SUCCESS = RGBColor(0x10, 0xb9, 0x81)
    SUCCESS_SOFT = RGBColor(0x6e, 0xe7, 0xb7)
    WARNING = RGBColor(0xf5, 0x9e, 0x0b)
    WARNING_SOFT = RGBColor(0xfc, 0xd3, 0x4d)
    DANGER = RGBColor(0xef, 0x44, 0x44)
    DANGER_SOFT = RGBColor(0xfc, 0xa5, 0xa5)
    
    WHITE = RGBColor(0xff, 0xff, 0xff)
    BLACK = RGBColor(0x00, 0x00, 0x00)


# =============================================================================
# DIMENSIONS & SCALING
# =============================================================================

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
SCALE_X = 13.333 / 960
SCALE_Y = 7.5 / 540

def px(val): 
    """Convert HTML-style pixels to PowerPoint inches (width-based)"""
    return Inches(val * SCALE_X)

def py(val): 
    """Convert HTML-style pixels to PowerPoint inches (height-based)"""
    return Inches(val * SCALE_Y)


# =============================================================================
# DATA STRUCTURES FOR ALL 20 TEMPLATE TYPES
# =============================================================================

# --- VOLUME 1: STANDARD TEMPLATES ---

@dataclass
class TitleContent:
    main_title: str = "Lecture Title"
    subtitle: str = "Course Name · Module Overview"
    instructor: str = "Instructor Name, RN, MSN"
    date: str = "Spring 2025"
    institution: str = "School of Nursing"

@dataclass
class ContentContent:
    title: str = "Slide Title"
    subtitle: str = "Module · Section"
    section_tag: str = "Learning Objectives"
    section_title: str = "Section Heading"
    intro_text: str = "Introduction text goes here."
    bullet_points: List[str] = field(default_factory=lambda: ["Point 1", "Point 2", "Point 3", "Point 4"])
    key_takeaways: List[str] = field(default_factory=lambda: ["Takeaway 1", "Takeaway 2", "Takeaway 3"])
    focus_phrase: str = "Key memorable phrase"
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class TwoColumnContent:
    title: str = "Comparison"
    subtitle: str = "Module · Section"
    left_title: str = "Column A"
    left_points: List[str] = field(default_factory=list)
    right_title: str = "Column B"
    right_points: List[str] = field(default_factory=list)
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class QuoteContent:
    quote: str = "Inspirational quote here."
    attribution: str = "— Author Name"
    context: str = ""
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class ProcessContent:
    title: str = "Process Overview"
    subtitle: str = "Module · Section"
    steps: List[Dict[str, str]] = field(default_factory=list)
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class StatisticsContent:
    title: str = "Key Statistics"
    subtitle: str = "Module · Section"
    stats: List[Dict[str, str]] = field(default_factory=list)
    source: str = ""
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class CaseStudyContent:
    title: str = "Case Study"
    subtitle: str = "Clinical Application"
    patient_info: str = ""
    presentation: str = ""
    question: str = ""
    considerations: List[str] = field(default_factory=list)
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class SummaryContent:
    title: str = "Key Takeaways"
    subtitle: str = "Module Summary"
    points: List[Dict[str, str]] = field(default_factory=list)
    next_topic: str = ""
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class QAContent:
    title: str = "Questions & Discussion"
    prompts: List[str] = field(default_factory=list)
    contact: str = ""
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class ImageContent:
    title: str = "Visual Focus"
    subtitle: str = "Module · Section"
    caption: str = "Image caption"
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

# --- VOLUME 2: PREMIUM TEMPLATES ---

@dataclass
class DecisionTreeContent:
    title: str = "Clinical Decision Pathway"
    subtitle: str = "Module · Decision Support"
    root_question: str = "Initial Assessment Question?"
    branches: List[Dict] = field(default_factory=lambda: [
        {"label": "Yes", "outcome": "Pathway A", "action": "Proceed with intervention"},
        {"label": "No", "outcome": "Pathway B", "action": "Consider alternatives"}
    ])
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class TimelineContent:
    title: str = "Clinical Timeline"
    subtitle: str = "Module · Progression"
    events: List[Dict] = field(default_factory=lambda: [
        {"time": "0h", "title": "Admission", "description": "Initial assessment"},
        {"time": "2h", "title": "Intervention", "description": "Treatment begins"},
        {"time": "6h", "title": "Reassessment", "description": "Evaluate response"},
        {"time": "24h", "title": "Discharge", "description": "Patient education"}
    ])
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class TableContent:
    title: str = "Data Overview"
    subtitle: str = "Module · Reference"
    headers: List[str] = field(default_factory=lambda: ["Category", "Value", "Range", "Status"])
    rows: List[List[str]] = field(default_factory=lambda: [
        ["Heart Rate", "72 bpm", "60-100", "Normal"],
        ["Blood Pressure", "120/80", "<120/80", "Normal"],
        ["Temperature", "98.6°F", "97.8-99.1", "Normal"]
    ])
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class CanvasContent:
    title: str = "Content Area"
    subtitle: str = "Module · Custom Layout"
    placeholder_text: str = "Insert your content here"
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class ChecklistContent:
    title: str = "Assessment Checklist"
    subtitle: str = "Module · Procedure"
    items: List[Dict] = field(default_factory=lambda: [
        {"text": "Verify patient identity", "checked": True},
        {"text": "Review allergies", "checked": True},
        {"text": "Assess vital signs", "checked": False},
        {"text": "Document findings", "checked": False}
    ])
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class MatrixContent:
    title: str = "Comparison Matrix"
    subtitle: str = "Module · Analysis"
    row_headers: List[str] = field(default_factory=lambda: ["Feature A", "Feature B", "Feature C"])
    col_headers: List[str] = field(default_factory=lambda: ["Option 1", "Option 2", "Option 3"])
    data: List[List[str]] = field(default_factory=lambda: [
        ["✓", "✓", "—"],
        ["✓", "—", "✓"],
        ["—", "✓", "✓"]
    ])
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class HierarchyContent:
    title: str = "Organizational Structure"
    subtitle: str = "Module · Framework"
    root: str = "Primary Concept"
    level1: List[str] = field(default_factory=lambda: ["Category A", "Category B", "Category C"])
    level2: List[List[str]] = field(default_factory=lambda: [
        ["Sub A1", "Sub A2"],
        ["Sub B1", "Sub B2"],
        ["Sub C1", "Sub C2"]
    ])
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class DefinitionContent:
    title: str = "Key Terminology"
    subtitle: str = "Module · Definitions"
    terms: List[Dict] = field(default_factory=lambda: [
        {"term": "Term One", "definition": "The detailed definition and explanation."},
        {"term": "Term Two", "definition": "Another important definition."},
        {"term": "Term Three", "definition": "Additional terminology explanation."}
    ])
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class BeforeAfterContent:
    title: str = "Transformation"
    subtitle: str = "Module · Comparison"
    before_title: str = "Before"
    before_points: List[str] = field(default_factory=lambda: ["Previous state", "Old approach", "Initial condition"])
    after_title: str = "After"
    after_points: List[str] = field(default_factory=lambda: ["Improved state", "New approach", "Desired outcome"])
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"

@dataclass
class AlertContent:
    title: str = "Critical Alert"
    subtitle: str = "Module · Safety"
    alert_type: str = "warning"  # warning, danger, success, info
    headline: str = "Important Notice"
    message: str = "Critical information that requires immediate attention."
    action_items: List[str] = field(default_factory=lambda: ["Action step one", "Action step two"])
    slide_number: str = "01"
    footer_text: str = "Nursing Program · 2025"


@dataclass
class NclexTipContent:
    """Template 21: Clean content slide with NCLEX tip banner"""
    title: str = "Slide Title"
    content_lines: List[str] = field(default_factory=lambda: [
        "Main content point or paragraph goes here.",
        "Additional information to support learning.",
        "Key concept explanation for student understanding."
    ])
    nclex_tip: str = "Remember: This is a high-yield NCLEX concept. Focus on the nursing priority and patient safety."
    slide_number: str = "01"


# =============================================================================
# BASE TEMPLATE CLASS
# =============================================================================

class TemplateBase:
    """Base class with all shared rendering methods"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.blank = self.prs.slide_layouts[6]
    
    def shape(self, slide, stype, l, t, w, h, fill=None, grad=None):
        """Add shape with styling"""
        s = slide.shapes.add_shape(stype, l, t, w, h)
        if fill:
            s.fill.solid()
            s.fill.fore_color.rgb = fill
        elif grad:
            s.fill.gradient()
            s.fill.gradient_angle = grad.get('angle', 135)
            stops = grad.get('stops', [])
            if len(stops) >= 2:
                s.fill.gradient_stops[0].color.rgb = stops[0]
                s.fill.gradient_stops[1].color.rgb = stops[1]
        else:
            s.fill.background()
        s.line.fill.background()
        return s
    
    def rounded(self, slide, l, t, w, h, fill=None, grad=None, r=0.08):
        """Add rounded rectangle"""
        s = self.shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill, grad)
        s.adjustments[0] = r
        return s
    
    def text(self, slide, l, t, w, h, txt, font="Arial", sz=12, clr=None,
             bold=False, italic=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
        """Add text box"""
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = txt
        p.alignment = align
        p.font.name = font
        p.font.size = Pt(sz)
        p.font.color.rgb = clr or Colors.TEXT_PRIMARY
        p.font.bold = bold
        p.font.italic = italic
        tf.vertical_anchor = valign
        return tb
    
    def bg(self, slide, c1, c2):
        """Add gradient background"""
        fill = slide.background.fill
        fill.gradient()
        fill.gradient_angle = 175
        fill.gradient_stops[0].color.rgb = c1
        fill.gradient_stops[1].color.rgb = c2
    
    def header(self, slide, title, subtitle, compact=False):
        """Standard header with medical icon"""
        h = py(65) if compact else py(85)
        
        # Main header gradient
        self.shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, h,
                   grad={'angle': 145, 'stops': [Colors.OCEAN_DEEP, Colors.OCEAN_MID]})
        
        # Top warm accent
        self.shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0), Inches(11.733), Pt(2),
                   fill=Colors.CORAL_BRIGHT)
        
        # Bottom cyan accent (layered glow)
        self.shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), h - Pt(4), SLIDE_WIDTH, Pt(4),
                   fill=Colors.CYAN_BRIGHT)
        self.shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), h - Pt(2), SLIDE_WIDTH, Pt(2),
                   fill=Colors.CYAN_SOFT)
        
        # Medical icon
        isz = py(42) if compact else py(50)
        itop = (h - isz) / 2
        ileft = px(28)
        
        # Icon shadow + main + highlight
        self.rounded(slide, ileft + Pt(2), itop + Pt(2), isz, isz, fill=Colors.CYAN_DARK, r=0.18)
        self.rounded(slide, ileft, itop, isz, isz, fill=Colors.CYAN_BRIGHT, r=0.18)
        self.shape(slide, MSO_SHAPE.RECTANGLE, ileft + Pt(4), itop + Pt(4), isz - Pt(8), Pt(2),
                   fill=Colors.CYAN_PALE)
        self.text(slide, ileft, itop, isz, isz, "+", "Arial", 22 if compact else 26,
                  Colors.OCEAN_DEEP, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Title
        tl = px(95)
        self.text(slide, tl, py(16) if compact else py(20), px(700), py(30), title,
                  "Georgia", 19 if compact else 21, Colors.WHITE, True)
        
        # Subtitle
        self.text(slide, tl, py(40) if compact else py(50), px(700), py(16), subtitle.upper(),
                  "Arial", 7 if compact else 8, Colors.CYAN_SOFT, True)
        
        # Decorative flow lines
        for i, (w, c) in enumerate([
            (80, Colors.CORAL_BRIGHT), (55, Colors.CYAN_SOFT),
            (35, Colors.CORAL_SOFT), (20, Colors.CYAN_PALE)
        ]):
            self.shape(slide, MSO_SHAPE.RECTANGLE, px(880 - w), py(20 + i * 9) if compact else py(26 + i * 11),
                       px(w), Pt(2), fill=c)
        
        return h
    
    def footer(self, slide, txt, num):
        """Standard footer with slide number badge"""
        top = SLIDE_HEIGHT - py(30)
        
        # Subtle divider line
        self.shape(slide, MSO_SHAPE.RECTANGLE, px(28), top - py(4), px(904), Pt(1),
                   fill=Colors.SAND_DEEP)
        
        # Footer text
        self.text(slide, px(32), top, px(300), py(22), txt.upper(),
                  "Arial", 7, Colors.TEXT_MUTED, True)
        
        # Slide number badge
        nsz = py(26)
        nleft = SLIDE_WIDTH - px(56)
        self.rounded(slide, nleft + Pt(1), top + Pt(1), nsz, nsz, fill=Colors.OCEAN_DEEP, r=0.2)
        self.rounded(slide, nleft, top, nsz, nsz,
                     grad={'angle': 145, 'stops': [Colors.OCEAN_MID, Colors.OCEAN_DEEP]}, r=0.2)
        self.text(slide, nleft, top, nsz, nsz, num,
                  "Arial", 9, Colors.CYAN_SOFT, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    
    def bullet(self, slide, l, t, clr=None):
        """Diamond bullet marker"""
        self.shape(slide, MSO_SHAPE.DIAMOND, l, t + Pt(4), Pt(10), Pt(10), fill=clr or Colors.CORAL_BRIGHT)
    
    def badge(self, slide, l, t, num, sz=None):
        """Numbered badge"""
        sz = sz or py(26)
        self.rounded(slide, l, t, sz, sz, fill=Colors.CYAN_BRIGHT, r=0.2)
        self.text(slide, l, t, sz, sz, str(num), "Arial", 10, Colors.OCEAN_DEEP, True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    
    def connector(self, slide, x1, y1, x2, y2, color=None, width=2):
        """Draw connector line"""
        if abs(x2 - x1) > abs(y2 - y1):
            self.shape(slide, MSO_SHAPE.RECTANGLE, x1, y1, x2 - x1, Pt(width), fill=color or Colors.CYAN_BRIGHT)
        else:
            self.shape(slide, MSO_SHAPE.RECTANGLE, x1, y1, Pt(width), y2 - y1, fill=color or Colors.CYAN_BRIGHT)
    
    def arrow_right(self, slide, l, t, w, h, color=None):
        """Right arrow shape"""
        return self.shape(slide, MSO_SHAPE.RIGHT_ARROW, l, t, w, h, fill=color or Colors.CYAN_BRIGHT)
    
    def arrow_down(self, slide, l, t, w, h, color=None):
        """Down arrow shape"""
        return self.shape(slide, MSO_SHAPE.DOWN_ARROW, l, t, w, h, fill=color or Colors.CYAN_BRIGHT)
    
    def save(self, fn):
        """Save presentation"""
        self.prs.save(fn)
        print(f"✅ Saved: {fn}")


# =============================================================================
# VOLUME 1: STANDARD TEMPLATE IMPLEMENTATIONS (1-10)
# =============================================================================

class TitleSlide(TemplateBase):
    """1. Title/Opening Slide"""
    
    def create(self, c: TitleContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.OCEAN_DEEP, Colors.OCEAN_MID)
        
        # Corner accents
        self.shape(s, MSO_SHAPE.RIGHT_TRIANGLE, SLIDE_WIDTH - px(200), Inches(0), px(200), py(150), fill=Colors.CYAN_MUTED)
        self.shape(s, MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), SLIDE_HEIGHT - py(120), px(180), py(120), fill=Colors.CYAN_MUTED)
        
        # Large medical icon
        isz = py(100)
        self.rounded(s, px(80), (SLIDE_HEIGHT - isz)/2 - py(20), isz, isz, fill=Colors.CYAN_BRIGHT, r=0.12)
        self.text(s, px(80), (SLIDE_HEIGHT - isz)/2 - py(20), isz, isz, "+", "Arial", 60, Colors.OCEAN_DEEP, True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Title
        tl = px(220)
        self.text(s, tl, py(180), px(680), py(80), c.main_title, "Georgia", 44, Colors.WHITE, True)
        self.shape(s, MSO_SHAPE.RECTANGLE, tl, py(270), px(120), Pt(4), fill=Colors.CYAN_BRIGHT)
        self.text(s, tl, py(290), px(600), py(36), c.subtitle, "Arial", 14, Colors.CYAN_SOFT)
        
        # Info
        it = SLIDE_HEIGHT - py(100)
        self.text(s, tl, it, px(300), py(24), c.instructor, "Arial", 11, Colors.TEXT_LIGHT)
        self.text(s, tl, it + py(28), px(300), py(24), f"{c.institution} · {c.date}", "Arial", 10, Colors.TEXT_MUTED)
        
        # Flow lines
        for i, (w, clr) in enumerate([(100, Colors.CORAL_BRIGHT), (70, Colors.CYAN_SOFT), (45, Colors.CORAL_SOFT)]):
            self.shape(s, MSO_SHAPE.RECTANGLE, SLIDE_WIDTH - px(40) - px(w), py(400 + i*16), px(w), Pt(2), fill=clr)


class ContentSlide(TemplateBase):
    """2. Standard Content Slide"""
    
    def create(self, c: ContentContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle)
        
        ct = hh + py(25)
        cl = px(28)
        mw, mh = px(680), py(365)
        
        # Main panel
        self.rounded(s, cl, ct, mw, mh, fill=Colors.WHITE, r=0.04)
        self.rounded(s, cl, ct + py(16), Pt(4), mh - py(32), fill=Colors.CYAN_BRIGHT, r=0.5)
        
        # Section tag + title
        tl, tt = cl + px(32), ct + py(28)
        self.shape(s, MSO_SHAPE.OVAL, tl, tt + Pt(3), Pt(6), Pt(6), fill=Colors.CYAN_BRIGHT)
        self.text(s, tl + Pt(12), tt, px(200), py(16), c.section_tag.upper(), "Arial", 8, Colors.CYAN_BRIGHT, True)
        self.text(s, tl, tt + py(18), px(600), py(30), c.section_title, "Georgia", 18, Colors.TEXT_PRIMARY, True)
        self.shape(s, MSO_SHAPE.RECTANGLE, tl, tt + py(52), px(600), Pt(1), fill=Colors.CYAN_SOFT)
        
        # Intro text
        if c.intro_text:
            self.text(s, tl, tt + py(64), px(600), py(48), c.intro_text, "Arial", 11, Colors.TEXT_SECONDARY)
        
        # Bullets
        bt = tt + py(120)
        for i, pt in enumerate(c.bullet_points[:4]):
            y = bt + py(32) * i
            self.bullet(s, tl, y)
            self.text(s, tl + px(28), y, px(580), py(28), pt, "Arial", 11, Colors.TEXT_PRIMARY)
        
        # Side panel
        sl = cl + mw + px(24)
        sw = px(200)
        self.rounded(s, sl, ct, sw, mh, grad={'angle': 165, 'stops': [Colors.OCEAN_MID, Colors.OCEAN_DEEP]}, r=0.04)
        
        spl, spt = sl + px(20), ct + py(24)
        self.text(s, spl, spt, px(160), py(24), "Key Takeaways", "Georgia", 14, Colors.WHITE, True)
        self.shape(s, MSO_SHAPE.RECTANGLE, spl, spt + py(28), px(160), Pt(1), fill=Colors.CYAN_SOFT)
        
        kt = spt + py(44)
        for i, item in enumerate(c.key_takeaways[:3]):
            y = kt + py(44) * i
            self.badge(s, spl, y, i + 1)
            self.text(s, spl + py(32), y, px(120), py(40), item, "Arial", 10, Colors.WHITE)
        
        # Focus anchor
        if c.focus_phrase:
            ft = kt + py(140)
            self.rounded(s, spl, ft, px(160), py(72), fill=RGBColor(0x2d, 0x1f, 0x14), r=0.08)
            self.shape(s, MSO_SHAPE.RECTANGLE, spl, ft, Pt(3), py(72), fill=Colors.CORAL_BRIGHT)
            self.text(s, spl + px(12), ft + py(10), px(136), py(12), "REMEMBER", "Arial", 7, Colors.CORAL_SOFT, True)
            self.text(s, spl + px(12), ft + py(24), px(136), py(44), f'"{c.focus_phrase}"', "Arial", 9, Colors.WHITE, italic=True)
        
        self.footer(s, c.footer_text, c.slide_number)


class TwoColumnSlide(TemplateBase):
    """3. Two-Column Comparison"""
    
    def create(self, c: TwoColumnContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle)
        
        ct = hh + py(25)
        cw, ch = px(435), py(365)
        gap = px(24)
        
        # Left column
        lx = px(28)
        self.rounded(s, lx, ct, cw, ch, fill=Colors.WHITE, r=0.04)
        self.shape(s, MSO_SHAPE.RECTANGLE, lx, ct, cw, Pt(4), fill=Colors.CYAN_BRIGHT)
        self.text(s, lx + px(24), ct + py(20), px(380), py(28), c.left_title, "Georgia", 18, Colors.TEXT_PRIMARY, True)
        for i, pt in enumerate(c.left_points[:5]):
            y = ct + py(60) + py(28) * i
            self.bullet(s, lx + px(24), y, Colors.CYAN_BRIGHT)
            self.text(s, lx + px(52), y, px(360), py(24), pt, "Arial", 11, Colors.TEXT_PRIMARY)
        
        # Right column
        rx = lx + cw + gap
        self.rounded(s, rx, ct, cw, ch, fill=Colors.WHITE, r=0.04)
        self.shape(s, MSO_SHAPE.RECTANGLE, rx, ct, cw, Pt(4), fill=Colors.CORAL_BRIGHT)
        self.text(s, rx + px(24), ct + py(20), px(380), py(28), c.right_title, "Georgia", 18, Colors.TEXT_PRIMARY, True)
        for i, pt in enumerate(c.right_points[:5]):
            y = ct + py(60) + py(28) * i
            self.bullet(s, rx + px(24), y, Colors.CORAL_BRIGHT)
            self.text(s, rx + px(52), y, px(360), py(24), pt, "Arial", 11, Colors.TEXT_PRIMARY)
        
        # VS badge
        vsz = py(40)
        self.shape(s, MSO_SHAPE.OVAL, lx + cw + gap/2 - vsz/2, ct + ch/2 - vsz/2, vsz, vsz, fill=Colors.OCEAN_DEEP)
        self.text(s, lx + cw + gap/2 - vsz/2, ct + ch/2 - vsz/2, vsz, vsz, "VS", "Arial", 10, Colors.CYAN_SOFT, True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        self.footer(s, c.footer_text, c.slide_number)


class QuoteSlide(TemplateBase):
    """4. Quote/Highlight Slide"""
    
    def create(self, c: QuoteContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.OCEAN_DEEP, Colors.OCEAN_MID)
        
        # Large quotation mark
        self.text(s, px(60), py(100), px(200), py(200), "\u201c", "Georgia", 180, Colors.CYAN_MUTED, True)
        
        # Quote text
        self.text(s, px(120), py(180), px(720), py(160), c.quote, "Georgia", 28, Colors.WHITE, italic=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Accent line
        self.shape(s, MSO_SHAPE.RECTANGLE, px(400), py(360), px(160), Pt(3), fill=Colors.CYAN_BRIGHT)
        
        # Attribution
        self.text(s, px(120), py(380), px(720), py(32), c.attribution, "Arial", 12, Colors.CYAN_SOFT, align=PP_ALIGN.CENTER)
        
        if c.context:
            self.text(s, px(120), py(420), px(720), py(48), c.context, "Arial", 11, Colors.TEXT_LIGHT, align=PP_ALIGN.CENTER)
        
        # Corner accents
        self.shape(s, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), px(80), Pt(4), fill=Colors.CORAL_BRIGHT)
        self.shape(s, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Pt(4), py(60), fill=Colors.CORAL_BRIGHT)
        self.shape(s, MSO_SHAPE.RECTANGLE, SLIDE_WIDTH - px(80), SLIDE_HEIGHT - Pt(4), px(80), Pt(4), fill=Colors.CYAN_BRIGHT)
        self.shape(s, MSO_SHAPE.RECTANGLE, SLIDE_WIDTH - Pt(4), SLIDE_HEIGHT - py(60), Pt(4), py(60), fill=Colors.CYAN_BRIGHT)
        
        self.footer(s, c.footer_text, c.slide_number)


class ProcessSlide(TemplateBase):
    """5. Process/Timeline Slide"""
    
    def create(self, c: ProcessContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        steps = c.steps[:5] or [{"title": f"Step {i+1}", "description": ""} for i in range(4)]
        n = len(steps)
        sw = px(160)
        tw = sw * n + px(40) * (n - 1)
        sx = (SLIDE_WIDTH - tw) / 2
        st = hh + py(60)
        
        for i, step in enumerate(steps):
            x = sx + (sw + px(40)) * i
            csz = py(50)
            
            # Step number circle
            self.shape(s, MSO_SHAPE.OVAL, x + sw/2 - csz/2, st, csz, csz, fill=Colors.CYAN_BRIGHT)
            self.text(s, x + sw/2 - csz/2, st, csz, csz, str(i+1), "Arial", 20, Colors.OCEAN_DEEP, True,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            
            # Arrow connector
            if i < n - 1:
                self.arrow_right(s, x + sw + px(5), st + csz/2 - py(8), px(30), py(16), Colors.CORAL_BRIGHT)
            
            # Step card
            cdt = st + csz + py(16)
            self.rounded(s, x, cdt, sw, py(200), fill=Colors.WHITE, r=0.06)
            self.text(s, x + px(12), cdt + py(16), sw - px(24), py(48), step.get("title", ""), "Georgia", 12, Colors.TEXT_PRIMARY, True, align=PP_ALIGN.CENTER)
            self.text(s, x + px(12), cdt + py(70), sw - px(24), py(120), step.get("description", ""), "Arial", 10, Colors.TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        
        self.footer(s, c.footer_text, c.slide_number)


class StatisticsSlide(TemplateBase):
    """6. Statistics/Data Slide"""
    
    def create(self, c: StatisticsContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        stats = c.stats[:4] or [{"value": "0", "label": "Metric"} for _ in range(3)]
        n = len(stats)
        sw, sh = px(200), py(180)
        gap = px(24)
        tw = sw * n + gap * (n - 1)
        sx = (SLIDE_WIDTH - tw) / 2
        st = hh + py(50)
        colors = [Colors.CYAN_BRIGHT, Colors.CORAL_BRIGHT, Colors.CYAN_DARK, Colors.CORAL_DEEP]
        
        for i, stat in enumerate(stats):
            x = sx + (sw + gap) * i
            ac = colors[i % len(colors)]
            
            self.rounded(s, x, st, sw, sh, fill=Colors.WHITE, r=0.06)
            self.shape(s, MSO_SHAPE.RECTANGLE, x, st, sw, Pt(5), fill=ac)
            self.text(s, x, st + py(30), sw, py(60), stat.get("value", "0"), "Georgia", 42, ac, True, align=PP_ALIGN.CENTER)
            self.text(s, x + px(12), st + py(100), sw - px(24), py(36), stat.get("label", ""), "Arial", 12, Colors.TEXT_PRIMARY, True, align=PP_ALIGN.CENTER)
            if stat.get("sublabel"):
                self.text(s, x + px(12), st + py(140), sw - px(24), py(32), stat.get("sublabel", ""), "Arial", 10, Colors.TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        
        if c.source:
            self.text(s, px(50), SLIDE_HEIGHT - py(60), px(860), py(24), f"Source: {c.source}", "Arial", 8, Colors.TEXT_MUTED, italic=True)
        
        self.footer(s, c.footer_text, c.slide_number)


class CaseStudySlide(TemplateBase):
    """7. Case Study Slide"""
    
    def create(self, c: CaseStudyContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        ct = hh + py(20)
        lm = px(28)
        
        # Patient info panel
        iw, ih = px(280), py(340)
        self.rounded(s, lm, ct, iw, ih, grad={'angle': 165, 'stops': [Colors.OCEAN_MID, Colors.OCEAN_DEEP]}, r=0.04)
        isz = py(40)
        self.shape(s, MSO_SHAPE.OVAL, lm + iw/2 - isz/2, ct + py(20), isz, isz, fill=Colors.CYAN_BRIGHT)
        self.text(s, lm + px(20), ct + py(75), iw - px(40), py(240), c.patient_info or "Patient info...", "Arial", 11, Colors.WHITE)
        
        # Presentation panel
        pl = lm + iw + px(20)
        pw = px(590)
        
        self.rounded(s, pl, ct, pw, py(160), fill=Colors.WHITE, r=0.04)
        self.text(s, pl + px(20), ct + py(12), px(200), py(20), "CLINICAL PRESENTATION", "Arial", 8, Colors.CYAN_BRIGHT, True)
        self.text(s, pl + px(20), ct + py(36), pw - px(40), py(110), c.presentation or "Details...", "Arial", 11, Colors.TEXT_PRIMARY)
        
        # Question panel
        qt = ct + py(175)
        self.rounded(s, pl, qt, pw, py(70), fill=RGBColor(0xfb, 0xf0, 0xe4), r=0.04)
        self.shape(s, MSO_SHAPE.RECTANGLE, pl, qt, Pt(4), py(70), fill=Colors.CORAL_BRIGHT)
        self.text(s, pl + px(20), qt + py(10), px(100), py(16), "KEY QUESTION", "Arial", 7, Colors.CORAL_DEEP, True)
        self.text(s, pl + px(20), qt + py(28), pw - px(40), py(36), c.question or "Priority?", "Arial", 12, Colors.TEXT_PRIMARY, True)
        
        # Considerations
        cst = qt + py(85)
        self.text(s, pl + px(20), cst, px(200), py(16), "CONSIDERATIONS", "Arial", 8, Colors.TEXT_MUTED, True)
        for i, item in enumerate(c.considerations[:3]):
            y = cst + py(20) + py(24) * i
            self.bullet(s, pl + px(20), y)
            self.text(s, pl + px(48), y, px(520), py(20), item, "Arial", 11, Colors.TEXT_PRIMARY)
        
        self.footer(s, c.footer_text, c.slide_number)


class SummarySlide(TemplateBase):
    """8. Summary/Recap Slide"""
    
    def create(self, c: SummaryContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        ct = hh + py(25)
        points = c.points[:4] or [{"title": f"Point {i+1}", "description": ""} for i in range(3)]
        
        cw, ch = px(420), py(130)
        gap = py(16)
        lx, rx = px(28), px(28) + cw + px(24)
        
        for i, pt in enumerate(points):
            col, row = i % 2, i // 2
            x = lx if col == 0 else rx
            y = ct + (ch + gap) * row
            
            self.rounded(s, x, y, cw, ch, fill=Colors.WHITE, r=0.04)
            self.badge(s, x + px(16), y + py(16), i + 1, py(32))
            self.text(s, x + px(60), y + py(16), px(340), py(28), pt.get("title", ""), "Georgia", 14, Colors.TEXT_PRIMARY, True)
            self.text(s, x + px(60), y + py(48), px(340), py(70), pt.get("description", ""), "Arial", 11, Colors.TEXT_SECONDARY)
        
        if c.next_topic:
            pt = SLIDE_HEIGHT - py(85)
            self.rounded(s, px(28), pt, px(904), py(45), fill=Colors.OCEAN_DEEP, r=0.04)
            self.text(s, px(48), pt + py(4), px(100), py(16), "NEXT UP", "Arial", 7, Colors.CYAN_SOFT, True)
            self.text(s, px(48), pt + py(20), px(860), py(20), c.next_topic, "Arial", 11, Colors.WHITE)
        
        self.footer(s, c.footer_text, c.slide_number)


class QASlide(TemplateBase):
    """9. Q&A/Discussion Slide"""
    
    def create(self, c: QAContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.OCEAN_DEEP, Colors.OCEAN_MID)
        
        self.text(s, px(80), py(120), px(800), py(100), c.title, "Georgia", 44, Colors.WHITE, True, align=PP_ALIGN.CENTER)
        self.shape(s, MSO_SHAPE.RECTANGLE, px(380), py(230), px(200), Pt(4), fill=Colors.CYAN_BRIGHT)
        
        prompts = c.prompts[:3]
        if prompts:
            pt = py(280)
            for i, prompt in enumerate(prompts):
                y = pt + py(50) * i
                self.text(s, px(250), y, py(30), py(40), "?", "Georgia", 24, Colors.CYAN_BRIGHT, True)
                self.text(s, px(290), y + py(5), px(420), py(40), prompt, "Arial", 12, Colors.WHITE)
        
        if c.contact:
            self.rounded(s, px(300), py(450), px(360), py(50), fill=Colors.CYAN_MUTED, r=0.06)
            self.text(s, px(300), py(450), px(360), py(50), c.contact, "Arial", 11, Colors.WHITE,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Decorative circles
        self.shape(s, MSO_SHAPE.OVAL, px(60), py(400), py(80), py(80), fill=Colors.CYAN_MUTED)
        self.shape(s, MSO_SHAPE.OVAL, SLIDE_WIDTH - px(140), py(100), py(60), py(60), fill=Colors.CORAL_BRIGHT)
        
        self.footer(s, c.footer_text, c.slide_number)


class ImageSlide(TemplateBase):
    """10. Image Focus Slide"""
    
    def create(self, c: ImageContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        ct = hh + py(20)
        iw, ih = px(600), py(320)
        il = (SLIDE_WIDTH - iw) / 2
        
        # Placeholder
        self.rounded(s, il, ct, iw, ih, fill=Colors.SAND_DARK, r=0.04)
        self.text(s, il, ct + ih/2 - py(30), iw, py(60), "\U0001F5BC\uFE0F", "Arial", 48, Colors.TEXT_MUTED, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        self.text(s, il, ct + ih/2 + py(20), iw, py(30), "Insert Image Here", "Arial", 12, Colors.TEXT_MUTED, align=PP_ALIGN.CENTER)
        
        # Caption
        capt = ct + ih + py(16)
        self.rounded(s, il, capt, iw, py(50), fill=Colors.WHITE, r=0.04)
        self.shape(s, MSO_SHAPE.RECTANGLE, il, capt, Pt(4), py(50), fill=Colors.CYAN_BRIGHT)
        self.text(s, il + px(20), capt + py(10), iw - px(40), py(30), c.caption, "Arial", 11, Colors.TEXT_SECONDARY, italic=True, align=PP_ALIGN.CENTER)
        
        self.footer(s, c.footer_text, c.slide_number)


# =============================================================================
# VOLUME 2: PREMIUM TEMPLATE IMPLEMENTATIONS (11-20)
# =============================================================================

class DecisionTreeSlide(TemplateBase):
    """11. Decision Tree Slide"""
    
    def create(self, c: DecisionTreeContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        ct = hh + py(20)
        center_x = SLIDE_WIDTH / 2
        
        # Root node
        root_w, root_h = px(380), py(70)
        root_x = center_x - root_w / 2
        root_y = ct
        
        self.rounded(s, root_x + Pt(3), root_y + Pt(3), root_w, root_h, fill=Colors.OCEAN_LIGHT, r=0.1)
        self.rounded(s, root_x, root_y, root_w, root_h,
                     grad={'angle': 135, 'stops': [Colors.OCEAN_MID, Colors.OCEAN_DEEP]}, r=0.1)
        self.shape(s, MSO_SHAPE.RECTANGLE, root_x, root_y, Pt(5), root_h, fill=Colors.CYAN_BRIGHT)
        self.text(s, root_x + px(16), root_y + py(10), root_w - px(32), root_h - py(20),
                  c.root_question, "Georgia", 15, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Connector
        conn_y = root_y + root_h
        self.connector(s, center_x - Pt(1), conn_y, center_x - Pt(1), conn_y + py(25), Colors.CYAN_BRIGHT, 3)
        
        split_y = conn_y + py(25)
        num_branches = len(c.branches)
        
        if num_branches >= 2:
            branch_spacing = px(300)
            total_width = branch_spacing * (num_branches - 1)
            start_x = center_x - total_width / 2
            
            self.connector(s, start_x, split_y, start_x + total_width, split_y, Colors.CYAN_BRIGHT, 3)
            
            for i, branch in enumerate(c.branches[:3]):
                bx = start_x + branch_spacing * i - px(130)
                by = split_y + py(20)
                bw, bh = px(260), py(50)
                
                self.connector(s, start_x + branch_spacing * i, split_y,
                               start_x + branch_spacing * i, by, Colors.CYAN_BRIGHT, 3)
                
                label_color = Colors.SUCCESS if branch.get("label", "").lower() in ["yes", "true", "positive"] else Colors.CORAL_BRIGHT
                self.rounded(s, bx + bw/2 - px(30), by - py(12), px(60), py(22), fill=label_color, r=0.3)
                self.text(s, bx + bw/2 - px(30), by - py(12), px(60), py(22),
                          branch.get("label", f"Option {i+1}"), "Arial", 9, Colors.WHITE, True,
                          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
                
                by += py(15)
                self.rounded(s, bx, by, bw, bh, fill=Colors.WHITE, r=0.08)
                self.shape(s, MSO_SHAPE.RECTANGLE, bx, by, bw, Pt(4), fill=label_color)
                self.text(s, bx + px(12), by + py(8), bw - px(24), py(20),
                          branch.get("outcome", "Outcome"), "Georgia", 12, Colors.TEXT_PRIMARY, True, align=PP_ALIGN.CENTER)
                
                action_y = by + bh + py(12)
                self.arrow_down(s, bx + bw/2 - py(10), by + bh, py(20), py(15), Colors.TEXT_MUTED)
                
                self.rounded(s, bx, action_y + py(10), bw, py(55), fill=Colors.SAND_WARM, r=0.08)
                self.shape(s, MSO_SHAPE.RECTANGLE, bx, action_y + py(10), Pt(4), py(55), fill=label_color)
                self.text(s, bx + px(14), action_y + py(18), bw - px(28), py(42),
                          branch.get("action", "Action"), "Arial", 10, Colors.TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        
        self.footer(s, c.footer_text, c.slide_number)


class TimelineSlide(TemplateBase):
    """12. Timeline Slide"""
    
    def create(self, c: TimelineContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        events = c.events[:6]
        n = len(events)
        
        tl_top = hh + py(100)
        tl_left = px(60)
        tl_right = px(900)
        tl_width = tl_right - tl_left
        
        # Timeline bar
        bar_h = py(8)
        self.rounded(s, tl_left, tl_top - bar_h/2, tl_width, bar_h, fill=Colors.SAND_DEEP, r=0.5)
        self.rounded(s, tl_left + Pt(2), tl_top - bar_h/2 + Pt(2), tl_width - Pt(4), bar_h - Pt(4),
                     grad={'angle': 90, 'stops': [Colors.CYAN_BRIGHT, Colors.CYAN_DARK]}, r=0.5)
        
        spacing = tl_width / (n - 1) if n > 1 else tl_width
        
        for i, event in enumerate(events):
            ex = tl_left + spacing * i if n > 1 else tl_left + tl_width / 2
            
            # Marker
            dot_sz = py(20)
            self.shape(s, MSO_SHAPE.OVAL, ex - dot_sz/2 + Pt(2), tl_top - dot_sz/2 + Pt(2), dot_sz, dot_sz, fill=Colors.CYAN_DARK)
            self.shape(s, MSO_SHAPE.OVAL, ex - dot_sz/2, tl_top - dot_sz/2, dot_sz, dot_sz, fill=Colors.CYAN_BRIGHT)
            self.shape(s, MSO_SHAPE.OVAL, ex - py(6), tl_top - py(6), py(12), py(12), fill=Colors.WHITE)
            
            # Time label
            self.rounded(s, ex - px(35), tl_top - py(50), px(70), py(28), fill=Colors.OCEAN_DEEP, r=0.15)
            self.text(s, ex - px(35), tl_top - py(50), px(70), py(28),
                      event.get("time", ""), "Arial", 10, Colors.CYAN_SOFT, True,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            
            # Event card
            card_top = tl_top + py(40)
            self.connector(s, ex, tl_top + dot_sz/2, ex, card_top, Colors.CYAN_MUTED, 2)
            
            card_w = px(140)
            card_h = py(130)
            card_x = ex - card_w / 2
            
            self.rounded(s, card_x + Pt(2), card_top + Pt(2), card_w, card_h, fill=Colors.SAND_DEEP, r=0.1)
            self.rounded(s, card_x, card_top, card_w, card_h, fill=Colors.WHITE, r=0.1)
            
            accent_color = Colors.CYAN_BRIGHT if i % 2 == 0 else Colors.CORAL_BRIGHT
            self.shape(s, MSO_SHAPE.RECTANGLE, card_x, card_top, card_w, Pt(4), fill=accent_color)
            
            self.text(s, card_x + px(10), card_top + py(14), card_w - px(20), py(28),
                      event.get("title", f"Event {i+1}"), "Georgia", 12, Colors.TEXT_PRIMARY, True, align=PP_ALIGN.CENTER)
            self.text(s, card_x + px(10), card_top + py(48), card_w - px(20), py(70),
                      event.get("description", ""), "Arial", 9, Colors.TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        
        self.footer(s, c.footer_text, c.slide_number)


class TableSlide(TemplateBase):
    """13. Table Slide"""
    
    def create(self, c: TableContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        headers = c.headers[:5]
        rows = c.rows[:6]
        
        num_cols = len(headers)
        num_rows = len(rows)
        
        tbl_left = px(50)
        tbl_top = hh + py(25)
        tbl_width = px(860)
        
        header_h = py(45)
        row_h = py(42)
        col_w = tbl_width / num_cols
        
        total_h = header_h + row_h * num_rows
        self.rounded(s, tbl_left + Pt(4), tbl_top + Pt(4), tbl_width, total_h, fill=Colors.OCEAN_LIGHT, r=0.06)
        
        # Header
        self.rounded(s, tbl_left, tbl_top, tbl_width, header_h,
                     grad={'angle': 180, 'stops': [Colors.OCEAN_MID, Colors.OCEAN_DEEP]}, r=0.06)
        
        for i, header in enumerate(headers):
            hx = tbl_left + col_w * i
            if i > 0:
                self.shape(s, MSO_SHAPE.RECTANGLE, hx, tbl_top + py(8), Pt(1), header_h - py(16), fill=Colors.OCEAN_LIGHT)
            self.text(s, hx + px(8), tbl_top + py(8), col_w - px(16), header_h - py(16),
                      header, "Arial", 11, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Data rows
        for ri, row in enumerate(rows):
            ry = tbl_top + header_h + row_h * ri
            row_color = Colors.WHITE if ri % 2 == 0 else Colors.SAND_WARM
            
            if ri == num_rows - 1:
                self.rounded(s, tbl_left, ry, tbl_width, row_h, fill=row_color, r=0.06)
                self.shape(s, MSO_SHAPE.RECTANGLE, tbl_left, ry, tbl_width, py(10), fill=row_color)
            else:
                self.shape(s, MSO_SHAPE.RECTANGLE, tbl_left, ry, tbl_width, row_h, fill=row_color)
            
            self.shape(s, MSO_SHAPE.RECTANGLE, tbl_left, ry + row_h - Pt(1), tbl_width, Pt(1), fill=Colors.SAND_DEEP)
            
            for ci, cell in enumerate(row[:num_cols]):
                cx = tbl_left + col_w * ci
                if ci > 0:
                    self.shape(s, MSO_SHAPE.RECTANGLE, cx, ry + py(6), Pt(1), row_h - py(12), fill=Colors.SAND_DEEP)
                
                # Status styling
                if ci == num_cols - 1 and cell.lower() in ['normal', 'good', 'pass', 'ok']:
                    self.rounded(s, cx + col_w/2 - px(40), ry + py(8), px(80), py(26), fill=Colors.SUCCESS, r=0.3)
                    self.text(s, cx + col_w/2 - px(40), ry + py(8), px(80), py(26),
                              cell, "Arial", 9, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
                elif ci == num_cols - 1 and cell.lower() in ['high', 'low', 'abnormal', 'warning']:
                    self.rounded(s, cx + col_w/2 - px(40), ry + py(8), px(80), py(26), fill=Colors.WARNING, r=0.3)
                    self.text(s, cx + col_w/2 - px(40), ry + py(8), px(80), py(26),
                              cell, "Arial", 9, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
                elif ci == num_cols - 1 and cell.lower() in ['critical', 'danger', 'fail']:
                    self.rounded(s, cx + col_w/2 - px(40), ry + py(8), px(80), py(26), fill=Colors.DANGER, r=0.3)
                    self.text(s, cx + col_w/2 - px(40), ry + py(8), px(80), py(26),
                              cell, "Arial", 9, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
                else:
                    self.text(s, cx + px(8), ry + py(6), col_w - px(16), row_h - py(12),
                              cell, "Arial", 10, Colors.TEXT_PRIMARY, ci == 0, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        self.footer(s, c.footer_text, c.slide_number)


class CanvasSlide(TemplateBase):
    """14. Canvas Slide (Freeform)"""
    
    def create(self, c: CanvasContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        canvas_top = hh + py(15)
        canvas_left = px(28)
        canvas_width = px(904)
        canvas_height = py(355)
        
        self.rounded(s, canvas_left, canvas_top, canvas_width, canvas_height, fill=Colors.WHITE, r=0.04)
        
        # Corner accents
        accent_size = py(30)
        for corner, (cx, cy, fill) in [
            ('tl', (canvas_left, canvas_top, Colors.CYAN_BRIGHT)),
            ('tr', (canvas_left + canvas_width - accent_size, canvas_top, Colors.CYAN_BRIGHT)),
            ('bl', (canvas_left, canvas_top + canvas_height - Pt(3), Colors.CORAL_BRIGHT)),
            ('br', (canvas_left + canvas_width - accent_size, canvas_top + canvas_height - Pt(3), Colors.CORAL_BRIGHT)),
        ]:
            if 't' in corner:
                self.shape(s, MSO_SHAPE.RECTANGLE, cx, cy, accent_size, Pt(3), fill=fill)
            if 'l' in corner:
                self.shape(s, MSO_SHAPE.RECTANGLE, canvas_left, cy if 'b' not in corner else cy - accent_size + Pt(3), Pt(3), accent_size, fill=fill)
            if 'r' in corner:
                self.shape(s, MSO_SHAPE.RECTANGLE, canvas_left + canvas_width - Pt(3), cy if 'b' not in corner else cy - accent_size + Pt(3), Pt(3), accent_size, fill=fill)
            if 'b' in corner:
                self.shape(s, MSO_SHAPE.RECTANGLE, cx, cy, accent_size, Pt(3), fill=fill)
        
        # Grid dots
        grid_spacing = px(80)
        for gx in range(1, 11):
            for gy in range(1, 5):
                dx = canvas_left + grid_spacing * gx
                dy = canvas_top + py(70) * gy
                if dx < canvas_left + canvas_width - px(40) and dy < canvas_top + canvas_height - py(40):
                    self.shape(s, MSO_SHAPE.OVAL, dx, dy, Pt(3), Pt(3), fill=Colors.SAND_DEEP)
        
        # Placeholder text
        self.text(s, canvas_left, canvas_top + canvas_height/2 - py(20),
                  canvas_width, py(40), c.placeholder_text,
                  "Arial", 14, Colors.TEXT_LIGHT, italic=True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        self.footer(s, c.footer_text, c.slide_number)


class ChecklistSlide(TemplateBase):
    """15. Checklist Slide"""
    
    def create(self, c: ChecklistContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        items = c.items[:8]
        
        list_left = px(80)
        list_top = hh + py(25)
        list_width = px(500)
        item_height = py(45)
        
        # Progress
        checked = sum(1 for item in items if item.get("checked", False))
        total = len(items)
        progress = checked / total if total > 0 else 0
        
        prog_top = list_top
        prog_width = px(280)
        prog_left = px(600)
        
        self.rounded(s, prog_left, prog_top, prog_width, py(12), fill=Colors.SAND_DEEP, r=0.5)
        if progress > 0:
            self.rounded(s, prog_left, prog_top, prog_width * progress, py(12),
                         fill=Colors.SUCCESS if progress == 1 else Colors.CYAN_BRIGHT, r=0.5)
        
        self.text(s, prog_left, prog_top + py(18), prog_width, py(20),
                  f"{checked}/{total} Complete ({int(progress*100)}%)", "Arial", 10, Colors.TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        
        # Summary card
        summary_top = prog_top + py(50)
        self.rounded(s, prog_left, summary_top, prog_width, py(180),
                     grad={'angle': 165, 'stops': [Colors.OCEAN_MID, Colors.OCEAN_DEEP]}, r=0.08)
        
        self.text(s, prog_left + px(20), summary_top + py(16), prog_width - px(40), py(24),
                  "Progress Summary", "Georgia", 14, Colors.WHITE, True)
        self.shape(s, MSO_SHAPE.RECTANGLE, prog_left + px(20), summary_top + py(44), prog_width - px(40), Pt(1), fill=Colors.CYAN_SOFT)
        
        self.text(s, prog_left + px(20), summary_top + py(60), px(80), py(50), str(checked), "Georgia", 32, Colors.CYAN_BRIGHT, True)
        self.text(s, prog_left + px(20), summary_top + py(105), px(80), py(20), "Done", "Arial", 9, Colors.TEXT_LIGHT)
        
        self.text(s, prog_left + px(120), summary_top + py(60), px(80), py(50), str(total - checked), "Georgia", 32, Colors.CORAL_SOFT, True)
        self.text(s, prog_left + px(120), summary_top + py(105), px(80), py(20), "Remaining", "Arial", 9, Colors.TEXT_LIGHT)
        
        # Checklist items
        for i, item in enumerate(items):
            iy = list_top + py(40) + item_height * i
            checked_item = item.get("checked", False)
            
            self.rounded(s, list_left, iy, list_width, item_height - py(6),
                         fill=Colors.WHITE if not checked_item else Colors.SAND_WARM, r=0.06)
            
            cb_size = py(24)
            cb_left = list_left + px(12)
            cb_top = iy + (item_height - py(6) - cb_size) / 2
            
            if checked_item:
                self.rounded(s, cb_left, cb_top, cb_size, cb_size, fill=Colors.SUCCESS, r=0.15)
                self.text(s, cb_left, cb_top, cb_size, cb_size, "✓", "Arial", 14, Colors.WHITE, True,
                          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            else:
                self.rounded(s, cb_left, cb_top, cb_size, cb_size, fill=Colors.SAND_DEEP, r=0.15)
                self.rounded(s, cb_left + Pt(2), cb_top + Pt(2), cb_size - Pt(4), cb_size - Pt(4), fill=Colors.WHITE, r=0.12)
            
            text_color = Colors.TEXT_MUTED if checked_item else Colors.TEXT_PRIMARY
            self.text(s, list_left + px(48), iy + py(8), list_width - px(60), item_height - py(22),
                      item.get("text", f"Item {i+1}"), "Arial", 12, text_color, valign=MSO_ANCHOR.MIDDLE)
            
            if checked_item:
                self.shape(s, MSO_SHAPE.RECTANGLE, list_left + px(48), iy + item_height/2 - Pt(1),
                           list_width - px(80), Pt(1), fill=Colors.TEXT_MUTED)
        
        self.footer(s, c.footer_text, c.slide_number)


class MatrixSlide(TemplateBase):
    """16. Comparison Matrix Slide"""
    
    def create(self, c: MatrixContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        row_headers = c.row_headers[:5]
        col_headers = c.col_headers[:4]
        data = c.data
        
        num_rows = len(row_headers)
        num_cols = len(col_headers)
        
        mtx_left = px(50)
        mtx_top = hh + py(25)
        
        row_label_w = px(180)
        col_w = px(160)
        row_h = py(50)
        header_h = py(55)
        
        mtx_width = row_label_w + col_w * num_cols
        mtx_height = header_h + row_h * num_rows
        
        self.rounded(s, mtx_left + Pt(4), mtx_top + Pt(4), mtx_width, mtx_height, fill=Colors.OCEAN_LIGHT, r=0.06)
        
        # Corner cell
        self.rounded(s, mtx_left, mtx_top, row_label_w, header_h,
                     grad={'angle': 135, 'stops': [Colors.OCEAN_DEEP, Colors.OCEAN_MID]}, r=0.06)
        
        # Column headers
        for ci, col_header in enumerate(col_headers):
            cx = mtx_left + row_label_w + col_w * ci
            self.shape(s, MSO_SHAPE.RECTANGLE, cx, mtx_top, col_w, header_h,
                       grad={'angle': 135, 'stops': [Colors.OCEAN_DEEP, Colors.OCEAN_MID]})
            if ci < num_cols - 1:
                self.shape(s, MSO_SHAPE.RECTANGLE, cx + col_w - Pt(1), mtx_top + py(10), Pt(1), header_h - py(20), fill=Colors.OCEAN_LIGHT)
            self.text(s, cx + px(8), mtx_top + py(10), col_w - px(16), header_h - py(20),
                      col_header, "Georgia", 12, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Rows
        for ri, row_header in enumerate(row_headers):
            ry = mtx_top + header_h + row_h * ri
            
            self.shape(s, MSO_SHAPE.RECTANGLE, mtx_left, ry, row_label_w, row_h, fill=Colors.OCEAN_SUBTLE)
            if ri < num_rows - 1:
                self.shape(s, MSO_SHAPE.RECTANGLE, mtx_left + px(12), ry + row_h - Pt(1), row_label_w - px(24), Pt(1), fill=Colors.OCEAN_LIGHT)
            self.text(s, mtx_left + px(12), ry + py(8), row_label_w - px(24), row_h - py(16),
                      row_header, "Arial", 11, Colors.WHITE, True, valign=MSO_ANCHOR.MIDDLE)
            
            for ci in range(num_cols):
                cx = mtx_left + row_label_w + col_w * ci
                cell_data = data[ri][ci] if ri < len(data) and ci < len(data[ri]) else "—"
                
                cell_bg = Colors.WHITE if ri % 2 == 0 else Colors.SAND_WARM
                self.shape(s, MSO_SHAPE.RECTANGLE, cx, ry, col_w, row_h, fill=cell_bg)
                self.shape(s, MSO_SHAPE.RECTANGLE, cx, ry + row_h - Pt(1), col_w, Pt(1), fill=Colors.SAND_DEEP)
                if ci < num_cols - 1:
                    self.shape(s, MSO_SHAPE.RECTANGLE, cx + col_w - Pt(1), ry + py(8), Pt(1), row_h - py(16), fill=Colors.SAND_DEEP)
                
                if cell_data in ["✓", "✔", "Yes"]:
                    self.shape(s, MSO_SHAPE.OVAL, cx + col_w/2 - py(14), ry + row_h/2 - py(14), py(28), py(28), fill=Colors.SUCCESS)
                    self.text(s, cx + col_w/2 - py(14), ry + row_h/2 - py(14), py(28), py(28),
                              "✓", "Arial", 14, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
                elif cell_data in ["✗", "—", "-", "No"]:
                    self.shape(s, MSO_SHAPE.OVAL, cx + col_w/2 - py(14), ry + row_h/2 - py(14), py(28), py(28), fill=Colors.SAND_DEEP)
                    self.text(s, cx + col_w/2 - py(14), ry + row_h/2 - py(14), py(28), py(28),
                              "—", "Arial", 14, Colors.TEXT_MUTED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
                else:
                    self.text(s, cx + px(8), ry + py(8), col_w - px(16), row_h - py(16),
                              cell_data, "Arial", 10, Colors.TEXT_PRIMARY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        self.footer(s, c.footer_text, c.slide_number)


class HierarchySlide(TemplateBase):
    """17. Hierarchy Slide"""
    
    def create(self, c: HierarchyContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        center_x = SLIDE_WIDTH / 2
        
        # Root
        root_w, root_h = px(280), py(55)
        root_x = center_x - root_w / 2
        root_y = hh + py(25)
        
        self.rounded(s, root_x + Pt(3), root_y + Pt(3), root_w, root_h, fill=Colors.OCEAN_LIGHT, r=0.12)
        self.rounded(s, root_x, root_y, root_w, root_h,
                     grad={'angle': 145, 'stops': [Colors.OCEAN_MID, Colors.OCEAN_DEEP]}, r=0.12)
        self.shape(s, MSO_SHAPE.RECTANGLE, root_x + px(10), root_y + root_h - Pt(4), root_w - px(20), Pt(4), fill=Colors.CYAN_BRIGHT)
        self.text(s, root_x + px(16), root_y + py(10), root_w - px(32), root_h - py(20),
                  c.root, "Georgia", 16, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Level 1
        l1_conn_y = root_y + root_h
        self.connector(s, center_x, l1_conn_y, center_x, l1_conn_y + py(20), Colors.CYAN_BRIGHT, 3)
        
        level1 = c.level1[:4]
        n1 = len(level1)
        
        l1_top = l1_conn_y + py(20)
        l1_node_w = px(180)
        l1_node_h = py(45)
        l1_spacing = px(220)
        l1_total = l1_spacing * (n1 - 1)
        l1_start = center_x - l1_total / 2
        
        self.connector(s, l1_start, l1_top, l1_start + l1_total, l1_top, Colors.CYAN_BRIGHT, 3)
        
        accent_colors = [Colors.CYAN_BRIGHT, Colors.CORAL_BRIGHT, Colors.SUCCESS, Colors.WARNING]
        
        for i, l1_item in enumerate(level1):
            l1_x = l1_start + l1_spacing * i - l1_node_w / 2
            
            self.connector(s, l1_start + l1_spacing * i, l1_top, l1_start + l1_spacing * i, l1_top + py(15), Colors.CYAN_BRIGHT, 3)
            
            l1_y = l1_top + py(15)
            
            self.rounded(s, l1_x + Pt(2), l1_y + Pt(2), l1_node_w, l1_node_h, fill=Colors.SAND_DEEP, r=0.1)
            self.rounded(s, l1_x, l1_y, l1_node_w, l1_node_h, fill=Colors.WHITE, r=0.1)
            self.shape(s, MSO_SHAPE.RECTANGLE, l1_x, l1_y, l1_node_w, Pt(4), fill=accent_colors[i % 4])
            self.text(s, l1_x + px(8), l1_y + py(8), l1_node_w - px(16), l1_node_h - py(16),
                      l1_item, "Arial", 11, Colors.TEXT_PRIMARY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            
            # Level 2
            if i < len(c.level2) and c.level2[i]:
                level2 = c.level2[i][:2]
                
                l2_conn_y = l1_y + l1_node_h
                self.connector(s, l1_start + l1_spacing * i, l2_conn_y, l1_start + l1_spacing * i, l2_conn_y + py(12), Colors.TEXT_MUTED, 2)
                
                l2_top = l2_conn_y + py(12)
                l2_node_w = px(130)
                l2_node_h = py(38)
                l2_spacing = px(140)
                l2_total = l2_spacing * (len(level2) - 1) if len(level2) > 1 else 0
                l2_start = l1_start + l1_spacing * i - l2_total / 2
                
                if len(level2) > 1:
                    self.connector(s, l2_start, l2_top, l2_start + l2_total, l2_top, Colors.TEXT_MUTED, 2)
                
                for j, l2_item in enumerate(level2):
                    l2_x = l2_start + l2_spacing * j - l2_node_w / 2 if len(level2) > 1 else l2_start - l2_node_w / 2
                    
                    if len(level2) > 1:
                        self.connector(s, l2_start + l2_spacing * j, l2_top, l2_start + l2_spacing * j, l2_top + py(10), Colors.TEXT_MUTED, 2)
                    
                    l2_y = l2_top + py(10)
                    
                    self.rounded(s, l2_x, l2_y, l2_node_w, l2_node_h, fill=Colors.SAND_WARM, r=0.08)
                    self.shape(s, MSO_SHAPE.RECTANGLE, l2_x, l2_y, Pt(3), l2_node_h, fill=accent_colors[i % 4])
                    self.text(s, l2_x + px(10), l2_y + py(6), l2_node_w - px(20), l2_node_h - py(12),
                              l2_item, "Arial", 9, Colors.TEXT_SECONDARY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        self.footer(s, c.footer_text, c.slide_number)


class DefinitionSlide(TemplateBase):
    """18. Definition Slide"""
    
    def create(self, c: DefinitionContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        terms = c.terms[:4]
        
        card_left = px(50)
        card_top = hh + py(20)
        card_width = px(860)
        card_height = py(85)
        card_gap = py(10)
        
        accent_colors = [Colors.CYAN_BRIGHT, Colors.CORAL_BRIGHT, Colors.SUCCESS, Colors.CYAN_DARK]
        
        for i, term_data in enumerate(terms):
            cy = card_top + (card_height + card_gap) * i
            accent = accent_colors[i % len(accent_colors)]
            
            self.rounded(s, card_left + Pt(3), cy + Pt(3), card_width, card_height, fill=Colors.SAND_DEEP, r=0.08)
            self.rounded(s, card_left, cy, card_width, card_height, fill=Colors.WHITE, r=0.08)
            self.shape(s, MSO_SHAPE.RECTANGLE, card_left, cy + py(12), Pt(5), card_height - py(24), fill=accent)
            
            # Term badge
            term_text = term_data.get("term", f"Term {i+1}")
            badge_w = px(160)
            self.rounded(s, card_left + px(20), cy + py(12), badge_w, py(28), fill=accent, r=0.15)
            self.text(s, card_left + px(20), cy + py(12), badge_w, py(28),
                      term_text, "Georgia", 12, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            
            # Definition
            def_text = term_data.get("definition", "Definition goes here.")
            self.text(s, card_left + px(20), cy + py(48), card_width - px(50), py(32),
                      def_text, "Arial", 11, Colors.TEXT_SECONDARY)
            
            # Index
            self.text(s, card_left + card_width - px(50), cy + py(12), px(40), py(28),
                      f"{i+1:02d}", "Georgia", 24, Colors.SAND_DEEP, bold=True, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)
        
        self.footer(s, c.footer_text, c.slide_number)


class BeforeAfterSlide(TemplateBase):
    """19. Before/After Slide"""
    
    def create(self, c: BeforeAfterContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        hh = self.header(s, c.title, c.subtitle, compact=True)
        
        ct = hh + py(20)
        col_width = px(400)
        col_height = py(340)
        gap = px(60)
        
        # BEFORE
        before_x = px(60)
        self.rounded(s, before_x, ct, px(100), py(30), fill=Colors.CORAL_BRIGHT, r=0.2)
        self.text(s, before_x, ct, px(100), py(30), c.before_title.upper(),
                  "Arial", 10, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        self.rounded(s, before_x, ct + py(40), col_width, col_height - py(40), fill=Colors.WHITE, r=0.08)
        self.shape(s, MSO_SHAPE.RECTANGLE, before_x, ct + py(40), col_width, Pt(4), fill=Colors.CORAL_BRIGHT)
        
        for i, point in enumerate(c.before_points[:5]):
            py_offset = ct + py(65) + py(50) * i
            self.shape(s, MSO_SHAPE.OVAL, before_x + px(20), py_offset, py(22), py(22), fill=Colors.CORAL_SOFT)
            self.text(s, before_x + px(20), py_offset, py(22), py(22), "✗",
                      "Arial", 12, Colors.CORAL_DEEP, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            self.text(s, before_x + px(54), py_offset, col_width - px(74), py(22),
                      point, "Arial", 11, Colors.TEXT_SECONDARY, valign=MSO_ANCHOR.MIDDLE)
        
        # Arrow
        arrow_x = before_x + col_width + gap / 2 - py(25)
        arrow_y = ct + col_height / 2
        
        self.shape(s, MSO_SHAPE.OVAL, arrow_x - py(5), arrow_y - py(5), py(60), py(60), fill=Colors.SAND_DARK)
        self.shape(s, MSO_SHAPE.OVAL, arrow_x, arrow_y, py(50), py(50), fill=Colors.OCEAN_DEEP)
        self.arrow_right(s, arrow_x + py(10), arrow_y + py(15), py(30), py(20), Colors.CYAN_BRIGHT)
        
        # AFTER
        after_x = before_x + col_width + gap
        self.rounded(s, after_x, ct, px(100), py(30), fill=Colors.SUCCESS, r=0.2)
        self.text(s, after_x, ct, px(100), py(30), c.after_title.upper(),
                  "Arial", 10, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        self.rounded(s, after_x, ct + py(40), col_width, col_height - py(40), fill=Colors.WHITE, r=0.08)
        self.shape(s, MSO_SHAPE.RECTANGLE, after_x, ct + py(40), col_width, Pt(4), fill=Colors.SUCCESS)
        
        for i, point in enumerate(c.after_points[:5]):
            py_offset = ct + py(65) + py(50) * i
            self.shape(s, MSO_SHAPE.OVAL, after_x + px(20), py_offset, py(22), py(22), fill=Colors.SUCCESS_SOFT)
            self.text(s, after_x + px(20), py_offset, py(22), py(22), "✓",
                      "Arial", 12, Colors.SUCCESS, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            self.text(s, after_x + px(54), py_offset, col_width - px(74), py(22),
                      point, "Arial", 11, Colors.TEXT_PRIMARY, valign=MSO_ANCHOR.MIDDLE)
        
        self.footer(s, c.footer_text, c.slide_number)


class AlertSlide(TemplateBase):
    """20. Alert Slide"""
    
    def create(self, c: AlertContent):
        s = self.prs.slides.add_slide(self.blank)
        
        alert_configs = {
            "danger": {"bg1": Colors.DANGER, "bg2": RGBColor(0xb9, 0x1c, 0x1c), "accent": Colors.DANGER_SOFT, "icon": "⚠"},
            "warning": {"bg1": Colors.WARNING, "bg2": RGBColor(0xd9, 0x77, 0x06), "accent": Colors.WARNING_SOFT, "icon": "⚡"},
            "success": {"bg1": Colors.SUCCESS, "bg2": RGBColor(0x05, 0x96, 0x69), "accent": Colors.SUCCESS_SOFT, "icon": "✓"},
            "info": {"bg1": Colors.CYAN_BRIGHT, "bg2": Colors.CYAN_DARK, "accent": Colors.CYAN_SOFT, "icon": "ℹ"}
        }
        
        config = alert_configs.get(c.alert_type, alert_configs["warning"])
        
        self.bg(s, config["bg1"], config["bg2"])
        
        # Large icon
        icon_size = py(120)
        icon_x = (SLIDE_WIDTH - icon_size) / 2
        icon_y = py(80)
        
        self.shape(s, MSO_SHAPE.OVAL, icon_x - py(10), icon_y - py(10), icon_size + py(20), icon_size + py(20), fill=config["accent"])
        self.shape(s, MSO_SHAPE.OVAL, icon_x, icon_y, icon_size, icon_size, fill=Colors.WHITE)
        self.text(s, icon_x, icon_y, icon_size, icon_size, config["icon"],
                  "Arial", 60, config["bg1"], True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Text
        self.text(s, px(80), icon_y + icon_size + py(30), px(800), py(50),
                  c.headline, "Georgia", 32, Colors.WHITE, True, align=PP_ALIGN.CENTER)
        
        self.text(s, px(100), icon_y + icon_size + py(90), px(760), py(60),
                  c.message, "Arial", 14, Colors.WHITE, align=PP_ALIGN.CENTER)
        
        # Actions
        if c.action_items:
            actions_top = icon_y + icon_size + py(170)
            action_w = px(500)
            action_h = py(30) * len(c.action_items) + py(30)
            action_x = (SLIDE_WIDTH - action_w) / 2
            
            self.rounded(s, action_x, actions_top, action_w, action_h, fill=Colors.WHITE, r=0.1)
            self.text(s, action_x + px(20), actions_top + py(10), action_w - px(40), py(20),
                      "ACTION REQUIRED", "Arial", 9, config["bg2"], True)
            
            for i, action in enumerate(c.action_items[:3]):
                ay = actions_top + py(35) + py(25) * i
                self.arrow_right(s, action_x + px(20), ay + py(3), py(14), py(14), config["bg1"])
                self.text(s, action_x + px(48), ay, action_w - px(68), py(20), action, "Arial", 11, Colors.TEXT_PRIMARY)
        
        self.footer(s, c.footer_text, c.slide_number)


class NclexTipSlide(TemplateBase):
    """21. NCLEX Tip Slide - Clean content with tip banner"""
    
    def create(self, c: NclexTipContent):
        s = self.prs.slides.add_slide(self.blank)
        self.bg(s, Colors.SAND_LIGHT, Colors.SAND_MID)
        
        # Custom header - 32pt title, no subtitle
        hh = self.nclex_header(s, c.title)
        
        # Calculate from bottom up for tight layout
        footer_height = py(24)
        tip_height = py(62)
        tip_bottom_gap = py(2)  # Minimal gap between tip and copyright
        
        # NCLEX tip positioned from bottom
        tip_top = SLIDE_HEIGHT - footer_height - tip_height - tip_bottom_gap
        
        # Content area - MAXIMIZED to fill space between header and tip
        content_top = hh + py(4)
        content_left = px(50)
        content_width = px(860)
        content_bottom_gap = py(4)  # Minimal gap between content and tip
        content_height = tip_top - content_top - content_bottom_gap  # Fill all available space
        
        # Main content panel - simple white background
        self.rounded(s, content_left, content_top, content_width, content_height, fill=Colors.WHITE, r=0.05)
        
        # Subtle left accent line (minimal embellishment for visual continuity)
        self.shape(s, MSO_SHAPE.RECTANGLE, content_left, content_top + py(12), Pt(3), content_height - py(24), fill=Colors.CYAN_BRIGHT)
        
        # Content text - clean paragraph format - 20pt font
        text_left = content_left + px(24)
        text_top = content_top + py(14)
        text_width = content_width - px(48)
        
        # Render content lines with 20pt font
        line_height = py(42)  # Spacing for 20pt font
        for i, line in enumerate(c.content_lines[:9]):  # Adjusted for expanded area
            y = text_top + line_height * i
            self.text(s, text_left, y, text_width, py(38), line, 
                      "Arial", 20, Colors.TEXT_PRIMARY, valign=MSO_ANCHOR.TOP)
        
        # NCLEX Tip Banner - positioned near bottom
        # Subtle shadow layer
        self.rounded(s, content_left + Pt(2), tip_top + Pt(2), content_width, tip_height,
                     fill=Colors.OCEAN_DEEP, r=0.06)
        # Main banner on top
        self.rounded(s, content_left, tip_top, content_width, tip_height,
                     grad={'angle': 135, 'stops': [Colors.OCEAN_MID, Colors.OCEAN_DEEP]}, r=0.06)
        
        # Left accent stripe (coral for attention/dopamine)
        self.shape(s, MSO_SHAPE.RECTANGLE, content_left, tip_top, Pt(4), tip_height, fill=Colors.CORAL_BRIGHT)
        
        # "NCLEX TIP" badge
        badge_w = px(85)
        badge_h = py(20)
        badge_x = content_left + px(16)
        badge_y = tip_top + py(8)
        
        # Badge glow effect
        self.rounded(s, badge_x - Pt(1), badge_y - Pt(1), badge_w + Pt(2), badge_h + Pt(2), 
                     fill=Colors.CORAL_SOFT, r=0.3)
        self.rounded(s, badge_x, badge_y, badge_w, badge_h, fill=Colors.CORAL_BRIGHT, r=0.3)
        self.text(s, badge_x, badge_y, badge_w, badge_h, "NCLEX TIP",
                  "Arial", 8, Colors.WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Tip icon (lightbulb)
        icon_x = content_left + px(112)
        icon_y = tip_top + py(6)
        icon_sz = py(24)
        
        # Glow behind icon
        self.shape(s, MSO_SHAPE.OVAL, icon_x - Pt(3), icon_y - Pt(3), icon_sz + Pt(6), icon_sz + Pt(6), 
                   fill=Colors.CYAN_MUTED)
        self.shape(s, MSO_SHAPE.OVAL, icon_x, icon_y, icon_sz, icon_sz, fill=Colors.CYAN_BRIGHT)
        self.text(s, icon_x, icon_y, icon_sz, icon_sz, "💡", "Arial", 10, Colors.OCEAN_DEEP,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Tip text - 20pt font
        tip_text_x = content_left + px(148)
        tip_text_y = tip_top + py(8)
        tip_text_w = content_width - px(168)
        
        self.text(s, tip_text_x, tip_text_y, tip_text_w, tip_height - py(16), c.nclex_tip,
                  "Arial", 20, Colors.WHITE, italic=True, valign=MSO_ANCHOR.TOP)
        
        # Pulsing dots (hypnotic visual rhythm)
        dot_y = tip_top + tip_height - py(12)
        dot_spacing = px(10)
        dot_start = content_left + content_width - px(65)
        
        for i in range(4):
            dot_size = Pt(3 + i)
            dot_x = dot_start + dot_spacing * i
            self.shape(s, MSO_SHAPE.OVAL, dot_x, dot_y - dot_size/2, dot_size, dot_size, 
                       fill=Colors.CYAN_SOFT if i % 2 == 0 else Colors.CORAL_SOFT)
        
        # Custom footer with copyright - tight to bottom
        self.nclex_footer(s, c.slide_number)
    
    def nclex_header(self, slide, title):
        """Custom header for NCLEX slides - 32pt title, no subtitle"""
        h = py(52)  # Compact header height
        
        # Header background
        self.rounded(slide, px(20), py(12), px(920), h,
                     grad={'angle': 135, 'stops': [Colors.OCEAN_MID, Colors.OCEAN_DEEP]}, r=0.06)
        
        # Shadow layer
        self.rounded(slide, px(20) + Pt(2), py(12) + Pt(2), px(920), h,
                     fill=Colors.OCEAN_DEEP, r=0.06)
        self.rounded(slide, px(20), py(12), px(920), h,
                     grad={'angle': 135, 'stops': [Colors.OCEAN_MID, Colors.OCEAN_DEEP]}, r=0.06)
        
        # Medical icon
        icon_sz = py(32)
        icon_x = px(36)
        icon_y = py(12) + (h - icon_sz) / 2
        self.shape(slide, MSO_SHAPE.OVAL, icon_x, icon_y, icon_sz, icon_sz, fill=Colors.CYAN_BRIGHT)
        self.text(slide, icon_x, icon_y, icon_sz, icon_sz, "+", "Arial", 18, Colors.OCEAN_DEEP, True,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        
        # Title - 32pt font
        self.text(slide, px(80), py(12), px(820), h, title,
                  "Georgia", 32, Colors.WHITE, True, valign=MSO_ANCHOR.MIDDLE)
        
        # Flow lines
        for i, offset in enumerate([0, 6, 12]):
            alpha = 0.3 - i * 0.1
            self.shape(slide, MSO_SHAPE.RECTANGLE,
                       px(850) + Pt(offset), py(12) + py(15), Pt(2), h - py(30),
                       fill=Colors.CYAN_SOFT)
        
        return py(12) + h
    
    def nclex_footer(self, slide, num):
        """Custom footer with copyright for NCLEX slides"""
        top = SLIDE_HEIGHT - py(20)  # Very tight to bottom
        
        # Subtle divider line
        self.shape(slide, MSO_SHAPE.RECTANGLE, px(28), top - py(1), px(904), Pt(0.5),
                   fill=Colors.SAND_DEEP)
        
        # Copyright text - 10pt font
        self.text(slide, px(32), top, px(900), py(18), 
                  "© 2026 Adaptive Academic Learning Solutions Inc | All Rights Reserved |",
                  "Arial", 10, Colors.TEXT_MUTED)


# =============================================================================
# TEMPLATE REGISTRY
# =============================================================================

ALL_TEMPLATES = {
    # Volume 1 - Standard
    'title': (TitleSlide, TitleContent),
    'content': (ContentSlide, ContentContent),
    'two_column': (TwoColumnSlide, TwoColumnContent),
    'quote': (QuoteSlide, QuoteContent),
    'process': (ProcessSlide, ProcessContent),
    'statistics': (StatisticsSlide, StatisticsContent),
    'case_study': (CaseStudySlide, CaseStudyContent),
    'summary': (SummarySlide, SummaryContent),
    'qa': (QASlide, QAContent),
    'image': (ImageSlide, ImageContent),
    # Volume 2 - Premium
    'decision_tree': (DecisionTreeSlide, DecisionTreeContent),
    'timeline': (TimelineSlide, TimelineContent),
    'table': (TableSlide, TableContent),
    'canvas': (CanvasSlide, CanvasContent),
    'checklist': (ChecklistSlide, ChecklistContent),
    'matrix': (MatrixSlide, MatrixContent),
    'hierarchy': (HierarchySlide, HierarchyContent),
    'definition': (DefinitionSlide, DefinitionContent),
    'before_after': (BeforeAfterSlide, BeforeAfterContent),
    'alert': (AlertSlide, AlertContent),
    # Volume 3 - Specialty
    'nclex_tip': (NclexTipSlide, NclexTipContent),
}


# =============================================================================
# SAMPLE DATA FOR ALL TEMPLATES
# =============================================================================

SAMPLE_DATA = {
    'two_column': {
        'left_points': ['Point A1', 'Point A2', 'Point A3', 'Point A4'],
        'right_points': ['Point B1', 'Point B2', 'Point B3', 'Point B4']
    },
    'process': {
        'steps': [
            {"title": "Assess", "description": "Gather patient data systematically"},
            {"title": "Diagnose", "description": "Identify nursing diagnoses"},
            {"title": "Plan", "description": "Set goals and interventions"},
            {"title": "Implement", "description": "Execute the care plan"},
            {"title": "Evaluate", "description": "Assess outcomes"}
        ]
    },
    'statistics': {
        'stats': [
            {"value": "95%", "label": "Patient Satisfaction", "sublabel": "Above national average"},
            {"value": "2.3M", "label": "Nurses in US", "sublabel": "As of 2024"},
            {"value": "6%", "label": "Job Growth", "sublabel": "Projected 2022-2032"}
        ],
        'source': "Bureau of Labor Statistics, 2024"
    },
    'case_study': {
        'patient_info': "68 y/o female\nPMH: HTN, DM Type 2\nAllergies: PCN\nAdmitted for: SOB, chest pain",
        'presentation': "Patient reports progressive shortness of breath over 3 days. VS: BP 158/92, HR 96, RR 24, SpO2 91% RA.",
        'question': "What are your priority nursing assessments?",
        'considerations': ["Cardiac vs pulmonary etiology", "Fluid volume status", "Pain characteristics"]
    },
    'summary': {
        'points': [
            {"title": "Assessment First", "description": "Always begin with thorough patient assessment"},
            {"title": "Evidence-Based", "description": "Use current research to guide practice"},
            {"title": "Patient-Centered", "description": "Involve patients in care decisions"},
            {"title": "Document Everything", "description": "Accurate records protect everyone"}
        ],
        'next_topic': "Module 2: Advanced Assessment Techniques"
    },
    'qa': {
        'prompts': ["What assessment findings concern you most?", "How would you prioritize interventions?", "What patient education is essential?"],
        'contact': "office.hours@nursing.edu"
    },
    'decision_tree': {
        'root_question': "Is the patient responsive?",
        'branches': [
            {"label": "Yes", "outcome": "Assess further", "action": "Check vital signs, symptoms, history"},
            {"label": "No", "outcome": "Emergency protocol", "action": "Call rapid response, begin BLS if needed"}
        ]
    },
    'timeline': {
        'events': [
            {"time": "0h", "title": "Admission", "description": "Initial assessment complete"},
            {"time": "2h", "title": "Labs", "description": "Blood work results received"},
            {"time": "4h", "title": "Treatment", "description": "Medication administered"},
            {"time": "8h", "title": "Reassess", "description": "Evaluate treatment response"},
            {"time": "24h", "title": "Discharge", "description": "Patient education complete"}
        ]
    },
    'table': {
        'headers': ["Parameter", "Value", "Normal Range", "Status"],
        'rows': [
            ["Heart Rate", "88 bpm", "60-100 bpm", "Normal"],
            ["Blood Pressure", "142/92", "<120/80", "High"],
            ["Temperature", "99.8°F", "97.8-99.1°F", "Warning"],
            ["SpO2", "97%", ">95%", "Normal"],
            ["Respiratory Rate", "18/min", "12-20/min", "Normal"]
        ]
    },
    'checklist': {
        'items': [
            {"text": "Verify patient identity (2 identifiers)", "checked": True},
            {"text": "Review medication allergies", "checked": True},
            {"text": "Check medication order accuracy", "checked": True},
            {"text": "Verify correct dose calculation", "checked": False},
            {"text": "Confirm administration route", "checked": False},
            {"text": "Document in medical record", "checked": False}
        ]
    },
    'matrix': {
        'row_headers': ["Effectiveness", "Side Effects", "Cost", "Availability"],
        'col_headers': ["Drug A", "Drug B", "Drug C"],
        'data': [["✓", "✓", "—"], ["—", "✓", "✓"], ["✓", "—", "✓"], ["✓", "✓", "—"]]
    },
    'hierarchy': {
        'root': "Patient Care Team",
        'level1': ["Physician", "Nursing", "Allied Health"],
        'level2': [["Attending", "Resident"], ["RN", "LPN"], ["PT", "OT"]]
    },
    'definition': {
        'terms': [
            {"term": "Tachycardia", "definition": "Heart rate >100 bpm at rest. May indicate fever, dehydration, anxiety, or cardiac pathology."},
            {"term": "Hypotension", "definition": "BP <90/60 mmHg. Can cause dizziness, fainting, and inadequate organ perfusion."},
            {"term": "Dyspnea", "definition": "Subjective breathing discomfort. Assess rate, depth, pattern, and oxygen saturation."},
            {"term": "Cyanosis", "definition": "Bluish skin/mucous membrane discoloration due to inadequate blood oxygenation."}
        ]
    },
    'before_after': {
        'before_title': "Before Intervention",
        'before_points': ["Uncontrolled pain (8/10)", "Limited mobility", "Poor sleep quality", "High anxiety"],
        'after_title': "After Intervention",
        'after_points': ["Pain managed (3/10)", "Improved mobility", "Restful sleep", "Reduced anxiety"]
    },
    'alert': {
        'alert_type': "danger",
        'headline': "Medication Safety Alert",
        'message': "High-alert medication requires independent double-check before administration.",
        'action_items': ["Perform independent double-check", "Verify patient allergies", "Document verification in MAR"]
    },
    'nclex_tip': {
        'title': 'Diuretic Therapy',
        'content_lines': [
            "Diuretics work by increasing urine output to reduce fluid volume overload.",
            "Loop diuretics (furosemide) are the most potent and act on the Loop of Henle.",
            "Monitor potassium levels closely - loop diuretics cause hypokalemia.",
            "Thiazide diuretics are less potent but useful for chronic hypertension management.",
            "Potassium-sparing diuretics (spironolactone) help prevent hypokalemia.",
            "Daily weights are the most accurate measure of fluid status changes.",
            "Report weight gain of >2 lbs/day or >5 lbs/week to the provider.",
            "Administer diuretics in the morning to prevent nocturia."
        ],
        'nclex_tip': "NCLEX Priority: When a question asks about diuretic therapy, look for answers that address potassium monitoring, daily weights, and I&O documentation."
    }
}


# =============================================================================
# GENERATOR FUNCTIONS
# =============================================================================

def generate_all_templates(output_dir: str = ".") -> List[str]:
    """Generate all 20 template types"""
    os.makedirs(output_dir, exist_ok=True)
    files = []
    
    for name, (cls, content_cls) in ALL_TEMPLATES.items():
        t = cls()
        c = content_cls()
        
        # Apply sample data if available
        if name in SAMPLE_DATA:
            for k, v in SAMPLE_DATA[name].items():
                setattr(c, k, v)
        
        t.create(c)
        fn = os.path.join(output_dir, f"template_{name}.pptx")
        t.save(fn)
        files.append(fn)
    
    return files


def generate_template(template_type: str, content_data: Dict = None, output_path: str = None) -> str:
    """Generate a single template"""
    if template_type not in ALL_TEMPLATES:
        raise ValueError(f"Unknown template type: {template_type}")
    
    cls, content_cls = ALL_TEMPLATES[template_type]
    t = cls()
    c = content_cls()
    
    # Apply provided data
    if content_data:
        for k, v in content_data.items():
            if hasattr(c, k):
                setattr(c, k, v)
    
    t.create(c)
    
    fn = output_path or f"template_{template_type}.pptx"
    t.save(fn)
    return fn


def generate_from_json(json_path: str, output_dir: str = ".") -> List[str]:
    """Generate templates from JSON file"""
    with open(json_path, 'r') as f:
        data = json.load(f)
    
    os.makedirs(output_dir, exist_ok=True)
    files = []
    
    slides = data.get('slides', [data])
    
    for i, slide_data in enumerate(slides):
        template_type = slide_data.get('type', 'content')
        content = slide_data.get('content', slide_data)
        
        fn = os.path.join(output_dir, f"slide_{i+1:02d}_{template_type}.pptx")
        generate_template(template_type, content, fn)
        files.append(fn)
    
    return files


# =============================================================================
# MAIN
# =============================================================================

def main():
    """Main entry point"""
    import argparse
    
    parser = argparse.ArgumentParser(description="Nursing Template Generator - 20 Slide Types")
    parser.add_argument('--type', help="Generate specific template type")
    parser.add_argument('--json', help="Generate from JSON file")
    parser.add_argument('--output', '-o', default="./all_templates", help="Output directory")
    parser.add_argument('--list', action='store_true', help="List all template types")
    parser.add_argument('--capacity', help="Show capacity info for template type")
    parser.add_argument('--auto-split', action='store_true', help="Auto-split overflow content")
    
    args = parser.parse_args()
    
    if args.list:
        print("\n📋 Available Template Types (21 total):\n")
        print("  VOLUME 1 - STANDARD:")
        for name in ['title', 'content', 'two_column', 'quote', 'process', 
                     'statistics', 'case_study', 'summary', 'qa', 'image']:
            spec = TEMPLATE_SPECS.get(name, {})
            cap = spec.get('capacity', {})
            cap_str = ', '.join(f"{k}:{v}" for k, v in cap.items() if isinstance(v, int))
            print(f"    • {name:15} [{cap_str}]")
        print("\n  VOLUME 2 - PREMIUM:")
        for name in ['decision_tree', 'timeline', 'table', 'canvas', 'checklist',
                     'matrix', 'hierarchy', 'definition', 'before_after', 'alert']:
            spec = TEMPLATE_SPECS.get(name, {})
            cap = spec.get('capacity', {})
            cap_str = ', '.join(f"{k}:{v}" for k, v in cap.items() if isinstance(v, int))
            print(f"    • {name:15} [{cap_str}]")
        print("\n  VOLUME 3 - SPECIALTY:")
        for name in ['nclex_tip']:
            spec = TEMPLATE_SPECS.get(name, {})
            cap = spec.get('capacity', {})
            cap_str = ', '.join(f"{k}:{v}" for k, v in cap.items() if isinstance(v, int))
            print(f"    • {name:15} [{cap_str}]")
        return
    
    if args.capacity:
        spec = TEMPLATE_SPECS.get(args.capacity)
        if spec:
            print(f"\n📊 Template: {args.capacity}")
            print(f"   Description: {spec.get('description', 'N/A')}")
            print(f"   Capacity: {spec.get('capacity', {})}")
            print(f"   Min items: {spec.get('min_content', 'N/A')}")
            print(f"   Max items: {spec.get('max_content', 'N/A')}")
            print(f"   Overflow: {spec.get('overflow_strategy', 'N/A')}")
            print(f"   Adaptive: {', '.join(spec.get('adaptive_features', []))}")
        else:
            print(f"Unknown template type: {args.capacity}")
        return
    
    print("\n" + "="*65)
    print("  NURSING LECTURE TEMPLATE SYSTEM")
    print("  21 Professional Slide Types")
    print("="*65 + "\n")
    
    if args.json:
        print(f"📄 Generating from JSON: {args.json}")
        files = generate_from_json(args.json, args.output)
    elif args.type:
        print(f"🎯 Generating single template: {args.type}")
        fn = generate_template(args.type, output_path=os.path.join(args.output, f"template_{args.type}.pptx"))
        files = [fn]
    else:
        print("📦 Generating all 20 templates...")
        files = generate_all_templates(args.output)
    
    print(f"\n✨ Generated {len(files)} template(s):")
    for f in files:
        print(f"   • {os.path.basename(f)}")
    
    print(f"\n📁 Output directory: {args.output}")
    print()


# =============================================================================
# ADAPTIVE CONTENT USAGE EXAMPLES
# =============================================================================

"""
EXAMPLE 1: Check if content fits before generating
------------------------------------------------------
from complete_template_system import check_content_fit, TimelineSlide, TimelineContent

# Content with 8 events (max is 6)
my_events = [{"time": f"Day {i}", "title": f"Event {i}"} for i in range(8)]

fit_check = check_content_fit('timeline', {'events': my_events})
print(fit_check)
# Output: {'fits': False, 'warnings': ['events: 8 items exceeds max 6'], 
#          'suggestions': ['Split into phases...']}


EXAMPLE 2: Auto-split overflow content
------------------------------------------------------
from complete_template_system import split_content_for_template, TimelineSlide, TimelineContent

# Content with 10 events
my_content = {
    'title': 'Patient Journey',
    'events': [{"time": f"Day {i}", "title": f"Milestone {i}"} for i in range(10)],
    'slide_number': '05'
}

# Split into multiple slides (6 + 4)
slide_contents = split_content_for_template('timeline', my_content)

for i, content in enumerate(slide_contents):
    slide = TimelineSlide()
    slide.create(TimelineContent(**content))
    slide.save(f"timeline_part{i+1}.pptx")
# Creates: timeline_part1.pptx (events 0-5), timeline_part2.pptx (events 6-9)


EXAMPLE 3: Get template suggestions based on content
------------------------------------------------------
from complete_template_system import suggest_template

# Describe your content
suggestions = suggest_template({
    'item_count': 5,
    'is_sequential': True,
    'content_type': 'process'
})

for template, confidence, reason in suggestions:
    print(f"{template}: {confidence:.0%} - {reason}")
# Output:
# process: 90% - Sequential steps fit process template
# timeline: 70% - More steps work well in timeline


EXAMPLE 4: Batch generate with auto-splitting
------------------------------------------------------
import json
from complete_template_system import split_content_for_template, ALL_TEMPLATES

with open('my_slides.json') as f:
    data = json.load(f)

slide_num = 1
for slide_data in data['slides']:
    template_type = slide_data['type']
    content = slide_data['content']
    
    # Auto-split if needed
    split_contents = split_content_for_template(template_type, content)
    
    for content_chunk in split_contents:
        content_chunk['slide_number'] = f"{slide_num:02d}"
        
        cls, content_cls = ALL_TEMPLATES[template_type]
        slide = cls()
        slide.create(content_cls(**content_chunk))
        slide.save(f"slide_{slide_num:02d}_{template_type}.pptx")
        slide_num += 1
"""


if __name__ == "__main__":
    main()
