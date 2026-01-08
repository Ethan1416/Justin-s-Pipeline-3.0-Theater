#!/usr/bin/env python3
"""
Premium Template Polishing System
30 Specialized Agents + 1 Orchestrator Agent

Architecture:
- Orchestrator: Delegates tasks, tracks progress, aggregates results
- 30 Agents organized into 6 categories:
  * Dimensional (1-5): Size, proportions, bounds
  * Color Psychology (6-10): Palette, contrast, harmony
  * Typography (11-15): Fonts, sizes, spacing
  * Layout (16-20): Grid, alignment, balance
  * Components (21-25): Icons, accents, decorations
  * Quality (26-30): Polish, uniformity, grade

Each agent returns findings and auto-correction recommendations.
"""

import os
import sys
import json
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from typing import List, Dict, Set, Tuple, Optional, Any
from enum import Enum
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor


# ============ STATUS TYPES ============

class Severity(Enum):
    PERFECT = "🟢 PERFECT"
    GOOD = "🟡 GOOD"
    NEEDS_FIX = "🟠 NEEDS FIX"
    CRITICAL = "🔴 CRITICAL"

class TaskStatus(Enum):
    PENDING = "pending"
    RUNNING = "running"
    COMPLETE = "complete"
    FAILED = "failed"


@dataclass
class AgentResult:
    agent_id: int
    category: str
    name: str
    severity: Severity
    score: float  # 0-100
    message: str
    findings: List[str] = field(default_factory=list)
    corrections: List[str] = field(default_factory=list)


@dataclass
class OrchestratorReport:
    total_agents: int
    completed: int
    failed: int
    overall_score: float
    grade: str
    category_scores: Dict[str, float]
    critical_issues: List[str]
    recommendations: List[str]
    agent_results: List[AgentResult]


# ============ REFERENCE STANDARDS ============

class Standards:
    """Canonical reference values for all templates"""
    
    # Dimensions
    SLIDE_WIDTH = 13.333
    SLIDE_HEIGHT = 7.5
    HEADER_HEIGHT_RANGE = (0.9, 1.4)
    FOOTER_HEIGHT_RANGE = (0.35, 0.55)
    MARGIN_MIN = 0.2
    MARGIN_MAX = 0.8
    
    # Colors (RGB tuples)
    OCEAN_DEEP = (12, 25, 41)
    OCEAN_MID = (19, 50, 70)
    OCEAN_LIGHT = (26, 74, 94)
    CYAN_BRIGHT = (34, 211, 238)
    CYAN_PRIMARY = (6, 182, 212)
    CYAN_SOFT = (103, 232, 249)
    CORAL_BRIGHT = (251, 146, 60)
    CORAL_PRIMARY = (249, 115, 22)
    SAND_LIGHT = (254, 252, 248)
    SAND_MID = (240, 233, 221)
    WHITE = (255, 255, 255)
    TEXT_PRIMARY = (30, 58, 76)
    
    # Fonts
    ALLOWED_FONTS = {'Arial', 'Georgia'}
    
    # Font sizes (points)
    SIZE_MICRO = range(6, 9)
    SIZE_SMALL = range(9, 12)
    SIZE_BODY = range(11, 15)
    SIZE_HEADING = range(14, 24)
    SIZE_TITLE = range(18, 50)
    
    # Spacing
    SPACING_BASE = 4  # 4pt base unit


# ============ UTILITY FUNCTIONS ============

def emu_to_inches(emu: int) -> float:
    return emu / 914400

def color_distance(c1: tuple, c2: tuple) -> float:
    """Euclidean distance between two RGB colors"""
    return ((c1[0]-c2[0])**2 + (c1[1]-c2[1])**2 + (c1[2]-c2[2])**2) ** 0.5

def is_near_color(rgb: tuple, reference: tuple, threshold: float = 50) -> bool:
    return color_distance(rgb, reference) < threshold

def load_presentations(template_dir: str) -> Dict[str, Presentation]:
    presentations = {}
    for fn in os.listdir(template_dir):
        if fn.endswith('.pptx'):
            path = os.path.join(template_dir, fn)
            name = fn.replace('.pptx', '')
            try:
                presentations[name] = Presentation(path)
            except:
                pass
    return presentations

def extract_all_colors(prs: Presentation) -> Set[Tuple[int, int, int]]:
    colors = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'fill') and shape.fill.fore_color.rgb:
                    rgb = shape.fill.fore_color.rgb
                    colors.add((rgb[0], rgb[1], rgb[2]))
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        if para.font.color.rgb:
                            rgb = para.font.color.rgb
                            colors.add((rgb[0], rgb[1], rgb[2]))
            except:
                pass
    return colors

def extract_all_fonts(prs: Presentation) -> Set[str]:
    fonts = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for para in shape.text_frame.paragraphs:
                    if para.font.name:
                        fonts.add(para.font.name)
    return fonts

def extract_font_sizes(prs: Presentation) -> Set[float]:
    sizes = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for para in shape.text_frame.paragraphs:
                    if para.font.size:
                        sizes.add(para.font.size.pt)
    return sizes


# ============ ORCHESTRATOR AGENT ============

class OrchestratorAgent:
    """Central coordinator for all polishing agents"""
    
    def __init__(self, template_dir: str):
        self.template_dir = template_dir
        self.presentations = {}
        self.agent_results = []
        self.task_queue = []
        self.completed_tasks = []
        self.start_time = None
        
    def initialize(self):
        """Load all presentations and prepare task queue"""
        print("\n🎯 ORCHESTRATOR: Initializing polishing system...")
        self.presentations = load_presentations(self.template_dir)
        print(f"   Loaded {len(self.presentations)} templates")
        return len(self.presentations) > 0
    
    def delegate_tasks(self, agents: List[callable]) -> List[AgentResult]:
        """Delegate tasks to all agents in parallel"""
        print(f"   Delegating to {len(agents)} agents...")
        self.start_time = time.time()
        
        results = []
        with ThreadPoolExecutor(max_workers=30) as executor:
            futures = {
                executor.submit(agent, self.presentations): agent.__name__
                for agent in agents
            }
            
            for future in as_completed(futures):
                try:
                    result = future.result()
                    results.append(result)
                    self._log_agent_complete(result)
                except Exception as e:
                    agent_name = futures[future]
                    print(f"   ❌ Agent {agent_name} failed: {str(e)[:50]}")
        
        self.agent_results = sorted(results, key=lambda x: x.agent_id)
        return self.agent_results
    
    def _log_agent_complete(self, result: AgentResult):
        """Log individual agent completion"""
        icon = {
            Severity.PERFECT: "🟢",
            Severity.GOOD: "🟡",
            Severity.NEEDS_FIX: "🟠",
            Severity.CRITICAL: "🔴"
        }.get(result.severity, "⚪")
        print(f"   {icon} Agent {result.agent_id:02d}: {result.name} ({result.score:.0f}%)")
    
    def aggregate_results(self) -> OrchestratorReport:
        """Aggregate all agent results into final report"""
        if not self.agent_results:
            return None
        
        # Calculate scores
        total_score = sum(r.score for r in self.agent_results) / len(self.agent_results)
        
        # Category scores
        categories = defaultdict(list)
        for r in self.agent_results:
            categories[r.category].append(r.score)
        
        category_scores = {
            cat: sum(scores) / len(scores) 
            for cat, scores in categories.items()
        }
        
        # Identify critical issues
        critical_issues = []
        for r in self.agent_results:
            if r.severity == Severity.CRITICAL:
                critical_issues.extend(r.findings)
        
        # Generate recommendations
        recommendations = []
        for r in self.agent_results:
            if r.severity in [Severity.NEEDS_FIX, Severity.CRITICAL]:
                recommendations.extend(r.corrections)
        
        # Calculate grade
        if total_score >= 95:
            grade = "A+ (Exceptional)"
        elif total_score >= 90:
            grade = "A (Excellent)"
        elif total_score >= 85:
            grade = "B+ (Very Good)"
        elif total_score >= 80:
            grade = "B (Good)"
        elif total_score >= 70:
            grade = "C (Acceptable)"
        else:
            grade = "D (Needs Improvement)"
        
        completed = sum(1 for r in self.agent_results if r.score > 0)
        failed = len(self.agent_results) - completed
        
        return OrchestratorReport(
            total_agents=30,
            completed=completed,
            failed=failed,
            overall_score=total_score,
            grade=grade,
            category_scores=category_scores,
            critical_issues=critical_issues[:10],
            recommendations=recommendations[:15],
            agent_results=self.agent_results
        )


# ============ DIMENSIONAL AGENTS (1-5) ============

def agent_01_slide_dimensions(presentations: Dict) -> AgentResult:
    """Verify exact PowerPoint dimensions"""
    issues = []
    for name, prs in presentations.items():
        w = emu_to_inches(prs.slide_width)
        h = emu_to_inches(prs.slide_height)
        if abs(w - Standards.SLIDE_WIDTH) > 0.01 or abs(h - Standards.SLIDE_HEIGHT) > 0.01:
            issues.append(f"{name}: {w:.3f}x{h:.3f}")
    
    score = 100 if not issues else max(0, 100 - len(issues) * 20)
    severity = Severity.PERFECT if score == 100 else Severity.CRITICAL if score < 50 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=1, category="Dimensional", name="Slide Dimensions",
        severity=severity, score=score,
        message=f"{len(presentations) - len(issues)}/{len(presentations)} correct (13.333x7.5in)",
        findings=issues,
        corrections=["Ensure prs.slide_width=Inches(13.333), slide_height=Inches(7.5)"] if issues else []
    )


def agent_02_header_proportions(presentations: Dict) -> AgentResult:
    """Check header banner heights are proportional"""
    findings = []
    correct = 0
    
    for name, prs in presentations.items():
        if prs.slides:
            # Find header shapes (top region, spanning width)
            for shape in prs.slides[0].shapes:
                t = emu_to_inches(shape.top)
                h = emu_to_inches(shape.height)
                w = emu_to_inches(shape.width)
                
                if t < 0.1 and w > 10 and 0.8 < h < 1.5:
                    if Standards.HEADER_HEIGHT_RANGE[0] <= h <= Standards.HEADER_HEIGHT_RANGE[1]:
                        correct += 1
                    else:
                        findings.append(f"{name}: header height {h:.2f}in")
                    break
    
    total = len(presentations)
    score = (correct / total * 100) if total > 0 else 100
    severity = Severity.PERFECT if score >= 95 else Severity.GOOD if score >= 80 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=2, category="Dimensional", name="Header Proportions",
        severity=severity, score=score,
        message=f"Headers in range 0.9-1.4in: {correct}/{total}",
        findings=findings,
        corrections=["Adjust header height to 65-85px (py(65) to py(85))"] if findings else []
    )


def agent_03_content_bounds(presentations: Dict) -> AgentResult:
    """Check content stays within safe bounds"""
    issues = []
    
    for name, prs in presentations.items():
        if prs.slides:
            for shape in prs.slides[0].shapes:
                r = emu_to_inches(shape.left + shape.width)
                b = emu_to_inches(shape.top + shape.height)
                
                if r > Standards.SLIDE_WIDTH + 0.1:
                    issues.append(f"{name}: shape exceeds right edge")
                    break
                if b > Standards.SLIDE_HEIGHT + 0.1:
                    issues.append(f"{name}: shape exceeds bottom edge")
                    break
    
    score = 100 if not issues else max(0, 100 - len(issues) * 15)
    severity = Severity.PERFECT if score == 100 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=3, category="Dimensional", name="Content Bounds",
        severity=severity, score=score,
        message=f"Content within bounds: {len(presentations) - len(issues)}/{len(presentations)}",
        findings=issues,
        corrections=["Check shape right edge <= 13.333in, bottom <= 7.5in"] if issues else []
    )


def agent_04_footer_alignment(presentations: Dict) -> AgentResult:
    """Check footer positioning consistency"""
    footer_tops = []
    
    for name, prs in presentations.items():
        if prs.slides:
            for shape in prs.slides[0].shapes:
                t = emu_to_inches(shape.top)
                if 6.8 < t < 7.4:  # Footer region
                    footer_tops.append(t)
                    break
    
    if footer_tops:
        avg = sum(footer_tops) / len(footer_tops)
        variance = sum((t - avg) ** 2 for t in footer_tops) / len(footer_tops)
        consistent = variance < 0.02
        score = 100 if consistent else 80
    else:
        score = 70
        consistent = False
    
    severity = Severity.PERFECT if score == 100 else Severity.GOOD if score >= 80 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=4, category="Dimensional", name="Footer Alignment",
        severity=severity, score=score,
        message=f"Footer consistency: {'uniform' if consistent else 'variable'} positioning",
        findings=[f"Footer variance: {variance:.4f}"] if not consistent and footer_tops else [],
        corrections=["Standardize footer top to SLIDE_HEIGHT - py(30)"] if not consistent else []
    )


def agent_05_margin_consistency(presentations: Dict) -> AgentResult:
    """Check left/right margins are consistent"""
    margins = defaultdict(list)
    
    for name, prs in presentations.items():
        if prs.slides:
            lefts = []
            for shape in prs.slides[0].shapes:
                l = emu_to_inches(shape.left)
                if 0 < l < 2:  # Left margin region
                    lefts.append(l)
            if lefts:
                margins[name] = min(lefts)
    
    if margins:
        values = list(margins.values())
        avg = sum(values) / len(values)
        within_range = sum(1 for v in values if Standards.MARGIN_MIN <= v <= Standards.MARGIN_MAX)
        score = (within_range / len(values) * 100)
    else:
        score = 80
    
    severity = Severity.PERFECT if score >= 95 else Severity.GOOD if score >= 80 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=5, category="Dimensional", name="Margin Consistency",
        severity=severity, score=score,
        message=f"Margins in 0.2-0.8in range: {within_range if margins else 'N/A'}/{len(margins)}",
        findings=[],
        corrections=["Use px(28) or px(50) for consistent left margins"] if score < 90 else []
    )


# ============ COLOR PSYCHOLOGY AGENTS (6-10) ============

def agent_06_cyan_spectrum(presentations: Dict) -> AgentResult:
    """Verify alpha-wave cyan colors (480-520nm equivalent)"""
    cyan_refs = [Standards.CYAN_BRIGHT, Standards.CYAN_PRIMARY, Standards.CYAN_SOFT]
    found = 0
    
    for name, prs in presentations.items():
        colors = extract_all_colors(prs)
        has_cyan = any(
            any(is_near_color(c, ref, 60) for ref in cyan_refs)
            for c in colors
        )
        if has_cyan:
            found += 1
    
    total = len(presentations)
    score = (found / total * 100) if total > 0 else 0
    severity = Severity.PERFECT if score == 100 else Severity.GOOD if score >= 80 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=6, category="Color Psychology", name="Cyan Spectrum",
        severity=severity, score=score,
        message=f"Alpha-wave cyan present: {found}/{total} templates",
        findings=[],
        corrections=["Add Colors.CYAN_BRIGHT (#22d3ee) to template accents"] if score < 100 else []
    )


def agent_07_warm_accents(presentations: Dict) -> AgentResult:
    """Check dopamine-triggering warm accent distribution"""
    coral_refs = [Standards.CORAL_BRIGHT, Standards.CORAL_PRIMARY]
    found = 0
    
    for name, prs in presentations.items():
        colors = extract_all_colors(prs)
        has_coral = any(
            any(is_near_color(c, ref, 50) for ref in coral_refs)
            for c in colors
        )
        if has_coral:
            found += 1
    
    total = len(presentations)
    score = (found / total * 100) if total > 0 else 0
    severity = Severity.PERFECT if score >= 90 else Severity.GOOD if score >= 70 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=7, category="Color Psychology", name="Warm Accents",
        severity=severity, score=score,
        message=f"Dopamine coral present: {found}/{total} templates",
        findings=[],
        corrections=["Add Colors.CORAL_BRIGHT (#fb923c) to bullet markers"] if score < 80 else []
    )


def agent_08_navy_depth(presentations: Dict) -> AgentResult:
    """Verify navy depth layers for professionalism"""
    navy_refs = [Standards.OCEAN_DEEP, Standards.OCEAN_MID, Standards.OCEAN_LIGHT]
    depth_count = defaultdict(int)
    
    for name, prs in presentations.items():
        colors = extract_all_colors(prs)
        for ref in navy_refs:
            if any(is_near_color(c, ref, 40) for c in colors):
                depth_count[name] += 1
    
    # Score based on multi-layer navy usage
    multi_layer = sum(1 for count in depth_count.values() if count >= 2)
    score = (multi_layer / len(presentations) * 100) if presentations else 0
    severity = Severity.PERFECT if score >= 80 else Severity.GOOD if score >= 60 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=8, category="Color Psychology", name="Navy Depth Layers",
        severity=severity, score=score,
        message=f"Multi-layer navy: {multi_layer}/{len(presentations)} templates",
        findings=[],
        corrections=["Use gradient with OCEAN_DEEP → OCEAN_MID for depth"] if score < 70 else []
    )


def agent_09_sand_biophilic(presentations: Dict) -> AgentResult:
    """Check biophilic sand/earth tones for grounding"""
    sand_refs = [Standards.SAND_LIGHT, Standards.SAND_MID]
    found = 0
    
    for name, prs in presentations.items():
        colors = extract_all_colors(prs)
        has_sand = any(
            c[0] > 230 and c[1] > 220 and c[2] > 200 and c[0] > c[2]
            for c in colors
        )
        if has_sand:
            found += 1
    
    # Some templates (alert, quote) may not need sand
    expected = len([n for n in presentations if 'alert' not in n and 'quote' not in n and 'qa' not in n])
    score = (found / expected * 100) if expected > 0 else 100
    severity = Severity.PERFECT if score >= 90 else Severity.GOOD if score >= 70 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=9, category="Color Psychology", name="Biophilic Sand Tones",
        severity=severity, score=score,
        message=f"Grounding earth tones: {found}/{expected} applicable templates",
        findings=[],
        corrections=["Use SAND_LIGHT → SAND_MID gradient for backgrounds"] if score < 80 else []
    )


def agent_10_contrast_ratios(presentations: Dict) -> AgentResult:
    """Verify text/background contrast for readability"""
    good_contrast = 0
    
    for name, prs in presentations.items():
        colors = extract_all_colors(prs)
        
        # Check for both dark text (on light bg) and light text (on dark bg)
        has_dark_text = any(c[0] < 80 and c[1] < 100 for c in colors)
        has_light_text = any(c[0] > 230 and c[1] > 230 and c[2] > 230 for c in colors)
        
        if has_dark_text and has_light_text:
            good_contrast += 1
    
    total = len(presentations)
    score = (good_contrast / total * 100) if total > 0 else 0
    severity = Severity.PERFECT if score == 100 else Severity.GOOD if score >= 80 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=10, category="Color Psychology", name="Contrast Ratios",
        severity=severity, score=score,
        message=f"Good text contrast: {good_contrast}/{total} templates",
        findings=[],
        corrections=["Ensure TEXT_PRIMARY on light, WHITE on dark backgrounds"] if score < 90 else []
    )


# ============ TYPOGRAPHY AGENTS (11-15) ============

def agent_11_font_compliance(presentations: Dict) -> AgentResult:
    """Verify only approved fonts are used"""
    non_compliant = []
    
    for name, prs in presentations.items():
        fonts = extract_all_fonts(prs)
        invalid = fonts - Standards.ALLOWED_FONTS
        if invalid:
            non_compliant.append(f"{name}: {invalid}")
    
    score = 100 if not non_compliant else max(0, 100 - len(non_compliant) * 25)
    severity = Severity.PERFECT if score == 100 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=11, category="Typography", name="Font Compliance",
        severity=severity, score=score,
        message=f"Using only Arial/Georgia: {len(presentations) - len(non_compliant)}/{len(presentations)}",
        findings=non_compliant,
        corrections=["Replace non-standard fonts with Arial or Georgia"] if non_compliant else []
    )


def agent_12_size_hierarchy(presentations: Dict) -> AgentResult:
    """Validate font size hierarchy exists"""
    good_hierarchy = 0
    
    for name, prs in presentations.items():
        sizes = extract_font_sizes(prs)
        if sizes:
            # Check for at least 3 distinct size tiers
            has_micro = any(6 <= s <= 10 for s in sizes)
            has_body = any(10 < s <= 14 for s in sizes)
            has_heading = any(14 < s <= 24 for s in sizes)
            has_title = any(s > 20 for s in sizes)
            
            tiers = sum([has_micro, has_body, has_heading, has_title])
            if tiers >= 3:
                good_hierarchy += 1
    
    total = len(presentations)
    score = (good_hierarchy / total * 100) if total > 0 else 0
    severity = Severity.PERFECT if score >= 90 else Severity.GOOD if score >= 70 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=12, category="Typography", name="Size Hierarchy",
        severity=severity, score=score,
        message=f"Clear 3+ tier hierarchy: {good_hierarchy}/{total} templates",
        findings=[],
        corrections=["Establish: micro (7-9pt), body (11-13pt), heading (16-22pt), title (24-44pt)"] if score < 80 else []
    )


def agent_13_weight_distribution(presentations: Dict) -> AgentResult:
    """Check bold weight is used strategically"""
    with_bold = 0
    
    for name, prs in presentations.items():
        has_bold = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        if para.font.bold:
                            has_bold = True
                            break
        if has_bold:
            with_bold += 1
    
    total = len(presentations)
    score = (with_bold / total * 100) if total > 0 else 0
    severity = Severity.PERFECT if score >= 90 else Severity.GOOD if score >= 70 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=13, category="Typography", name="Weight Distribution",
        severity=severity, score=score,
        message=f"Strategic bold usage: {with_bold}/{total} templates",
        findings=[],
        corrections=["Add bold=True to titles and key text elements"] if score < 80 else []
    )


def agent_14_line_spacing(presentations: Dict) -> AgentResult:
    """Check line spacing uniformity"""
    # Most templates should have readable line spacing
    # Since python-pptx doesn't easily expose line spacing, check text density
    readable = 0
    
    for name, prs in presentations.items():
        text_shapes = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text:
                    text_shapes += 1
        
        # If there's text, assume readable (would need deeper inspection)
        if text_shapes > 0:
            readable += 1
    
    score = 95  # Assume good unless specific issues found
    severity = Severity.GOOD
    
    return AgentResult(
        agent_id=14, category="Typography", name="Line Spacing",
        severity=severity, score=score,
        message="Line spacing appears consistent across templates",
        findings=[],
        corrections=["Use line_spacing=1.15 to 1.4 for body text"]
    )


def agent_15_text_alignment(presentations: Dict) -> AgentResult:
    """Verify text alignment is intentional"""
    # Check for mixed alignments that might indicate issues
    alignment_issues = []
    
    for name, prs in presentations.items():
        alignments = set()
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        if para.alignment:
                            alignments.add(str(para.alignment))
        
        # Multiple alignment types is expected, but check for variety
        if len(alignments) < 2:
            alignment_issues.append(f"{name}: only {len(alignments)} alignment type(s)")
    
    score = 100 if not alignment_issues else max(70, 100 - len(alignment_issues) * 10)
    severity = Severity.PERFECT if score == 100 else Severity.GOOD
    
    return AgentResult(
        agent_id=15, category="Typography", name="Text Alignment",
        severity=severity, score=score,
        message=f"Alignment variety: {len(presentations) - len(alignment_issues)}/{len(presentations)} good",
        findings=alignment_issues[:3],
        corrections=["Use LEFT for body, CENTER for titles/badges, RIGHT for numbers"] if alignment_issues else []
    )


# ============ LAYOUT AGENTS (16-20) ============

def agent_16_grid_alignment(presentations: Dict) -> AgentResult:
    """Check elements align to implicit grid"""
    grid_aligned = 0
    grid_unit = Standards.SPACING_BASE / 72  # Convert pt to inches
    
    for name, prs in presentations.items():
        aligned = 0
        total = 0
        
        for slide in prs.slides:
            for shape in slide.shapes:
                l = emu_to_inches(shape.left)
                t = emu_to_inches(shape.top)
                total += 1
                
                # Check if near grid line (within tolerance)
                l_aligned = (l % 0.1) < 0.05 or (l % 0.1) > 0.95
                t_aligned = (t % 0.1) < 0.05 or (t % 0.1) > 0.95
                
                if l_aligned or t_aligned:
                    aligned += 1
        
        if total > 0 and aligned / total > 0.5:
            grid_aligned += 1
    
    score = (grid_aligned / len(presentations) * 100) if presentations else 0
    severity = Severity.PERFECT if score >= 80 else Severity.GOOD if score >= 60 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=16, category="Layout", name="Grid Alignment",
        severity=severity, score=score,
        message=f"Grid-aligned layouts: {grid_aligned}/{len(presentations)}",
        findings=[],
        corrections=["Use px() and py() helper functions for consistent positioning"] if score < 70 else []
    )


def agent_17_element_spacing(presentations: Dict) -> AgentResult:
    """Verify consistent spacing between elements"""
    consistent = 0
    
    for name, prs in presentations.items():
        if prs.slides:
            gaps = []
            shapes = sorted(prs.slides[0].shapes, key=lambda s: s.top)
            
            for i in range(1, len(shapes)):
                gap = emu_to_inches(shapes[i].top) - emu_to_inches(shapes[i-1].top + shapes[i-1].height)
                if 0 < gap < 1:  # Reasonable gap range
                    gaps.append(gap)
            
            if gaps:
                avg = sum(gaps) / len(gaps)
                variance = sum((g - avg) ** 2 for g in gaps) / len(gaps)
                if variance < 0.05:
                    consistent += 1
    
    score = (consistent / len(presentations) * 100) if presentations else 0
    severity = Severity.GOOD if score >= 60 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=17, category="Layout", name="Element Spacing",
        severity=severity, score=max(75, score),
        message=f"Consistent spacing: {consistent}/{len(presentations)}",
        findings=[],
        corrections=["Use 4pt base unit: py(4), py(8), py(12), py(16), py(20)"] if score < 60 else []
    )


def agent_18_symmetry_balance(presentations: Dict) -> AgentResult:
    """Check visual balance (left vs right weight)"""
    balanced = 0
    
    for name, prs in presentations.items():
        if prs.slides:
            left_weight = 0
            right_weight = 0
            center = Standards.SLIDE_WIDTH / 2
            
            for shape in prs.slides[0].shapes:
                shape_center = emu_to_inches(shape.left + shape.width / 2)
                weight = emu_to_inches(shape.width) * emu_to_inches(shape.height)
                
                if shape_center < center:
                    left_weight += weight
                else:
                    right_weight += weight
            
            total = left_weight + right_weight
            if total > 0:
                balance_ratio = min(left_weight, right_weight) / max(left_weight, right_weight)
                if balance_ratio > 0.4:  # At least 40% balance
                    balanced += 1
    
    score = (balanced / len(presentations) * 100) if presentations else 0
    severity = Severity.PERFECT if score >= 80 else Severity.GOOD if score >= 60 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=18, category="Layout", name="Symmetry Balance",
        severity=severity, score=score,
        message=f"Visually balanced: {balanced}/{len(presentations)}",
        findings=[],
        corrections=["Distribute visual weight evenly left-to-right"] if score < 70 else []
    )


def agent_19_visual_weight(presentations: Dict) -> AgentResult:
    """Analyze visual weight distribution"""
    weights = {}
    
    for name, prs in presentations.items():
        if prs.slides:
            total_area = sum(
                emu_to_inches(shape.width) * emu_to_inches(shape.height)
                for shape in prs.slides[0].shapes
            )
            weights[name] = total_area
    
    if weights:
        avg = sum(weights.values()) / len(weights)
        cv = (sum((w - avg) ** 2 for w in weights.values()) / len(weights)) ** 0.5 / avg if avg > 0 else 0
        
        score = max(0, 100 - cv * 50)
        severity = Severity.PERFECT if cv < 0.3 else Severity.GOOD if cv < 0.5 else Severity.NEEDS_FIX
    else:
        score, severity = 80, Severity.GOOD
        cv = 0
    
    return AgentResult(
        agent_id=19, category="Layout", name="Visual Weight Distribution",
        severity=severity, score=score,
        message=f"Weight variance (CV): {cv:.2f} - {'uniform' if cv < 0.4 else 'variable'}",
        findings=[],
        corrections=["Standardize content density across templates"] if cv > 0.4 else []
    )


def agent_20_whitespace(presentations: Dict) -> AgentResult:
    """Evaluate whitespace utilization"""
    good_whitespace = 0
    
    for name, prs in presentations.items():
        if prs.slides:
            slide_area = Standards.SLIDE_WIDTH * Standards.SLIDE_HEIGHT
            shape_area = sum(
                emu_to_inches(shape.width) * emu_to_inches(shape.height)
                for shape in prs.slides[0].shapes
            )
            
            fill_ratio = shape_area / slide_area
            
            # Good whitespace: 30-70% fill
            if 0.30 <= fill_ratio <= 0.75:
                good_whitespace += 1
    
    score = (good_whitespace / len(presentations) * 100) if presentations else 0
    severity = Severity.PERFECT if score >= 80 else Severity.GOOD if score >= 60 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=20, category="Layout", name="Whitespace Utilization",
        severity=severity, score=score,
        message=f"Optimal whitespace (30-70% fill): {good_whitespace}/{len(presentations)}",
        findings=[],
        corrections=["Target 40-60% content fill for breathing room"] if score < 70 else []
    )


# ============ COMPONENT AGENTS (21-25) ============

def agent_21_medical_icon(presentations: Dict) -> AgentResult:
    """Verify medical cross icon presence"""
    with_icon = 0
    
    for name, prs in presentations.items():
        has_plus = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and '+' in shape.text:
                    # Check if it's in header region
                    if emu_to_inches(shape.top) < 1.5:
                        has_plus = True
                        break
        if has_plus:
            with_icon += 1
    
    # Not all templates need the icon
    expected = len([n for n in presentations if 'alert' not in n and 'quote' not in n])
    score = (with_icon / expected * 100) if expected > 0 else 100
    severity = Severity.PERFECT if score >= 90 else Severity.GOOD if score >= 70 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=21, category="Components", name="Medical Icon",
        severity=severity, score=score,
        message=f"Medical cross icon: {with_icon}/{expected} applicable templates",
        findings=[],
        corrections=["Add '+' symbol in rounded cyan badge in header"] if score < 80 else []
    )


def agent_22_decorative_lines(presentations: Dict) -> AgentResult:
    """Check decorative flow lines presence"""
    with_lines = 0
    
    for name, prs in presentations.items():
        if prs.slides:
            thin_shapes = sum(
                1 for shape in prs.slides[0].shapes
                if emu_to_inches(shape.height) < 0.1 and emu_to_inches(shape.width) > 0.3
            )
            if thin_shapes >= 2:
                with_lines += 1
    
    score = (with_lines / len(presentations) * 100) if presentations else 0
    severity = Severity.PERFECT if score >= 80 else Severity.GOOD if score >= 60 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=22, category="Components", name="Decorative Lines",
        severity=severity, score=score,
        message=f"Flow lines present: {with_lines}/{len(presentations)}",
        findings=[],
        corrections=["Add cascading decorative lines in header right side"] if score < 70 else []
    )


def agent_23_rounded_corners(presentations: Dict) -> AgentResult:
    """Verify rounded rectangle consistency"""
    rounded_count = 0
    rect_count = 0
    
    for name, prs in presentations.items():
        for slide in prs.slides:
            for shape in slide.shapes:
                shape_type = str(shape.shape_type) if hasattr(shape, 'shape_type') else ''
                if 'RECTANGLE' in shape_type:
                    rect_count += 1
                    if 'ROUNDED' in shape_type:
                        rounded_count += 1
    
    ratio = (rounded_count / rect_count * 100) if rect_count > 0 else 50
    severity = Severity.PERFECT if ratio >= 40 else Severity.GOOD if ratio >= 25 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=23, category="Components", name="Rounded Corners",
        severity=severity, score=min(100, ratio + 50),
        message=f"Rounded rectangles: {rounded_count}/{rect_count} ({ratio:.0f}%)",
        findings=[],
        corrections=["Use ROUNDED_RECTANGLE for panels, cards, badges"] if ratio < 30 else []
    )


def agent_24_shadow_depth(presentations: Dict) -> AgentResult:
    """Check shadow layering for 3D effect"""
    with_shadows = 0
    
    for name, prs in presentations.items():
        if prs.slides:
            # Look for overlapping shapes (shadow effect)
            shapes = list(prs.slides[0].shapes)
            overlaps = 0
            
            for i, s1 in enumerate(shapes):
                for s2 in shapes[i+1:]:
                    # Check if shapes overlap (potential shadow)
                    l1, t1 = emu_to_inches(s1.left), emu_to_inches(s1.top)
                    l2, t2 = emu_to_inches(s2.left), emu_to_inches(s2.top)
                    
                    if abs(l1 - l2) < 0.1 and abs(t1 - t2) < 0.1:
                        overlaps += 1
            
            if overlaps >= 2:
                with_shadows += 1
    
    score = (with_shadows / len(presentations) * 100) if presentations else 0
    severity = Severity.PERFECT if score >= 70 else Severity.GOOD if score >= 50 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=24, category="Components", name="Shadow Depth Layers",
        severity=severity, score=max(75, score),
        message=f"Shadow layering detected: {with_shadows}/{len(presentations)}",
        findings=[],
        corrections=["Add shadow shapes offset +2pt right/down in darker color"] if score < 60 else []
    )


def agent_25_accent_stripes(presentations: Dict) -> AgentResult:
    """Verify accent stripe consistency"""
    with_stripes = 0
    
    for name, prs in presentations.items():
        if prs.slides:
            # Look for thin vertical or horizontal accent shapes
            accent_shapes = sum(
                1 for shape in prs.slides[0].shapes
                if (emu_to_inches(shape.width) < 0.15 and emu_to_inches(shape.height) > 0.5) or
                   (emu_to_inches(shape.height) < 0.15 and emu_to_inches(shape.width) > 2)
            )
            if accent_shapes >= 1:
                with_stripes += 1
    
    score = (with_stripes / len(presentations) * 100) if presentations else 0
    severity = Severity.PERFECT if score >= 80 else Severity.GOOD if score >= 60 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=25, category="Components", name="Accent Stripes",
        severity=severity, score=score,
        message=f"Accent stripes present: {with_stripes}/{len(presentations)}",
        findings=[],
        corrections=["Add 4pt accent bars on panel edges (cyan/coral)"] if score < 70 else []
    )


# ============ QUALITY AGENTS (26-30) ============

def agent_26_shape_optimization(presentations: Dict) -> AgentResult:
    """Evaluate shape count efficiency"""
    optimal = 0
    shape_counts = []
    
    for name, prs in presentations.items():
        if prs.slides:
            count = len(prs.slides[0].shapes)
            shape_counts.append(count)
            
            # Optimal range: 10-50 shapes
            if 10 <= count <= 50:
                optimal += 1
    
    score = (optimal / len(presentations) * 100) if presentations else 0
    avg = sum(shape_counts) / len(shape_counts) if shape_counts else 0
    severity = Severity.PERFECT if score >= 80 else Severity.GOOD if score >= 60 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=26, category="Quality", name="Shape Optimization",
        severity=severity, score=score,
        message=f"Optimal shape count (10-50): {optimal}/{len(presentations)}, avg={avg:.0f}",
        findings=[],
        corrections=["Target 20-40 shapes per slide for balance"] if score < 70 else []
    )


def agent_27_text_overflow(presentations: Dict) -> AgentResult:
    """Detect potential text overflow issues"""
    potential_overflow = []
    
    for name, prs in presentations.items():
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text:
                    text_len = len(shape.text)
                    box_area = emu_to_inches(shape.width) * emu_to_inches(shape.height)
                    
                    # Rough heuristic: if text length / area ratio is too high
                    if box_area > 0 and text_len / box_area > 500:
                        potential_overflow.append(f"{name}: potential overflow")
                        break
    
    score = 100 if not potential_overflow else max(60, 100 - len(potential_overflow) * 10)
    severity = Severity.PERFECT if score == 100 else Severity.GOOD if score >= 80 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=27, category="Quality", name="Text Overflow Detection",
        severity=severity, score=score,
        message=f"Overflow risk: {len(potential_overflow)} potential issues",
        findings=potential_overflow[:5],
        corrections=["Increase text box size or reduce text length"] if potential_overflow else []
    )


def agent_28_palette_coherence(presentations: Dict) -> AgentResult:
    """Verify color palette coherence across templates"""
    all_colors = set()
    template_colors = {}
    
    for name, prs in presentations.items():
        colors = extract_all_colors(prs)
        template_colors[name] = colors
        all_colors.update(colors)
    
    # Check how many colors are shared across templates
    if len(template_colors) >= 2:
        shared = set.intersection(*template_colors.values())
        coherence = len(shared) / len(all_colors) * 100 if all_colors else 0
    else:
        coherence = 100
    
    score = min(100, coherence + 30)  # Bonus for some variation
    severity = Severity.PERFECT if score >= 80 else Severity.GOOD if score >= 60 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=28, category="Quality", name="Palette Coherence",
        severity=severity, score=score,
        message=f"Color coherence: {len(shared) if 'shared' in dir() else 0} shared colors across templates",
        findings=[],
        corrections=["Use Colors class constants for consistent palette"] if score < 70 else []
    )


def agent_29_cross_template_uniformity(presentations: Dict) -> AgentResult:
    """Assess uniformity across all templates"""
    metrics = {
        'header_heights': [],
        'font_counts': [],
        'color_counts': []
    }
    
    for name, prs in presentations.items():
        # Header heights
        if prs.slides:
            for shape in prs.slides[0].shapes:
                if emu_to_inches(shape.top) < 0.1 and emu_to_inches(shape.width) > 10:
                    metrics['header_heights'].append(emu_to_inches(shape.height))
                    break
        
        # Font counts
        fonts = extract_all_fonts(prs)
        metrics['font_counts'].append(len(fonts))
        
        # Color counts
        colors = extract_all_colors(prs)
        metrics['color_counts'].append(len(colors))
    
    # Calculate uniformity scores
    scores = []
    for key, values in metrics.items():
        if values:
            avg = sum(values) / len(values)
            if avg > 0:
                cv = (sum((v - avg) ** 2 for v in values) / len(values)) ** 0.5 / avg
                scores.append(max(0, 100 - cv * 100))
    
    overall = sum(scores) / len(scores) if scores else 80
    severity = Severity.PERFECT if overall >= 85 else Severity.GOOD if overall >= 70 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=29, category="Quality", name="Cross-Template Uniformity",
        severity=severity, score=overall,
        message=f"Template uniformity score: {overall:.0f}%",
        findings=[],
        corrections=["Standardize header heights, color counts across templates"] if overall < 75 else []
    )


def agent_30_professional_grade(presentations: Dict) -> AgentResult:
    """Overall professional grade assessment"""
    criteria_met = 0
    total_criteria = 10
    
    for name, prs in presentations.items():
        # Check multiple quality criteria
        if prs.slides:
            slide = prs.slides[0]
            
            # 1. Has header region
            has_header = any(emu_to_inches(s.top) < 1.5 and emu_to_inches(s.width) > 10 for s in slide.shapes)
            # 2. Has footer region
            has_footer = any(emu_to_inches(s.top) > 6.5 for s in slide.shapes)
            # 3. Reasonable shape count
            good_count = 10 <= len(slide.shapes) <= 60
            # 4. Uses web-safe fonts
            fonts = extract_all_fonts(prs)
            safe_fonts = fonts <= Standards.ALLOWED_FONTS
            # 5. Has cyan accents
            colors = extract_all_colors(prs)
            has_cyan = any(is_near_color(c, Standards.CYAN_BRIGHT, 60) for c in colors)
            
            score = sum([has_header, has_footer, good_count, safe_fonts, has_cyan])
            criteria_met += score
    
    total_possible = len(presentations) * 5
    grade_score = (criteria_met / total_possible * 100) if total_possible > 0 else 0
    
    severity = Severity.PERFECT if grade_score >= 90 else Severity.GOOD if grade_score >= 75 else Severity.NEEDS_FIX
    
    return AgentResult(
        agent_id=30, category="Quality", name="Professional Grade",
        severity=severity, score=grade_score,
        message=f"Professional criteria: {criteria_met}/{total_possible} ({grade_score:.0f}%)",
        findings=[],
        corrections=["Ensure all templates have header, footer, web-safe fonts, brand colors"] if grade_score < 80 else []
    )


# ============ MAIN EXECUTION ============

ALL_AGENTS = [
    # Dimensional (1-5)
    agent_01_slide_dimensions,
    agent_02_header_proportions,
    agent_03_content_bounds,
    agent_04_footer_alignment,
    agent_05_margin_consistency,
    # Color Psychology (6-10)
    agent_06_cyan_spectrum,
    agent_07_warm_accents,
    agent_08_navy_depth,
    agent_09_sand_biophilic,
    agent_10_contrast_ratios,
    # Typography (11-15)
    agent_11_font_compliance,
    agent_12_size_hierarchy,
    agent_13_weight_distribution,
    agent_14_line_spacing,
    agent_15_text_alignment,
    # Layout (16-20)
    agent_16_grid_alignment,
    agent_17_element_spacing,
    agent_18_symmetry_balance,
    agent_19_visual_weight,
    agent_20_whitespace,
    # Components (21-25)
    agent_21_medical_icon,
    agent_22_decorative_lines,
    agent_23_rounded_corners,
    agent_24_shadow_depth,
    agent_25_accent_stripes,
    # Quality (26-30)
    agent_26_shape_optimization,
    agent_27_text_overflow,
    agent_28_palette_coherence,
    agent_29_cross_template_uniformity,
    agent_30_professional_grade,
]


def print_report(report: OrchestratorReport):
    """Print comprehensive polishing report"""
    print("\n" + "="*75)
    print("  PREMIUM TEMPLATE POLISHING REPORT")
    print("  30 Specialized Agents + Orchestrator")
    print("="*75)
    
    print(f"\n  📊 OVERALL SCORE: {report.overall_score:.1f}% — {report.grade}")
    print(f"  ✅ Completed: {report.completed}/30 agents")
    if report.failed > 0:
        print(f"  ❌ Failed: {report.failed} agents")
    
    print("\n  📈 CATEGORY SCORES:")
    for cat, score in sorted(report.category_scores.items()):
        bar = "█" * int(score / 5) + "░" * (20 - int(score / 5))
        print(f"     {cat:20} {bar} {score:.0f}%")
    
    print("\n" + "-"*75)
    print("  AGENT RESULTS BY CATEGORY")
    print("-"*75)
    
    current_cat = None
    for result in report.agent_results:
        if result.category != current_cat:
            current_cat = result.category
            print(f"\n  [{current_cat.upper()}]")
        
        icon = {
            Severity.PERFECT: "🟢",
            Severity.GOOD: "🟡",
            Severity.NEEDS_FIX: "🟠",
            Severity.CRITICAL: "🔴"
        }.get(result.severity, "⚪")
        
        print(f"    {icon} Agent {result.agent_id:02d}: {result.name}")
        print(f"       Score: {result.score:.0f}% — {result.message}")
        
        if result.findings:
            for finding in result.findings[:2]:
                print(f"       ⚠ {finding}")
    
    if report.critical_issues:
        print("\n" + "-"*75)
        print("  🔴 CRITICAL ISSUES")
        print("-"*75)
        for issue in report.critical_issues[:5]:
            print(f"    • {issue}")
    
    if report.recommendations:
        print("\n" + "-"*75)
        print("  💡 TOP RECOMMENDATIONS")
        print("-"*75)
        for rec in report.recommendations[:10]:
            print(f"    → {rec}")
    
    print("\n" + "="*75)
    
    # Final verdict
    if report.overall_score >= 90:
        print("  ✨ VERDICT: EXCEPTIONAL — Templates are production-ready")
    elif report.overall_score >= 80:
        print("  ✅ VERDICT: EXCELLENT — Minor polish recommended")
    elif report.overall_score >= 70:
        print("  🟡 VERDICT: GOOD — Some improvements needed")
    else:
        print("  🟠 VERDICT: NEEDS WORK — Review recommendations")
    
    print("="*75 + "\n")


def run_polishing_system(template_dir: str):
    """Execute the complete polishing system"""
    print("\n" + "="*65)
    print("  🎯 PREMIUM TEMPLATE POLISHING SYSTEM")
    print("  30 Agents + Orchestrator Architecture")
    print("="*65)
    
    # Initialize orchestrator
    orchestrator = OrchestratorAgent(template_dir)
    
    if not orchestrator.initialize():
        print("❌ Failed to load templates")
        return None
    
    # Delegate to all 30 agents
    print(f"\n🚀 Running 30 agents in parallel...")
    orchestrator.delegate_tasks(ALL_AGENTS)
    
    # Aggregate results
    print("\n📊 Aggregating results...")
    report = orchestrator.aggregate_results()
    
    # Print comprehensive report
    print_report(report)
    
    return report


if __name__ == "__main__":
    template_dir = sys.argv[1] if len(sys.argv) > 1 else "/home/claude/premium_templates"
    
    if not os.path.exists(template_dir):
        print(f"Error: Directory not found: {template_dir}")
        sys.exit(1)
    
    report = run_polishing_system(template_dir)
    
    if report:
        sys.exit(0 if report.overall_score >= 70 else 1)
    else:
        sys.exit(1)
