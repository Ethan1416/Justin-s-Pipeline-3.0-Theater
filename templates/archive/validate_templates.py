#!/usr/bin/env python3
"""
=============================================================================
NURSING TEMPLATE VALIDATION SYSTEM
Complete validation for all 20 slide types
=============================================================================

Validates:
- Slide dimensions (16:9 aspect ratio)
- Color palette compliance
- Font compliance (Arial/Georgia only)
- Typography hierarchy
- Required components (header, footer, medical icon)
- Layout bounds
- Shape optimization

Usage:
    python validate_templates.py                      # Validate all in current dir
    python validate_templates.py /path/to/templates   # Validate specific directory
    python validate_templates.py --strict             # Fail on warnings
"""

import os
import sys
from dataclasses import dataclass, field
from typing import List, Dict, Set, Tuple, Optional
from collections import defaultdict
from pptx import Presentation
from pptx.util import Emu


# =============================================================================
# REFERENCE STANDARDS
# =============================================================================

class Standards:
    """Canonical reference values for validation"""
    
    # Dimensions
    SLIDE_WIDTH_INCHES = 13.333
    SLIDE_HEIGHT_INCHES = 7.5
    ASPECT_RATIO = 16 / 9
    
    # Colors (RGB tuples) - Primary palette
    REFERENCE_COLORS = {
        'ocean_deep': (12, 25, 41),
        'ocean_mid': (19, 50, 70),
        'ocean_light': (26, 74, 94),
        'cyan_bright': (34, 211, 238),
        'cyan_primary': (6, 182, 212),
        'cyan_soft': (103, 232, 249),
        'coral_bright': (251, 146, 60),
        'coral_primary': (249, 115, 22),
        'sand_light': (254, 252, 248),
        'sand_warm': (250, 245, 237),
        'sand_mid': (240, 233, 221),
        'white': (255, 255, 255),
        'text_primary': (30, 58, 76),
        'text_secondary': (61, 90, 108),
        'success': (16, 185, 129),
        'warning': (245, 158, 11),
        'danger': (239, 68, 68),
    }
    
    # Fonts
    ALLOWED_FONTS = {'Arial', 'Georgia'}
    
    # Font size ranges (points)
    SIZE_RANGES = {
        'micro': (6, 9),
        'small': (9, 12),
        'body': (10, 14),
        'heading': (14, 24),
        'title': (18, 50),
        'hero': (40, 70)
    }
    
    # Shape counts
    MIN_SHAPES = 5
    MAX_SHAPES = 80
    OPTIMAL_MIN = 10
    OPTIMAL_MAX = 50


# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================

def emu_to_inches(emu: int) -> float:
    """Convert EMU to inches"""
    return emu / 914400

def rgb_to_tuple(rgb) -> Tuple[int, int, int]:
    """Convert RGBColor to tuple"""
    try:
        return (rgb[0], rgb[1], rgb[2])
    except:
        return None

def color_distance(c1: tuple, c2: tuple) -> float:
    """Euclidean distance between RGB colors"""
    if not c1 or not c2:
        return float('inf')
    return ((c1[0]-c2[0])**2 + (c1[1]-c2[1])**2 + (c1[2]-c2[2])**2) ** 0.5

def is_approved_color(rgb: tuple, threshold: float = 60) -> bool:
    """Check if color is within threshold of any approved color"""
    if not rgb:
        return True
    for ref in Standards.REFERENCE_COLORS.values():
        if color_distance(rgb, ref) < threshold:
            return True
    return False


# =============================================================================
# VALIDATION RESULTS
# =============================================================================

@dataclass
class ValidationIssue:
    level: str  # 'ERROR', 'WARNING', 'INFO'
    category: str
    message: str
    details: str = ""

@dataclass
class TemplateValidation:
    filename: str
    is_valid: bool
    score: float
    issues: List[ValidationIssue] = field(default_factory=list)
    metrics: Dict = field(default_factory=dict)

@dataclass
class ValidationReport:
    templates_checked: int
    templates_passed: int
    templates_failed: int
    overall_score: float
    results: List[TemplateValidation] = field(default_factory=list)
    summary: Dict = field(default_factory=dict)


# =============================================================================
# VALIDATOR CLASS
# =============================================================================

class TemplateValidator:
    """Comprehensive template validator"""
    
    def __init__(self, strict: bool = False):
        self.strict = strict
        self.results = []
    
    def validate_file(self, filepath: str) -> TemplateValidation:
        """Validate a single PPTX file"""
        filename = os.path.basename(filepath)
        issues = []
        metrics = {}
        
        try:
            prs = Presentation(filepath)
        except Exception as e:
            return TemplateValidation(
                filename=filename,
                is_valid=False,
                score=0,
                issues=[ValidationIssue('ERROR', 'File', f'Cannot open file: {str(e)}')]
            )
        
        # Run all validations
        issues.extend(self._validate_dimensions(prs, metrics))
        issues.extend(self._validate_colors(prs, metrics))
        issues.extend(self._validate_fonts(prs, metrics))
        issues.extend(self._validate_typography(prs, metrics))
        issues.extend(self._validate_components(prs, metrics))
        issues.extend(self._validate_layout(prs, metrics))
        issues.extend(self._validate_shapes(prs, metrics))
        
        # Calculate score
        errors = sum(1 for i in issues if i.level == 'ERROR')
        warnings = sum(1 for i in issues if i.level == 'WARNING')
        
        score = 100 - (errors * 15) - (warnings * 5)
        score = max(0, min(100, score))
        
        is_valid = errors == 0 and (not self.strict or warnings == 0)
        
        return TemplateValidation(
            filename=filename,
            is_valid=is_valid,
            score=score,
            issues=issues,
            metrics=metrics
        )
    
    def _validate_dimensions(self, prs: Presentation, metrics: Dict) -> List[ValidationIssue]:
        """Validate slide dimensions"""
        issues = []
        
        width = emu_to_inches(prs.slide_width)
        height = emu_to_inches(prs.slide_height)
        ratio = width / height if height > 0 else 0
        
        metrics['width'] = width
        metrics['height'] = height
        metrics['aspect_ratio'] = ratio
        
        # Check dimensions
        if abs(width - Standards.SLIDE_WIDTH_INCHES) > 0.01:
            issues.append(ValidationIssue(
                'ERROR', 'Dimensions',
                f'Incorrect width: {width:.3f}in (expected {Standards.SLIDE_WIDTH_INCHES}in)'
            ))
        
        if abs(height - Standards.SLIDE_HEIGHT_INCHES) > 0.01:
            issues.append(ValidationIssue(
                'ERROR', 'Dimensions',
                f'Incorrect height: {height:.3f}in (expected {Standards.SLIDE_HEIGHT_INCHES}in)'
            ))
        
        # Check aspect ratio
        if abs(ratio - Standards.ASPECT_RATIO) > 0.01:
            issues.append(ValidationIssue(
                'WARNING', 'Dimensions',
                f'Non-standard aspect ratio: {ratio:.3f} (expected {Standards.ASPECT_RATIO:.3f})'
            ))
        
        return issues
    
    def _validate_colors(self, prs: Presentation, metrics: Dict) -> List[ValidationIssue]:
        """Validate color palette usage"""
        issues = []
        colors_used = set()
        non_standard_colors = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                # Check fill colors
                try:
                    if hasattr(shape, 'fill') and shape.fill.fore_color and shape.fill.fore_color.rgb:
                        rgb = rgb_to_tuple(shape.fill.fore_color.rgb)
                        if rgb:
                            colors_used.add(rgb)
                            if not is_approved_color(rgb):
                                non_standard_colors.append(rgb)
                except:
                    pass
                
                # Check text colors
                try:
                    if hasattr(shape, 'text_frame'):
                        for para in shape.text_frame.paragraphs:
                            if para.font.color and para.font.color.rgb:
                                rgb = rgb_to_tuple(para.font.color.rgb)
                                if rgb:
                                    colors_used.add(rgb)
                except:
                    pass
        
        metrics['colors_used'] = len(colors_used)
        metrics['non_standard_colors'] = len(set(non_standard_colors))
        
        # Check for required color families
        has_cyan = any(is_approved_color(c, 80) and c[1] > 150 and c[2] > 150 for c in colors_used)
        has_coral = any(c[0] > 200 and 50 < c[1] < 180 and c[2] < 100 for c in colors_used)
        has_navy = any(c[0] < 50 and c[1] < 80 and c[2] < 120 for c in colors_used)
        
        if not has_cyan:
            issues.append(ValidationIssue(
                'WARNING', 'Colors',
                'Missing cyan accent colors (alpha-wave spectrum)'
            ))
        
        if not has_coral:
            issues.append(ValidationIssue(
                'INFO', 'Colors',
                'No coral accent detected (dopamine trigger)'
            ))
        
        if not has_navy:
            issues.append(ValidationIssue(
                'WARNING', 'Colors',
                'Missing navy/dark colors for depth'
            ))
        
        if len(set(non_standard_colors)) > 5:
            issues.append(ValidationIssue(
                'WARNING', 'Colors',
                f'{len(set(non_standard_colors))} non-standard colors detected'
            ))
        
        return issues
    
    def _validate_fonts(self, prs: Presentation, metrics: Dict) -> List[ValidationIssue]:
        """Validate font usage"""
        issues = []
        fonts_used = set()
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        if para.font.name:
                            fonts_used.add(para.font.name)
        
        metrics['fonts_used'] = list(fonts_used)
        
        # Check for non-approved fonts
        non_approved = fonts_used - Standards.ALLOWED_FONTS
        if non_approved:
            issues.append(ValidationIssue(
                'ERROR', 'Typography',
                f'Non-approved fonts: {", ".join(non_approved)}',
                'Only Arial and Georgia are allowed'
            ))
        
        # Check font variety
        if len(fonts_used) > 3:
            issues.append(ValidationIssue(
                'WARNING', 'Typography',
                f'Too many fonts ({len(fonts_used)}). Use 2-3 maximum.'
            ))
        
        return issues
    
    def _validate_typography(self, prs: Presentation, metrics: Dict) -> List[ValidationIssue]:
        """Validate typography hierarchy"""
        issues = []
        sizes_used = set()
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        if para.font.size:
                            sizes_used.add(para.font.size.pt)
        
        metrics['font_sizes'] = sorted(sizes_used)
        
        # Check for hierarchy
        if sizes_used:
            min_size = min(sizes_used)
            max_size = max(sizes_used)
            
            if max_size - min_size < 8:
                issues.append(ValidationIssue(
                    'WARNING', 'Typography',
                    'Limited size hierarchy. Consider more variation for visual hierarchy.'
                ))
            
            # Check for extremely small text
            if min_size < 6:
                issues.append(ValidationIssue(
                    'WARNING', 'Typography',
                    f'Text too small ({min_size}pt). Minimum recommended: 7pt'
                ))
            
            # Check for extremely large text
            if max_size > 60:
                issues.append(ValidationIssue(
                    'INFO', 'Typography',
                    f'Very large text ({max_size}pt). Ensure readability.'
                ))
        
        return issues
    
    def _validate_components(self, prs: Presentation, metrics: Dict) -> List[ValidationIssue]:
        """Validate required components"""
        issues = []
        
        if not prs.slides:
            issues.append(ValidationIssue('ERROR', 'Structure', 'No slides in presentation'))
            return issues
        
        slide = prs.slides[0]
        
        # Check for header (shape at top spanning width)
        has_header = False
        for shape in slide.shapes:
            t = emu_to_inches(shape.top)
            w = emu_to_inches(shape.width)
            h = emu_to_inches(shape.height)
            if t < 0.2 and w > 10 and 0.5 < h < 2:
                has_header = True
                break
        
        # Check for footer (shape at bottom)
        has_footer = False
        for shape in slide.shapes:
            t = emu_to_inches(shape.top)
            if t > 6.5:
                has_footer = True
                break
        
        # Check for medical icon (+ symbol)
        has_medical_icon = False
        for shape in slide.shapes:
            if hasattr(shape, 'text') and '+' in (shape.text or ''):
                t = emu_to_inches(shape.top)
                if t < 2:  # In header region
                    has_medical_icon = True
                    break
        
        metrics['has_header'] = has_header
        metrics['has_footer'] = has_footer
        metrics['has_medical_icon'] = has_medical_icon
        
        if not has_header:
            issues.append(ValidationIssue(
                'WARNING', 'Components',
                'No header detected'
            ))
        
        if not has_footer:
            issues.append(ValidationIssue(
                'WARNING', 'Components',
                'No footer detected'
            ))
        
        return issues
    
    def _validate_layout(self, prs: Presentation, metrics: Dict) -> List[ValidationIssue]:
        """Validate layout bounds"""
        issues = []
        
        if not prs.slides:
            return issues
        
        slide = prs.slides[0]
        out_of_bounds = []
        
        for shape in slide.shapes:
            right = emu_to_inches(shape.left + shape.width)
            bottom = emu_to_inches(shape.top + shape.height)
            
            if right > Standards.SLIDE_WIDTH_INCHES + 0.1:
                out_of_bounds.append(f'right edge: {right:.2f}in')
            if bottom > Standards.SLIDE_HEIGHT_INCHES + 0.1:
                out_of_bounds.append(f'bottom edge: {bottom:.2f}in')
        
        metrics['out_of_bounds'] = len(out_of_bounds)
        
        if out_of_bounds:
            issues.append(ValidationIssue(
                'ERROR', 'Layout',
                f'Content exceeds slide bounds: {", ".join(out_of_bounds[:3])}'
            ))
        
        return issues
    
    def _validate_shapes(self, prs: Presentation, metrics: Dict) -> List[ValidationIssue]:
        """Validate shape count and optimization"""
        issues = []
        
        if not prs.slides:
            return issues
        
        shape_count = len(prs.slides[0].shapes)
        metrics['shape_count'] = shape_count
        
        if shape_count < Standards.MIN_SHAPES:
            issues.append(ValidationIssue(
                'WARNING', 'Optimization',
                f'Very few shapes ({shape_count}). Slide may appear sparse.'
            ))
        
        if shape_count > Standards.MAX_SHAPES:
            issues.append(ValidationIssue(
                'ERROR', 'Optimization',
                f'Too many shapes ({shape_count}). May affect performance.'
            ))
        elif shape_count > Standards.OPTIMAL_MAX:
            issues.append(ValidationIssue(
                'WARNING', 'Optimization',
                f'High shape count ({shape_count}). Consider simplification.'
            ))
        
        return issues
    
    def validate_directory(self, dirpath: str) -> ValidationReport:
        """Validate all PPTX files in directory"""
        results = []
        
        for filename in os.listdir(dirpath):
            if filename.endswith('.pptx'):
                filepath = os.path.join(dirpath, filename)
                result = self.validate_file(filepath)
                results.append(result)
        
        passed = sum(1 for r in results if r.is_valid)
        failed = len(results) - passed
        avg_score = sum(r.score for r in results) / len(results) if results else 0
        
        return ValidationReport(
            templates_checked=len(results),
            templates_passed=passed,
            templates_failed=failed,
            overall_score=avg_score,
            results=results
        )


# =============================================================================
# REPORT PRINTER
# =============================================================================

def print_report(report: ValidationReport):
    """Print formatted validation report"""
    print("\n" + "="*70)
    print("  TEMPLATE VALIDATION REPORT")
    print("="*70)
    
    print(f"\n  📊 SUMMARY")
    print(f"     Templates Checked: {report.templates_checked}")
    print(f"     Passed: {report.templates_passed}")
    print(f"     Failed: {report.templates_failed}")
    print(f"     Overall Score: {report.overall_score:.1f}%")
    
    print("\n" + "-"*70)
    print("  INDIVIDUAL RESULTS")
    print("-"*70)
    
    for result in sorted(report.results, key=lambda x: x.filename):
        status = "✅" if result.is_valid else "❌"
        print(f"\n  {status} {result.filename}")
        print(f"     Score: {result.score:.0f}%")
        
        if result.metrics:
            m = result.metrics
            details = []
            if 'shape_count' in m:
                details.append(f"shapes={m['shape_count']}")
            if 'colors_used' in m:
                details.append(f"colors={m['colors_used']}")
            if 'fonts_used' in m:
                details.append(f"fonts={m['fonts_used']}")
            if details:
                print(f"     Metrics: {', '.join(details)}")
        
        errors = [i for i in result.issues if i.level == 'ERROR']
        warnings = [i for i in result.issues if i.level == 'WARNING']
        
        if errors:
            for issue in errors:
                print(f"     🔴 ERROR: {issue.message}")
        if warnings:
            for issue in warnings[:3]:
                print(f"     🟡 WARNING: {issue.message}")
    
    print("\n" + "="*70)
    
    # Final verdict
    if report.overall_score >= 90:
        print("  ✨ EXCELLENT - Templates meet all standards")
    elif report.overall_score >= 80:
        print("  ✅ GOOD - Minor issues detected")
    elif report.overall_score >= 70:
        print("  🟡 ACCEPTABLE - Some improvements recommended")
    else:
        print("  🔴 NEEDS WORK - Review errors above")
    
    print("="*70 + "\n")


# =============================================================================
# MAIN
# =============================================================================

def main():
    import argparse
    
    parser = argparse.ArgumentParser(description="Validate nursing lecture templates")
    parser.add_argument('path', nargs='?', default='.', help="Directory or file to validate")
    parser.add_argument('--strict', action='store_true', help="Treat warnings as errors")
    parser.add_argument('--quiet', action='store_true', help="Only show summary")
    
    args = parser.parse_args()
    
    validator = TemplateValidator(strict=args.strict)
    
    if os.path.isfile(args.path):
        result = validator.validate_file(args.path)
        report = ValidationReport(
            templates_checked=1,
            templates_passed=1 if result.is_valid else 0,
            templates_failed=0 if result.is_valid else 1,
            overall_score=result.score,
            results=[result]
        )
    elif os.path.isdir(args.path):
        report = validator.validate_directory(args.path)
    else:
        print(f"Error: Path not found: {args.path}")
        sys.exit(1)
    
    if not args.quiet:
        print_report(report)
    else:
        print(f"Score: {report.overall_score:.1f}% | Passed: {report.templates_passed}/{report.templates_checked}")
    
    sys.exit(0 if report.templates_failed == 0 else 1)


if __name__ == "__main__":
    main()
